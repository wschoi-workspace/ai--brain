"""
RXR 2-Layer Case Study: 가나초콜릿 vs 가나초콜릿하우스
HTML 1:1 완전 재현 — 모든 텍스트·테이블·박스 빠짐없이 반영
R디자인가이드 (#6666FF, Pretendard, 다크모드, 16:9)
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── 컬러 ───
P  = RGBColor(0x66,0x66,0xFF)  # Primary
PD = RGBColor(0x53,0x53,0xFF)  # Deep
BK = RGBColor(0x00,0x00,0x00)
D1 = RGBColor(0x1A,0x1A,0x1A)
D2 = RGBColor(0x26,0x26,0x26)
D3 = RGBColor(0x30,0x30,0x30)
W  = RGBColor(0xFF,0xFF,0xFF)
G1 = RGBColor(0xA1,0xA1,0xAA)
G2 = RGBColor(0x71,0x71,0x80)
GR = RGBColor(0x10,0xB9,0x81)
WN = RGBColor(0xF5,0x9E,0x0B)
RD = RGBColor(0xC0,0x00,0x00)
CO = RGBColor(0xFF,0x50,0x50)
BL = RGBColor(0x4D,0x93,0xF7)
PK = RGBColor(0xD9,0x46,0xEF)
CH = RGBColor(0x8B,0x5E,0x3C)
SP = RGBColor(0x7C,0x3A,0xED)
F  = "Pretendard"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def bg(s, c=BK):
    f=s.background.fill; f.solid(); f.fore_color.rgb=c
def rr(s, l, t, w, h, c=D2):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,l,t,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()
    try: sh.adjustments[0]=0.06
    except: pass
    return sh
def ln(s, l, t, w, h=Inches(0.04), c=P):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb=c; sh.line.fill.background()
def bx(s, l, t, w, h, txt="", sz=14, c=W, b=False, a=PP_ALIGN.LEFT):
    tb=s.shapes.add_textbox(l,t,w,h)
    tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.text=txt; p.font.size=Pt(sz); p.font.color.rgb=c; p.font.bold=b; p.font.name=F; p.alignment=a
    return tb
def pa(tf, txt, sz=14, c=W, b=False, a=PP_ALIGN.LEFT, sp=0):
    p=tf.add_paragraph(); p.text=txt; p.font.size=Pt(sz); p.font.color.rgb=c; p.font.bold=b; p.font.name=F; p.alignment=a
    if sp: p.space_before=Pt(sp)
def btm(s): ln(s, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08))
def stag(s, t): bx(s, Inches(0.8), Inches(0.5), Inches(4), Inches(0.3), t, 11, P, True)
def stit(s, t, sz=28): bx(s, Inches(0.8), Inches(0.85), Inches(11.5), Inches(0.6), t, sz, W, True)
def tb(s, l, t, w, h, r, c):
    return s.shapes.add_table(r,c,l,t,w,h).table
def hd(t, txts, bg_c=P):
    for i,txt in enumerate(txts):
        c=t.rows[0].cells[i]; c.text=txt; c.fill.solid(); c.fill.fore_color.rgb=bg_c
        for p in c.text_frame.paragraphs:
            p.font.color.rgb=W; p.font.size=Pt(9); p.font.bold=True; p.font.name=F
def cl(c, txt, sz=9, b=False, fc=W, bg_c=D2):
    c.text=txt; c.fill.solid(); c.fill.fore_color.rgb=bg_c
    for p in c.text_frame.paragraphs:
        p.font.size=Pt(sz); p.font.bold=b; p.font.color.rgb=fc; p.font.name=F
def gb(s, l, t, pct, mw, h, c, label=""):
    rr(s, Inches(l), Inches(t), Inches(mw), Inches(h), D3)
    bw=max(mw*pct/100, 0.1)
    rr(s, Inches(l), Inches(t), Inches(bw), Inches(h), c)
    if label: bx(s, Inches(l+mw+0.1), Inches(t-0.01), Inches(0.8), Inches(h+0.02), label, 9, G1)

# ══════════════════════════════════════
# S1: 표지
# ══════════════════════════════════════
s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
ln(s,Inches(0),Inches(0),Inches(13.333),Inches(0.06))
bx(s,Inches(1),Inches(1.2),Inches(6),Inches(0.4),"RXR 2-LAYER CASE STUDY",13,P,True)
bx(s,Inches(1),Inches(2.0),Inches(11),Inches(1.5),"가나초콜릿 vs 가나초콜릿하우스\n동일 브랜드, 완전히 다른 심리 구조",36,W,True)
bx(s,Inches(1),Inches(4.0),Inches(9),Inches(0.8),"제품 리뷰와 공간 리뷰를 RXR 2-Layer 프레임워크(Content Layer + Psyche Layer)로 분석하여,\n단일 레이어 분석이 놓치는 심리적 전환을 실증한다.",14,G1)
bx(s,Inches(1),Inches(5.5),Inches(5),Inches(0.4),"Project RENT  ·  R-lab  ·  2026.04",12,G2)
btm(s)

# ══════════════════════════════════════
# S2: 분석 대상 & 데이터 개요
# ══════════════════════════════════════
s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
stag(s,"SECTION 1"); stit(s,"분석 대상 & 데이터 개요")
# 제품 카드
rr(s,Inches(0.8),Inches(1.8),Inches(5.6),Inches(3.2),D2); ln(s,Inches(0.8),Inches(1.8),Inches(5.6),Inches(0.05),CH)
bx(s,Inches(1.1),Inches(2.0),Inches(5),Inches(0.35),"A  가나초콜릿 (제품)",16,CH,True)
t=bx(s,Inches(1.1),Inches(2.5),Inches(5),Inches(2.3),"유형: 판형 초콜릿 (FMCG)",11,G1)
pa(t.text_frame,"역사: 1975~ (51년), 누적 매출 1.3조 원",11,G1,sp=3)
pa(t.text_frame,"누적 판매량: 66억 갑 (초당 4개)",11,G1,sp=3)
pa(t.text_frame,"포지셔닝: \"국민 초콜릿\"",11,G1,sp=3)
pa(t.text_frame,"수집 리뷰: 커뮤니티/블로그/쇼핑몰 기반",11,G1,sp=3)
# 공간 카드
rr(s,Inches(6.9),Inches(1.8),Inches(5.6),Inches(3.2),D2); ln(s,Inches(6.9),Inches(1.8),Inches(5.6),Inches(0.05),SP)
bx(s,Inches(7.2),Inches(2.0),Inches(5),Inches(0.35),"B  가나초콜릿하우스 (공간)",16,SP,True)
t=bx(s,Inches(7.2),Inches(2.5),Inches(5),Inches(2.3),"유형: 체험형 팝업스토어 / 디저트 카페",11,G1)
pa(t.text_frame,"역사: 2022~ (시즌3까지), 150평 규모",11,G1,sp=3)
pa(t.text_frame,"구매전환율: 95%+, 예약 1분 만에 마감",11,G1,sp=3)
pa(t.text_frame,"포지셔닝: \"가나, 디저트가 되다\"",11,G1,sp=3)
pa(t.text_frame,"수집 리뷰: 블로그/SNS/리뷰플랫폼 기반",11,G1,sp=3)
# 분석 프레임 박스
rr(s,Inches(0.8),Inches(5.2),Inches(11.7),Inches(2.0),D3); ln(s,Inches(0.8),Inches(5.2),Inches(0.05),Inches(2.0),P)
bx(s,Inches(1.1),Inches(5.3),Inches(11),Inches(0.3),"분석 프레임",13,P,True)
t=bx(s,Inches(1.1),Inches(5.7),Inches(11),Inches(0.3),"Layer 1 (Content) - n-gram, 감성분류, 토픽 클러스터링으로 \"무엇을 말했는가\" 분석",11,G1)
pa(t.text_frame,"Layer 2 (Psyche) - LIWC 4대 변수 + 대명사/조사/어미 패턴으로 \"어떻게 말했는가\" 분석",11,G1,sp=3)
pa(t.text_frame,"Layer 3 (Synthesis) - What x How 교차 분석으로 \"왜 다른 반응인가\" 진단",11,G1,sp=3)
btm(s)

# ══════════════════════════════════════
# S3: L1 감성 분포 비교
# ══════════════════════════════════════
s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
stag(s,"SECTION 2 — LAYER 1"); stit(s,"2-1. 감성 분포 비교: \"무엇을 말했는가\"")
# 제품
rr(s,Inches(0.8),Inches(1.8),Inches(5.6),Inches(5.4),D2)
bx(s,Inches(1.1),Inches(1.95),Inches(5),Inches(0.3),"A  가나초콜릿 - 감성 분포",14,CH,True)
for lbl,y,pct,c,lb in [("긍정",2.5,48,GR,"48%"),("중립",2.9,22,G2,"22%"),("부정",3.3,30,CO,"30%")]:
    bx(s,Inches(1.1),Inches(y),Inches(0.8),Inches(0.25),lbl,10,c,True); gb(s,2.0,y+0.02,pct,3.8,0.22,c,lb)
bx(s,Inches(1.1),Inches(3.75),Inches(5),Inches(0.2),"부정 30% - \"팜유\", \"싸구려 기름맛\", \"맛없다\" 직접 비판",9,G2)
t=bx(s,Inches(1.1),Inches(4.2),Inches(5),Inches(0.2),"대표 긍정:",10,G1,True)
pa(t.text_frame,"\"입에서 사르르 녹는 달콤한 맛이 특징이다\"",9,G2,sp=2)
pa(t.text_frame,"\"난 아직도 갈색이랑 빨간색 맛의 차이를 모르겠음\"",9,G2,sp=2)
pa(t.text_frame,"",4,sp=4)
pa(t.text_frame,"대표 부정:",10,CO,True,sp=2)
pa(t.text_frame,"\"카카오버터 대신 싸구려 팜유로 만드니까\"",9,G2,sp=2)
pa(t.text_frame,"\"가나초콜렛은 가나초콜렛이지 초콜렛은 아니야\"",9,G2,sp=2)
# 공간
rr(s,Inches(6.9),Inches(1.8),Inches(5.6),Inches(5.4),D2)
bx(s,Inches(7.2),Inches(1.95),Inches(5),Inches(0.3),"B  가나초콜릿하우스 - 감성 분포",14,SP,True)
for lbl,y,pct,c,lb in [("긍정",2.5,86,GR,"86%"),("중립",2.9,10,G2,"10%"),("부정",3.3,4,CO,"4%")]:
    bx(s,Inches(7.2),Inches(y),Inches(0.8),Inches(0.25),lbl,10,c,True); gb(s,8.1,y+0.02,pct,3.8,0.22,c,lb)
bx(s,Inches(7.2),Inches(3.75),Inches(5),Inches(0.2),"부정은 품절/웨이팅 불만에 집중, 공간 자체 비판은 극소",9,G2)
t=bx(s,Inches(7.2),Inches(4.2),Inches(5),Inches(0.2),"대표 긍정:",10,G1,True)
pa(t.text_frame,"\"결론부터 이야기하자면 생각보다 더 좋았습니다!\"",9,G2,sp=2)
pa(t.text_frame,"\"처음에는 홍보 팝업으로 생각했는데 완성도에 감동\"",9,G2,sp=2)
pa(t.text_frame,"\"다른 카페 가지 말고 여기 있다 가자\"",9,G2,sp=2)
pa(t.text_frame,"",4,sp=4)
pa(t.text_frame,"대표 부정:",10,CO,True,sp=2)
pa(t.text_frame,"\"별거 볼 건 없더라고요ㅠㅠ 초콜릿 특색이 덜 보이는 점 아쉬움\"",9,G2,sp=2)
btm(s)

# ══════════════════════════════════════
# S4: L1 토픽 클러스터 & 키워드 비교
# ══════════════════════════════════════
s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
stag(s,"SECTION 2 — LAYER 1"); stit(s,"2-2. 토픽 클러스터 & 키워드 비교")
t=tb(s,Inches(0.8),Inches(1.7),Inches(11.7),Inches(5.0),8,5)
hd(t,["토픽","제품 비중","제품 대표 키워드","공간 비중","공간 대표 키워드"])
data=[
    ("맛/식감","42%","\"달콤\", \"사르르 녹는\", \"부드럽다\", \"쌉싸름\"","18%","\"진한 농도\", \"조화\", \"단짠단짠\""),
    ("공간/분위기","0%","-","36%","\"고급스럽다\", \"빈티지\", \"포토존\", \"초콜릿 왕국\""),
    ("가격/가성비","18%","\"납득할 만한\", \"가격 인상\", \"저가형\"","8%","\"합리적\", \"35,000원 페어링\""),
    ("품질/원료","25%","\"팜유\", \"카카오 함량\", \"카카오버터\"","2%","-"),
    ("추억/브랜드","15%","\"어릴 때\", \"국민 초콜릿\", \"추억\"","6%","\"50주년\", \"헤리티지\""),
    ("관계/사회적 경험","0%","-","22%","\"데이트\", \"같이 가자\", \"가족\", \"눈치 싸움\""),
    ("SNS/포토","0%","-","8%","\"사진\", \"인스타\", \"포토부스\""),
]
hi_p={"42%","25%","18%","15%"}; hi_s={"36%","22%","18%"}
for i,(tp,pv,pk,sv,sk) in enumerate(data):
    r=i+1
    cl(t.rows[r].cells[0],tp,9,True)
    cl(t.rows[r].cells[1],pv,10,True,CH if pv in hi_p else W)
    cl(t.rows[r].cells[2],pk,8)
    cl(t.rows[r].cells[3],sv,10,True,SP if sv in hi_s else W)
    cl(t.rows[r].cells[4],sk,8)
for i,w in enumerate([1.5,1.0,3.5,1.0,4.7]): t.columns[i].width=Inches(w)
rr(s,Inches(0.8),Inches(6.9),Inches(11.7),Inches(0.4),D3)
bx(s,Inches(1.0),Inches(6.92),Inches(11.3),Inches(0.35),"핵심:  제품은 맛/원료(42%+25%=67%)에 집중  |  공간은 분위기+관계(36%+22%=58%)에 집중  →  소비 경험의 질적 전환",9,G1)
btm(s)

# ══════════════════════════════════════
# S5: L2 LIWC 4대 변수 — 게이지 + 테이블
# ══════════════════════════════════════
s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
stag(s,"SECTION 3 — LAYER 2"); stit(s,"3-1. LIWC 4대 변수 비교: \"어떻게 말했는가\"",26)
t=tb(s,Inches(0.8),Inches(1.7),Inches(11.7),Inches(5.5),5,5)
hd(t,["변수","제품","공간","변화","기능어 근거"])
liwc=[
    ("Authenticity\n(진정성)","65 (높음)","45 (중간)","↓ 20",
     "제품: \"난 아직도\", \"솔직히\", 부정 표현 포함 → 자기 노출 높음\n공간: 감탄사 과다, 1인칭 부재, 부정 표현 거의 없음 → 사회적 바람직성 작용"),
    ("Clout\n(추천 확신도)","35 (낮음)","65 (높음)","↑ 30",
     "제품: \"맛있다\" 단순 평가, 추천 발화 없음\n공간: \"꼭 가보세요\", \"여기 있다 가자\", \"-자\" 청유형 → 타인에 대한 영향력 행사"),
    ("Analytical\n(인식 방식)","65 (분석적)","45 (감성적)","↓ 20",
     "제품: \"카카오 함량\", \"팜유\", \"가격 대비\" → 인과적/비교적 서술\n공간: \"대박\", \"미쳤다\", \"환상\" → 감탄형/직관적 반응"),
    ("Emotional Tone\n(감정 강도)","55 (중립)","80 (고강도 긍정)","↑ 25",
     "제품: 긍정/부정 혼재 → 톤 중화\n공간: \"감동\", \"왕국\", \"환상\" + 부정어 거의 없음 → 순수 긍정 집중"),
]
pc=[GR,WN,BL,G1]; sc=[WN,GR,WN,GR]; dc=[CO,GR,CO,GR]
for i,(v,pv,sv,dv,ev) in enumerate(liwc):
    r=i+1
    cl(t.rows[r].cells[0],v,9,True,P)
    cl(t.rows[r].cells[1],pv,10,True,pc[i])
    cl(t.rows[r].cells[2],sv,10,True,sc[i])
    cl(t.rows[r].cells[3],dv,12,True,dc[i])
    cl(t.rows[r].cells[4],ev,8,False,G1)
for i,w in enumerate([1.8,1.5,1.5,0.9,6.0]): t.columns[i].width=Inches(w)
btm(s)

# ══════════════════════════════════════
# S6: L2 추가 심리 지표 (RXR 고유)
# ══════════════════════════════════════
s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
stag(s,"SECTION 3 — LAYER 2 (RXR EXCLUSIVE)"); stit(s,"3-2. 추가 심리 지표 (RXR 고유)",26)
t=tb(s,Inches(0.8),Inches(1.7),Inches(11.7),Inches(4.6),5,5)
hd(t,["지표","제품","공간","기능어 근거","해석"])
psy=[
    ("Freshness Index\n(\"이/가\" vs \"은/는\")","0.35\n\"은/는\" 우세","0.78\n\"이/가\" 우세",
     "제품: \"가나초콜릿은\", \"맛은\" → 이미 알려진 주제\n공간: \"분위기가\", \"조화가\", \"공간이\" → 새로운 발견의 초점화",
     "제품은 51년 된 익숙한 존재\n공간은 새로운 발견으로 체험"),
    ("Psych. Distance\n(대명사 패턴)","0.82\n\"나/내\" 지배적","0.41\n\"여기/우리\" 증가",
     "제품: \"난 아직도\", \"내 입속으로\" → 개인 경험 서술\n공간: \"여기 있다 가자\", \"우리 같이\" → 장소 귀속+집단 경험",
     "제품은 관찰자 시점\n공간은 참여자/귀속 시점"),
    ("시제 프레임\n(과거/현재/미래)","과거 회고 우세\n\"-었-\", \"어릴 때\"","현재+미래 우세\n\"-ㄴ다\", \"-자\"",
     "제품: \"먹었다\", \"어릴 때부터\" → 기억 속의 브랜드\n공간: \"펼쳐진다\", \"가자\" → 지금-여기 + 행동 촉구",
     "제품은 노스탤지어\n공간은 현재 몰입+확산 의지"),
    ("감각 축\n(지배 감각)","미각+촉각\n근접 감각(Proximal)","시각+공간\n원격 감각(Distal)",
     "제품: \"녹는\", \"부드럽다\", \"달콤한\" → 입안 감각\n공간: \"고급스럽다\", \"포토존\", \"분위기\" → 눈과 몸의 공간 경험",
     "제품→공간 전환 시\n감각 축 자체가 이동"),
]
for i,(nm,pv,sv,ev,ip) in enumerate(psy):
    r=i+1
    cl(t.rows[r].cells[0],nm,9,True,P)
    cl(t.rows[r].cells[1],pv,9,True,CH)
    cl(t.rows[r].cells[2],sv,9,True,SP)
    cl(t.rows[r].cells[3],ev,8,False,G1)
    cl(t.rows[r].cells[4],ip,8,False,W)
for i,w in enumerate([1.8,1.5,1.5,4.5,2.4]): t.columns[i].width=Inches(w)
# L2 결론 박스
rr(s,Inches(0.8),Inches(6.4),Inches(11.7),Inches(0.9),D3); ln(s,Inches(0.8),Inches(6.4),Inches(0.05),Inches(0.9),PK)
bx(s,Inches(1.1),Inches(6.45),Inches(11.2),Inches(0.2),"Layer 2 결론:",11,PK,True)
t2=bx(s,Inches(1.1),Inches(6.7),Inches(11.2),Inches(0.2),"Auth ↓20: 솔직한 개인 평가 → 사회적 공유 발화  |  Clout ↑30: 수동적 소비자 → 능동적 추천자",9,G1)
pa(t2.text_frame,"Freshness 0.35→0.78: \"이미 아는 것\" → \"새로운 발견\"  |  이 변화는 Layer 1(감성 점수)만으로는 절대 포착할 수 없다.",9,G1,sp=2)
btm(s)

# ══════════════════════════════════════
# S7: 패턴 ID 레퍼런스 (6컬럼 — 의미 포함)
# ══════════════════════════════════════
s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
stag(s,"SECTION 4 — PATTERN DICTIONARY"); stit(s,"4-0. 2-Layer 교차 분석 패턴 사전",26)
bx(s,Inches(0.8),Inches(1.3),Inches(11),Inches(0.3),"Layer 1(감성)과 Layer 2(심리)를 교차하면 아래 6개 패턴으로 분류된다. 각 패턴은 전혀 다른 전략적 의미를 갖는다.",10,G1)
t=tb(s,Inches(0.8),Inches(1.7),Inches(11.7),Inches(5.5),7,6)
hd(t,["패턴 ID","이름","L1 (Content) 신호","L2 (Psyche) 신호","의미","대응 전략"])
pdata=[
    ("P-01","진짜 팬","긍정 80%+","Auth 높음\nClout 높음","진심으로 좋아하고,\n자신 있게 추천까지 한다","VIP 관리, 레퍼럴 앰배서더 전환,\n이 세그먼트의 리뷰를 마케팅에 활용"),
    ("P-02","인스타 소비","긍정 80%+\n\"사진\", \"예쁘다\"","Auth 낮음\nClout 낮음","바이럴은 되지만 진심이 아님.\n\"올리고 싶었을 뿐\"\n재방문 의향 낮음","체험 깊이 강화(시식/DIY),\n포토존 뒤에 제품 터치포인트 배치"),
    ("P-03","관심은 있지만\n확신 부족","긍정 70%+\n\"괜찮다\", \"좋다\"","Auth 높음\nClout 낮음","솔직하게 좋다고 느끼지만,\n남에게 추천까지는 안 함.\n\"나는 좋은데 네가 좋아할지는...\"","제품 체험 터치포인트 강화,\n샘플링 확대로 확신 전환 유도"),
    ("P-04","신선함 소진","키워드 변화 없음\n감성 소폭 하락","Freshness 급락\n(\"이/가\"→\"은/는\")","내용은 같지만 \"새로운 발견\"\n느낌이 사라짐.\n2~3주차부터 진행되는 체험 피로","콘텐츠 리프레시 (신규 메뉴/존 오픈),\n주간 리프레시 캘린더 적용"),
    ("P-05","탐색은 길지만\n결정 못 함","체류시간 ↑\n구매 ↓","Analytical 높음\nClout 낮음","비교/분석 모드에 갇혀 있음.\n\"좋은 건 알겠는데 뭘 사야 할지...\"\n선택 과부하(Choice Overload)","선택지 축소, 큐레이션 동선,\n직원 추천 / 베스트 메뉴 강조"),
    ("P-06","브랜드 피로도\n(조용한 이탈)","긍정 유지\nbut 언급량 ↓","Tone ↓ 추세\nClout ↓ 추세\nDistance 변화 없음","불만은 아니지만 관심이 사라짐.\n\"싫어서 안 가는 게 아니라\n그냥 생각이 안 나서 안 감\"","긴급 리프레시 필요,\n이벤트성 콘텐츠 투입,\n기존 고객 재활성화 프로그램"),
]
pclr=[GR,WN,P,PK,G2,CO]
for i,(pid,nm,l1,l2,mn,st) in enumerate(pdata):
    r=i+1
    cl(t.rows[r].cells[0],pid,10,True,pclr[i])
    cl(t.rows[r].cells[1],nm,8,True)
    cl(t.rows[r].cells[2],l1,7)
    cl(t.rows[r].cells[3],l2,7)
    cl(t.rows[r].cells[4],mn,7)
    cl(t.rows[r].cells[5],st,7)
for i,w in enumerate([0.8,1.2,1.6,1.8,2.8,3.5]): t.columns[i].width=Inches(w)
btm(s)

# ══════════════════════════════════════
# S8: Zone 매트릭스 + 읽는 법
# ══════════════════════════════════════
s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
stag(s,"SECTION 4 — ZONE MATRIX"); stit(s,"Zone 분류 (감성 x 진정성 매트릭스)",26)
zones=[
    (0.8,1.8,"Zone A: 진짜 팬","감성 긍정 + Auth 높음","진심으로 좋아한다. P-01과 동일 영역.\n재방문/재구매 확률 가장 높은 세그먼트.",GR),
    (6.9,1.8,"Zone B: 인스타 소비","감성 긍정 + Auth 낮음","\"예쁘니까 올린 것.\" P-02와 동일 영역.\n바이럴 기여는 하지만 재방문 가치 낮음.",WN),
    (0.8,4.1,"Zone C: 진심 불만","감성 부정 + Auth 높음","실제 문제가 존재한다는 신호. 무시 금지.\n\"팜유 싸구려\", \"특색 없다\" = 진짜 개선점.",CO),
    (6.9,4.1,"Zone D: 노이즈","감성 부정 + Auth 낮음","악의적/경쟁사 공격성 리뷰 가능성.\n분석 대상에서 필터링 후 제외.",G2),
]
for x,y,title,cond,desc,c in zones:
    rr(s,Inches(x),Inches(y),Inches(5.6),Inches(2.0),D2); ln(s,Inches(x),Inches(y),Inches(5.6),Inches(0.05),c)
    bx(s,Inches(x+0.2),Inches(y+0.2),Inches(5.2),Inches(0.3),title,15,c,True)
    bx(s,Inches(x+0.2),Inches(y+0.55),Inches(5.2),Inches(0.2),cond,10,P)
    bx(s,Inches(x+0.2),Inches(y+0.9),Inches(5.2),Inches(0.9),desc,10,G1)
# 읽는 법
rr(s,Inches(0.8),Inches(6.3),Inches(11.7),Inches(0.9),D3); ln(s,Inches(0.8),Inches(6.3),Inches(0.05),Inches(0.9),P)
bx(s,Inches(1.1),Inches(6.35),Inches(11.2),Inches(0.2),"읽는 법:",11,P,True)
bx(s,Inches(1.1),Inches(6.6),Inches(11.2),Inches(0.5),"아래 교차 분석에서 각 리뷰 유형에 패턴 ID가 부여된다. 패턴 ID를 보면 \"이 리뷰가 어떤 심리 상태이고,\n어떤 전략으로 대응해야 하는지\"가 즉시 파악된다. 이것이 단순 감성 분석(긍정/부정)과 2-Layer 분석의 결정적 차이.",9,G1)
btm(s)

# ══════════════════════════════════════
# S9: 제품 리뷰 교차 분석
# ══════════════════════════════════════
s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
stag(s,"SECTION 4 — CROSS ANALYSIS (PRODUCT)"); stit(s,"4-1. 제품 리뷰 - 패턴 매칭 결과",26)
t=tb(s,Inches(0.8),Inches(1.7),Inches(11.7),Inches(3.6),4,5)
hd(t,["리뷰 유형","L1 (Content)","L2 (Psyche)","패턴 ID","2-Layer 인사이트"])
rd=[
    ("솔직한 팬","긍정 / \"맛있다\", \"맞다\"","Auth 높음, Clout 낮음","P-03","관심은 있지만 추천 의지 없음.\n\"난 좋아하는데 남한테 권하기는...\" - 개인적 만족에 머무름"),
    ("분석적 비판자","부정 / \"팜유\", \"싸구려\"","Auth 높음, Analytical 높음","P-03 변형\n(Zone C)","진심 불만. 실제 품질 문제 인지.\n무시하면 안 됨"),
    ("추억형 소비자","긍정 / \"어릴 때\", \"추억\"","Auth 높음, Tone 중립,\n시제: 과거","P-06 전조","노스탤지어 의존. 현재 브랜드 경험이 아닌\n과거 기억에 의존 - 새 경험 없이 점진적 이탈 가능"),
]
pc2=[P,CO,WN]
for i,(typ,l1,l2,pat,ins) in enumerate(rd):
    r=i+1; cl(t.rows[r].cells[0],typ,9,True); cl(t.rows[r].cells[1],l1,8); cl(t.rows[r].cells[2],l2,8)
    cl(t.rows[r].cells[3],pat,10,True,pc2[i]); cl(t.rows[r].cells[4],ins,8)
for i,w in enumerate([1.3,2.2,2.5,1.3,4.4]): t.columns[i].width=Inches(w)
# 진단 박스
rr(s,Inches(0.8),Inches(5.5),Inches(11.7),Inches(1.7),D3); ln(s,Inches(0.8),Inches(5.5),Inches(0.05),Inches(1.7),CO)
bx(s,Inches(1.1),Inches(5.6),Inches(11.2),Inches(0.3),"제품 진단:",13,CO,True)
bx(s,Inches(1.1),Inches(5.95),Inches(11.2),Inches(1.1),"L1만 보면 \"긍정 48%, 그럭저럭\" - L2를 더하면 긍정 리뷰조차 추천 의지가 없고(Clout 35),\n과거 추억에 의존하고 있음이 드러남. 브랜드가 현재 시점에서 새로운 긍정 경험을 생산하지 못하고 있다.",11,G1)
btm(s)

# ══════════════════════════════════════
# S10: 공간 리뷰 교차 분석
# ══════════════════════════════════════
s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
stag(s,"SECTION 4 — CROSS ANALYSIS (SPACE)"); stit(s,"4-2. 공간 리뷰 - 패턴 매칭 결과",26)
t=tb(s,Inches(0.8),Inches(1.7),Inches(11.7),Inches(3.6),5,5)
hd(t,["리뷰 유형","L1 (Content)","L2 (Psyche)","패턴 ID","2-Layer 인사이트"])
rd2=[
    ("인스타 소비형","긍정 90%+ / \"예쁘다\", \"사진\"","Auth 낮음, Clout 낮음,\nTone 높음","P-02","바이럴 =/= 충성. SNS 공유는 하지만\n재방문 의향 약함. 인스타에 올리고 싶었을 뿐"),
    ("기대 위반형","긍정 / \"생각보다 좋았다\",\n\"감동\"","Auth 높음, Clout 높음,\nFreshness 높음","P-01","진짜 팬 전환. \"제품 팝업인 줄 알았는데 진심인 공간\"\n스키마 위반이 진정한 긍정 반응 생성"),
    ("관계 소비형","긍정 / \"같이 가자\",\n\"데이트\"","Clout 높음, Distance 낮음\n(우리/여기), 시제: 미래","P-01 변형","행동 전환 완료. 단순 감상을 넘어 동반자에게\n행동을 촉구 - 가장 높은 전환 가치"),
    ("실망형","부정 / \"별거 없다\",\n\"아쉽다\"","Auth 높음, Tone 낮음","Zone C","진심 불만. 소수지만 Auth 높아 실제 문제 반영.\n특히 \"초콜릿 특색 부족\"은 제품-공간 연결 약점 시사"),
]
pc3=[WN,GR,GR,CO]
for i,(typ,l1,l2,pat,ins) in enumerate(rd2):
    r=i+1; cl(t.rows[r].cells[0],typ,9,True); cl(t.rows[r].cells[1],l1,8); cl(t.rows[r].cells[2],l2,8)
    cl(t.rows[r].cells[3],pat,10,True,pc3[i]); cl(t.rows[r].cells[4],ins,8)
for i,w in enumerate([1.3,2.2,2.5,1.3,4.4]): t.columns[i].width=Inches(w)
# 진단 박스
rr(s,Inches(0.8),Inches(5.5),Inches(11.7),Inches(1.7),D3); ln(s,Inches(0.8),Inches(5.5),Inches(0.05),Inches(1.7),GR)
bx(s,Inches(1.1),Inches(5.6),Inches(11.2),Inches(0.3),"공간 진단: 긍정 86%를 분해하면",13,GR,True)
t2=bx(s,Inches(1.1),Inches(5.95),Inches(11.2),Inches(0.3),"- 인스타 소비 (P-02): 약 40~45% - Auth 낮음, Clout 낮음. 바이럴은 되지만 재방문 가치 낮음",10,G1)
pa(t2.text_frame,"- 진짜 팬 (P-01): 약 30~35% - Auth+Clout 모두 높음. 진정한 브랜드 옹호자",10,G1,sp=2)
pa(t2.text_frame,"- 관계 소비 (P-01 변형): 약 15~20% - 행동 전환까지 완료. 최고 가치 세그먼트",10,G1,sp=2)
pa(t2.text_frame,"Layer 1의 \"긍정 86%\"는 하나의 숫자지만, Layer 2를 추가하면 3개의 전혀 다른 세그먼트가 드러난다.",10,P,True,sp=6)
btm(s)

# ══════════════════════════════════════
# S11: 5대 심리축 전환 테이블 (6컬럼)
# ══════════════════════════════════════
s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
stag(s,"SECTION 5 — KEY FINDING"); stit(s,"핵심 발견: 5대 심리 축의 체계적 전환",26)
bx(s,Inches(0.8),Inches(1.3),Inches(11.5),Inches(0.3),"동일 브랜드 \"가나\"가 제품→공간으로 확장될 때, 소비자의 언어 구조에서 5개 심리 축이 동시에 전환된다.",10,G1)
t=tb(s,Inches(0.8),Inches(1.7),Inches(11.7),Inches(4.8),6,6)
hd(t,["심리 축","제품","→","공간","기능어 증거","비즈니스 의미"])
ax=[
    ("1. 자아 축\n\"나\"→\"여기/우리\"","개인적 경험","→","사회적 경험","1인칭 단수 ↓\n장소 지시어+복수형 ↑","공간이 소비를 \"공유 가능한 경험\"으로\n전환 - 바이럴 구조 형성"),
    ("2. 시간 축\n과거→현재+미래","노스탤지어","→","현재 몰입","과거 시제 ↓\n현재형+청유형(-자) ↑","브랜드가 \"기억\"에서 \"현재 경험\"으로\n리포지셔닝 - 노스탤지어 의존 탈피"),
    ("3. 감각 축\n미각→시각","입안 감각","→","공간 감각","미각/촉각 단어 ↓\n시각/공간 단어 ↑","근접(Proximal)→원격(Distal) 전환\nSNS 공유에 최적화된 감각 구조"),
    ("4. 인지 축\n분석→감동","이성적 평가","→","감성적 반응","Analytical 65→45\n인과 접속사 ↓, 감탄사 ↑","분석적 비교 대상→감정적 몰입 대상\n가격 경쟁에서 가치 경쟁으로"),
    ("5. 행동 축\n평가→추천","수동적 소비자","→","능동적 추천자","Clout 35→65\n청유형 어미 ↑","소비자가 자발적 영업사원으로 전환\n마케팅 비용 구조 변화"),
]
for i,(nm,pv,ar,sv,ev,bm) in enumerate(ax):
    r=i+1
    cl(t.rows[r].cells[0],nm,8,True,P)
    cl(t.rows[r].cells[1],pv,9,False,CH)
    cl(t.rows[r].cells[2],ar,14,True,P)
    cl(t.rows[r].cells[3],sv,9,True,SP)
    cl(t.rows[r].cells[4],ev,7,False,PK)
    cl(t.rows[r].cells[5],bm,7,False,G1)
for i,w in enumerate([1.8,1.3,0.5,1.3,2.8,4.0]): t.columns[i].width=Inches(w)
# 인사이트 박스
rr(s,Inches(0.8),Inches(6.6),Inches(11.7),Inches(0.7),D3); ln(s,Inches(0.8),Inches(6.6),Inches(0.05),Inches(0.7),GR)
bx(s,Inches(1.1),Inches(6.65),Inches(11.2),Inches(0.6),"이 전환 데이터의 가치: \"공간이 좋다\"는 누구나 알 수 있지만, \"5개 심리 축이 동시에 전환된다\"는 정량 데이터는 RXR 2-Layer만 제공할 수 있다.\nIR/제안서에서 \"Clout 35→65, Freshness 0.35→0.78, 감각 축 Proximal→Distal 전환\"이라고 말할 수 있는 근거.",9,G1)
btm(s)

# ══════════════════════════════════════
# S12: L1만 vs 2-Layer 비교
# ══════════════════════════════════════
s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
stag(s,"SECTION 6 — COMPARISON"); stit(s,"6-1. Layer 1만 vs 2-Layer 통합 - 인사이트 품질 비교",24)
t=tb(s,Inches(0.8),Inches(1.7),Inches(11.7),Inches(5.5),6,3)
hd(t,["질문","Layer 1만 답변","2-Layer 통합 답변"])
qa=[
    ("공간이 성공했나?","\"긍정 86%, 성공\"","\"긍정 86% 중 진짜 팬은 30~35%. 나머지 40~45%는\n인스타 소비형(Auth 낮음)으로 재방문 가치 낮음\""),
    ("제품 브랜드는 건강한가?","\"긍정 48%, 부정 30%, 보통\"","\"긍정 리뷰도 Clout 35 - 추천 의지 없음. 긍정의 대부분이\n과거 추억 기반(시제 과거 우세)이라 현재 경험 생산력 부족\""),
    ("공간이 제품 대비 나은 이유?","\"감성 점수가 더 높다\"","\"5개 심리 축(자아/시간/감각/인지/행동)이 동시에 전환되며\n소비자를 평가자에서 추천자로 변환\""),
    ("다음 팝업에서 뭘 바꿔야?","\"긍정 유지하면 된다\"","\"인스타 소비(P-02)를 진짜 팬(P-01)으로 전환하려면 체험 깊이\n강화 필요. 포토존만으로는 Auth 올릴 수 없음. DIY/테이스팅 배치\""),
    ("시즌4를 해야 하나?","\"긍정이니 해도 됨\"","\"Freshness 0.78은 '아직 신선'이지만, 시즌 반복 시 하락 필연.\n주간 리프레시 캘린더와 시즌별 차별화 콘텐츠 필수\""),
]
for i,(q,a1,a2) in enumerate(qa):
    r=i+1; cl(t.rows[r].cells[0],q,9,True); cl(t.rows[r].cells[1],a1,8,False,G1); cl(t.rows[r].cells[2],a2,8,False,W,D3)
for i,w in enumerate([2.5,3.0,6.2]): t.columns[i].width=Inches(w)
btm(s)

# ══════════════════════════════════════
# S13: 세그먼트별 전략
# ══════════════════════════════════════
s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
stag(s,"SECTION 6 — STRATEGY"); stit(s,"6-2. 가나초콜릿하우스를 위한 세그먼트별 전략",24)
segs=[
    ("인스타 소비형 (40~45%)",WN,
     "L1: 긍정 / \"예쁘다\", \"사진\"\nL2: Auth 낮음, Clout 낮음",
     "포토존 뒤에 제품 체험 동선 배치.\n\"사진 찍고 → 시식 → 구매\"로 연결.\nSNS 공유 시 체험 쿠폰 제공하여\n바이럴을 체험으로 전환"),
    ("진짜 팬 (30~35%)",GR,
     "L1: 긍정 / \"감동\", \"진심\"\nL2: Auth 높음, Clout 높음",
     "VIP 프로그램(시즌 프리오픈, 한정 메뉴).\n리뷰어 앰배서더 프로그램.\n이 세그먼트의 리뷰가\n다른 소비자의 P-02→P-01 전환을 유도"),
    ("관계 소비형 (15~20%)",P,
     "L1: 긍정 / \"같이\", \"데이트\"\nL2: Clout 높음, Distance 낮음",
     "2인/그룹 전용 메뉴, 커플/가족 체험 코스.\n\"함께 만드는 초콜릿\" DIY를 핵심 프로그램.\n관계 기억에 브랜드 각인"),
]
for i,(title,c,profile,strategy) in enumerate(segs):
    x=0.8+i*4.1
    rr(s,Inches(x),Inches(1.8),Inches(3.8),Inches(5.4),D2); ln(s,Inches(x),Inches(1.8),Inches(3.8),Inches(0.05),c)
    bx(s,Inches(x+0.2),Inches(1.95),Inches(3.4),Inches(0.35),title,14,c,True)
    bx(s,Inches(x+0.2),Inches(2.5),Inches(3.4),Inches(0.25),"프로필",10,G2,True)
    rr(s,Inches(x+0.2),Inches(2.8),Inches(3.4),Inches(1.3),D3)
    bx(s,Inches(x+0.35),Inches(2.85),Inches(3.1),Inches(1.2),profile,9,G1)
    bx(s,Inches(x+0.2),Inches(4.2),Inches(3.4),Inches(0.25),"전략",10,c,True)
    rr(s,Inches(x+0.2),Inches(4.5),Inches(3.4),Inches(2.5),D3)
    bx(s,Inches(x+0.35),Inches(4.55),Inches(3.1),Inches(2.4),strategy,10,W)
btm(s)

# ══════════════════════════════════════
# S14: 핵심 요약
# ══════════════════════════════════════
s=prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
ln(s,Inches(0),Inches(0),Inches(13.333),Inches(0.06))
bx(s,Inches(1),Inches(0.8),Inches(11.3),Inches(1.0),"같은 \"긍정 86%\"도, Layer 2를 추가하면\n3개의 전혀 다른 세그먼트가 보인다.",32,W,True,PP_ALIGN.CENTER)
# L1 vs 2-Layer 대비
rr(s,Inches(0.8),Inches(2.1),Inches(5.5),Inches(2.8),D2); ln(s,Inches(0.8),Inches(2.1),Inches(5.5),Inches(0.04),WN)
bx(s,Inches(1.1),Inches(2.2),Inches(5),Inches(0.25),"Layer 1만의 결론",12,WN,True)
t=bx(s,Inches(1.1),Inches(2.55),Inches(5),Inches(0.2),"- 제품: 긍정 48%, 보통",10,G1)
pa(t.text_frame,"- 공간: 긍정 86%, 대성공",10,G1,sp=2)
pa(t.text_frame,"- \"공간이 제품보다 좋다\"",10,G1,sp=2)
pa(t.text_frame,"- 전략? \"계속 하면 된다\"",10,G1,sp=2)

rr(s,Inches(6.9),Inches(2.1),Inches(5.6),Inches(2.8),D2); ln(s,Inches(6.9),Inches(2.1),Inches(5.6),Inches(0.04),P)
bx(s,Inches(7.2),Inches(2.2),Inches(5),Inches(0.25),"2-Layer 통합의 결론",12,P,True)
t=bx(s,Inches(7.2),Inches(2.55),Inches(5),Inches(0.2),"- 제품: Clout 35, 추천 의지 없음. 과거 추억 의존",10,W)
pa(t.text_frame,"- 공간: 진짜 팬 30~35%, 인스타 소비 40~45%",10,W,sp=2)
pa(t.text_frame,"- 제품→공간 전환 시 5개 심리 축 동시 전환 실증",10,W,sp=2)
pa(t.text_frame,"- 전략: 세그먼트별 차별화 + Freshness 모니터링",10,W,sp=2)

# 핵심 메시지
rr(s,Inches(0.8),Inches(5.1),Inches(11.7),Inches(0.6),D3)
bx(s,Inches(1.0),Inches(5.15),Inches(11.3),Inches(0.5),"Layer 1은 \"좋다/나쁘다\"를 말해준다. Layer 2는 \"왜 좋은지, 얼마나 진심인지, 언제까지 지속될지\"를 말해준다.",13,P,False,PP_ALIGN.CENTER)

# 5축 미니 카드
for i,(nm,val) in enumerate([("자아 축","\"나\"→\"우리\""),("시간 축","과거→현재"),("감각 축","미각→시각"),("인지 축","분석→감동"),("행동 축","평가→추천")]):
    x=1.2+i*2.3
    rr(s,Inches(x),Inches(5.9),Inches(2.0),Inches(1.0),D2); ln(s,Inches(x),Inches(5.9),Inches(2.0),Inches(0.04),P)
    bx(s,Inches(x),Inches(6.0),Inches(2.0),Inches(0.2),nm,9,P,True,PP_ALIGN.CENTER)
    bx(s,Inches(x),Inches(6.3),Inches(2.0),Inches(0.3),val,13,W,True,PP_ALIGN.CENTER)
bx(s,Inches(1),Inches(7.05),Inches(11.3),Inches(0.3),"Project RENT  ·  R-lab  ·  RXR 2-Layer Case Study  ·  2026.04",10,G2,False,PP_ALIGN.CENTER)
btm(s)

# ─── 저장 ───
out=os.path.join(os.path.dirname(__file__),"rxr-2layer-case-ghana-v3.pptx")
prs.save(out)
print("Done:",out)
print("Slides:",len(prs.slides))
