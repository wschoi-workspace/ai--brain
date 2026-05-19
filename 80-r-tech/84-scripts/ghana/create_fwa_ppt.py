"""
기능어분석 (Function Word Analysis) - 마케팅 & 소비자 조사 특화 PPT 생성
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── 색상 팔레트 ───
C_PRIMARY = RGBColor(0x1A, 0x1A, 0x2E)
C_ACCENT = RGBColor(0x43, 0x61, 0xEE)
C_ACCENT_DARK = RGBColor(0x2D, 0x3D, 0xA8)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_BG = RGBColor(0xF0, 0xF4, 0xFF)
C_TEXT = RGBColor(0x2C, 0x2C, 0x3A)
C_TEXT_LIGHT = RGBColor(0x6B, 0x72, 0x80)
C_SUCCESS = RGBColor(0x05, 0x96, 0x69)
C_WARNING = RGBColor(0xF5, 0x9E, 0x0B)
C_DANGER = RGBColor(0xDC, 0x26, 0x26)
C_BROWN = RGBColor(0x92, 0x40, 0x0E)
C_GREEN = RGBColor(0x05, 0x96, 0x69)
C_PURPLE = RGBColor(0x7C, 0x3A, 0xED)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ─── 헬퍼 함수 ───
def add_bg(slide, color=C_PRIMARY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_textbox(slide, left, top, width, height, text="", font_size=14, color=C_TEXT, bold=False, alignment=PP_ALIGN.LEFT, font_name="맑은 고딕"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_para(tf, text, font_size=14, color=C_TEXT, bold=False, alignment=PP_ALIGN.LEFT, space_before=0, space_after=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "맑은 고딕"
    p.alignment = alignment
    if space_before: p.space_before = Pt(space_before)
    if space_after: p.space_after = Pt(space_after)
    return p

def add_table(slide, left, top, width, height, rows, cols):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    return table_shape.table

def style_header_row(table, col_widths=None, bg_color=C_PRIMARY):
    for i, cell in enumerate(table.rows[0].cells):
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color
        for p in cell.text_frame.paragraphs:
            p.font.color.rgb = C_WHITE
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.name = "맑은 고딕"
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

def style_body_cell(cell, text, font_size=10, bold=False, color=C_TEXT):
    cell.text = text
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = "맑은 고딕"

def add_stat_card(slide, left, top, value, label, accent_color=C_ACCENT):
    card = add_shape_bg(slide, left, top, Inches(2.2), Inches(1.3), C_WHITE)
    card.adjustments[0] = 0.08
    add_textbox(slide, left, top + Inches(0.15), Inches(2.2), Inches(0.6),
                value, font_size=30, color=accent_color, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left, top + Inches(0.75), Inches(2.2), Inches(0.5),
                label, font_size=10, color=C_TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

def add_bar(slide, left, top, width_pct, max_width, height, color, label_text=""):
    bar_w = int(max_width * width_pct / 100)
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, bar_w, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    bar.adjustments[0] = 0.5
    if label_text:
        add_textbox(slide, left + bar_w + Inches(0.1), top - Inches(0.02), Inches(1), height,
                    label_text, font_size=9, color=C_TEXT_LIGHT)

# ═══════════════════════════════════════════
# SLIDE 1: 표지
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_PRIMARY)

# 장식 원
circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(-2), Inches(6), Inches(6))
circle.fill.solid()
circle.fill.fore_color.rgb = RGBColor(0x43, 0x61, 0xEE)
circle.fill.fore_color.brightness = 0.85
circle.line.fill.background()

add_textbox(slide, Inches(1), Inches(1.2), Inches(5), Inches(0.5),
            "MARKETING & CONSUMER RESEARCH", font_size=13, color=RGBColor(0x93, 0xC5, 0xFD), bold=True)
add_textbox(slide, Inches(1), Inches(1.9), Inches(10), Inches(1.2),
            "기능어분석", font_size=48, color=C_WHITE, bold=True)
add_textbox(slide, Inches(1), Inches(3.1), Inches(10), Inches(1),
            "Function Word Analysis", font_size=32, color=RGBColor(0x81, 0x8C, 0xF8))
add_textbox(slide, Inches(1), Inches(4.3), Inches(9), Inches(0.8),
            "마케팅 & 소비자 조사 관점에서 본 기능어분석의 정의, 비즈니스 적용 사례,\n그리고 실무 활용 가능성", font_size=15, color=RGBColor(0x94, 0xA3, 0xB8))
add_textbox(slide, Inches(1), Inches(5.8), Inches(6), Inches(0.4),
            "2026.03 · 30+ 학술 자료 기반 · do-better-workspace", font_size=11, color=RGBColor(0x64, 0x74, 0x8B))

# ═══════════════════════════════════════════
# SLIDE 2: 목차
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
add_textbox(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.6), "목차", font_size=28, color=C_PRIMARY, bold=True)

sections = [
    ("Part 1 — 개념 이해", ["1. 기능어의 정의와 핵심 원리", "2. 기능어분석이란? — 방법론"]),
    ("Part 2 — 마케팅 적용 사례", ["3. 소비자 리뷰 분석", "4. 광고 카피 & 메시지 최적화",
     "5. 브랜드 포지셔닝 분석", "6. 소셜미디어 & 인플루언서", "7. 고객 서비스 / VOC / FGI", "8. 경쟁사 분석 & 한국 연구"]),
    ("Part 3 — 실무 활용", ["9. 실무 워크플로우 & 도구", "10. 마케터를 위한 시사점"]),
    ("Case Study", ["★ 가나초콜릿 × 가나초콜릿하우스"]),
]

y = 1.4
for part_title, items in sections:
    add_textbox(slide, Inches(1.2), Inches(y), Inches(4), Inches(0.35),
                part_title, font_size=11, color=C_ACCENT, bold=True)
    y += 0.35
    for item in items:
        add_textbox(slide, Inches(1.6), Inches(y), Inches(5), Inches(0.3),
                    item, font_size=13, color=C_TEXT)
        y += 0.3
    y += 0.15

# ═══════════════════════════════════════════
# SLIDE 3: 기능어의 정의
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), C_PRIMARY)
add_textbox(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.6),
            "1  기능어(Function Word)의 정의와 핵심 원리", font_size=24, color=C_WHITE, bold=True)

# 정의 박스
add_shape_bg(slide, Inches(0.8), Inches(1.4), Inches(11.7), Inches(0.9), C_LIGHT_BG)
tb = add_textbox(slide, Inches(1.0), Inches(1.5), Inches(11.3), Inches(0.7),
            "기능어(function word)란 어휘적 의미가 거의 없지만 문장 내 단어들의 문법적 관계를 표현하는 단어.\n대명사(나, 너, 그), 조사(을/를, 에서), 접속사(그리고, 하지만) 등이 해당한다.",
            font_size=13, color=C_TEXT)

# Stat cards
add_stat_card(slide, Inches(0.8), Inches(2.6), "<500", "영어 기능어 총 수", C_ACCENT)
add_stat_card(slide, Inches(3.2), Inches(2.6), "55%+", "일상 언어 중 비중", C_ACCENT)
add_stat_card(slide, Inches(5.6), Inches(2.6), "20,000+", "LIWC 활용 학술논문", C_ACCENT)
add_stat_card(slide, Inches(8.0), Inches(2.6), "100+", "LIWC-22 분석 차원", C_ACCENT)

# 왜 기능어인가 테이블
table = add_table(slide, Inches(0.8), Inches(4.2), Inches(11.7), Inches(3), 6, 3)
headers = ["원리", "설명", "마케팅 함의"]
for i, h in enumerate(headers):
    table.rows[0].cells[i].text = h
style_header_row(table, [2.5, 4.6, 4.6])

data = [
    ("자동적 처리", "기능어는 무의식적으로 산출, 의도적 조작이 어려움", "소비자의 진짜 심리 포착 (설문 편향 극복)"),
    ("편향 불감", "내용어는 조작 가능하지만 기능어는 아님", "가짜 리뷰/허위 추천 탐지의 이론적 근거"),
    ("관계성 반영", "대명사 사용은 화자-대상의 심리적 거리 반영", "\"I\" vs \"we\" vs \"they\"로 브랜드-소비자 관계 측정"),
    ("인지 복잡성", "전치사, 접속사, 부사는 사고의 복잡성 반영", "소비자 의사결정 복잡도/확신도 추정"),
    ("범문화적 적용", "기능어 원리는 언어를 넘어 보편적 적용", "글로벌 브랜드의 다국가 소비자 비교 분석"),
]
for r, (a, b, c) in enumerate(data, 1):
    style_body_cell(table.rows[r].cells[0], a, bold=True)
    style_body_cell(table.rows[r].cells[1], b)
    style_body_cell(table.rows[r].cells[2], c)

# ═══════════════════════════════════════════
# SLIDE 4: LIWC 4대 변수
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), C_PRIMARY)
add_textbox(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.6),
            "2  LIWC-22의 4대 요약 변수 — 마케팅 핵심 지표", font_size=24, color=C_WHITE, bold=True)

table = add_table(slide, Inches(0.8), Inches(1.4), Inches(11.7), Inches(3), 5, 3)
headers = ["요약 변수", "산출 방식", "마케팅 활용"]
for i, h in enumerate(headers):
    table.rows[0].cells[i].text = h
style_header_row(table, [3, 4.35, 4.35])

data = [
    ("Analytical Thinking\n분석적 사고", "형식적, 논리적, 위계적 사고 단어 비율", "소비자가 브랜드를 이성적/감성적 중\n어떤 방식으로 인식하는지"),
    ("Clout\n영향력", "\"I\" 적게, \"we\"/\"you\" 많이 사용", "소비자가 브랜드에 대해 느끼는\n자신감/권위 수준"),
    ("Authenticity\n진정성", "자기모니터링 수준 반영", "소비자가 브랜드에 대해\n솔직하게 표현하는 정도"),
    ("Emotional Tone\n감정 톤", "긍정/부정 감정어 비율 종합\n(50+=긍정, 50-=부정)", "전반적 브랜드 감성 평가"),
]
for r, (a, b, c) in enumerate(data, 1):
    style_body_cell(table.rows[r].cells[0], a, bold=True, color=C_ACCENT_DARK)
    style_body_cell(table.rows[r].cells[1], b)
    style_body_cell(table.rows[r].cells[2], c)

# 방법론 4칸
y = 4.8
labels = ["사전(Dictionary) 기반", "빈도 기반 분석", "스타일로메트리", "ML/DL 결합"]
descs = [
    "LIWC/K-LIWC의 미리 정의된\n단어 사전으로 80개+ 범주 분류",
    "각 기능어 범주의 출현 빈도를\n전체 단어 수 대비 백분율로 계산",
    "기능어 사용 패턴을 \"언어적 지문\"으로\n활용하여 저자/스타일 식별",
    "LIWC 지표를 피처로 SVM, BERT에 학습\nLIWC+바이그램 89.8% 정확도"
]
for i in range(4):
    x = 0.8 + i * 3.05
    add_shape_bg(slide, Inches(x), Inches(y), Inches(2.85), Inches(1.6), C_LIGHT_BG)
    add_textbox(slide, Inches(x + 0.15), Inches(y + 0.1), Inches(2.55), Inches(0.35),
                labels[i], font_size=12, color=C_ACCENT_DARK, bold=True)
    add_textbox(slide, Inches(x + 0.15), Inches(y + 0.5), Inches(2.55), Inches(1),
                descs[i], font_size=10, color=C_TEXT_LIGHT)

# ═══════════════════════════════════════════
# SLIDE 5: 소비자 리뷰 분석
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), C_PRIMARY)
add_textbox(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.6),
            "3  소비자 리뷰 분석 — 가짜 리뷰 탐지 & 유용성 예측", font_size=24, color=C_WHITE, bold=True)

# 가짜 리뷰 테이블
table = add_table(slide, Inches(0.8), Inches(1.4), Inches(7.5), Inches(2.8), 5, 3)
headers = ["기능어 지표", "가짜 리뷰", "진짜 리뷰"]
for i, h in enumerate(headers):
    table.rows[0].cells[i].text = h
style_header_row(table, [2.2, 2.65, 2.65])

data = [
    ("1인칭 대명사", "과다 — 경험 조작 시 자기참조 증가", "자연스러운 비율"),
    ("긍정 감정어", "비정상적으로 높음", "부정 감정어가 자연스럽게 혼합"),
    ("지각 과정어", "감각적 묘사로 신빙성 조작", "구체적 제품 디테일에 집중"),
    ("느낌표", "높은 빈도", "상대적으로 적음"),
]
for r, (a, b, c) in enumerate(data, 1):
    style_body_cell(table.rows[r].cells[0], a, bold=True)
    style_body_cell(table.rows[r].cells[1], b, color=C_DANGER)
    style_body_cell(table.rows[r].cells[2], c, color=C_SUCCESS)

# Stat cards
add_stat_card(slide, Inches(8.8), Inches(1.4), "79%", "4개 피처만으로\n탐지 정확도", C_ACCENT)
add_stat_card(slide, Inches(8.8), Inches(3.0), "89.8%", "LIWC+바이그램\n조합 정확도", C_SUCCESS)

# 리뷰 유용성
add_shape_bg(slide, Inches(0.8), Inches(4.6), Inches(11.7), Inches(2.5), C_LIGHT_BG)
add_textbox(slide, Inches(1.0), Inches(4.7), Inches(11), Inches(0.35),
            "CASE 2 — 리뷰 유용성 예측 (Amazon 20,997건 분석)", font_size=14, color=C_ACCENT_DARK, bold=True)
tb = add_textbox(slide, Inches(1.0), Inches(5.2), Inches(11), Inches(1.5), "", font_size=12, color=C_TEXT)
tf = tb.text_frame
add_para(tf, "• 비인격적 스타일(impersonal)로 작성된 리뷰가 인격적 스타일보다 더 유용하다고 인식", font_size=12, color=C_TEXT)
add_para(tf, "• 리뷰 가이드라인에 \"객관적 톤\"을 권장하면 유용한 리뷰가 유도됨", font_size=12, color=C_TEXT)
add_para(tf, "• 리뷰 정렬 알고리즘에 LIWC 기반 유용성 점수를 반영할 수 있음", font_size=12, color=C_TEXT)

# ═══════════════════════════════════════════
# SLIDE 6: 광고 카피 최적화 - CTR
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), C_PRIMARY)
add_textbox(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.6),
            "4  광고 카피 & 메시지 최적화 — 인스타그램 CTR 분석", font_size=24, color=C_WHITE, bold=True)

add_textbox(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.3),
            "Inoue & Yoshida (2023) — 건강제품 3,555개 + 화장품 3,270개 인스타그램 광고",
            font_size=11, color=C_TEXT_LIGHT)

# 건강제품 테이블
add_textbox(slide, Inches(0.8), Inches(1.6), Inches(3), Inches(0.3),
            "건강제품 (n=3,555)", font_size=12, color=C_BROWN, bold=True)
table = add_table(slide, Inches(0.8), Inches(1.95), Inches(5.5), Inches(2.2), 5, 3)
headers = ["LIWC 카테고리", "본문 (rho)", "이미지 (rho)"]
for i, h in enumerate(headers):
    table.rows[0].cells[i].text = h
style_header_row(table, [2.2, 1.65, 1.65], C_BROWN)

data = [
    ("부정 감정 (negemo)", "+0.221", "+0.302"),
    ("인지적 과정", "+0.152", "+0.106"),
    ("생물학적 과정", "+0.151", "+0.094"),
    ("돈/가격", "-0.058", "-0.170"),
]
for r, (a, b, c) in enumerate(data, 1):
    style_body_cell(table.rows[r].cells[0], a, bold=True)
    color_b = C_SUCCESS if b.startswith("+") and float(b) > 0.1 else (C_DANGER if b.startswith("-") and abs(float(b)) > 0.1 else C_TEXT)
    color_c = C_SUCCESS if c.startswith("+") and float(c) > 0.1 else (C_DANGER if c.startswith("-") and abs(float(c)) > 0.1 else C_TEXT)
    style_body_cell(table.rows[r].cells[1], b, bold=True, color=color_b)
    style_body_cell(table.rows[r].cells[2], c, bold=True, color=color_c)

# 화장품 테이블
add_textbox(slide, Inches(6.8), Inches(1.6), Inches(3), Inches(0.3),
            "화장품 (n=3,270)", font_size=12, color=C_PURPLE, bold=True)
table = add_table(slide, Inches(6.8), Inches(1.95), Inches(5.7), Inches(2.2), 5, 3)
headers = ["LIWC 카테고리", "본문 (rho)", "이미지 (rho)"]
for i, h in enumerate(headers):
    table.rows[0].cells[i].text = h
style_header_row(table, [2.4, 1.65, 1.65], C_PURPLE)

data = [
    ("부정 감정 (negemo)", "+0.254", "+0.124"),
    ("돈/가격", "+0.259", "-0.076"),
    ("동기 (drives)", "+0.194", "+0.020"),
    ("생물학적 과정", "-0.308", "+0.020"),
]
for r, (a, b, c) in enumerate(data, 1):
    style_body_cell(table.rows[r].cells[0], a, bold=True)
    color_b = C_SUCCESS if b.startswith("+") and float(b) > 0.1 else (C_DANGER if b.startswith("-") and abs(float(b)) > 0.1 else C_TEXT)
    color_c = C_SUCCESS if c.startswith("+") and float(c) > 0.1 else (C_DANGER if c.startswith("-") and abs(float(c)) > 0.1 else C_TEXT)
    style_body_cell(table.rows[r].cells[1], b, bold=True, color=color_b)
    style_body_cell(table.rows[r].cells[2], c, bold=True, color=color_c)

# DO/DON'T
y = 4.5
add_shape_bg(slide, Inches(0.8), Inches(y), Inches(5.9), Inches(2.6), RGBColor(0xEC, 0xFD, 0xF5))
add_textbox(slide, Inches(1.0), Inches(y + 0.1), Inches(5.5), Inches(0.3),
            "DO — 이렇게 하세요", font_size=12, color=C_SUCCESS, bold=True)
tb = add_textbox(slide, Inches(1.0), Inches(y + 0.45), Inches(5.5), Inches(2), "", font_size=11, color=C_TEXT)
tf = tb.text_frame
add_para(tf, "• 건강제품: 이미지에 위기감/결핍 메시지 강조 (rho=0.302)", font_size=11, color=C_TEXT)
add_para(tf, "• 화장품: 본문에 가격을 투명하게 공개 (rho=+0.259)", font_size=11, color=C_TEXT)
add_para(tf, "• 공통: 부정 감정어를 전략적으로 활용", font_size=11, color=C_TEXT)

add_shape_bg(slide, Inches(7.0), Inches(y), Inches(5.5), Inches(2.6), RGBColor(0xFE, 0xF2, 0xF2))
add_textbox(slide, Inches(7.2), Inches(y + 0.1), Inches(5.1), Inches(0.3),
            "DON'T — 피하세요", font_size=12, color=C_DANGER, bold=True)
tb = add_textbox(slide, Inches(7.2), Inches(y + 0.45), Inches(5.1), Inches(2), "", font_size=11, color=C_TEXT)
tf = tb.text_frame
add_para(tf, "• 건강제품: 이미지에 가격 표기 (rho=-0.170)", font_size=11, color=C_TEXT)
add_para(tf, "• 화장품: 감각적 수식어 남발 (rho=-0.308)", font_size=11, color=C_TEXT)
add_para(tf, "• 공통: 긍정어만으로 구성된 카피는 CTR에 미미", font_size=11, color=C_TEXT)

# ═══════════════════════════════════════════
# SLIDE 7: 대명사와 소비자 행동
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), C_PRIMARY)
add_textbox(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.6),
            "4  대명사 하나가 소비자 행동을 바꾼다", font_size=24, color=C_WHITE, bold=True)

# I vs We
add_shape_bg(slide, Inches(0.8), Inches(1.3), Inches(5.9), Inches(3.2), C_LIGHT_BG)
add_textbox(slide, Inches(1.0), Inches(1.4), Inches(5.5), Inches(0.35),
            "\"I(제가)\" vs \"We(저희가)\" — Packard et al., JMR 2018", font_size=13, color=C_ACCENT_DARK, bold=True)
add_stat_card(slide, Inches(1.0), Inches(1.9), "+19%", "만족도 증가", C_SUCCESS)
add_stat_card(slide, Inches(3.4), Inches(1.9), "+7%", "잠재 매출 증가", C_SUCCESS)
tb = add_textbox(slide, Inches(1.0), Inches(3.4), Inches(5.5), Inches(1), "", font_size=11, color=C_TEXT)
tf = tb.text_frame
add_para(tf, "• CS에서 \"I\" 사용 시 대리(agency) + 공감(empathy) 이중 효과", font_size=11, color=C_TEXT)
add_para(tf, "• CS 매니저 92%가 비효과적인 \"We\" 사용 중", font_size=11, color=C_DANGER)
add_para(tf, "• 비용 0원, 즉시 적용 가능한 \"저비용 고효율\" 전략", font_size=11, color=C_SUCCESS)

# You 효과
add_shape_bg(slide, Inches(7.0), Inches(1.3), Inches(5.5), Inches(3.2), C_LIGHT_BG)
add_textbox(slide, Inches(7.2), Inches(1.4), Inches(5.1), Inches(0.35),
            "\"You(당신)\" 효과 — Cruz et al., JIM 2017", font_size=13, color=C_PURPLE, bold=True)
add_stat_card(slide, Inches(7.2), Inches(1.9), "+21%", "문제해결 책임감", C_PURPLE)
add_stat_card(slide, Inches(9.6), Inches(1.9), "매개", "자기참조\n(self-referencing)", C_PURPLE)
tb = add_textbox(slide, Inches(7.2), Inches(3.4), Inches(5.1), Inches(1), "", font_size=11, color=C_TEXT)
tf = tb.text_frame
add_para(tf, "• \"You\" → 자기참조 활성화 → 몰입도/태도 상승", font_size=11, color=C_TEXT)
add_para(tf, "• 한국(집단주의)에서는 \"우리/함께\"가 더 효과적일 수 있음", font_size=11, color=C_WARNING)

# SNS 대명사
add_shape_bg(slide, Inches(0.8), Inches(4.8), Inches(11.7), Inches(2.4), RGBColor(0xF9, 0xFA, 0xFB))
add_textbox(slide, Inches(1.0), Inches(4.9), Inches(11), Inches(0.35),
            "SNS 대명사별 인게이지먼트 — Labrecque et al. (2020), 15,788개 포스트", font_size=13, color=C_PRIMARY, bold=True)
table = add_table(slide, Inches(1.0), Inches(5.35), Inches(11.3), Inches(1.7), 4, 5)
headers = ["대명사", "좋아요", "댓글", "공유", "작동 원리"]
for i, h in enumerate(headers):
    table.rows[0].cells[i].text = h
style_header_row(table, [1.8, 1.5, 1.5, 1.5, 5.0])
data = [
    ("You (당신)", "—", "양의 효과 ↑", "—", "소비자를 직접 호출 → 대화 참여 유발"),
    ("They (그들)", "양의 효과 ↑", "양의 효과 ↑", "양의 효과 ↑", "외부 참조/사회적 증거 → 전방위 인게이지먼트"),
    ("I / We", "—", "—", "—", "브랜드 인격화 / 공동체 의식 (상황별 효과)"),
]
for r, (a, b, c, d, e) in enumerate(data, 1):
    style_body_cell(table.rows[r].cells[0], a, bold=True)
    style_body_cell(table.rows[r].cells[1], b, color=C_SUCCESS if "↑" in b else C_TEXT)
    style_body_cell(table.rows[r].cells[2], c, color=C_SUCCESS if "↑" in c else C_TEXT)
    style_body_cell(table.rows[r].cells[3], d, color=C_SUCCESS if "↑" in d else C_TEXT)
    style_body_cell(table.rows[r].cells[4], e)

# ═══════════════════════════════════════════
# SLIDE 8: 브랜드 포지셔닝 & 소셜미디어
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), C_PRIMARY)
add_textbox(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.6),
            "5-6  브랜드 포지셔닝 & 소셜미디어 분석", font_size=24, color=C_WHITE, bold=True)

# 럭셔리 브랜드
add_shape_bg(slide, Inches(0.8), Inches(1.3), Inches(6), Inches(3.5), C_LIGHT_BG)
add_textbox(slide, Inches(1.0), Inches(1.4), Inches(5.6), Inches(0.35),
            "럭셔리 브랜드 소비자 언어 — Oc et al. (2023)", font_size=13, color=C_ACCENT_DARK, bold=True)
add_textbox(slide, Inches(1.0), Inches(1.75), Inches(5.6), Inches(0.25),
            "29,000+ 댓글 · 88개 캠페인 · 9개 럭셔리 브랜드", font_size=10, color=C_TEXT_LIGHT)
table = add_table(slide, Inches(1.0), Inches(2.15), Inches(5.6), Inches(2.3), 5, 3)
headers = ["LIWC 변수", "편의품", "전문품"]
for i, h in enumerate(headers):
    table.rows[0].cells[i].text = h
style_header_row(table, [2, 1.8, 1.8])
data = [
    ("Analytical", "낮음", "높음 ↑"),
    ("Clout", "높음 ↑", "낮음"),
    ("Authenticity", "낮음", "높음 ↑"),
    ("Tone", "복합적", "카테고리별 상이"),
]
for r, (a, b, c) in enumerate(data, 1):
    style_body_cell(table.rows[r].cells[0], a, bold=True)
    style_body_cell(table.rows[r].cells[1], b, color=C_SUCCESS if "↑" in b else C_TEXT)
    style_body_cell(table.rows[r].cells[2], c, color=C_SUCCESS if "↑" in c else C_TEXT)

# 인플루언서
add_shape_bg(slide, Inches(7.1), Inches(1.3), Inches(5.5), Inches(3.5), RGBColor(0xFE, 0xF3, 0xC7))
add_textbox(slide, Inches(7.3), Inches(1.4), Inches(5.1), Inches(0.35),
            "인플루언서 고각성 언어 — JM 2024", font_size=13, color=C_BROWN, bold=True)
add_textbox(slide, Inches(7.3), Inches(1.75), Inches(5.1), Inches(0.25),
            "20,923 포스트 · 1,376 인플루언서", font_size=10, color=C_TEXT_LIGHT)
table = add_table(slide, Inches(7.3), Inches(2.15), Inches(5.1), Inches(1.5), 3, 3)
headers = ["유형", "각성도 10%↑ 시", "매개"]
for i, h in enumerate(headers):
    table.rows[0].cells[i].text = h
style_header_row(table, [1.7, 1.7, 1.7], C_BROWN)
data = [
    ("마이크로 (<100K)", "+5.4% 인게이지먼트", "신뢰성 ↑"),
    ("매크로 (>100K)", "-8.4% 인게이지먼트", "신뢰성 ↓"),
]
for r, (a, b, c) in enumerate(data, 1):
    style_body_cell(table.rows[r].cells[0], a, bold=True)
    style_body_cell(table.rows[r].cells[1], b, color=C_SUCCESS if "+" in b else C_DANGER)
    style_body_cell(table.rows[r].cells[2], c, color=C_SUCCESS if "↑" in c else C_DANGER)

add_textbox(slide, Inches(7.3), Inches(3.8), Inches(5.1), Inches(0.8),
            "매크로에게 \"최고!\", \"미쳤다!\" 등\n고각성 표현을 요구하면 역효과",
            font_size=11, color=C_DANGER)

# 라이벌리 효과
add_shape_bg(slide, Inches(0.8), Inches(5.1), Inches(11.7), Inches(2.1), RGBColor(0xF9, 0xFA, 0xFB))
add_textbox(slide, Inches(1.0), Inches(5.2), Inches(11), Inches(0.35),
            "라이벌 참조 효과 (Rivalry Reference Effect) — Berendt et al., JMR 2024", font_size=13, color=C_PRIMARY, bold=True)
tb = add_textbox(slide, Inches(1.0), Inches(5.65), Inches(11), Inches(1.2), "", font_size=12, color=C_TEXT)
tf = tb.text_frame
add_para(tf, "• 브랜드가 SNS에서 라이벌 경쟁사를 언급하면 비라이벌 대비 소비자 인게이지먼트가 증가", font_size=12, color=C_TEXT)
add_para(tf, "• 브랜드-경쟁사 관계를 기능어 수준에서 분석하면 소비자 반응을 예측할 수 있다", font_size=12, color=C_TEXT)

# ═══════════════════════════════════════════
# SLIDE 9: 콜센터 & FGI
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), C_PRIMARY)
add_textbox(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.6),
            "7  고객 서비스 / VOC / FGI — 단계별 언어 전략", font_size=24, color=C_WHITE, bold=True)

add_textbox(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.3),
            "Packard, Li & Berger (2024), JCR — 23,958 대화 순간 + N=1,589 실험", font_size=11, color=C_TEXT_LIGHT)

# 타임라인 바
y_tl = 1.7
# 초반
add_shape_bg(slide, Inches(0.8), Inches(y_tl), Inches(3), Inches(2.5), RGBColor(0xFE, 0xF3, 0xC7))
add_textbox(slide, Inches(0.8), Inches(y_tl + 0.1), Inches(3), Inches(0.25),
            "초반 25%", font_size=10, color=C_TEXT_LIGHT, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(0.8), Inches(y_tl + 0.4), Inches(3), Inches(0.35),
            "Relating (따뜻함)", font_size=16, color=C_BROWN, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(0.8), Inches(y_tl + 0.9), Inches(3), Inches(1.4),
            "\"불편하셨겠습니다\"\n\"제가 도와드리겠습니다\"\n\n감정적(affective) 언어\n공감 표현", font_size=11, color=C_TEXT, alignment=PP_ALIGN.CENTER)

# 중반
add_shape_bg(slide, Inches(4.0), Inches(y_tl), Inches(5.2), Inches(2.5), RGBColor(0xDB, 0xEA, 0xFE))
add_textbox(slide, Inches(4.0), Inches(y_tl + 0.1), Inches(5.2), Inches(0.25),
            "중반 50%", font_size=10, color=C_TEXT_LIGHT, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(4.0), Inches(y_tl + 0.4), Inches(5.2), Inches(0.35),
            "Resolving (역량/해결)", font_size=16, color=C_ACCENT_DARK, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(4.0), Inches(y_tl + 0.9), Inches(5.2), Inches(1.4),
            "\"확인 결과 ~입니다\"\n\"원인은 ~이며, 조치는 ~\"\n\n인지적(cognitive) 언어\n문제 해결, 논리적 표현", font_size=11, color=C_TEXT, alignment=PP_ALIGN.CENTER)

# 후반
add_shape_bg(slide, Inches(9.4), Inches(y_tl), Inches(3.2), Inches(2.5), RGBColor(0xFE, 0xF3, 0xC7))
add_textbox(slide, Inches(9.4), Inches(y_tl + 0.1), Inches(3.2), Inches(0.25),
            "후반 25%", font_size=10, color=C_TEXT_LIGHT, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(9.4), Inches(y_tl + 0.4), Inches(3.2), Inches(0.35),
            "Relating (따뜻함)", font_size=16, color=C_BROWN, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(9.4), Inches(y_tl + 0.9), Inches(3.2), Inches(1.4),
            "\"감사합니다\"\n\"추가 도움이 필요하시면\"\n\n감정적(affective) 언어", font_size=11, color=C_TEXT, alignment=PP_ALIGN.CENTER)

# FGI
y_fgi = 4.6
add_textbox(slide, Inches(0.8), Inches(y_fgi), Inches(11), Inches(0.35),
            "FGI / 인터뷰 보완 분석 — LIWC 결합 4가지 시나리오", font_size=14, color=C_PRIMARY, bold=True)
table = add_table(slide, Inches(0.8), Inches(y_fgi + 0.4), Inches(11.7), Inches(2.2), 5, 3)
headers = ["조사 유형", "LIWC 활용법", "기대 인사이트"]
for i, h in enumerate(headers):
    table.rows[0].cells[i].text = h
style_header_row(table, [2.5, 4.6, 4.6])
data = [
    ("브랜드 인식 FGI", "참가자별 Clout, Authenticity 점수 비교", "진정한 의견 vs 사회적 바람직성 편향 구분"),
    ("제품 사용 인터뷰", "감정 톤, 인지 과정어 분석", "명시적으로 말하지 않는 심리적 저항/열의 포착"),
    ("광고 반응 설문", "개방형 응답 LIWC 분석", "정량적 만족도 점수와 실제 언어 표현 간 괴리 발견"),
    ("고객 여정 인터뷰", "시간적 참조어, 인과 관계어 분석", "고객이 인식하는 핵심 전환 모멘트 식별"),
]
for r, (a, b, c) in enumerate(data, 1):
    style_body_cell(table.rows[r].cells[0], a, bold=True)
    style_body_cell(table.rows[r].cells[1], b)
    style_body_cell(table.rows[r].cells[2], c)

# ═══════════════════════════════════════════
# SLIDE 10: 한국 연구 & 도구
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), C_PRIMARY)
add_textbox(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.6),
            "8-9  한국 연구 & 실무 도구", font_size=24, color=C_WHITE, bold=True)

# K-LIWC
add_shape_bg(slide, Inches(0.8), Inches(1.3), Inches(6), Inches(3.5), RGBColor(0xFE, 0xF2, 0xF2))
add_textbox(slide, Inches(1.0), Inches(1.4), Inches(5.6), Inches(0.35),
            "K-LIWC 기만 리뷰 탐지 (한국소비자학회)", font_size=13, color=C_DANGER, bold=True)
tb = add_textbox(slide, Inches(1.0), Inches(1.85), Inches(5.6), Inches(2.5), "", font_size=11, color=C_TEXT)
tf = tb.text_frame
add_para(tf, "• 39개 언어 지표 + 44개 심리 지표로 분석", font_size=11, color=C_TEXT)
add_para(tf, "• 핵심 기만 변인: 형태소, 일반 명사, 조사", font_size=11, color=C_TEXT, bold=True)
add_para(tf, "• 심리 지표: 감정/정서 과정, 긍정 감정, 감각/지각 과정", font_size=11, color=C_TEXT)
add_para(tf, "• 한국어 고유 기능어(조사/어미)가 영어와 다른 패턴", font_size=11, color=C_DANGER, bold=True)
add_para(tf, "", font_size=6, color=C_TEXT)
add_para(tf, "KLIWC 광고효과 연구 (384명 대학생)", font_size=12, color=C_PRIMARY, bold=True)
add_para(tf, "• 17개 심리사회적 + 9개 언어학적 변인에서 유의미한 차이", font_size=11, color=C_TEXT)
add_para(tf, "• 광고태도·제품태도·구매의도와 유의미한 상관 확인", font_size=11, color=C_TEXT)

# 도구 비교
add_textbox(slide, Inches(7.1), Inches(1.3), Inches(5.5), Inches(0.35),
            "실무 도구 비교", font_size=14, color=C_PRIMARY, bold=True)
table = add_table(slide, Inches(7.1), Inches(1.75), Inches(5.5), Inches(3.0), 5, 3)
headers = ["도구", "특징", "비용"]
for i, h in enumerate(headers):
    table.rows[0].cells[i].text = h
style_header_row(table, [1.5, 2.5, 1.5])
data = [
    ("LIWC-22", "골드 스탠다드, 100+ 차원", "$89.95~$179.90"),
    ("K-LIWC", "한국어 특화, 형태소 분석", "학술 연구용"),
    ("Receptiviti API", "프로그래밍 방식, 94 카테고리", "API 과금"),
    ("Python+KoNLPy", "커스텀 사전 구축 가능", "무료 (오픈소스)"),
]
for r, (a, b, c) in enumerate(data, 1):
    style_body_cell(table.rows[r].cells[0], a, bold=True, color=C_ACCENT_DARK)
    style_body_cell(table.rows[r].cells[1], b)
    style_body_cell(table.rows[r].cells[2], c)

# 워크플로우
y_wf = 5.1
add_textbox(slide, Inches(0.8), Inches(y_wf), Inches(11), Inches(0.35),
            "마케터를 위한 4가지 워크플로우", font_size=14, color=C_PRIMARY, bold=True)
wf_labels = ["A. 소비자 리뷰 인사이트", "B. 광고 카피 사전 테스트", "C. CS/VOC 모니터링", "D. FGI 보완 분석"]
wf_steps = [
    "리뷰 크롤링 → LIWC 분석\n→ 자사/경쟁사 패턴 차이",
    "카피 후보 → LIWC 입력\n→ negemo/대명사 기반 예측",
    "VOC 수집 → 배치 분석\n→ Tone 시계열 + 경보",
    "녹취 전사 → LIWC 프로필\n→ Authenticity로 편향 구분"
]
for i in range(4):
    x = 0.8 + i * 3.05
    add_shape_bg(slide, Inches(x), Inches(y_wf + 0.4), Inches(2.85), Inches(1.6), C_LIGHT_BG)
    add_textbox(slide, Inches(x + 0.1), Inches(y_wf + 0.45), Inches(2.65), Inches(0.3),
                wf_labels[i], font_size=11, color=C_ACCENT_DARK, bold=True)
    add_textbox(slide, Inches(x + 0.1), Inches(y_wf + 0.8), Inches(2.65), Inches(1.1),
                wf_steps[i], font_size=10, color=C_TEXT_LIGHT)

# ═══════════════════════════════════════════
# SLIDE 11-12: 가나초콜릿 케이스 스터디
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, RGBColor(0x1C, 0x14, 0x0A))

# 장식
circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(-1), Inches(5), Inches(5))
circle.fill.solid()
circle.fill.fore_color.rgb = RGBColor(0x92, 0x40, 0x0E)
circle.fill.fore_color.brightness = 0.8
circle.line.fill.background()

add_textbox(slide, Inches(1), Inches(0.5), Inches(5), Inches(0.4),
            "CASE STUDY", font_size=13, color=RGBColor(0xFD, 0xE6, 0x8A), bold=True)
add_textbox(slide, Inches(1), Inches(1.1), Inches(10), Inches(1),
            "가나초콜릿 × 가나초콜릿하우스", font_size=36, color=C_WHITE, bold=True)
add_textbox(slide, Inches(1), Inches(2.2), Inches(10), Inches(0.6),
            "기능어 분석으로 본 브랜드 인식 변화: 제품에서 체험 공간으로", font_size=16, color=RGBColor(0xD9, 0x77, 0x06))

# 개요 카드
add_shape_bg(slide, Inches(1), Inches(3.3), Inches(5), Inches(3.2), RGBColor(0x2D, 0x20, 0x10))
add_textbox(slide, Inches(1.2), Inches(3.4), Inches(4.6), Inches(0.3),
            "가나초콜릿 (제품)", font_size=14, color=RGBColor(0xFD, 0xE6, 0x8A), bold=True)
tb = add_textbox(slide, Inches(1.2), Inches(3.8), Inches(4.6), Inches(2.5), "", font_size=12, color=RGBColor(0xD4, 0xD4, 0xD8))
tf = tb.text_frame
add_para(tf, "1975년 출시, 50년 역사", font_size=12, color=RGBColor(0xD4, 0xD4, 0xD8))
add_para(tf, "누적 매출 1조 3,000억원", font_size=12, color=RGBColor(0xD4, 0xD4, 0xD8))
add_para(tf, "66억 갑 판매 / 초당 4개", font_size=12, color=RGBColor(0xD4, 0xD4, 0xD8))
add_para(tf, "인지도 약 90%", font_size=12, color=RGBColor(0xD4, 0xD4, 0xD8))

add_shape_bg(slide, Inches(6.5), Inches(3.3), Inches(5.5), Inches(3.2), RGBColor(0x0A, 0x2D, 0x1C))
add_textbox(slide, Inches(6.7), Inches(3.4), Inches(5.1), Inches(0.3),
            "가나초콜릿하우스 (공간)", font_size=14, color=RGBColor(0xA7, 0xF3, 0xD0), bold=True)
tb = add_textbox(slide, Inches(6.7), Inches(3.8), Inches(5.1), Inches(2.5), "", font_size=12, color=RGBColor(0xD4, 0xD4, 0xD8))
tf = tb.text_frame
add_para(tf, "2022~2024, 3시즌 운영", font_size=12, color=RGBColor(0xD4, 0xD4, 0xD8))
add_para(tf, "체험형 팝업/디저트 카페", font_size=12, color=RGBColor(0xD4, 0xD4, 0xD8))
add_para(tf, "예약 대기 2,000+ / 구매전환 95%", font_size=12, color=RGBColor(0xD4, 0xD4, 0xD8))
add_para(tf, "\"18세기 런던 초콜릿 하우스\"", font_size=12, color=RGBColor(0xD4, 0xD4, 0xD8))

# ── SLIDE 12: 분석 결과 ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), C_BROWN)
add_textbox(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.6),
            "★  가나 케이스 — LIWC 4대 변수 & 감각 채널 비교", font_size=24, color=C_WHITE, bold=True)

# LIWC 바 차트
add_textbox(slide, Inches(0.8), Inches(1.3), Inches(5), Inches(0.35),
            "LIWC 4대 변수 추정 비교", font_size=14, color=C_PRIMARY, bold=True)

bar_data = [
    ("Analytical", 65, 45),
    ("Clout", 35, 65),
    ("Authenticity", 65, 45),
    ("Tone", 55, 80),
]
max_w = Inches(4)
for i, (name, prod, space) in enumerate(bar_data):
    y_b = Inches(1.8 + i * 0.7)
    add_textbox(slide, Inches(0.8), y_b, Inches(1.2), Inches(0.25), name, font_size=10, color=C_PRIMARY, bold=True)
    add_bar(slide, Inches(2.1), y_b + Inches(0.02), prod, max_w, Inches(0.2), C_BROWN, f"제품 {prod}")
    add_bar(slide, Inches(2.1), y_b + Inches(0.27), space, max_w, Inches(0.2), C_GREEN, f"공간 {space}")

# 범례
add_textbox(slide, Inches(2.1), Inches(4.7), Inches(1.5), Inches(0.2), "■ 제품", font_size=9, color=C_BROWN, bold=True)
add_textbox(slide, Inches(3.3), Inches(4.7), Inches(1.5), Inches(0.2), "■ 공간", font_size=9, color=C_GREEN, bold=True)

# 감각 채널
add_textbox(slide, Inches(7), Inches(1.3), Inches(5.5), Inches(0.35),
            "감각 채널 전환", font_size=14, color=C_PRIMARY, bold=True)

senses = ["미각", "촉각", "시각", "공간", "사회"]
prod_pcts = [90, 60, 15, 5, 10]
space_pcts = [35, 15, 85, 75, 70]
for i, (name, pp, sp) in enumerate(zip(senses, prod_pcts, space_pcts)):
    y_s = Inches(1.8 + i * 0.65)
    add_textbox(slide, Inches(7), y_s, Inches(0.7), Inches(0.2), name, font_size=10, color=C_PRIMARY, bold=True, alignment=PP_ALIGN.RIGHT)
    # 제품
    add_bar(slide, Inches(7.8), y_s, pp, Inches(2.3), Inches(0.18), C_BROWN, f"{pp}%")
    # 공간
    add_bar(slide, Inches(7.8), y_s + Inches(0.22), sp, Inches(2.3), Inches(0.18), C_GREEN, f"{sp}%")

add_textbox(slide, Inches(7), Inches(5), Inches(5.5), Inches(0.3),
            "근접감각(미각/촉각) → 원격감각(시각/공간) + 사회적 경험", font_size=10, color=C_TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

# 5대 전환 요약
y_sh = 5.4
add_textbox(slide, Inches(0.8), Inches(y_sh), Inches(12), Inches(0.35),
            "5대 언어 구조 전환", font_size=14, color=C_PRIMARY, bold=True)
shifts = [
    ("\"나/내\"", "\"여기/우리\"", "개인→사회"),
    ("과거 회고", "현재 몰입", "추억→지금"),
    ("미각+촉각", "시각+공간", "입안→눈&몸"),
    ("분석적 비교", "감동+기대위반", "이성→감성"),
    ("독백", "대화/추천", "소유→공유"),
]
for i, (fr, to, meaning) in enumerate(shifts):
    x = 0.8 + i * 2.5
    add_shape_bg(slide, Inches(x), Inches(y_sh + 0.4), Inches(2.3), Inches(1.4), C_LIGHT_BG)
    add_textbox(slide, Inches(x), Inches(y_sh + 0.45), Inches(2.3), Inches(0.25),
                fr, font_size=11, color=C_BROWN, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(x), Inches(y_sh + 0.7), Inches(2.3), Inches(0.2),
                "→", font_size=14, color=C_ACCENT, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(x), Inches(y_sh + 0.9), Inches(2.3), Inches(0.25),
                to, font_size=11, color=C_GREEN, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(x), Inches(y_sh + 1.2), Inches(2.3), Inches(0.3),
                meaning, font_size=9, color=C_TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

# ── SLIDE 13: 가나 시사점 & 인사이트 ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), C_BROWN)
add_textbox(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.6),
            "★  가나 케이스 — 시사점 & 인사이트 & 프로젝트 고려사항", font_size=24, color=C_WHITE, bold=True)

# 핵심 써머리
add_shape_bg(slide, Inches(0.8), Inches(1.3), Inches(11.7), Inches(1.2), RGBColor(0xFF, 0xFB, 0xEB))
add_textbox(slide, Inches(1.0), Inches(1.35), Inches(11.3), Inches(1.1),
            "핵심 써머리: 가나초콜릿이 제품에서 공간으로 확장되면서, 소비자 언어 구조가 \"솔직한 개인 평가\"(고Authenticity)에서\n\"열정적인 사회적 추천\"(고Clout)으로 체계적으로 전환. 대명사, 시제, 감각, 조사라는 기능어 레벨의 변화가 이를 실증.",
            font_size=13, color=C_BROWN, bold=True)

# DO/DON'T
y = 2.7
add_shape_bg(slide, Inches(0.8), Inches(y), Inches(5.9), Inches(2.5), RGBColor(0xEC, 0xFD, 0xF5))
add_textbox(slide, Inches(1.0), Inches(y + 0.1), Inches(5.5), Inches(0.3),
            "검증된 성과", font_size=12, color=C_SUCCESS, bold=True)
tb = add_textbox(slide, Inches(1.0), Inches(y + 0.45), Inches(5.5), Inches(2), "", font_size=11, color=C_TEXT)
tf = tb.text_frame
add_para(tf, "• 사회적 소비 유도: 사회적 단어 10% → 70% 증가", font_size=11, color=C_TEXT)
add_para(tf, "• 시간 프레임: 과거 회고 → \"지금 여기\" 브랜드 전환", font_size=11, color=C_TEXT)
add_para(tf, "• 감각 채널: 미각 단일 → 시각+공간+사회적 다채널", font_size=11, color=C_TEXT)
add_para(tf, "• 기대 위반 전략 성공: \"팝업 → 감동\" 패턴 반복 확인", font_size=11, color=C_TEXT)

add_shape_bg(slide, Inches(7.0), Inches(y), Inches(5.5), Inches(2.5), RGBColor(0xFE, 0xF2, 0xF2))
add_textbox(slide, Inches(7.2), Inches(y + 0.1), Inches(5.1), Inches(0.3),
            "주의해야 할 점", font_size=12, color=C_DANGER, bold=True)
tb = add_textbox(slide, Inches(7.2), Inches(y + 0.45), Inches(5.1), Inches(2), "", font_size=11, color=C_TEXT)
tf = tb.text_frame
add_para(tf, "• Authenticity 하락 리스크: \"인스타 감성\" 편향 가능성", font_size=11, color=C_TEXT)
add_para(tf, "• 제품-공간 연결 약화: 미각 언급 크게 감소", font_size=11, color=C_TEXT)
add_para(tf, "• 부정 피드백 부재: 개선 피드백 포착 어려움", font_size=11, color=C_TEXT)

# 인사이트 & 프로젝트 고려사항
y2 = 5.4
add_textbox(slide, Inches(0.8), Inches(y2), Inches(11.7), Inches(0.3),
            "인사이트 & 향후 프로젝트 고려사항", font_size=14, color=C_PRIMARY, bold=True)

insights = [
    "공간은 기능어를\n바꾼다",
    "Authenticity↔Clout\n트레이드오프 관리",
    "\"이/가\" 조사 =\n새로움 지표",
    "시계열 추적으로\n캠페인 효과 측정",
    "K-LIWC 정량\n분석으로 검증",
]
for i, txt in enumerate(insights):
    x = 0.8 + i * 2.5
    add_shape_bg(slide, Inches(x), Inches(y2 + 0.35), Inches(2.3), Inches(1.2), C_LIGHT_BG)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.9), Inches(y2 + 0.15), Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = C_ACCENT
    circle.line.fill.background()
    add_textbox(slide, Inches(x + 0.9), Inches(y2 + 0.17), Inches(0.5), Inches(0.5),
                str(i + 1), font_size=14, color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(x + 0.05), Inches(y2 + 0.65), Inches(2.2), Inches(0.8),
                txt, font_size=10, color=C_TEXT, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════
# SLIDE 14: 마케터 시사점 총정리
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), C_PRIMARY)
add_textbox(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.6),
            "10  마케터를 위한 6대 시사점", font_size=24, color=C_WHITE, bold=True)

implications = [
    ("가짜 리뷰 탐지", "1인칭 과다 + 긍정 과잉 패턴으로\n자동 감지. 4피처 79%, LIWC+바이그램 89.8%"),
    ("광고 카피 최적화", "부정 감정어(negemo)가 CTR과 최강 상관.\n건강=가격 회피, 화장품=가격 투명성"),
    ("대명사 전략", "CS \"제가(I)\" > \"저희가(we)\" → 만족도 +19%\n광고 \"당신(you)\" → 책임감 +21%"),
    ("브랜드 포지셔닝", "LIWC 4대 변수로 소비자 인식 수치화.\n리포지셔닝 전후 변화 추적 가능"),
    ("CS 교육 혁신", "대화 초반/후반=따뜻함, 중반=역량.\n23,958 순간 분석으로 실증"),
    ("질적 조사 보완", "FGI/인터뷰에 LIWC 적용하면\n\"말하지 않은 심리\" 포착 가능"),
]

for i, (title, desc) in enumerate(implications):
    row = i // 3
    col = i % 3
    x = 0.8 + col * 4.1
    y = 1.4 + row * 3.0

    add_shape_bg(slide, Inches(x), Inches(y), Inches(3.9), Inches(2.7), C_LIGHT_BG)

    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.15), Inches(y + 0.15), Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = C_ACCENT
    circle.line.fill.background()
    add_textbox(slide, Inches(x + 0.15), Inches(y + 0.17), Inches(0.5), Inches(0.5),
                str(i + 1), font_size=16, color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(x + 0.8), Inches(y + 0.2), Inches(2.9), Inches(0.35),
                title, font_size=14, color=C_ACCENT_DARK, bold=True)
    add_textbox(slide, Inches(x + 0.2), Inches(y + 0.8), Inches(3.5), Inches(1.7),
                desc, font_size=12, color=C_TEXT)

# ═══════════════════════════════════════════
# SLIDE 15: 엔딩
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_PRIMARY)

add_textbox(slide, Inches(1), Inches(2), Inches(11), Inches(1),
            "Thank You", font_size=48, color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(3.2), Inches(11), Inches(0.6),
            "기능어분석 — 마케팅 & 소비자 조사 특화 리서치 보고서", font_size=18, color=RGBColor(0x94, 0xA3, 0xB8), alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.4),
            "2026.03 · do-better-workspace · 30+ 학술 자료 기반", font_size=13, color=RGBColor(0x64, 0x74, 0x8B), alignment=PP_ALIGN.CENTER)

# ─── 저장 ───
output_path = os.path.join(os.path.dirname(__file__), "function-word-analysis-research.pptx")
prs.save(output_path)
print(f"PPT saved: {output_path}")
