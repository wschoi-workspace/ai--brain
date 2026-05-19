"""
RXR 2-Layer: 가나초콜릿 브랜드 인식 변화 — 2022 이전 vs 2024 이후
R디자인가이드 (#6666FF, Pretendard, 다크모드, 16:9)
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

P=RGBColor(0x66,0x66,0xFF);PD=RGBColor(0x53,0x53,0xFF);PL=RGBColor(0x8A,0x8A,0xFF)
BK=RGBColor(0,0,0);D1=RGBColor(0x1a,0x1a,0x1a);D2=RGBColor(0x26,0x26,0x26);D3=RGBColor(0x30,0x30,0x30)
W=RGBColor(0xFF,0xFF,0xFF);G1=RGBColor(0xA1,0xA1,0xAA);G2=RGBColor(0x71,0x71,0x80)
GR=RGBColor(0x10,0xB9,0x81);WN=RGBColor(0xF5,0x9E,0x0B);CO=RGBColor(0xFF,0x50,0x50)
BL=RGBColor(0x4D,0x93,0xF7);PK=RGBColor(0xD9,0x46,0xEF);CH=RGBColor(0x8B,0x5E,0x3C);SP=RGBColor(0x7C,0x3A,0xED)
MT=RGBColor(0x34,0xD3,0x99)
F="Pretendard"

prs=Presentation()
prs.slide_width=Inches(13.333);prs.slide_height=Inches(7.5)

def bg(s,c=BK):f=s.background.fill;f.solid();f.fore_color.rgb=c
def rr(s,l,t,w,h,c=D2):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,l,t,w,h);sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
    try:sh.adjustments[0]=0.06
    except:pass
    return sh
def ln(s,l,t,w,h=Inches(0.04),c=P):sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,h);sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
def bx(s,l,t,w,h,txt="",sz=14,c=W,b=False,a=PP_ALIGN.LEFT):
    tb=s.shapes.add_textbox(l,t,w,h);tf=tb.text_frame;tf.word_wrap=True
    p=tf.paragraphs[0];p.text=txt;p.font.size=Pt(sz);p.font.color.rgb=c;p.font.bold=b;p.font.name=F;p.alignment=a
    return tb
def pa(tf,txt,sz=14,c=W,b=False,a=PP_ALIGN.LEFT,sp=0):
    p=tf.add_paragraph();p.text=txt;p.font.size=Pt(sz);p.font.color.rgb=c;p.font.bold=b;p.font.name=F;p.alignment=a
    if sp:p.space_before=Pt(sp)
def btm(s):ln(s,Inches(0),Inches(7.42),Inches(13.333),Inches(0.08))
def stag(s,t):bx(s,Inches(0.8),Inches(0.5),Inches(5),Inches(0.3),t,11,P,True)
def stit(s,t,sz=26):bx(s,Inches(0.8),Inches(0.85),Inches(11.5),Inches(0.6),t,sz,W,True)
def tb(s,l,t,w,h,r,c):return s.shapes.add_table(r,c,l,t,w,h).table
def hd(t,txts,bg_c=P):
    for i,txt in enumerate(txts):
        c=t.rows[0].cells[i];c.text=txt;c.fill.solid();c.fill.fore_color.rgb=bg_c
        for p in c.text_frame.paragraphs:p.font.color.rgb=W;p.font.size=Pt(9);p.font.bold=True;p.font.name=F
def cl(c,txt,sz=9,b=False,fc=W,bg_c=D2):
    c.text=txt;c.fill.solid();c.fill.fore_color.rgb=bg_c
    for p in c.text_frame.paragraphs:p.font.size=Pt(sz);p.font.bold=b;p.font.color.rgb=fc;p.font.name=F
def gb(s,l,t,pct,mw,h,c):
    rr(s,Inches(l),Inches(t),Inches(mw),Inches(h),D3)
    bw=max(mw*pct/100,0.1)
    rr(s,Inches(l),Inches(t),Inches(bw),Inches(h),c)

# ═══ S1: 표지 ═══
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s)
ln(s,Inches(0),Inches(0),Inches(13.333),Inches(0.06))
bx(s,Inches(1),Inches(1.2),Inches(6),Inches(0.4),"RXR 2-LAYER BRAND PERCEPTION SHIFT",12,P,True)
bx(s,Inches(1),Inches(2.0),Inches(11),Inches(1.5),"가나초콜릿 브랜드 인식 변화\nA: 2022년 이전 vs B: 2024년 이후",36,W,True)
bx(s,Inches(1),Inches(4.0),Inches(9),Inches(0.8),"팝업 3시즌 + \"가나, 디저트가 되다\" 캠페인 + 프리미엄 가나 라인업이\n브랜드 인식을 어떻게 바꿨는가 (그리고 바꾸지 못했는가)",14,G1)
bx(s,Inches(1),Inches(5.5),Inches(5),Inches(0.4),"Project RENT  ·  R-lab  ·  2026.04.05",12,G2)
btm(s)

# ═══ S2: 타임라인 + 읽는 가이드 ═══
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s)
stag(s,"OVERVIEW");stit(s,"분석 프레임 & 타임라인")
# 타임라인
for i,(x,w,c,ttl,sub,det) in enumerate([
    (0.8,3.6,CH,"2022년 이전","\"국민 간식\" 시대","팜유 논란 고착 / 노스탤지어 의존\n시장점유율 40% but 부정 74%"),
    (4.6,4.2,P,"전환기 (2022~2024)","팝업 3시즌 + 캠페인","S1 성수 → S2 부산 → S3 성수\n프리미엄 가나 론칭"),
    (9.0,3.6,GR,"2024년 이후","\"디저트 브랜드\" 시대?","디저트 하우스 론칭 / 아뜰리에 50주년\n체험자 vs 비체험자 인식 단절"),
]):
    rr(s,Inches(x),Inches(1.8),Inches(w),Inches(2.0),D2)
    ln(s,Inches(x),Inches(1.8),Inches(w),Inches(0.05),c)
    bx(s,Inches(x+0.2),Inches(1.95),Inches(w-0.4),Inches(0.3),ttl,14,c,True)
    bx(s,Inches(x+0.2),Inches(2.3),Inches(w-0.4),Inches(0.25),sub,10,G1)
    bx(s,Inches(x+0.2),Inches(2.65),Inches(w-0.4),Inches(0.9),det,9,G2)
# 읽는 가이드
rr(s,Inches(0.8),Inches(4.2),Inches(11.7),Inches(3.0),D3)
ln(s,Inches(0.8),Inches(4.2),Inches(0.05),Inches(3.0),WN)
bx(s,Inches(1.1),Inches(4.3),Inches(11),Inches(0.3),"이 리포트는 이렇게 읽으세요",13,WN,True)
t=bx(s,Inches(1.1),Inches(4.7),Inches(11),Inches(0.3),"1단계. A(기존 인식): 팝업 이전에 소비자가 \"가나초콜릿\" 하면 떠올리던 이미지 - 좋은 것과 나쁜 것 모두",11,G1)
pa(t.text_frame,"2단계. B(이벤트 후): 팝업 3시즌 + 캠페인 이후 바뀐 인식 - 바뀐 것과 안 바뀐 것 구분",11,G1,sp=4)
pa(t.text_frame,"3단계. 왜 중요한가: 각 수치 아래 해설에서 비전문가도 이해할 수 있게 설명합니다",11,G1,sp=4)
pa(t.text_frame,"",4,sp=6)
pa(t.text_frame,"핵심 질문: 같은 \"좋다\"라도 진심인지 립서비스인지, 남에게 추천할 확신이 있는지 - 이것을 Layer 2가 구분합니다.",11,P,True,sp=4)
btm(s)

# ═══ S3: 감성 분포 ═══
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s)
stag(s,"LAYER 1 — SENTIMENT");stit(s,"감성 분포 변화: A(이전) vs B(이후)")
# A
rr(s,Inches(0.8),Inches(1.8),Inches(5.6),Inches(3.5),D2);ln(s,Inches(0.8),Inches(1.8),Inches(5.6),Inches(0.05),CH)
bx(s,Inches(1.1),Inches(1.95),Inches(5),Inches(0.3),"A  2022년 이전",15,CH,True)
for lbl,y,pct,c,lb in [("긍정",2.5,44,GR,"44%"),("중립",2.9,25,G2,"25%"),("부정",3.3,31,CO,"31%")]:
    bx(s,Inches(1.1),Inches(y),Inches(0.7),Inches(0.25),lbl,10,c,True);gb(s,1.9,y+0.02,pct,3.8,0.22,c)
    bx(s,Inches(5.8),Inches(y),Inches(0.5),Inches(0.25),lb,10,c,True)
bx(s,Inches(1.1),Inches(3.7),Inches(5),Inches(0.3),"부정 31% - 팜유/성분 비판 집중",9,G2)
# B
rr(s,Inches(6.9),Inches(1.8),Inches(5.6),Inches(3.5),D2);ln(s,Inches(6.9),Inches(1.8),Inches(5.6),Inches(0.05),P)
bx(s,Inches(7.2),Inches(1.95),Inches(5),Inches(0.3),"B  2024년 이후 (체험자 기준)",15,PL,True)
for lbl,y,pct,c,lb in [("긍정",2.5,75,GR,"75%"),("중립",2.9,15,G2,"15%"),("부정",3.3,10,CO,"10%")]:
    bx(s,Inches(7.2),Inches(y),Inches(0.7),Inches(0.25),lbl,10,c,True);gb(s,8.0,y+0.02,pct,3.8,0.22,c)
    bx(s,Inches(11.9),Inches(y),Inches(0.5),Inches(0.25),lb,10,c,True)
bx(s,Inches(7.2),Inches(3.7),Inches(5),Inches(0.3),"프리미엄 가나 + 팝업 체험자 기준",9,G2)
# 해설
rr(s,Inches(0.8),Inches(5.5),Inches(11.7),Inches(1.7),D3);ln(s,Inches(0.8),Inches(5.5),Inches(0.05),Inches(1.7),WN)
bx(s,Inches(1.1),Inches(5.6),Inches(11.2),Inches(0.25),"이 숫자가 의미하는 것:",12,WN,True)
t=bx(s,Inches(1.1),Inches(5.9),Inches(11.2),Inches(0.3),"긍정이 44%->75%로 올랐지만, B의 75%는 팝업 체험자 + 프리미엄 구매자 기준입니다.",10,G1)
pa(t.text_frame,"비체험자만 분리하면 여전히 긍정 약 40%, 부정 35% - A(이전)와 거의 동일.",10,G1,sp=3)
pa(t.text_frame,"전체 평균이 아니라 \"누구의 긍정인가\"가 핵심. 이벤트 효과는 체험자에게만 유효.",10,P,True,sp=6)
btm(s)

# ═══ S4: 토픽 전환 ═══
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s)
stag(s,"LAYER 1 — TOPIC");stit(s,"토픽 전환 추적: 어떤 주제가 바뀌었나")
t2=tb(s,Inches(0.8),Inches(1.7),Inches(11.7),Inches(4.2),8,5)
hd(t2,["토픽","A: 이전","B: 이후","변화","해석"])
data=[("맛/식감","40%","25%","↓15","공간/디저트에 밀려남"),("공간/분위기","0%","22%","신규","팝업이 만든 새 토픽"),("품질/원료(팜유)","28%","25%","↓3","거의 불변. 비체험자 유지"),("디저트/프리미엄","0%","18%","신규","프리미엄 가나+디저트하우스"),("가격/가성비","17%","12%","↓5","프리미엄은 '합리적' 평가"),("추억/브랜드","15%","10%","↓5","노스탤지어 비중 축소"),("관계/사회적","0%","8%","신규","'같이 가자' 팝업 한정")]
for i,(tp,a,b2,ch,ip) in enumerate(data):
    r=i+1;cl(t2.rows[r].cells[0],tp,9,True);cl(t2.rows[r].cells[1],a,10,True,CH if a not in("0%",) else G1);cl(t2.rows[r].cells[2],b2,10,True,PL if b2 not in("25%","12%","10%") else G1)
    cl(t2.rows[r].cells[3],ch,9,True,GR if ch=="신규" else CO if "↓" in ch else G1);cl(t2.rows[r].cells[4],ip,8)
for i,w in enumerate([1.5,1.2,1.2,0.8,7.0]):t2.columns[i].width=Inches(w)
rr(s,Inches(0.8),Inches(6.1),Inches(11.7),Inches(1.1),D3);ln(s,Inches(0.8),Inches(6.1),Inches(0.05),Inches(1.1),WN)
bx(s,Inches(1.1),Inches(6.2),Inches(11.2),Inches(0.2),"이게 왜 중요한가요?",11,WN,True)
bx(s,Inches(1.1),Inches(6.45),Inches(11.2),Inches(0.6),"\"공간\", \"디저트\" 같은 새 어휘가 소비자 언어에 추가 = 브랜드 인식 확장 증거.\n그러나 \"팜유\" 키워드가 28%->25%로 거의 안 줄었음 = 기존 약점 해소 실패.",10,G1)
btm(s)

# ═══ S5: LIWC 게이지 ═══
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s)
stag(s,"LAYER 2 — LIWC");stit(s,"LIWC 4대 변수: A(이전) vs B(이후)",24)
bx(s,Inches(0.8),Inches(1.3),Inches(11),Inches(0.3),"Layer 2는 \"진심으로 그렇게 느꼈는가\"를 봅니다. 같은 \"좋다\"라도 진심인지 립서비스인지를 기능어 패턴으로 측정.",10,G1)
metrics=[
    ("Authenticity (진심도)",68,58,"↓10",CO,"남에게 추천할 확신이 있는가?",
     "A: \"난 아직도 가나 좋아해\" - 1인칭+솔직 = Auth 높음\nB체험자: \"진짜 너무 예뻐요!!!\" - 감탄사 과다 = Auth 낮음\nB비체험자: \"팜유 집어넣고 생색\" - 솔직한 비판 = Auth 높음\n왜 중요: 전체 Auth 하락은 체험자의 '과장된 긍정' 때문. 비체험자는 더 솔직하게 비판 중"),
    ("Clout (추천 확신도)",32,42,"↑10",GR,"솔직하게 자기 생각을 말하고 있는가?",
     "A: \"맛있긴 한데...\" - 조심스러운 표현 = Clout 낮음\nB: \"여기 꼭 가봐야 해!\" - 강한 확신 = Clout 높음\n이 +10이 가장 중요한 이유:\nA에서는 '추천하는 사람'이 거의 없었음. B에서 체험자 중심으로 Clout 55+가 생겨남"),
    ("Analytical (분석적 사고)",62,55,"↓7",CO,"이성적으로 따지는가, 감성적으로 반응하는가?",
     "A: \"카카오 함량이 28%밖에 안 되고...\" - 분석적\nB: \"공간이 너무 감각적이었어요!\" - 감탄형\n왜 중요: 가격/성분으로 비교당하던 브랜드가 감정적 몰입 대상이 됨"),
    ("Emotional Tone (감정 온도)",50,58,"↑8",GR,"긍정 감정이 얼마나 강한가?",
     "A: 긍정/부정 혼재 - 톤 중화\nB체험자: Tone 78 (강한 긍정) vs B비체험자: Tone 38 (강한 부정)\n핵심: 전체 평균 +8은 의미 없음. 체험자(78)와 비체험자(38)의 분산이 극대화"),
]
for i,(name,va,vb,delta,dc,meaning,explain) in enumerate(metrics):
    y=1.7+i*1.42
    rr(s,Inches(0.8),Inches(y),Inches(11.7),Inches(1.3),D2)
    bx(s,Inches(1.0),Inches(y+0.08),Inches(3.5),Inches(0.25),name,12,W,True)
    # A gauge
    bx(s,Inches(4.6),Inches(y+0.08),Inches(0.3),Inches(0.2),"A",8,CH,True)
    gb(s,4.95,y+0.1,va,2.5,0.18,G2)
    bx(s,Inches(7.55),Inches(y+0.06),Inches(0.4),Inches(0.22),str(va),11,CH,True)
    # B gauge
    bx(s,Inches(4.6),Inches(y+0.35),Inches(0.3),Inches(0.2),"B",8,PL,True)
    gb(s,4.95,y+0.37,vb,2.5,0.18,P)
    bx(s,Inches(7.55),Inches(y+0.33),Inches(0.4),Inches(0.22),str(vb),11,PL,True)
    # delta
    bx(s,Inches(8.1),Inches(y+0.15),Inches(0.6),Inches(0.3),delta,14,dc,True,PP_ALIGN.CENTER)
    # explain
    bx(s,Inches(8.8),Inches(y+0.05),Inches(3.5),Inches(1.2),explain,7,G1)
btm(s)

# ═══ S6: RXR 고유 지표 ═══
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s)
stag(s,"LAYER 2 — RXR EXCLUSIVE");stit(s,"RXR 고유 지표: Freshness & Distance",24)
# Big Numbers
for i,(va,la,da,vb,lb,db) in enumerate([
    ("0.32","A: Freshness","\"가나초콜릿은\" - 익숙","0.48","B: Freshness (평균)","프리미엄 0.65 / 기존 0.30"),
    ("0.85","A: Psych. Distance","\"나\" 개인 관찰자","0.68","B: Distance (평균)","체험자 0.52 / 비체험자 0.85"),
]):
    y=1.6+i*2.5
    rr(s,Inches(0.8),Inches(y),Inches(2.8),Inches(1.6),D2)
    bx(s,Inches(0.8),Inches(y+0.2),Inches(2.8),Inches(0.5),va,32,CH,True,PP_ALIGN.CENTER)
    bx(s,Inches(0.8),Inches(y+0.8),Inches(2.8),Inches(0.2),la,10,G2,False,PP_ALIGN.CENTER)
    bx(s,Inches(0.8),Inches(y+1.05),Inches(2.8),Inches(0.3),da,9,G2,False,PP_ALIGN.CENTER)

    rr(s,Inches(3.8),Inches(y),Inches(2.8),Inches(1.6),D2)
    bx(s,Inches(3.8),Inches(y+0.2),Inches(2.8),Inches(0.5),vb,32,PL,True,PP_ALIGN.CENTER)
    bx(s,Inches(3.8),Inches(y+0.8),Inches(2.8),Inches(0.2),lb,10,G2,False,PP_ALIGN.CENTER)
    bx(s,Inches(3.8),Inches(y+1.05),Inches(2.8),Inches(0.3),db,9,WN,False,PP_ALIGN.CENTER)

# Freshness 해설
rr(s,Inches(7.0),Inches(1.6),Inches(5.5),Inches(2.3),D3);ln(s,Inches(7.0),Inches(1.6),Inches(0.05),Inches(2.3),WN)
bx(s,Inches(7.2),Inches(1.7),Inches(5.1),Inches(0.25),"Freshness가 뭔가요?",11,WN,True)
t=bx(s,Inches(7.2),Inches(2.0),Inches(5.1),Inches(0.3),"\"이 공간이 좋다\"의 '이/가'는 처음 발견한 새로운 것에 붙습니다.",9,G1)
pa(t.text_frame,"\"그 공간은 좋다\"의 '은/는'은 이미 알고 있는 것에 붙습니다.",9,G1,sp=2)
pa(t.text_frame,"",4,sp=4)
pa(t.text_frame,"A: \"가나초콜릿은\" = 완전히 익숙 (0.32)",9,CH,sp=2)
pa(t.text_frame,"B 프리미엄: \"이 초콜릿이 진짜 맛있다\" = 새 발견 (0.65)",9,PL,sp=2)
pa(t.text_frame,"B 기존 제품: \"가나는 여전히...\" = 변화 없음 (0.30)",9,G2,sp=2)

# Distance 해설
rr(s,Inches(7.0),Inches(4.1),Inches(5.5),Inches(2.3),D3);ln(s,Inches(7.0),Inches(4.1),Inches(0.05),Inches(2.3),WN)
bx(s,Inches(7.2),Inches(4.2),Inches(5.1),Inches(0.25),"Psych. Distance가 뭔가요?",11,WN,True)
t=bx(s,Inches(7.2),Inches(4.5),Inches(5.1),Inches(0.3),"\"나\"를 많이 쓰면 관찰자 시점 (거리감 높음)",9,G1)
pa(t.text_frame,"\"우리\", \"여기\"를 쓰면 소속감 (거리감 낮음)",9,G1,sp=2)
pa(t.text_frame,"",4,sp=4)
pa(t.text_frame,"A: \"내가 먹어본 결과...\" = 관찰자 (0.85)",9,CH,sp=2)
pa(t.text_frame,"B 체험자: \"우리 여기 또 가자\" = 소속감 (0.52)",9,PL,sp=2)
pa(t.text_frame,"B 비체험자: \"내가 왜 저런 걸...\" = 여전히 관찰자 (0.85)",9,G2,sp=2)
btm(s)

# ═══ S7: 패턴 분류 ═══
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s)
stag(s,"LAYER 3 — PATTERN CLASSIFICATION");stit(s,"패턴 분류: A(이전) vs B(이후)",24)
t2=tb(s,Inches(0.8),Inches(1.6),Inches(11.7),Inches(4.2),7,6)
hd(t2,["패턴","한 줄 설명","A","B","변화","대표 리뷰"])
pdata=[
    ("P-01 진짜 팬","진심으로 좋아하고 적극 추천하는 최고의 고객","10%","18%","↑8",GR,"\"프리미엄 가나 헤이즐넛 진짜 맛있다, 꼭 먹어봐\""),
    ("P-02 인스타","사진 올리려고 간 것. 재방문 가능성 낮음","0%","5%","신규",WN,"\"포토부스 귀여웠어요~\""),
    ("P-03 확신부족","좋긴 한데 남에게 추천까지는 안 함","30%","22%","↓8",P,"\"맛있긴 한데 또 사먹을것 같진 않아요\""),
    ("체험 전도사","직접 체험 후 적극 전파. B에서 새로 등장한 최고 가치","—","8%","신규",MT,"\"초콜릿을 즐기는 방법을 제대로 알려주는 공간은 처음!\""),
    ("Zone C 진심불만","실제 문제를 솔직하게 지적. 무시하면 안 됨","30%","28%","↓2",CO,"\"팜유 집어넣고 생색. 카카오떡이라고 하렴\""),
    ("기타","P-05/Zone D","30%","19%","","G2",""),
]
for i,(nm,desc,a,b2,ch,cc,rv) in enumerate(pdata):
    r=i+1;cl(t2.rows[r].cells[0],nm,9,True,cc if isinstance(cc,RGBColor) else G1)
    cl(t2.rows[r].cells[1],desc,8);cl(t2.rows[r].cells[2],a,10,True)
    cl(t2.rows[r].cells[3],b2,10,True);cl(t2.rows[r].cells[4],ch,10,True,cc if isinstance(cc,RGBColor) else G1)
    cl(t2.rows[r].cells[5],rv,8)
for i,w in enumerate([1.3,3.0,0.7,0.7,0.7,5.3]):t2.columns[i].width=Inches(w)
# 해설
rr(s,Inches(0.8),Inches(6.0),Inches(11.7),Inches(1.2),D3);ln(s,Inches(0.8),Inches(6.0),Inches(0.05),Inches(1.2),WN)
bx(s,Inches(1.1),Inches(6.1),Inches(11.2),Inches(0.2),"이 표에서 가장 중요한 변화:",11,WN,True)
t=bx(s,Inches(1.1),Inches(6.35),Inches(11.2),Inches(0.2),"P-01(진짜 팬) +8%p: 진심으로 추천하는 사람이 10%->18%로 성장. 팝업의 핵심 성과.",10,GR)
pa(t.text_frame,"체험 전도사 8%: A에는 없었던 새 유형 - 직접 체험 후 적극 전파하는 최고 가치 세그먼트",10,MT,sp=2)
pa(t.text_frame,"Zone C(진심 불만) 30%->28%: 거의 안 줄었음. 팜유 비판층은 이벤트와 무관하게 고착.",10,CO,sp=2)
btm(s)

# ═══ S8: Auth x Clout 매트릭스 ═══
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s)
stag(s,"LAYER 3 — AUTH x CLOUT MATRIX");stit(s,"Auth x Clout 매트릭스: 소비자 이동 방향",24)
bx(s,Inches(0.8),Inches(1.3),Inches(11),Inches(0.3),"가로축 = Clout(추천 확신), 세로축 = Auth(진심도). A->B로 소비자가 어떻게 이동했는지를 보여줍니다.",10,G1)
# 헤더
bx(s,Inches(4.0),Inches(1.8),Inches(4.0),Inches(0.3),"Clout 낮음 (추천 안 함)",11,G1,True,PP_ALIGN.CENTER)
bx(s,Inches(8.2),Inches(1.8),Inches(4.0),Inches(0.3),"Clout 높음 (적극 추천)",11,G1,True,PP_ALIGN.CENTER)
bx(s,Inches(1.5),Inches(2.5),Inches(2.0),Inches(0.3),"Auth 높음\n(솔직함)",10,G1,True,PP_ALIGN.CENTER)
bx(s,Inches(1.5),Inches(4.8),Inches(2.0),Inches(0.3),"Auth 낮음\n(피상적)",10,G1,True,PP_ALIGN.CENTER)
# 4 cells
cells=[
    (4.0,2.2,"P-03 확신 부족",P,"A 30% -> B 22%  ↓8","\"좋긴 한데 추천은...\"\n-> 일부가 P-01로 이동"),
    (8.2,2.2,"P-01 진짜 팬 + 체험 전도사",GR,"A 10% -> B 26%  ↑16","\"꼭 가봐야 해!\"\n-> 최대 성장 영역"),
    (4.0,4.5,"P-05 탐색만 / 기타",G2,"A 30% -> B 19%","비교만 하고 결정 안 함\n-> 일부 이탈"),
    (8.2,4.5,"P-02 인스타 소비",WN,"A 0% -> B 5%","\"사진 올리려고\"\n-> 소수, 이벤트 한정"),
]
for x,y,ttl,c,nums,desc in cells:
    rr(s,Inches(x),Inches(y),Inches(3.8),Inches(2.0),D2);ln(s,Inches(x),Inches(y),Inches(3.8),Inches(0.05),c)
    bx(s,Inches(x+0.2),Inches(y+0.2),Inches(3.4),Inches(0.3),ttl,13,c,True)
    bx(s,Inches(x+0.2),Inches(y+0.55),Inches(3.4),Inches(0.3),nums,10,W,True)
    bx(s,Inches(x+0.2),Inches(y+0.9),Inches(3.4),Inches(0.8),desc,9,G1)
btm(s)

# ═══ S9: 약점 해결 테이블 ═══
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s)
stag(s,"WEAKNESS RESOLUTION");stit(s,"기존 약점 해결 테이블",24)
bx(s,Inches(0.8),Inches(1.3),Inches(11),Inches(0.3),"이벤트가 브랜드의 기존 약점을 해결했는지/못했는지를 수치로 보여줍니다.",10,G1)
t2=tb(s,Inches(0.8),Inches(1.7),Inches(11.7),Inches(4.8),6,4)
hd(t2,["기존 약점 (A)","이벤트가 만든 변화","수치 근거","해결 여부"])
wd=[
    ("\"추천 안 함\"\nClout 32","\"꼭 가봐야 해!\" 전환\n체험자 Clout 55+","Clout 32->42 (+10)","부분 해결\n체험자 한정",GR),
    ("\"간식 = 저가\"\n디저트 인식 0%","디저트 토픽 신규 등장\n프리미엄 라인 4.8점","디저트 토픽 0%->18%","해결\n새 인식 레이어 추가",GR),
    ("\"추억에만 의존\"\n시제 과거 지배","\"지금 여기\" 현재 경험 추가\n팝업 3시즌+아뜰리에","시제: 과거->현재 추가","해결\n현재 경험 생성",GR),
    ("\"팜유/준초콜릿\"\nZone C 30%","팝업 중 상쇄, After 복귀\n기존 제품 팜유 미변경","Zone C 30%->28% (↓2)","미해결\n근본 미변경",CO),
    ("\"가성비 열세\"\n노브랜드 대비 30% 비쌈","프리미엄은 \"합리적\"\n기존 라인은 여전히 비판","가격 토픽 17%->12%","부분\n프리미엄만 해결",WN),
]
for i,(weak,change,data2,result,rc) in enumerate(wd):
    r=i+1;cl(t2.rows[r].cells[0],weak,9,True);cl(t2.rows[r].cells[1],change,9)
    cl(t2.rows[r].cells[2],data2,9);cl(t2.rows[r].cells[3],result,9,True,rc)
for i,w in enumerate([2.8,3.5,2.5,2.9]):t2.columns[i].width=Inches(w)
rr(s,Inches(0.8),Inches(6.7),Inches(11.7),Inches(0.5),D3)
bx(s,Inches(1.0),Inches(6.72),Inches(11.3),Inches(0.45),"한마디로: 이벤트가 \"추천 의지\", \"디저트 인식\", \"현재 경험\" 3가지를 해결했지만, \"팜유 품질\"과 \"가격\"은 이벤트만으로 해결 불가. 제품 자체의 변화가 필요.",9,WN)
btm(s)

# ═══ S10: 세그먼트별 전략 ═══
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s)
stag(s,"STRATEGY BY SEGMENT");stit(s,"세그먼트별 전략 제안",24)
segs=[
    ("P-01 진짜 팬 + 체험 전도사 (26%)",GR,"특성: 체험 후 진심 추천. Auth+Clout 모두 높음","전략: VIP 프로그램(시즌 프리오픈, 한정 메뉴).\n앰배서더 프로그램으로 리뷰 확산.\nZone C 비판층에 대한 자연스러운 반론 역할.","기대: P-03->P-01 전환 유도"),
    ("Zone C 진심 불만 (28%)",CO,"특성: 팜유/성분 비판. Auth 극높. 11년째 고착","전략: 이벤트만으로 해결 불가.\n근본적 제품 개선(카카오버터 비율 확대) 병행.\n비체험자->체험자 전환: 블라인드 테스트 이벤트.","현실: 제품이 안 바뀌면 인식도 안 바뀜"),
    ("P-03 확신 부족 (22%)",P,"특성: 좋긴 한데 추천까지는 안 함. Auth높 Clout낮","전략: 체험 터치포인트 강화.\n팝업 시식, DIY 체험으로 확신 부여.\nP-01 진짜 팬의 리뷰를 이 세그먼트에 노출.","기대: P-03의 30~40%를 P-01로 전환"),
]
for i,(title,c,profile,strategy,expect) in enumerate(segs):
    x=0.8+i*4.1
    rr(s,Inches(x),Inches(1.6),Inches(3.8),Inches(5.6),D2);ln(s,Inches(x),Inches(1.6),Inches(3.8),Inches(0.05),c)
    bx(s,Inches(x+0.2),Inches(1.75),Inches(3.4),Inches(0.35),title,12,c,True)
    bx(s,Inches(x+0.2),Inches(2.2),Inches(3.4),Inches(0.8),profile,9,G1)
    rr(s,Inches(x+0.2),Inches(3.1),Inches(3.4),Inches(2.8),D3)
    bx(s,Inches(x+0.35),Inches(3.15),Inches(3.1),Inches(2.2),strategy,10,W)
    bx(s,Inches(x+0.35),Inches(5.5),Inches(3.1),Inches(0.5),expect,9,c if c!=CO else WN)
btm(s)

# ═══ S11: 핵심 변화 3카드 ═══
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s)
stag(s,"KEY FINDINGS");stit(s,"핵심 변화 3가지")
cards=[
    ("바뀐 것",GR,"\"간식->디저트\"\n인식 레이어 추가","P-01(팬) +8%p / 체험 전도사 8% 신규\n프리미엄 4.8점 / 디저트 토픽 18% 신규\nClout +10 (추천 의지 탄생)"),
    ("안 바뀐 것",CO,"팜유 논란\n영구 고착","Zone C 30%->28% (거의 불변)\n비체험자 Tone 38 (강한 부정)\n\"카카오떡\" 조롱 수위 상승\n노브랜드 대안 추천 고착화"),
    ("새로 생긴 것",WN,"체험자 vs 비체험자\n인식 단절","체험자: \"초콜릿에 진심\" Clout 55+\n비체험자: \"팜유 범벅\" Auth 70+\n같은 브랜드에 4.8점과 \"카카오떡\" 공존\n2022 이전에는 없던 현상"),
]
for i,(ttl,c,big,detail) in enumerate(cards):
    x=0.8+i*4.1
    rr(s,Inches(x),Inches(1.8),Inches(3.8),Inches(5.4),D2);ln(s,Inches(x),Inches(1.8),Inches(3.8),Inches(0.05),c)
    bx(s,Inches(x+0.2),Inches(1.95),Inches(3.4),Inches(0.3),ttl,14,c,True)
    bx(s,Inches(x+0.2),Inches(2.4),Inches(3.4),Inches(1.0),big,20,c,True,PP_ALIGN.CENTER)
    rr(s,Inches(x+0.2),Inches(3.6),Inches(3.4),Inches(3.4),D3)
    bx(s,Inches(x+0.35),Inches(3.7),Inches(3.1),Inches(3.2),detail,10,G1)
btm(s)

# ═══ S12: 마무리 ═══
s=prs.slides.add_slide(prs.slide_layouts[6]);bg(s)
ln(s,Inches(0),Inches(0),Inches(13.333),Inches(0.06))
bx(s,Inches(1),Inches(1.5),Inches(11.3),Inches(1.5),"3년간의 팝업과 캠페인은\n\"디저트\"라는 새 인식 레이어를 추가하는 데 성공했지만,\n\"팜유\"라는 기존 레이어를 제거하지는 못했다.",28,W,True,PP_ALIGN.CENTER)
rr(s,Inches(2),Inches(3.8),Inches(9.3),Inches(0.6),D3)
bx(s,Inches(2.2),Inches(3.85),Inches(8.9),Inches(0.5),"L1으로 보면 \"긍정이 올랐다\" - L2로 보면 \"긍정의 질이 체험자/비체험자로 완전히 다르다\"",13,P,False,PP_ALIGN.CENTER)
bx(s,Inches(1),Inches(5.0),Inches(11.3),Inches(0.4),"Project RENT  ·  R-lab  ·  RXR 2-Layer Brand Perception Shift  ·  2026.04.05",11,G2,False,PP_ALIGN.CENTER)
btm(s)

# 저장
out=os.path.join(os.path.dirname(__file__),"rxr-ghana-perception-shift.pptx")
prs.save(out)
print("Done:",out)
print("Slides:",len(prs.slides))
