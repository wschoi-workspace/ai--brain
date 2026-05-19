"""
RXR 2-Layer Case Study: 가나초콜릿 vs 가나초콜릿하우스
16:9 / R디자인가이드 (Primary: #6666FF, Pretendard, 다크모드)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── R디자인가이드 컬러 ───
C_PRIMARY   = RGBColor(0x66, 0x66, 0xFF)  # #6666FF
C_DEEP      = RGBColor(0x53, 0x53, 0xFF)  # #5353FF
C_LIGHT     = RGBColor(0x8A, 0x8A, 0xFF)  # #8A8AFF
C_PALE      = RGBColor(0xE1, 0xE1, 0xFF)  # #E1E1FF
C_BLACK     = RGBColor(0x00, 0x00, 0x00)
C_DARK      = RGBColor(0x1A, 0x1A, 0x1A)
C_DARK2     = RGBColor(0x26, 0x26, 0x26)
C_DARK3     = RGBColor(0x30, 0x30, 0x30)
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_LGRAY     = RGBColor(0xF6, 0xF6, 0xF6)
C_BODY      = RGBColor(0x42, 0x4D, 0x61)
C_GRAY      = RGBColor(0xA1, 0xA1, 0xAA)
C_GRAY2     = RGBColor(0x71, 0x71, 0x80)
C_RED       = RGBColor(0xC0, 0x00, 0x00)
C_CORAL     = RGBColor(0xFF, 0x50, 0x50)
C_BLUE      = RGBColor(0x4D, 0x93, 0xF7)
C_GREEN     = RGBColor(0x10, 0xB9, 0x81)
C_WARN      = RGBColor(0xF5, 0x9E, 0x0B)
C_PINK      = RGBColor(0xD9, 0x46, 0xEF)
C_CHOCO     = RGBColor(0x8B, 0x5E, 0x3C)
C_SPACE     = RGBColor(0x7C, 0x3A, 0xED)

FONT = "Pretendard"
FONT_EN = "Poppins"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ─── 헬퍼 ───
def bg(slide, color=C_BLACK):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = color

def rect(slide, l, t, w, h, color=C_DARK2, radius=0.06):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    try: s.adjustments[0] = radius
    except: pass
    return s

def line_shape(slide, l, t, w, h=Inches(0.04), color=C_PRIMARY):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s

def box(slide, l, t, w, h, txt="", sz=14, color=C_WHITE, bold=False, align=PP_ALIGN.LEFT, font=FONT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = txt
    p.font.size = Pt(sz); p.font.color.rgb = color; p.font.bold = bold; p.font.name = font; p.alignment = align
    return tb

def para(tf, txt, sz=14, color=C_WHITE, bold=False, align=PP_ALIGN.LEFT, sp_before=0, font=FONT):
    p = tf.add_paragraph(); p.text = txt
    p.font.size = Pt(sz); p.font.color.rgb = color; p.font.bold = bold; p.font.name = font; p.alignment = align
    if sp_before: p.space_before = Pt(sp_before)
    return p

def tag(slide, l, t, txt, bg_color=C_PRIMARY, txt_color=C_WHITE, sz=9):
    r = rect(slide, l, t, Inches(len(txt)*0.12+0.3), Inches(0.28), bg_color, 0.5)
    box(slide, l, t+Inches(0.02), Inches(len(txt)*0.12+0.3), Inches(0.24), txt, sz=sz, color=txt_color, bold=True, align=PP_ALIGN.CENTER, font=FONT_EN)

def stat_card(slide, l, t, val, label, vc=C_PRIMARY, w=2.5, h=1.4):
    rect(slide, Inches(l), Inches(t), Inches(w), Inches(h), C_DARK2)
    box(slide, Inches(l), Inches(t+0.15), Inches(w), Inches(0.6), val, sz=32, color=vc, bold=True, align=PP_ALIGN.CENTER)
    box(slide, Inches(l), Inches(t+0.85), Inches(w), Inches(0.4), label, sz=10, color=C_GRAY, align=PP_ALIGN.CENTER)

def gauge_bar(slide, l, t, pct, max_w, h, color, label=""):
    # track
    rect(slide, Inches(l), Inches(t), Inches(max_w), Inches(h), C_DARK3, 0.5)
    # bar
    bw = max(max_w * pct / 100, 0.15)
    r = rect(slide, Inches(l), Inches(t), Inches(bw), Inches(h), color, 0.5)
    # label
    if label:
        box(slide, Inches(l + max_w + 0.15), Inches(t-0.02), Inches(1.2), Inches(h+0.04), label, sz=10, color=C_GRAY)

def bottom_line(slide):
    line_shape(slide, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), C_PRIMARY)

def slide_tag(slide, txt):
    box(slide, Inches(0.8), Inches(0.5), Inches(3), Inches(0.3), txt, sz=11, color=C_PRIMARY, bold=True, font=FONT_EN)

def slide_title(slide, txt, sz=30):
    box(slide, Inches(0.8), Inches(0.85), Inches(11), Inches(0.7), txt, sz=sz, color=C_WHITE, bold=True)

def tbl(slide, l, t, w, h, rows, cols):
    return slide.shapes.add_table(rows, cols, l, t, w, h).table

def hdr(table, texts, bg_color=C_PRIMARY):
    for i, txt in enumerate(texts):
        c = table.rows[0].cells[i]; c.text = txt; c.fill.solid(); c.fill.fore_color.rgb = bg_color
        for p in c.text_frame.paragraphs:
            p.font.color.rgb = C_WHITE; p.font.size = Pt(10); p.font.bold = True; p.font.name = FONT

def cell(c, txt, sz=10, bold=False, color=C_WHITE):
    c.text = txt; c.fill.solid(); c.fill.fore_color.rgb = C_DARK2
    for p in c.text_frame.paragraphs:
        p.font.size = Pt(sz); p.font.bold = bold; p.font.color.rgb = color; p.font.name = FONT

def cell_highlight(c, txt, sz=10, bold=True, color=C_PRIMARY, bg=C_DARK3):
    c.text = txt; c.fill.solid(); c.fill.fore_color.rgb = bg
    for p in c.text_frame.paragraphs:
        p.font.size = Pt(sz); p.font.bold = bold; p.font.color.rgb = color; p.font.name = FONT


# ═══════════════════════════════════════════════════
# SLIDE 1: 표지
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
line_shape(s, Inches(0), Inches(0), Inches(13.333), Inches(0.06))

box(s, Inches(1), Inches(1.2), Inches(6), Inches(0.4), "RXR 2-LAYER CASE STUDY", sz=13, color=C_PRIMARY, bold=True, font=FONT_EN)
box(s, Inches(1), Inches(2.0), Inches(11), Inches(1.5),
    "가나초콜릿 vs 가나초콜릿하우스\n동일 브랜드, 완전히 다른 심리 구조", sz=38, color=C_WHITE, bold=True)
box(s, Inches(1), Inches(4.0), Inches(9), Inches(0.8),
    "제품 리뷰와 공간 리뷰를 RXR 2-Layer 프레임워크로 분석하여,\n단일 레이어 분석이 놓치는 심리적 전환을 실증한다.", sz=15, color=C_GRAY)
box(s, Inches(1), Inches(5.5), Inches(5), Inches(0.4), "Project RENT  ·  R-lab  ·  2026.04", sz=12, color=C_GRAY2)
bottom_line(s)


# ═══════════════════════════════════════════════════
# SLIDE 2: 분석 대상
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
slide_tag(s, "OVERVIEW")
slide_title(s, "분석 대상 & 프레임워크")

# 가나초콜릿 카드
rect(s, Inches(0.8), Inches(2.0), Inches(5.6), Inches(4.5), C_DARK2)
line_shape(s, Inches(0.8), Inches(2.0), Inches(5.6), Inches(0.05), C_CHOCO)
box(s, Inches(1.1), Inches(2.3), Inches(5), Inches(0.4), "A  가나초콜릿 (제품)", sz=18, color=C_CHOCO, bold=True)
tb = box(s, Inches(1.1), Inches(2.9), Inches(5), Inches(3.3), "유형: 판형 초콜릿 (FMCG)", sz=12, color=C_GRAY)
para(tb.text_frame, "역사: 1975~ (51년), 누적 매출 1.3조 원", sz=12, color=C_GRAY, sp_before=4)
para(tb.text_frame, "누적 판매량: 66억 갑 (초당 4개)", sz=12, color=C_GRAY, sp_before=4)
para(tb.text_frame, "포지셔닝: \"국민 초콜릿\"", sz=12, color=C_GRAY, sp_before=4)
para(tb.text_frame, "", sz=6, sp_before=8)
para(tb.text_frame, "리뷰 특성", sz=13, color=C_WHITE, bold=True, sp_before=6)
para(tb.text_frame, "· 맛/식감 중심 평가 (미각+촉각)", sz=11, color=C_GRAY, sp_before=4)
para(tb.text_frame, "· 긍정/부정 양분 (팜유 논란)", sz=11, color=C_GRAY, sp_before=2)
para(tb.text_frame, "· 추억/노스탤지어 기반 긍정", sz=11, color=C_GRAY, sp_before=2)

# 가나초콜릿하우스 카드
rect(s, Inches(6.9), Inches(2.0), Inches(5.6), Inches(4.5), C_DARK2)
line_shape(s, Inches(6.9), Inches(2.0), Inches(5.6), Inches(0.05), C_SPACE)
box(s, Inches(7.2), Inches(2.3), Inches(5), Inches(0.4), "B  가나초콜릿하우스 (공간)", sz=18, color=C_SPACE, bold=True)
tb = box(s, Inches(7.2), Inches(2.9), Inches(5), Inches(3.3), "유형: 체험형 팝업스토어 / 디저트 카페", sz=12, color=C_GRAY)
para(tb.text_frame, "역사: 2022~ (시즌3까지), 150평 규모", sz=12, color=C_GRAY, sp_before=4)
para(tb.text_frame, "구매전환율: 95%+, 예약 1분 만에 마감", sz=12, color=C_GRAY, sp_before=4)
para(tb.text_frame, "포지셔닝: \"가나, 디저트가 되다\"", sz=12, color=C_GRAY, sp_before=4)
para(tb.text_frame, "", sz=6, sp_before=8)
para(tb.text_frame, "리뷰 특성", sz=13, color=C_WHITE, bold=True, sp_before=6)
para(tb.text_frame, "· 공간/분위기 중심 (시각+공간감각)", sz=11, color=C_GRAY, sp_before=4)
para(tb.text_frame, "· 압도적 긍정 (86%)", sz=11, color=C_GRAY, sp_before=2)
para(tb.text_frame, "· 관계적 소비 (\"같이 가자\")", sz=11, color=C_GRAY, sp_before=2)
bottom_line(s)


# ═══════════════════════════════════════════════════
# SLIDE 3: Layer 1 감성 분포
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
slide_tag(s, "LAYER 1 — CONTENT")
slide_title(s, "감성 분포 비교: \"무엇을 말했는가\"")

# 제품 감성
rect(s, Inches(0.8), Inches(2.0), Inches(5.6), Inches(4.8), C_DARK2)
box(s, Inches(1.1), Inches(2.2), Inches(5), Inches(0.35), "A  가나초콜릿 — 감성 분포", sz=15, color=C_CHOCO, bold=True)

box(s, Inches(1.1), Inches(2.8), Inches(0.9), Inches(0.3), "긍정", sz=11, color=C_GREEN, bold=True)
gauge_bar(s, 2.1, 2.82, 48, 4.0, 0.28, C_GREEN, "48%")
box(s, Inches(1.1), Inches(3.2), Inches(0.9), Inches(0.3), "중립", sz=11, color=C_GRAY, bold=True)
gauge_bar(s, 2.1, 3.22, 22, 4.0, 0.28, C_GRAY2, "22%")
box(s, Inches(1.1), Inches(3.6), Inches(0.9), Inches(0.3), "부정", sz=11, color=C_CORAL, bold=True)
gauge_bar(s, 2.1, 3.62, 30, 4.0, 0.28, C_CORAL, "30%")

tb = box(s, Inches(1.1), Inches(4.2), Inches(5), Inches(0.3), "대표 긍정:", sz=11, color=C_GRAY)
para(tb.text_frame, "\"입에서 사르르 녹는 달콤한 맛이 특징이다\"", sz=10, color=C_GRAY2, sp_before=2)
para(tb.text_frame, "", sz=4, sp_before=4)
para(tb.text_frame, "대표 부정:", sz=11, color=C_CORAL, sp_before=2)
para(tb.text_frame, "\"카카오버터 대신 싸구려 팜유로 만드니까\"", sz=10, color=C_GRAY2, sp_before=2)
para(tb.text_frame, "\"가나초콜렛은 가나초콜렛이지 초콜렛은 아니야\"", sz=10, color=C_GRAY2, sp_before=2)

# 공간 감성
rect(s, Inches(6.9), Inches(2.0), Inches(5.6), Inches(4.8), C_DARK2)
box(s, Inches(7.2), Inches(2.2), Inches(5), Inches(0.35), "B  가나초콜릿하우스 — 감성 분포", sz=15, color=C_SPACE, bold=True)

box(s, Inches(7.2), Inches(2.8), Inches(0.9), Inches(0.3), "긍정", sz=11, color=C_GREEN, bold=True)
gauge_bar(s, 8.2, 2.82, 86, 4.0, 0.28, C_GREEN, "86%")
box(s, Inches(7.2), Inches(3.2), Inches(0.9), Inches(0.3), "중립", sz=11, color=C_GRAY, bold=True)
gauge_bar(s, 8.2, 3.22, 10, 4.0, 0.28, C_GRAY2, "10%")
box(s, Inches(7.2), Inches(3.6), Inches(0.9), Inches(0.3), "부정", sz=11, color=C_CORAL, bold=True)
gauge_bar(s, 8.2, 3.62, 4, 4.0, 0.28, C_CORAL, "4%")

tb = box(s, Inches(7.2), Inches(4.2), Inches(5), Inches(0.3), "대표 긍정:", sz=11, color=C_GRAY)
para(tb.text_frame, "\"생각보다 더 좋았습니다! 공간이 초콜릿 컬러에 맞게 구성\"", sz=10, color=C_GRAY2, sp_before=2)
para(tb.text_frame, "\"다른 카페 가지 말고 여기 있다 가자\"", sz=10, color=C_GRAY2, sp_before=2)
para(tb.text_frame, "", sz=4, sp_before=4)
para(tb.text_frame, "대표 부정:", sz=11, color=C_CORAL, sp_before=2)
para(tb.text_frame, "\"별거 볼 건 없더라고요ㅠㅠ 초콜릿 특색이 덜 보이는 점 아쉬움\"", sz=10, color=C_GRAY2, sp_before=2)

# L1 결론 바
rect(s, Inches(0.8), Inches(7.0), Inches(11.7), Inches(0.35), C_DARK3)
box(s, Inches(1.0), Inches(7.02), Inches(11.3), Inches(0.3),
    "Layer 1 결론:  제품은 긍정/부정 양분(48:30),  공간은 압도적 긍정(86:4).  하지만 그 86%가 진심인지는 모른다.",
    sz=10, color=C_GRAY)
bottom_line(s)


# ═══════════════════════════════════════════════════
# SLIDE 3-B: 토픽 클러스터 & 키워드 비교
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
slide_tag(s, "LAYER 1 — TOPIC CLUSTER")
slide_title(s, "토픽 클러스터 & 키워드 비교")

t = tbl(s, Inches(0.8), Inches(1.9), Inches(11.7), Inches(4.8), 8, 5)
hdr(t, ["토픽", "제품 비중", "제품 키워드", "공간 비중", "공간 키워드"])
topics = [
    ("맛/식감",        "42%", "\"달콤\", \"사르르 녹는\", \"부드럽다\"",    "18%", "\"진한 농도\", \"조화\", \"단짠단짠\""),
    ("공간/분위기",    "0%",  "-",                                           "36%", "\"고급스럽다\", \"빈티지\", \"포토존\""),
    ("가격/가성비",    "18%", "\"납득할 만한\", \"가격 인상\"",              "8%",  "\"합리적\", \"35,000원 페어링\""),
    ("품질/원료",      "25%", "\"팜유\", \"카카오 함량\", \"카카오버터\"",   "2%",  "-"),
    ("추억/브랜드",    "15%", "\"어릴 때\", \"국민 초콜릿\", \"추억\"",     "6%",  "\"50주년\", \"헤리티지\""),
    ("관계/사회적",    "0%",  "-",                                           "22%", "\"데이트\", \"같이 가자\", \"가족\""),
    ("SNS/포토",       "0%",  "-",                                           "8%",  "\"사진\", \"인스타\", \"포토부스\""),
]
for i, (topic, pv, pk, sv, sk) in enumerate(topics):
    r = i + 1
    cell(t.rows[r].cells[0], topic, sz=10, bold=True)
    # 제품 비중 - 높으면 강조
    if pv in ("42%", "25%", "18%", "15%"):
        cell_highlight(t.rows[r].cells[1], pv, sz=11, bold=True, color=C_CHOCO)
    else:
        cell(t.rows[r].cells[1], pv, sz=10)
    cell(t.rows[r].cells[2], pk, sz=9)
    # 공간 비중 - 높으면 강조
    if sv in ("36%", "22%", "18%"):
        cell_highlight(t.rows[r].cells[3], sv, sz=11, bold=True, color=C_SPACE)
    else:
        cell(t.rows[r].cells[3], sv, sz=10)
    cell(t.rows[r].cells[4], sk, sz=9)

t.columns[0].width = Inches(1.5)
t.columns[1].width = Inches(1.2)
t.columns[2].width = Inches(3.5)
t.columns[3].width = Inches(1.2)
t.columns[4].width = Inches(4.3)

rect(s, Inches(0.8), Inches(6.9), Inches(11.7), Inches(0.4), C_DARK3)
box(s, Inches(1.0), Inches(6.92), Inches(11.3), Inches(0.35),
    "핵심 차이:  제품은 맛/원료(42%+25%=67%)에 집중  |  공간은 분위기+관계(36%+22%=58%)에 집중  →  소비 경험의 질적 전환",
    sz=10, color=C_GRAY)
bottom_line(s)


# ═══════════════════════════════════════════════════
# SLIDE 3-C: 대표 리뷰 샘플
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
slide_tag(s, "LAYER 1 — REVIEW SAMPLES")
slide_title(s, "대표 리뷰 텍스트 비교")

# 제품 리뷰
rect(s, Inches(0.8), Inches(2.0), Inches(5.6), Inches(5.0), C_DARK2)
line_shape(s, Inches(0.8), Inches(2.0), Inches(5.6), Inches(0.05), C_CHOCO)
box(s, Inches(1.1), Inches(2.2), Inches(5), Inches(0.35), "A  가나초콜릿 리뷰", sz=15, color=C_CHOCO, bold=True)

prod_reviews = [
    ("\"입에서 사르르 녹는 달콤한 맛이 특징이다\"", "블로그 리뷰", C_GREEN),
    ("\"난 아직도 갈색이랑 빨간색 맛의 차이를 모르겠음.\n그냥 둘 다 맛있으니 내 입속으로\"", "블로그 리뷰", C_GREEN),
    ("\"어릴 때부터 가나만 먹어서 그런가\n가나가 입에 젤 맞다\"", "블로그 리뷰", C_WARN),
    ("\"카카오 함량도 낮고 카카오버터 대신\n싸구려 팜유로 만드니까\"", "클리앙", C_CORAL),
    ("\"가나초콜렛은 가나초콜렛이지 초콜렛은 아니야\"", "루리웹", C_CORAL),
]
for i, (text, src, clr) in enumerate(prod_reviews):
    y = 2.7 + i * 0.8
    line_shape(s, Inches(1.1), Inches(y), Inches(0.04), Inches(0.6), clr)
    box(s, Inches(1.3), Inches(y), Inches(4.8), Inches(0.55), text, sz=9, color=C_GRAY)
    box(s, Inches(1.3), Inches(y+0.55), Inches(2), Inches(0.2), src, sz=7, color=C_GRAY2)

# 공간 리뷰
rect(s, Inches(6.9), Inches(2.0), Inches(5.6), Inches(5.0), C_DARK2)
line_shape(s, Inches(6.9), Inches(2.0), Inches(5.6), Inches(0.05), C_SPACE)
box(s, Inches(7.2), Inches(2.2), Inches(5), Inches(0.35), "B  가나초콜릿하우스 리뷰", sz=15, color=C_SPACE, bold=True)

space_reviews = [
    ("\"결론부터 이야기하자면 생각보다 더 좋았습니다!\n공간의 톤이 초콜릿 컬러에 맞게 구성\"", "Polle", C_GREEN),
    ("\"처음에는 단순한 제품 홍보 팝업으로 생각했는데\n공간과 서비스 완성도에 감동\"", "heyPOP", C_GREEN),
    ("\"다른 카페 가지 말고 여기 있다 가자\"", "현장 대화 / 시사저널e", C_PRIMARY),
    ("\"여기는 초콜릿에 진심이구나\"", "heyPOP", C_GREEN),
    ("\"생각보다 별거 볼 건 없더라고요ㅠㅠ\n초콜릿의 특색이 덜 보이는 점도 아쉬움\"", "Polle", C_CORAL),
]
for i, (text, src, clr) in enumerate(space_reviews):
    y = 2.7 + i * 0.8
    line_shape(s, Inches(7.2), Inches(y), Inches(0.04), Inches(0.6), clr)
    box(s, Inches(7.4), Inches(y), Inches(4.8), Inches(0.55), text, sz=9, color=C_GRAY)
    box(s, Inches(7.4), Inches(y+0.55), Inches(2), Inches(0.2), src, sz=7, color=C_GRAY2)

bottom_line(s)


# ═══════════════════════════════════════════════════
# SLIDE 4: Layer 2 LIWC 4대 변수
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
slide_tag(s, "LAYER 2 — PSYCHE")
slide_title(s, "LIWC 4대 변수 비교: \"어떻게 말했는가\"")

metrics = [
    ("Authenticity (진정성)",  65, 45, "↓20", C_CORAL,
     "제품: \"난 아직도\", \"솔직히\", 부정 표현 포함 → 자기 노출 높음\n공간: 감탄사 과다, 1인칭 부재, 부정 거의 없음 → 사회적 바람직성"),
    ("Clout (추천 확신도)",    35, 65, "↑30", C_GREEN,
     "제품: \"맛있다\" 단순 평가, 추천 발화 없음\n공간: \"꼭 가보세요\", \"여기 있다 가자\", -자 청유형 → 영향력 행사"),
    ("Analytical (인식 방식)", 65, 45, "↓20", C_CORAL,
     "제품: \"카카오 함량\", \"팜유\", \"가격 대비\" → 인과적/비교적 서술\n공간: \"대박\", \"미쳤다\", \"환상\" → 감탄형/직관적 반응"),
    ("Emotional Tone (감정)",  55, 80, "↑25", C_GREEN,
     "제품: 긍정/부정 혼재 → 톤 중화\n공간: \"감동\", \"왕국\", \"환상\" + 부정어 거의 없음 → 순수 긍정 집중"),
]

for i, (name, v_prod, v_space, delta, dc, note) in enumerate(metrics):
    y = 2.0 + i * 1.3
    rect(s, Inches(0.8), Inches(y), Inches(11.7), Inches(1.15), C_DARK2)
    box(s, Inches(1.0), Inches(y+0.1), Inches(2.8), Inches(0.3), name, sz=13, color=C_WHITE, bold=True)

    # 제품 게이지
    box(s, Inches(3.9), Inches(y+0.08), Inches(0.8), Inches(0.25), "제품", sz=9, color=C_CHOCO, bold=True)
    gauge_bar(s, 4.7, y+0.1, v_prod, 2.8, 0.22, C_CHOCO, str(v_prod))

    # 공간 게이지
    box(s, Inches(3.9), Inches(y+0.42), Inches(0.8), Inches(0.25), "공간", sz=9, color=C_SPACE, bold=True)
    gauge_bar(s, 4.7, y+0.44, v_space, 2.8, 0.22, C_SPACE, str(v_space))

    # 변화
    box(s, Inches(8.0), Inches(y+0.15), Inches(0.7), Inches(0.4), delta, sz=16, color=dc, bold=True, align=PP_ALIGN.CENTER)

    # 해석
    box(s, Inches(8.8), Inches(y+0.1), Inches(3.5), Inches(0.9), note, sz=9, color=C_GRAY)

bottom_line(s)


# ═══════════════════════════════════════════════════
# SLIDE 5: Layer 2 추가 심리 지표
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
slide_tag(s, "LAYER 2 — RXR EXCLUSIVE")
slide_title(s, "RXR 고유 심리 지표: Freshness · Distance · 시제 · 감각")

indicators = [
    ("Freshness Index", "\"이/가\" vs \"은/는\"", "0.35", "0.78",
     "제품: \"가나초콜릿은\", \"맛은\" → 이미 알려진 주제\n공간: \"분위기가\", \"조화가\", \"공간이\" → 새로운 발견의 초점화",
     "제품은 51년 된 익숙한 존재\n공간은 새로운 발견으로 체험"),
    ("Psych. Distance", "대명사 패턴", "0.82", "0.41",
     "제품: \"난 아직도\", \"내 입속으로\" → 개인 경험 서술\n공간: \"여기 있다 가자\", \"우리 같이\" → 장소 귀속+집단 경험",
     "제품은 관찰자 시점\n공간은 참여자/귀속 시점"),
    ("시제 프레임", "과거/현재/미래", "과거 회고", "현재+미래",
     "제품: \"먹었다\", \"어릴 때부터\" → 기억 속의 브랜드\n공간: \"펼쳐진다\", \"가자\" → 지금-여기 + 행동 촉구",
     "제품은 노스탤지어\n공간은 현재 몰입+확산 의지"),
    ("감각 축", "지배 감각", "미각+촉각", "시각+공간",
     "제품: \"녹는\", \"부드럽다\", \"달콤\" → 입안 감각\n공간: \"고급스럽다\", \"포토존\", \"분위기\" → 눈과 몸의 공간 경험",
     "제품→공간 전환 시\n감각 축 자체가 이동\nProximal→Distal"),
]

t = tbl(s, Inches(0.8), Inches(1.9), Inches(11.7), Inches(5.2), 5, 5)
hdr(t, ["지표", "제품", "공간", "기능어 근거", "해석"])

for i, (name_val, sub, v1, v2, n1, n2) in enumerate(indicators):
    r = i + 1
    cell_highlight(t.rows[r].cells[0], f"{name_val}\n({sub})", sz=9, bold=True, color=C_PRIMARY)
    cell_highlight(t.rows[r].cells[1], v1, sz=10, bold=True, color=C_CHOCO)
    cell_highlight(t.rows[r].cells[2], v2, sz=10, bold=True, color=C_SPACE)
    cell(t.rows[r].cells[3], n1, sz=8)
    cell(t.rows[r].cells[4], n2, sz=8)

t.columns[0].width = Inches(1.8)
t.columns[1].width = Inches(1.3)
t.columns[2].width = Inches(1.3)
t.columns[3].width = Inches(4.5)
t.columns[4].width = Inches(2.8)

# L2 결론 박스
rect(s, Inches(0.8), Inches(7.1), Inches(11.7), Inches(0.3), C_DARK3)
box(s, Inches(1.0), Inches(7.12), Inches(11.3), Inches(0.25),
    "L2 결론:  Auth↓20 (개인평가→사회공유)  |  Clout↑30 (소비자→추천자)  |  Freshness 0.35→0.78 (익숙→새발견)  → L1만으로는 절대 포착 불가",
    sz=9, color=C_PINK)
bottom_line(s)


# ═══════════════════════════════════════════════════
# SLIDE 6: 패턴 ID 레퍼런스
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
slide_tag(s, "PATTERN DICTIONARY")
slide_title(s, "2-Layer 교차 분석 패턴 사전")

t = tbl(s, Inches(0.8), Inches(1.9), Inches(11.7), Inches(5.2), 7, 6)
hdr(t, ["패턴 ID", "이름", "L1 신호", "L2 신호", "의미", "대응 전략"])

data = [
    ("P-01", "진짜 팬", "긍정 80%+",
     "Auth 높음\nClout 높음",
     "진심으로 좋아하고,\n자신 있게 추천까지 한다",
     "VIP 관리,\n앰배서더 전환"),
    ("P-02", "인스타 소비", "긍정 80%+\n\"사진\", \"예쁘다\"",
     "Auth 낮음\nClout 낮음",
     "바이럴은 되지만 진심 아님.\n\"올리고 싶었을 뿐\"\n재방문 의향 낮음",
     "체험 깊이 강화,\n제품 터치포인트 배치"),
    ("P-03", "관심 있지만\n확신 부족", "긍정 70%+\n\"괜찮다\", \"좋다\"",
     "Auth 높음\nClout 낮음",
     "솔직하게 좋다고 느끼지만\n추천까지는 안 함.\n\"나는 좋은데 네가 좋아할지는...\"",
     "샘플링 확대,\n확신 전환 유도"),
    ("P-04", "신선함 소진", "키워드 변화 없음\n감성 소폭 하락",
     "Freshness 급락\n(이/가→은/는)",
     "내용은 같지만 \"새로운 발견\"\n느낌 사라짐.\n2~3주차부터 체험 피로",
     "콘텐츠 리프레시,\n신규 존 오픈"),
    ("P-05", "탐색은 길지만\n결정 못 함", "체류시간↑\n구매↓",
     "Analytical 높음\nClout 낮음",
     "비교/분석 모드에 갇힘.\n\"좋은 건 알겠는데 뭘 사야 할지\"\n선택 과부하(Choice Overload)",
     "선택지 축소,\n큐레이션 동선"),
    ("P-06", "브랜드 피로도\n(조용한 이탈)", "긍정 유지\nbut 언급량↓",
     "Tone↓ 추세\nClout↓ 추세\nDistance 변화 없음",
     "불만은 아니지만 관심 사라짐.\n\"싫어서 안 가는 게 아니라\n그냥 생각이 안 나서 안 감\"",
     "긴급 리프레시,\n이벤트성 콘텐츠,\n재활성화 프로그램"),
]
colors = [C_GREEN, C_WARN, C_PRIMARY, C_PINK, C_GRAY, C_CORAL]

for i, (pid, name, l1, l2, meaning, strategy) in enumerate(data):
    row = i + 1
    cell_highlight(t.rows[row].cells[0], pid, sz=11, bold=True, color=colors[i])
    cell(t.rows[row].cells[1], name, sz=9, bold=True)
    cell(t.rows[row].cells[2], l1, sz=8)
    cell(t.rows[row].cells[3], l2, sz=8)
    cell(t.rows[row].cells[4], meaning, sz=8)
    cell(t.rows[row].cells[5], strategy, sz=8)

for col_i, w in enumerate([0.9, 1.4, 1.8, 2.0, 3.2, 2.4]):
    t.columns[col_i].width = Inches(w)
bottom_line(s)


# ═══════════════════════════════════════════════════
# SLIDE 7: Zone 매트릭스
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
slide_tag(s, "ZONE MATRIX")
slide_title(s, "감성 × 진정성 매트릭스 — 4개 Zone 분류")

zones = [
    (0.8, 2.2, "Zone A: 진짜 팬", "감성 긍정 + Auth 높음", "진심으로 좋아한다.\n재방문/재구매 확률 가장 높은 세그먼트.", C_GREEN),
    (6.9, 2.2, "Zone B: 인스타 소비", "감성 긍정 + Auth 낮음", "\"예쁘니까 올린 것.\"\n바이럴 기여는 하지만 재방문 가치 낮음.", C_WARN),
    (0.8, 4.6, "Zone C: 진심 불만", "감성 부정 + Auth 높음", "실제 문제 존재. 무시 금지.\n\"팜유\", \"특색 없다\" = 진짜 개선점.", C_CORAL),
    (6.9, 4.6, "Zone D: 노이즈", "감성 부정 + Auth 낮음", "악의적/경쟁사 공격성 리뷰.\n분석 대상에서 필터링 후 제외.", C_GRAY2),
]

# 축 레이블
box(s, Inches(3.2), Inches(1.9), Inches(3), Inches(0.3), "Auth 높음  ←───────→  Auth 낮음", sz=10, color=C_GRAY, align=PP_ALIGN.CENTER)

for x, y, title, cond, desc, color in zones:
    rect(s, Inches(x), Inches(y), Inches(5.6), Inches(2.1), C_DARK2)
    line_shape(s, Inches(x), Inches(y), Inches(5.6), Inches(0.05), color)
    box(s, Inches(x+0.2), Inches(y+0.2), Inches(5.2), Inches(0.35), title, sz=16, color=color, bold=True)
    box(s, Inches(x+0.2), Inches(y+0.6), Inches(5.2), Inches(0.25), cond, sz=10, color=C_PRIMARY)
    box(s, Inches(x+0.2), Inches(y+1.0), Inches(5.2), Inches(1.0), desc, sz=11, color=C_GRAY)

bottom_line(s)


# ═══════════════════════════════════════════════════
# SLIDE 8: 교차 분석 — 제품
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
slide_tag(s, "CROSS ANALYSIS — PRODUCT")
slide_title(s, "제품 리뷰 교차 분석: 패턴 매칭 결과")

t = tbl(s, Inches(0.8), Inches(1.9), Inches(11.7), Inches(3.6), 4, 5)
hdr(t, ["리뷰 유형", "L1 (Content)", "L2 (Psyche)", "패턴", "2-Layer 인사이트"])

rows_data = [
    ("솔직한 팬", "긍정 / \"맛있다\"", "Auth 높음, Clout 낮음", "P-03", "관심은 있지만 추천 의지 없음.\n\"난 좋은데 남한테 권하기는...\""),
    ("분석적 비판자", "부정 / \"팜유\", \"싸구려\"", "Auth 높음, Analytical 높음", "Zone C", "진심 불만. 실제 품질 문제.\n무시하면 안 됨"),
    ("추억형 소비자", "긍정 / \"어릴 때\"", "Auth 높음, 시제: 과거", "P-06 전조", "노스탤지어 의존.\n새 경험 없이 점진적 이탈 가능"),
]
p_colors = [C_PRIMARY, C_CORAL, C_WARN]
for i, (typ, l1, l2, pat, insight) in enumerate(rows_data):
    r = i + 1
    cell(t.rows[r].cells[0], typ, sz=10, bold=True)
    cell(t.rows[r].cells[1], l1, sz=9)
    cell(t.rows[r].cells[2], l2, sz=9)
    cell_highlight(t.rows[r].cells[3], pat, sz=11, bold=True, color=p_colors[i])
    cell(t.rows[r].cells[4], insight, sz=9)

for col_i, w in enumerate([1.5, 2.2, 2.5, 1.2, 4.3]):
    t.columns[col_i].width = Inches(w)

# 진단 박스
rect(s, Inches(0.8), Inches(5.8), Inches(11.7), Inches(1.2), C_DARK3)
line_shape(s, Inches(0.8), Inches(5.8), Inches(0.05), Inches(1.2), C_CORAL)
tb = box(s, Inches(1.1), Inches(5.95), Inches(11.2), Inches(0.3), "제품 진단", sz=13, color=C_CORAL, bold=True)
box(s, Inches(1.1), Inches(6.3), Inches(11.2), Inches(0.6),
    "L1만 보면 \"긍정 48%, 그럭저럭\" — L2를 더하면 긍정 리뷰조차 추천 의지 없고(Clout 35), 과거 추억에 의존.\n브랜드가 현재 시점에서 새로운 긍정 경험을 생산하지 못하고 있다.",
    sz=11, color=C_GRAY)
bottom_line(s)


# ═══════════════════════════════════════════════════
# SLIDE 9: 교차 분석 — 공간
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
slide_tag(s, "CROSS ANALYSIS — SPACE")
slide_title(s, "공간 리뷰 교차 분석: 패턴 매칭 결과")

t = tbl(s, Inches(0.8), Inches(1.9), Inches(11.7), Inches(3.6), 5, 5)
hdr(t, ["리뷰 유형", "L1 (Content)", "L2 (Psyche)", "패턴", "2-Layer 인사이트"])

rows_data = [
    ("인스타 소비형", "긍정 90%+, \"사진\"", "Auth 낮음, Clout 낮음", "P-02", "바이럴 ≠ 충성.\n인스타에 올리고 싶었을 뿐"),
    ("기대 위반형", "긍정 / \"생각보다\"", "Auth 높음, Clout 높음", "P-01", "진짜 팬 전환.\n스키마 위반이 진정 긍정 생성"),
    ("관계 소비형", "긍정 / \"같이 가자\"", "Clout 높음, Distance 낮음", "P-01+", "행동 전환 완료.\n가장 높은 전환 가치"),
    ("실망형", "부정 / \"별거 없다\"", "Auth 높음, Tone 낮음", "Zone C", "진심 불만. 제품-공간 연결 약점"),
]
p_colors = [C_WARN, C_GREEN, C_GREEN, C_CORAL]
for i, (typ, l1, l2, pat, insight) in enumerate(rows_data):
    r = i + 1
    cell(t.rows[r].cells[0], typ, sz=10, bold=True)
    cell(t.rows[r].cells[1], l1, sz=9)
    cell(t.rows[r].cells[2], l2, sz=9)
    cell_highlight(t.rows[r].cells[3], pat, sz=11, bold=True, color=p_colors[i])
    cell(t.rows[r].cells[4], insight, sz=9)

for col_i, w in enumerate([1.5, 2.2, 2.5, 1.2, 4.3]):
    t.columns[col_i].width = Inches(w)

# 세그먼트 분해
rect(s, Inches(0.8), Inches(5.8), Inches(11.7), Inches(1.2), C_DARK3)
line_shape(s, Inches(0.8), Inches(5.8), Inches(0.05), Inches(1.2), C_GREEN)
box(s, Inches(1.1), Inches(5.95), Inches(11.2), Inches(0.3), "공간 진단: 긍정 86%를 분해하면", sz=13, color=C_GREEN, bold=True)

# 3 mini stats
stat_card(s, 1.1, 6.35, "40~45%", "인스타 소비 (P-02)\n바이럴O 재방문X", C_WARN, 3.4, 0.7)
stat_card(s, 4.8, 6.35, "30~35%", "진짜 팬 (P-01)\nAuth+Clout 높음", C_GREEN, 3.4, 0.7)
stat_card(s, 8.5, 6.35, "15~20%", "관계 소비 (P-01+)\n행동 전환 완료", C_PRIMARY, 3.4, 0.7)
bottom_line(s)


# ═══════════════════════════════════════════════════
# SLIDE 10: 5대 심리축 전환
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
slide_tag(s, "KEY FINDING")
slide_title(s, "5대 심리 축의 체계적 전환: 제품 → 공간")

axes = [
    ("자아 축",  "\"나\"",       "\"여기/우리\"", "개인적 → 사회적 경험\n바이럴 구조 형성"),
    ("시간 축",  "과거 회고",    "현재+미래",     "노스탤지어 → 현재 몰입\n브랜드 리포지셔닝"),
    ("감각 축",  "미각+촉각",    "시각+공간",     "Proximal → Distal\nSNS 공유 최적화"),
    ("인지 축",  "분석적 평가",  "감성적 반응",   "가격 경쟁 → 가치 경쟁"),
    ("행동 축",  "수동적 소비",  "능동적 추천",   "평가자 → 자발적 영업사원\n마케팅 비용 구조 변화"),
]

for i, (name, v1, v2, note) in enumerate(axes):
    y = 2.0 + i * 1.05
    rect(s, Inches(0.8), Inches(y), Inches(11.7), Inches(0.9), C_DARK2)
    box(s, Inches(1.0), Inches(y+0.15), Inches(1.6), Inches(0.35), name, sz=14, color=C_PRIMARY, bold=True)

    # 제품
    rect(s, Inches(2.8), Inches(y+0.12), Inches(2.2), Inches(0.55), C_DARK3)
    box(s, Inches(2.8), Inches(y+0.2), Inches(2.2), Inches(0.35), v1, sz=14, color=C_CHOCO, bold=True, align=PP_ALIGN.CENTER)

    # 화살표
    box(s, Inches(5.15), Inches(y+0.15), Inches(0.6), Inches(0.4), "→", sz=22, color=C_PRIMARY, bold=True, align=PP_ALIGN.CENTER)

    # 공간
    rect(s, Inches(5.9), Inches(y+0.12), Inches(2.2), Inches(0.55), C_DARK3)
    box(s, Inches(5.9), Inches(y+0.2), Inches(2.2), Inches(0.35), v2, sz=14, color=C_SPACE, bold=True, align=PP_ALIGN.CENTER)

    # 해석
    box(s, Inches(8.4), Inches(y+0.1), Inches(3.9), Inches(0.7), note, sz=10, color=C_GRAY)

bottom_line(s)


# ═══════════════════════════════════════════════════
# SLIDE 10-B: 5대 심리축 상세 — 기능어 증거 + 비즈니스 의미
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
slide_tag(s, "KEY FINDING — DETAIL")
slide_title(s, "5대 심리축 전환: 기능어 증거 & 비즈니스 의미")

t = tbl(s, Inches(0.8), Inches(1.9), Inches(11.7), Inches(5.2), 6, 4)
hdr(t, ["심리 축", "제품 → 공간", "기능어 증거", "비즈니스 의미"])

axis_detail = [
    ("1. 자아 축\n\"나\"→\"여기/우리\"",
     "개인적 경험\n→ 사회적 경험",
     "1인칭 단수 감소\n장소 지시어+복수형 증가",
     "공간이 소비를 \"공유 가능한 경험\"으로\n전환 → 바이럴 구조 형성"),
    ("2. 시간 축\n과거→현재+미래",
     "노스탤지어\n→ 현재 몰입",
     "과거 시제(-었-) 감소\n현재형+청유형(-자) 증가",
     "브랜드가 \"기억\"에서 \"현재 경험\"으로\n리포지셔닝 → 노스탤지어 의존 탈피"),
    ("3. 감각 축\n미각→시각",
     "입안 감각\n→ 공간 감각",
     "미각/촉각 단어 감소\n시각/공간 단어 증가",
     "Proximal→Distal 전환\nSNS 공유에 최적화된 감각 구조"),
    ("4. 인지 축\n분석→감동",
     "이성적 평가\n→ 감성적 반응",
     "Analytical 65→45\n인과 접속사 감소, 감탄사 증가",
     "분석적 비교 대상→감정적 몰입 대상\n가격 경쟁에서 가치 경쟁으로"),
    ("5. 행동 축\n평가→추천",
     "수동적 소비자\n→ 능동적 추천자",
     "Clout 35→65\n청유형 어미(-자, -세요) 증가",
     "소비자가 자발적 영업사원으로 전환\n마케팅 비용 구조 변화"),
]
for i, (axis, change, evidence, biz) in enumerate(axis_detail):
    r = i + 1
    cell_highlight(t.rows[r].cells[0], axis, sz=9, bold=True, color=C_PRIMARY)
    cell(t.rows[r].cells[1], change, sz=9)
    cell_highlight(t.rows[r].cells[2], evidence, sz=9, bold=False, color=C_PINK, bg=C_DARK3)
    cell(t.rows[r].cells[3], biz, sz=9)

t.columns[0].width = Inches(2.2)
t.columns[1].width = Inches(2.0)
t.columns[2].width = Inches(3.5)
t.columns[3].width = Inches(4.0)

rect(s, Inches(0.8), Inches(7.1), Inches(11.7), Inches(0.3), C_DARK3)
box(s, Inches(1.0), Inches(7.12), Inches(11.3), Inches(0.25),
    "이 전환 데이터의 가치:  \"Clout 35→65, Freshness 0.35→0.78\" = IR/제안서에서 경쟁사가 제출할 수 없는 정량 데이터",
    sz=9, color=C_PRIMARY)
bottom_line(s)


# ═══════════════════════════════════════════════════
# SLIDE 11: L1만 vs 2-Layer 비교
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
slide_tag(s, "COMPARISON")
slide_title(s, "Layer 1만 vs 2-Layer 통합 — 인사이트 품질 비교")

t = tbl(s, Inches(0.8), Inches(2.0), Inches(11.7), Inches(5.0), 6, 3)
hdr(t, ["질문", "Layer 1만", "2-Layer 통합"])

qa = [
    ("공간이 성공했나?", "\"긍정 86%, 성공\"", "진짜 팬은 30~35%.\n40~45%는 인스타 소비형"),
    ("제품 브랜드 건강한가?", "\"긍정 48%, 보통\"", "긍정도 Clout 35, 추천 의지 없음.\n과거 추억 기반"),
    ("공간이 나은 이유?", "\"감성 점수가 높다\"", "5개 심리 축 동시 전환.\n평가자→추천자 변환"),
    ("다음 팝업 개선점?", "\"긍정 유지하면 된다\"", "P-02→P-01 전환 필요.\n체험 깊이 강화"),
    ("시즌4 해야 하나?", "\"긍정이니 해도 됨\"", "Freshness 모니터링 필수.\n주간 리프레시 캘린더 적용"),
]
for i, (q, a1, a2) in enumerate(qa):
    r = i + 1
    cell(t.rows[r].cells[0], q, sz=10, bold=True)
    cell(t.rows[r].cells[1], a1, sz=10, color=C_GRAY)
    cell_highlight(t.rows[r].cells[2], a2, sz=10, bold=False, color=C_WHITE, bg=C_DARK3)

t.columns[0].width = Inches(2.8)
t.columns[1].width = Inches(3.2)
t.columns[2].width = Inches(5.7)
bottom_line(s)


# ═══════════════════════════════════════════════════
# SLIDE 11-B: 세그먼트별 전략
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
slide_tag(s, "STRATEGY BY SEGMENT")
slide_title(s, "가나초콜릿하우스를 위한 세그먼트별 전략")

segs = [
    ("인스타 소비형 (40~45%)", C_WARN,
     "L1: 긍정 / \"예쁘다\", \"사진\"\nL2: Auth 낮음, Clout 낮음",
     "포토존 뒤에 제품 체험 동선 배치.\n\"사진 찍고 → 시식 → 구매\"로 연결.\nSNS 공유 시 체험 쿠폰 제공하여\n바이럴을 체험으로 전환"),
    ("진짜 팬 (30~35%)", C_GREEN,
     "L1: 긍정 / \"감동\", \"진심\"\nL2: Auth 높음, Clout 높음",
     "VIP 프로그램(시즌 프리오픈, 한정 메뉴).\n리뷰어 앰배서더 프로그램.\n이 세그먼트의 리뷰가 다른 소비자의\nP-02→P-01 전환을 유도"),
    ("관계 소비형 (15~20%)", C_PRIMARY,
     "L1: 긍정 / \"같이\", \"데이트\"\nL2: Clout 높음, Distance 낮음",
     "2인/그룹 전용 메뉴, 커플/가족 체험 코스.\n\"함께 만드는 초콜릿\" DIY를 핵심 프로그램.\n관계 기억에 브랜드 각인"),
]

for i, (title, color, profile, strategy) in enumerate(segs):
    x = 0.8 + i * 4.1
    rect(s, Inches(x), Inches(2.0), Inches(3.8), Inches(5.2), C_DARK2)
    line_shape(s, Inches(x), Inches(2.0), Inches(3.8), Inches(0.05), color)
    box(s, Inches(x+0.2), Inches(2.2), Inches(3.4), Inches(0.4), title, sz=14, color=color, bold=True)

    # 프로필
    box(s, Inches(x+0.2), Inches(2.8), Inches(3.4), Inches(0.3), "프로필", sz=10, color=C_GRAY2, bold=True)
    rect(s, Inches(x+0.2), Inches(3.15), Inches(3.4), Inches(1.2), C_DARK3)
    box(s, Inches(x+0.35), Inches(3.2), Inches(3.1), Inches(1.1), profile, sz=9, color=C_GRAY)

    # 전략
    box(s, Inches(x+0.2), Inches(4.5), Inches(3.4), Inches(0.3), "전략", sz=10, color=color, bold=True)
    rect(s, Inches(x+0.2), Inches(4.85), Inches(3.4), Inches(2.2), C_DARK3)
    box(s, Inches(x+0.35), Inches(4.9), Inches(3.1), Inches(2.1), strategy, sz=10, color=C_WHITE)

bottom_line(s)


# ═══════════════════════════════════════════════════
# SLIDE 12: 핵심 요약 (L1 vs 2-Layer 대비 포함)
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
line_shape(s, Inches(0), Inches(0), Inches(13.333), Inches(0.06))

box(s, Inches(1), Inches(1.0), Inches(11), Inches(1.2),
    "같은 \"긍정 86%\"도, Layer 2를 추가하면\n3개의 전혀 다른 세그먼트가 보인다.", sz=32, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)

# L1 vs 2-Layer 대비 불릿
rect(s, Inches(0.8), Inches(2.5), Inches(5.5), Inches(3.0), C_DARK2)
line_shape(s, Inches(0.8), Inches(2.5), Inches(5.5), Inches(0.04), C_WARN)
box(s, Inches(1.1), Inches(2.6), Inches(5), Inches(0.3), "Layer 1만의 결론", sz=13, color=C_WARN, bold=True)
tb = box(s, Inches(1.1), Inches(3.0), Inches(5), Inches(0.3), "  - 제품: 긍정 48%, 보통", sz=11, color=C_GRAY)
para(tb.text_frame, "  - 공간: 긍정 86%, 대성공", sz=11, color=C_GRAY, sp_before=3)
para(tb.text_frame, "  - \"공간이 제품보다 좋다\"", sz=11, color=C_GRAY, sp_before=3)
para(tb.text_frame, "  - 전략? \"계속 하면 된다\"", sz=11, color=C_GRAY, sp_before=3)

rect(s, Inches(6.9), Inches(2.5), Inches(5.6), Inches(3.0), C_DARK2)
line_shape(s, Inches(6.9), Inches(2.5), Inches(5.6), Inches(0.04), C_PRIMARY)
box(s, Inches(7.2), Inches(2.6), Inches(5), Inches(0.3), "2-Layer 통합의 결론", sz=13, color=C_PRIMARY, bold=True)
tb = box(s, Inches(7.2), Inches(3.0), Inches(5), Inches(0.3), "  - 제품: Clout 35, 추천 의지 없음. 과거 추억 의존", sz=11, color=C_WHITE)
para(tb.text_frame, "  - 공간: 진짜 팬 30~35%, 인스타 소비 40~45%", sz=11, color=C_WHITE, sp_before=3)
para(tb.text_frame, "  - 제품→공간 전환 시 5개 심리 축 동시 전환 실증", sz=11, color=C_WHITE, sp_before=3)
para(tb.text_frame, "  - 전략: 세그먼트별 차별화 + Freshness 모니터링", sz=11, color=C_WHITE, sp_before=3)

# 핵심 메시지
rect(s, Inches(0.8), Inches(5.7), Inches(11.7), Inches(0.6), C_DARK3)
box(s, Inches(1.0), Inches(5.75), Inches(11.3), Inches(0.5),
    "Layer 1은 \"좋다/나쁘다\"를 말해준다.  Layer 2는 \"왜 좋은지, 얼마나 진심인지, 언제까지 지속될지\"를 말해준다.",
    sz=13, color=C_PRIMARY, align=PP_ALIGN.CENTER)

# 5축 미니 카드
axes_mini = [
    ("자아 축", "\"나\"→\"우리\""),
    ("시간 축", "과거→현재"),
    ("감각 축", "미각→시각"),
    ("인지 축", "분석→감동"),
    ("행동 축", "평가→추천"),
]
for i, (name, val) in enumerate(axes_mini):
    x = 1.2 + i * 2.3
    rect(s, Inches(x), Inches(6.4), Inches(2.0), Inches(0.8), C_DARK2)
    line_shape(s, Inches(x), Inches(6.4), Inches(2.0), Inches(0.04), C_PRIMARY)
    box(s, Inches(x), Inches(6.5), Inches(2.0), Inches(0.2), name, sz=9, color=C_PRIMARY, bold=True, align=PP_ALIGN.CENTER)
    box(s, Inches(x), Inches(6.75), Inches(2.0), Inches(0.3), val, sz=12, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
bottom_line(s)


# ─── 저장 ───
out = os.path.join(os.path.dirname(__file__), "rxr-2layer-case-ghana-v2.pptx")
prs.save(out)
print("Done:", out)
print("Slides:", len(prs.slides))
