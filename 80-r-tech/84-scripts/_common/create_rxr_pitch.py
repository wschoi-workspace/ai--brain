"""
RXR 4축 프레임워크 — VC/주니어용 비즈니스 피치덱
16:9 / R스타일가이드 적용 (Primary: #FF5050, Font: 맑은 고딕)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── R 스타일가이드 컬러 ───
R_RED = RGBColor(0xFF, 0x50, 0x50)
R_DARK = RGBColor(0x1A, 0x1A, 0x1A)
R_DARK2 = RGBColor(0x2A, 0x2A, 0x2A)
R_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
R_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)
R_GRAY = RGBColor(0x99, 0x99, 0x99)
R_GRAY2 = RGBColor(0x66, 0x66, 0x66)
R_ACCENT = RGBColor(0xFF, 0x50, 0x50)  # 레드
R_BLUE = RGBColor(0x3B, 0x82, 0xF6)
R_GREEN = RGBColor(0x10, 0xB9, 0x81)
R_PINK = RGBColor(0xD9, 0x46, 0xEF)
R_ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
R_CYAN = RGBColor(0x0E, 0xA5, 0xE9)

FONT_KR = "맑은 고딕"
FONT_EN = "Arial"  # ADAM.CG PRO 대체

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ─── 헬퍼 ───
def bg(slide, color=R_DARK):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = color

def rect(slide, l, t, w, h, color=R_DARK2):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    try: s.adjustments[0] = 0.06
    except: pass
    return s

def box(slide, l, t, w, h, txt="", sz=14, color=R_WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = txt
    p.font.size = Pt(sz); p.font.color.rgb = color; p.font.bold = bold; p.font.name = FONT_KR; p.alignment = align
    return tb

def para(tf, txt, sz=14, color=R_WHITE, bold=False, align=PP_ALIGN.LEFT, sp_before=0):
    p = tf.add_paragraph(); p.text = txt
    p.font.size = Pt(sz); p.font.color.rgb = color; p.font.bold = bold; p.font.name = FONT_KR; p.alignment = align
    if sp_before: p.space_before = Pt(sp_before)
    return p

def stat(slide, l, t, val, label, vc=R_RED):
    rect(slide, l, t, Inches(2.4), Inches(1.4), R_DARK2)
    box(slide, l, t+Inches(0.15), Inches(2.4), Inches(0.7), val, sz=32, color=vc, bold=True, align=PP_ALIGN.CENTER)
    box(slide, l, t+Inches(0.85), Inches(2.4), Inches(0.5), label, sz=10, color=R_GRAY, align=PP_ALIGN.CENTER)

def tbl(slide, l, t, w, h, rows, cols):
    return slide.shapes.add_table(rows, cols, l, t, w, h).table

def hdr(table, texts, widths=None, bg_color=R_RED):
    for i, txt in enumerate(texts):
        c = table.rows[0].cells[i]; c.text = txt; c.fill.solid(); c.fill.fore_color.rgb = bg_color
        for p in c.text_frame.paragraphs:
            p.font.color.rgb = R_WHITE; p.font.size = Pt(10); p.font.bold = True; p.font.name = FONT_KR
    if widths:
        for i, w in enumerate(widths): table.columns[i].width = Inches(w)

def cell(c, txt, sz=10, bold=False, color=R_WHITE):
    c.text = txt
    for p in c.text_frame.paragraphs:
        p.font.size = Pt(sz); p.font.bold = bold; p.font.color.rgb = color; p.font.name = FONT_KR

def bar(slide, l, t, pct, max_w, h, color, label=""):
    bw = int(max_w * pct / 100)
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, bw, h)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    try: s.adjustments[0] = 0.5
    except: pass
    if label:
        box(slide, l + bw + Inches(0.1), t - Inches(0.02), Inches(1.5), h, label, sz=9, color=R_GRAY)


# ═══════════════════════════════════════
# SLIDE 1: 표지
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, R_DARK)

# 레드 라인
slide_shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.06))
slide_shape.fill.solid(); slide_shape.fill.fore_color.rgb = R_RED; slide_shape.line.fill.background()

box(s, Inches(1), Inches(1.5), Inches(6), Inches(0.4), "RXR 4-AXIS DATA FRAMEWORK", sz=13, color=R_RED, bold=True)
box(s, Inches(1), Inches(2.3), Inches(11), Inches(1.2), "오프라인 매장의 모든 것을\n숫자로 바꾸고, 고객의 진심까지 읽는다", sz=40, color=R_WHITE, bold=True)
box(s, Inches(1), Inches(4.0), Inches(9), Inches(0.8), "기존 솔루션은 퍼즐의 한 조각만 봅니다.\nRXR은 전체 그림을 완성하고, 고객의 무의식까지 데이터화합니다.", sz=16, color=R_GRAY)
box(s, Inches(1), Inches(5.5), Inches(5), Inches(0.4), "Project RENT  ·  R-lab  ·  2026", sz=12, color=R_GRAY2)

# ═══════════════════════════════════════
# SLIDE 2: PROBLEM
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, R_DARK)
box(s, Inches(0.8), Inches(0.5), Inches(3), Inches(0.3), "PROBLEM", sz=11, color=R_RED, bold=True)
box(s, Inches(0.8), Inches(0.9), Inches(10), Inches(0.7), "지금 시장의 블라인드 스팟", sz=32, color=R_WHITE, bold=True)

cards = [
    ("POS만 보는 세상", "\"얼마 팔렸는지\"만 안다.\n1,000명이 왔는데 10명만 샀다면,\n990명은 왜 안 샀는지 모른다.", "유입 · 행동 · 인식 = 전부 사각지대"),
    ("카메라만 보는 세상", "\"어디서 멈췄는지\"까지.\n메이아이는 동선과 체류를 잡지만\n그게 매출이 됐는지는 모른다.", "매출 연결 · 고객 심리 = 사각지대"),
    ("설문만 보는 세상", "\"좋았다고 말했다.\"\nNPS 점수는 높은데\n그게 진심인지 모른다.", "응답 편향 · 무의식 심리 = 사각지대"),
]
for i, (title, desc, blind) in enumerate(cards):
    x = 0.8 + i * 4.0
    rect(s, Inches(x), Inches(2.0), Inches(3.7), Inches(4.8), R_DARK2)
    box(s, Inches(x+0.2), Inches(2.2), Inches(3.3), Inches(0.4), title, sz=18, color=R_WHITE, bold=True, align=PP_ALIGN.CENTER)
    box(s, Inches(x+0.2), Inches(2.8), Inches(3.3), Inches(2.5), desc, sz=13, color=R_GRAY, align=PP_ALIGN.CENTER)
    # blind 박스
    rect(s, Inches(x+0.15), Inches(5.6), Inches(3.4), Inches(0.9), RGBColor(0x3A, 0x1A, 0x1A))
    box(s, Inches(x+0.3), Inches(5.7), Inches(3.1), Inches(0.7), blind, sz=11, color=R_RED, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════
# SLIDE 3: SOLUTION — 4축 퍼널
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, R_DARK)
box(s, Inches(0.8), Inches(0.5), Inches(3), Inches(0.3), "SOLUTION", sz=11, color=R_RED, bold=True)
box(s, Inches(0.8), Inches(0.9), Inches(10), Inches(0.7), "RXR 4축 프레임워크", sz=32, color=R_WHITE, bold=True)
box(s, Inches(0.8), Inches(1.5), Inches(10), Inches(0.4), "봤다 → 들어왔다 → 머물렀다 → 샀다 → 공유했다 → 진짜 느꼈다", sz=13, color=R_GRAY)

# 퍼널 단계
steps = [
    ("AXIS 1", "유입\nTraffic", "\"몇 명이 봤고\n들어왔는가?\"", "방문자 수\n시선 인식률\n입장 전환률", R_BLUE),
    ("AXIS 2", "행동\nBehavior", "\"어디서 머물고\n뭘 봤는가?\"", "동선 히트맵\n체류 시간\n성별·연령", RGBColor(0x63, 0x66, 0xF1)),
    ("AXIS 3", "매출\nConversion", "\"누가 얼마나\n샀는가?\"", "구매 전환률\n객단가\n타겟별 구매", R_GREEN),
    ("AXIS 4", "인식\nFeedback", "\"어떻게 느끼고\n퍼뜨렸는가?\"", "SNS 언급\n감성 분석\nNPS", R_PINK),
    ("+ FWA", "진심\n읽기", "\"정말 진심으로\n좋았는가?\"", "진정성\n추천 확신도\n신선함", R_RED),
]
for i, (num, title, q, kpi, color) in enumerate(steps):
    x = 0.5 + i * 2.5
    rect(s, Inches(x), Inches(2.2), Inches(2.3), Inches(4.8), R_DARK2)
    # 상단 색상 바
    top_bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(2.2), Inches(2.3), Inches(0.08))
    top_bar.fill.solid(); top_bar.fill.fore_color.rgb = color; top_bar.line.fill.background()
    box(s, Inches(x), Inches(2.4), Inches(2.3), Inches(0.3), num, sz=10, color=color, bold=True, align=PP_ALIGN.CENTER)
    box(s, Inches(x), Inches(2.75), Inches(2.3), Inches(0.7), title, sz=16, color=R_WHITE, bold=True, align=PP_ALIGN.CENTER)
    box(s, Inches(x+0.1), Inches(3.5), Inches(2.1), Inches(0.8), q, sz=11, color=R_GRAY, align=PP_ALIGN.CENTER)
    box(s, Inches(x+0.1), Inches(4.5), Inches(2.1), Inches(1.5), kpi, sz=11, color=RGBColor(0xCC, 0xCC, 0xCC), align=PP_ALIGN.CENTER)
    if i < 4:
        box(s, Inches(x+2.3), Inches(3.8), Inches(0.2), Inches(0.4), ">", sz=16, color=R_GRAY2, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════
# SLIDE 4: SECRET WEAPON — FWA Before/After
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, R_DARK)
box(s, Inches(0.8), Inches(0.5), Inches(3), Inches(0.3), "SECRET WEAPON", sz=11, color=R_RED, bold=True)
box(s, Inches(0.8), Inches(0.9), Inches(11), Inches(0.7), "같은 별점 5점이라도 진심의 깊이가 다르다", sz=28, color=R_WHITE, bold=True)

# Before
rect(s, Inches(0.8), Inches(2.0), Inches(5.2), Inches(4.5), R_DARK2)
box(s, Inches(0.8), Inches(2.1), Inches(5.2), Inches(0.3), "BEFORE — 기존 분석", sz=11, color=R_GRAY, bold=True, align=PP_ALIGN.CENTER)
box(s, Inches(0.8), Inches(2.6), Inches(5.2), Inches(0.8), "81%", sz=48, color=R_GRAY, bold=True, align=PP_ALIGN.CENTER)
box(s, Inches(0.8), Inches(3.4), Inches(5.2), Inches(0.3), "긍정 리뷰 비율", sz=12, color=R_GRAY2, align=PP_ALIGN.CENTER)
rect(s, Inches(1.2), Inches(4.0), Inches(4.4), Inches(1.5), RGBColor(0x22, 0x22, 0x22))
box(s, Inches(1.4), Inches(4.1), Inches(4), Inches(1.3), "\"긍정이 81%니까 좋은 거네요\"\n\n→ 끝. 더 이상 알 수 없음.", sz=13, color=R_GRAY)

# VS
box(s, Inches(6.1), Inches(3.3), Inches(1.1), Inches(0.6), "VS", sz=28, color=R_GRAY2, bold=True, align=PP_ALIGN.CENTER)

# After
rect(s, Inches(7.3), Inches(2.0), Inches(5.2), Inches(4.5), R_DARK2)
# 레드 보더
top_bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.3), Inches(2.0), Inches(5.2), Inches(0.06))
top_bar.fill.solid(); top_bar.fill.fore_color.rgb = R_RED; top_bar.line.fill.background()
box(s, Inches(7.3), Inches(2.15), Inches(5.2), Inches(0.3), "AFTER — RXR + FWA", sz=11, color=R_RED, bold=True, align=PP_ALIGN.CENTER)
box(s, Inches(7.3), Inches(2.6), Inches(5.2), Inches(0.8), "81%", sz=48, color=R_RED, bold=True, align=PP_ALIGN.CENTER)
box(s, Inches(7.3), Inches(3.4), Inches(5.2), Inches(0.3), "긍정 리뷰 비율", sz=12, color=R_GRAY, align=PP_ALIGN.CENTER)
rect(s, Inches(7.7), Inches(4.0), Inches(4.4), Inches(1.5), RGBColor(0x2A, 0x1A, 0x1A))
tb = box(s, Inches(7.9), Inches(4.1), Inches(4), Inches(1.3), "", sz=13, color=R_GRAY)
tf = tb.text_frame
para(tf, "Authenticity 42 → 과장/편향 의심", sz=13, color=R_RED, bold=True)
para(tf, "Clout 38 → 단순 감상, 추천 확신 약함", sz=13, color=R_RED, bold=True)
para(tf, "→ \"화제성은 있으나 진정한 옹호는 아직\"", sz=12, color=R_GREEN, sp_before=6)

# 하단 4 메트릭
metrics = [
    ("Authenticity", "진심도 측정기", "\"진짜인가, 과장인가\""),
    ("Clout", "추천 확신도", "\"자신 있게 권하는가\""),
    ("Analytical", "인식 방식", "\"이성적 vs 감성적\""),
    ("Emotional Tone", "감정 심층 구조", "\"표면 너머 진짜 감정\""),
]
for i, (name, sub, desc) in enumerate(metrics):
    x = 0.8 + i * 3.1
    rect(s, Inches(x), Inches(6.0), Inches(2.9), Inches(1.1), R_DARK2)
    box(s, Inches(x+0.1), Inches(6.05), Inches(2.7), Inches(0.3), name, sz=13, color=R_RED, bold=True, align=PP_ALIGN.CENTER)
    box(s, Inches(x+0.1), Inches(6.35), Inches(2.7), Inches(0.3), sub, sz=11, color=R_WHITE, align=PP_ALIGN.CENTER)
    box(s, Inches(x+0.1), Inches(6.65), Inches(2.7), Inches(0.3), desc, sz=10, color=R_GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════
# SLIDE 5: CASE — 가나초콜릿
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, R_DARK)
box(s, Inches(0.8), Inches(0.5), Inches(3), Inches(0.3), "PROVEN CASE", sz=11, color=R_ORANGE, bold=True)
box(s, Inches(0.8), Inches(0.9), Inches(10), Inches(0.7), "가나초콜릿이 증명한 것", sz=28, color=R_WHITE, bold=True)
box(s, Inches(0.8), Inches(1.5), Inches(10), Inches(0.3), "같은 브랜드가 제품→공간으로 확장됐을 때, 소비자의 언어 구조가 완전히 달라졌습니다", sz=12, color=R_GRAY)

# 감각 바
senses = ["미각", "촉각", "시각", "공간", "사회"]
prod = [90, 60, 15, 5, 10]
space = [35, 15, 85, 75, 70]

# 제품
rect(s, Inches(0.8), Inches(2.1), Inches(5.5), Inches(3.3), R_DARK2)
box(s, Inches(1.0), Inches(2.2), Inches(5), Inches(0.3), "가나초콜릿 (제품)", sz=13, color=R_ORANGE, bold=True)
for i, (nm, p) in enumerate(zip(senses, prod)):
    y = 2.65 + i * 0.55
    box(s, Inches(1.0), Inches(y), Inches(0.6), Inches(0.2), nm, sz=10, color=R_ORANGE, bold=True, align=PP_ALIGN.RIGHT)
    bar(s, Inches(1.7), Inches(y), p, Inches(3.2), Inches(0.22), RGBColor(0xB4, 0x53, 0x09), f"{p}%")

# 공간
rect(s, Inches(7.0), Inches(2.1), Inches(5.5), Inches(3.3), R_DARK2)
box(s, Inches(7.2), Inches(2.2), Inches(5), Inches(0.3), "가나초콜릿하우스 (공간)", sz=13, color=R_GREEN, bold=True)
for i, (nm, sp) in enumerate(zip(senses, space)):
    y = 2.65 + i * 0.55
    box(s, Inches(7.2), Inches(y), Inches(0.6), Inches(0.2), nm, sz=10, color=R_GREEN, bold=True, align=PP_ALIGN.RIGHT)
    bar(s, Inches(7.9), Inches(y), sp, Inches(3.2), Inches(0.22), R_GREEN, f"{sp}%")

# LIWC 변수
liwc_data = [("Analytical", "65", "45", "이성적→감성적"), ("Clout", "35", "65", "소극→적극추천"), ("Authenticity", "65", "45", "솔직→인스타감성"), ("Tone", "55", "80", "중립→압도적 긍정")]
for i, (nm, v1, v2, desc) in enumerate(liwc_data):
    x = 0.8 + i * 3.1
    rect(s, Inches(x), Inches(5.6), Inches(2.9), Inches(1.5), R_DARK2)
    box(s, Inches(x), Inches(5.65), Inches(2.9), Inches(0.25), nm, sz=10, color=R_GRAY, align=PP_ALIGN.CENTER)
    box(s, Inches(x), Inches(5.95), Inches(1.3), Inches(0.4), v1, sz=22, color=R_ORANGE, bold=True, align=PP_ALIGN.CENTER)
    box(s, Inches(x+1.15), Inches(6.05), Inches(0.6), Inches(0.3), "→", sz=14, color=R_GRAY2, align=PP_ALIGN.CENTER)
    box(s, Inches(x+1.6), Inches(5.95), Inches(1.3), Inches(0.4), v2, sz=22, color=R_GREEN, bold=True, align=PP_ALIGN.CENTER)
    box(s, Inches(x), Inches(6.45), Inches(2.9), Inches(0.3), desc, sz=10, color=R_GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════
# SLIDE 6: WHY EMOTION MATTERS
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, R_DARK)
box(s, Inches(0.8), Inches(0.5), Inches(3), Inches(0.3), "WHY IT MATTERS", sz=11, color=R_RED, bold=True)
box(s, Inches(0.8), Inches(0.9), Inches(11), Inches(0.7), "소비자의 감정을 읽는 것이 왜 비즈니스의 핵심인가", sz=26, color=R_WHITE, bold=True)

reasons = [
    ("매출은 후행지표다", "감정 데이터는 선행지표.\nClout 점수가 하락하면\n2~3주 뒤 재방문율이 떨어짐.\n매출이 떨어지기 전에\n운영 전략을 바꿀 수 있다.", "-2~3주\n매출 하락 전\n조기 경고"),
    ("설문은 거짓말한다", "NPS \"추천하겠다\" 답변 중\n실제 추천 비율은 20% 미만.\n기능어분석은 무의식적\n언어 패턴을 읽어\n이 편향을 뚫고 진심을 포착.", "FWA\n무의식 패턴\n= 진짜 의도"),
    ("진짜 팬 vs 구경꾼", "방문자 1만 명 중\n진짜 브랜드를 사랑하는\n사람은 몇 명인가?\nClout 높은 리뷰어 =\n자발적 추천 핵심 팬.", "ROI\n핵심 팬 식별\n→ 타겟 집중"),
]
for i, (title, desc, kpi) in enumerate(reasons):
    x = 0.8 + i * 4.0
    rect(s, Inches(x), Inches(1.9), Inches(3.7), Inches(5.2), R_DARK2)
    # 상단 레드 바
    top_bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.9), Inches(3.7), Inches(0.06))
    top_bar.fill.solid(); top_bar.fill.fore_color.rgb = R_RED; top_bar.line.fill.background()
    box(s, Inches(x+0.2), Inches(2.1), Inches(3.3), Inches(0.35), title, sz=17, color=R_WHITE, bold=True, align=PP_ALIGN.CENTER)
    box(s, Inches(x+0.2), Inches(2.6), Inches(3.3), Inches(2.8), desc, sz=12, color=R_GRAY, align=PP_ALIGN.CENTER)
    # KPI 박스
    rect(s, Inches(x+0.3), Inches(5.7), Inches(3.1), Inches(1.1), RGBColor(0x3A, 0x1A, 0x1A))
    box(s, Inches(x+0.3), Inches(5.8), Inches(3.1), Inches(0.9), kpi, sz=12, color=R_RED, bold=True, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════
# SLIDE 7: HOW IT WORKS
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, R_DARK)
box(s, Inches(0.8), Inches(0.5), Inches(3), Inches(0.3), "HOW IT WORKS", sz=11, color=R_RED, bold=True)
box(s, Inches(0.8), Inches(0.9), Inches(10), Inches(0.7), "어떻게 감정을 분석하는가", sz=28, color=R_WHITE, bold=True)

# 3단계
steps_hw = [
    ("STEP 1. 수집", "SNS 리뷰, 설문 개방형 응답,\n블로그 후기 등 텍스트 수집"),
    ("STEP 2. 기능어 추출", "대명사(나/우리/여기),\n조사(은/는 vs 이/가),\n감정어, 시제 패턴 분리"),
    ("STEP 3. 심리 수치화", "LIWC 알고리즘으로\nAuthenticity, Clout 등\n0~100점으로 변환"),
]
for i, (title, desc) in enumerate(steps_hw):
    x = 0.8 + i * 4.0
    rect(s, Inches(x), Inches(1.9), Inches(3.7), Inches(2.3), R_DARK2)
    box(s, Inches(x+0.2), Inches(2.0), Inches(3.3), Inches(0.35), title, sz=14, color=R_RED, bold=True, align=PP_ALIGN.CENTER)
    box(s, Inches(x+0.2), Inches(2.5), Inches(3.3), Inches(1.5), desc, sz=12, color=R_GRAY, align=PP_ALIGN.CENTER)
    if i < 2:
        box(s, Inches(x+3.7), Inches(2.7), Inches(0.3), Inches(0.5), ">", sz=18, color=R_GRAY2, align=PP_ALIGN.CENTER)

# 리뷰 비교
rect(s, Inches(0.8), Inches(4.5), Inches(5.8), Inches(2.7), R_DARK2)
# 레드 헤더
h_bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.5), Inches(5.8), Inches(0.4))
h_bar.fill.solid(); h_bar.fill.fore_color.rgb = RGBColor(0x3A, 0x1A, 0x1A); h_bar.line.fill.background()
box(s, Inches(1.0), Inches(4.52), Inches(5.4), Inches(0.35), "리뷰 A — 과장/조작 의심  |  Authenticity: 28", sz=11, color=R_RED, bold=True)
box(s, Inches(1.0), Inches(5.0), Inches(5.4), Inches(1.3), "\"나는 정말 너무너무 좋았고 진짜 최고의\n경험이었어요!! 분위기도 완전 대박이고\n모든 게 완벽했습니다!!\"\n\n→ 1인칭 과다 + 극단적 긍정어 반복 + 느낌표 과잉", sz=11, color=R_GRAY)

rect(s, Inches(6.9), Inches(4.5), Inches(5.8), Inches(2.7), R_DARK2)
h_bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.9), Inches(4.5), Inches(5.8), Inches(0.4))
h_bar.fill.solid(); h_bar.fill.fore_color.rgb = RGBColor(0x0A, 0x2D, 0x1C); h_bar.line.fill.background()
box(s, Inches(7.1), Inches(4.52), Inches(5.4), Inches(0.35), "리뷰 B — 진정한 경험 공유  |  Authenticity: 78", sz=11, color=R_GREEN, bold=True)
box(s, Inches(7.1), Inches(5.0), Inches(5.4), Inches(1.3), "\"처음엔 단순 팝업인 줄 알았는데, 공간이\n생각보다 정성스럽게 구성돼 있었어요.\n테린느가 괜찮았고, 한번쯤 가볼 만해요.\"\n\n→ 초점 조사 + 구체적 디테일 + 균형 잡힌 평가", sz=11, color=R_GRAY)

# ═══════════════════════════════════════
# SLIDE 8: BUSINESS IMPACT
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, R_DARK)
box(s, Inches(0.8), Inches(0.5), Inches(3), Inches(0.3), "BUSINESS IMPACT", sz=11, color=R_GREEN, bold=True)
box(s, Inches(0.8), Inches(0.9), Inches(10), Inches(0.7), "���정 데이터가 만드는 5가지 비즈니스 효과", sz=26, color=R_WHITE, bold=True)

impacts = [
    ("브랜드 피로도 조기 감지", "Freshness Index 하락 → 매출 하락 전 운영 전략 변경", "-2~3주", "조기 경고"),
    ("가짜 바이럴 vs 진짜 옹호 구분", "Authenticity 기반 캠페인 ROI 정밀 측정", "79%", "탐지 정확도"),
    ("체류↑ 구매↓ 원인 규명", "Analytical↑ Clout↓ = 관심은 있으나 확신 부족", "+19%", "전환율 개선 가능"),
    ("다음 프로젝트 재현 근거", "Clout 최고 요소를 식별 → 데이터 기반 기획", "DATA", "감 아닌 근거"),
    ("IR/제안서 독점 지표", "경쟁사가 말할 수 없는 심리 데이터 = 가치 경쟁", "ONLY", "독점 영역"),
]
for i, (title, desc, val, vlabel) in enumerate(impacts):
    y = 1.8 + i * 1.08
    rect(s, Inches(0.8), Inches(y), Inches(11.7), Inches(0.95), R_DARK2)
    box(s, Inches(1.0), Inches(y+0.05), Inches(3.5), Inches(0.35), title, sz=14, color=R_WHITE, bold=True)
    box(s, Inches(1.0), Inches(y+0.45), Inches(7.5), Inches(0.4), desc, sz=11, color=R_GRAY)
    # KPI
    rect(s, Inches(10.4), Inches(y+0.1), Inches(1.8), Inches(0.75), RGBColor(0x3A, 0x1A, 0x1A))
    box(s, Inches(10.4), Inches(y+0.1), Inches(1.8), Inches(0.45), val, sz=18, color=R_RED, bold=True, align=PP_ALIGN.CENTER)
    box(s, Inches(10.4), Inches(y+0.5), Inches(1.8), Inches(0.3), vlabel, sz=9, color=R_GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════
# SLIDE 9: GLOBAL LANDSCAPE
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, R_DARK)
box(s, Inches(0.8), Inches(0.5), Inches(3), Inches(0.3), "GLOBAL LANDSCAPE", sz=11, color=R_CYAN, bold=True)
box(s, Inches(0.8), Inches(0.9), Inches(10), Inches(0.7), "해외 리테일테크와 무엇이 다른가", sz=28, color=R_WHITE, bold=True)

solutions = [
    ("RetailNext", "미국 · 세계 1위", "센서+비디오\n히트맵, 동선, 전환율", "감정·인식·SNS = 전부 불가"),
    ("Sensormatic", "미국 · 연 400억건", "최대 트래픽 네트워크\n직원 제외, Re-ID", "\"왜 샀는가\" 답할 수 없음"),
    ("Placer.ai", "미국 · 모바일GPS", "센서 없이 위치 분석\n경쟁 매장 이동 패턴", "매장 안 행동 모름. 감정 제로"),
    ("Quividi", "프랑스 · 표정분석", "유일한 표정 mood 측정\n사이니지 앞 주목도", "표정 읽기 ≠ 진심 읽기"),
    ("V-Count", "터키 · 130국", "99% 정확도, 200+ KPI\n가성비 히트맵", "행동 \"너머\" 데이터 없음"),
    ("메이아이", "한국 · 300+ 매장", "기존 CCTV 활용\nRe-ID 92%", "POS 미연동. 감정 불가"),
]
for i, (name, loc, feat, blind) in enumerate(solutions):
    row = i // 3; col = i % 3
    x = 0.8 + col * 4.0; y = 1.8 + row * 2.7
    rect(s, Inches(x), Inches(y), Inches(3.7), Inches(2.4), R_DARK2)
    box(s, Inches(x+0.15), Inches(y+0.1), Inches(3.4), Inches(0.3), name, sz=15, color=R_WHITE, bold=True)
    box(s, Inches(x+0.15), Inches(y+0.4), Inches(3.4), Inches(0.2), loc, sz=9, color=R_GRAY)
    box(s, Inches(x+0.15), Inches(y+0.7), Inches(3.4), Inches(0.8), feat, sz=10, color=R_GRAY)
    rect(s, Inches(x+0.1), Inches(y+1.6), Inches(3.5), Inches(0.65), RGBColor(0x3A, 0x1A, 0x1A))
    box(s, Inches(x+0.2), Inches(y+1.65), Inches(3.3), Inches(0.55), blind, sz=10, color=R_RED)

# ═══════════════════════════════════════
# SLIDE 10: VS CONSULTING
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, R_DARK)
box(s, Inches(0.8), Inches(0.5), Inches(5), Inches(0.3), "VS CONSULTING FRAMEWORKS", sz=11, color=R_RED, bold=True)
box(s, Inches(0.8), Inches(0.9), Inches(11), Inches(0.7), "글로벌 리서치 프레임워크와도 차원이 다르다", sz=26, color=R_WHITE, bold=True)

consults = [
    ("McKinsey", "CMAC · Retail Analytics", "SKU 경제성, 가격 탄력성,\n고객 CLV, 미디어 믹스 최적화", "수억 원/프로젝트. 일회성.\n비정형 데이터 약함. 감정=설문"),
    ("Bain NPS", "Net Promoter Score", "\"추천하시겠습니까?\" 단일 질문.\nNPS = 성장률의 20~60% 설명.", "\"왜\"인지 모름. 문화 편향.\n조작 가능. 무의식 포착 불가."),
    ("Kantar / Ipsos", "Brand Health Tracking", "인지도, 선호도, 구매의향,\n브랜드 자산 종합. 21,000+ DB.", "연간 수억~수십억 원.\n실시간 불가. 의식적 답변만."),
]
for i, (name, sub, feat, limit) in enumerate(consults):
    x = 0.8 + i * 4.0
    rect(s, Inches(x), Inches(1.8), Inches(3.7), Inches(3.5), R_DARK2)
    box(s, Inches(x+0.15), Inches(1.9), Inches(3.4), Inches(0.35), name, sz=17, color=R_WHITE, bold=True)
    box(s, Inches(x+0.15), Inches(2.25), Inches(3.4), Inches(0.2), sub, sz=10, color=R_GRAY)
    box(s, Inches(x+0.15), Inches(2.6), Inches(3.4), Inches(1), feat, sz=11, color=R_GRAY)
    rect(s, Inches(x+0.1), Inches(3.8), Inches(3.5), Inches(1.2), RGBColor(0x3A, 0x1A, 0x1A))
    box(s, Inches(x+0.2), Inches(3.85), Inches(3.3), Inches(1.1), limit, sz=10, color=R_RED)

# 하단 매트릭스 요약
rect(s, Inches(0.8), Inches(5.6), Inches(11.7), Inches(1.5), RGBColor(0x1A, 0x1A, 0x2E))
box(s, Inches(0.8), Inches(5.7), Inches(11.7), Inches(0.35), "RXR은 이 모든 것의 교차점에 있다", sz=15, color=R_WHITE, bold=True, align=PP_ALIGN.CENTER)

labels = [
    ("리테일테크\n행동 데이터", R_CYAN),
    ("컨설팅\n전략 프레임", R_ORANGE),
    ("리서치\n브랜드 측정", RGBColor(0x8B, 0x5C, 0xF6)),
    ("아무도 못하는\n무의식 심리 = RXR", R_RED),
]
for i, (txt, color) in enumerate(labels):
    x = 1.0 + i * 3.0
    box(s, Inches(x), Inches(6.15), Inches(2.7), Inches(0.8), txt, sz=12, color=color, bold=True, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════
# SLIDE 11: DELIVERABLES
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, R_DARK)
box(s, Inches(0.8), Inches(0.5), Inches(3), Inches(0.3), "HOW WE DELIVER", sz=11, color=R_ORANGE, bold=True)
box(s, Inches(0.8), Inches(0.9), Inches(10), Inches(0.7), "실제로 무엇을 받게 되는가", sz=28, color=R_WHITE, bold=True)

# Daily
rect(s, Inches(0.8), Inches(1.9), Inches(5.8), Inches(5.2), R_DARK2)
top_bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.9), Inches(5.8), Inches(0.06))
top_bar.fill.solid(); top_bar.fill.fore_color.rgb = R_RED; top_bar.line.fill.background()
box(s, Inches(1.0), Inches(2.05), Inches(5.4), Inches(0.35), "Daily AI 인사이트 (매일 자동)", sz=15, color=R_WHITE, bold=True)
rect(s, Inches(1.0), Inches(2.6), Inches(5.4), Inches(1.8), RGBColor(0x22, 0x22, 0x22))
box(s, Inches(1.1), Inches(2.65), Inches(5.2), Inches(0.25), "Alert", sz=11, color=R_RED, bold=True)
box(s, Inches(1.1), Inches(2.95), Inches(5.2), Inches(1.3), "\"유입 +26.5% 증가했지만 구매 전환\n-1.8%p 하락. 리뷰 Clout 점수도 38로\n낮아 체험 중심 방문 비율이 높은 것으로\n보입니다.\"", sz=11, color=R_GRAY)
rect(s, Inches(1.0), Inches(4.6), Inches(5.4), Inches(1.4), RGBColor(0x0A, 0x2D, 0x1C))
box(s, Inches(1.1), Inches(4.65), Inches(5.2), Inches(0.25), "Action", sz=11, color=R_GREEN, bold=True)
box(s, Inches(1.1), Inches(4.95), Inches(5.2), Inches(0.9), "\"피크 시간대(18~20시) 상품 설명형\nPOP 강화 + 한정 프로모션 테스트 필요\"", sz=11, color=R_GRAY)

# Monthly
rect(s, Inches(6.9), Inches(1.9), Inches(5.8), Inches(5.2), R_DARK2)
top_bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.9), Inches(1.9), Inches(5.8), Inches(0.06))
top_bar.fill.solid(); top_bar.fill.fore_color.rgb = R_PINK; top_bar.line.fill.background()
box(s, Inches(7.1), Inches(2.05), Inches(5.4), Inches(0.35), "Monthly 전략 리포트 (프로젝트 종료)", sz=15, color=R_WHITE, bold=True)
rect(s, Inches(7.1), Inches(2.6), Inches(5.4), Inches(1.8), RGBColor(0x22, 0x22, 0x22))
box(s, Inches(7.2), Inches(2.65), Inches(5.2), Inches(0.25), "감정 추이 그��프", sz=11, color=R_PINK, bold=True)
box(s, Inches(7.2), Inches(2.95), Inches(5.2), Inches(1.3), "주차별 Authenticity · Clout · Freshness\nIndex 변화 추이\n\"2주차 Clout 피크(72) → 4주차 45 하락\n= 신선함 소진 감지\"", sz=11, color=R_GRAY)
rect(s, Inches(7.1), Inches(4.6), Inches(5.4), Inches(1.4), RGBColor(0x22, 0x22, 0x22))
box(s, Inches(7.2), Inches(4.65), Inches(5.2), Inches(0.25), "핵심 인사이트", sz=11, color=R_ORANGE, bold=True)
box(s, Inches(7.2), Inches(4.95), Inches(5.2), Inches(0.9), "\"20대 여성 방문 48%이지만 구매 기여 31%.\n30대 고객이 구매 44% — 진짜 매출은 여기.\n다음 프로젝트에서 30대 시간대 운영 강화\"", sz=11, color=R_GRAY)

# ═══════════════════════════════════════
# SLIDE 12: BIZ MODEL
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, R_DARK)
box(s, Inches(0.8), Inches(0.5), Inches(3), Inches(0.3), "BUSINESS MODEL", sz=11, color=R_GREEN, bold=True)
box(s, Inches(0.8), Inches(0.9), Inches(10), Inches(0.7), "어디서 돈이 되는가", sz=32, color=R_WHITE, bold=True)

models = [
    ("프로젝트 분석 리포트", "브랜드 팝업/매장 프로젝트별\n4축 통합 데이터 프로덕트 납품\n\n\"이 매장이 왜 성공했는가\"를\n숫자로 증명하는 보고서"),
    ("SaaS 대시보드", "월간 구독형\n실시간 AI 인사이트 엔진\n\n매일 자동 생성되는\n\"오늘 뭘 해야 하는가\"\n운영 의사결정 지원"),
    ("IR / 제안서 데이터", "브랜드 클라이언트가\n투자자/본사에 보고할 수 있는\n정량 데이터 패키지\n\n메이아이가 못하는\n\"심리적 성과 데��터\""),
]
for i, (title, desc) in enumerate(models):
    x = 0.8 + i * 4.0
    rect(s, Inches(x), Inches(2.0), Inches(3.7), Inches(4.8), R_DARK2)
    top_bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(2.0), Inches(3.7), Inches(0.06))
    top_bar.fill.solid(); top_bar.fill.fore_color.rgb = R_RED; top_bar.line.fill.background()
    box(s, Inches(x+0.2), Inches(2.2), Inches(3.3), Inches(0.4), title, sz=17, color=R_WHITE, bold=True, align=PP_ALIGN.CENTER)
    box(s, Inches(x+0.2), Inches(2.8), Inches(3.3), Inches(3.5), desc, sz=13, color=R_GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════
# SLIDE 13: MARKET
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, R_DARK)
box(s, Inches(0.8), Inches(0.5), Inches(3), Inches(0.3), "MARKET OPPORTUNITY", sz=11, color=R_CYAN, bold=True)
box(s, Inches(0.8), Inches(0.9), Inches(10), Inches(0.7), "왜 지금인가", sz=32, color=R_WHITE, bold=True)

stat(s, Inches(0.8), Inches(2.0), "1.8조+", "국내 팝업스토어\n시장 규모 (2024)")
stat(s, Inches(3.5), Inches(2.0), "300+", "메이아이 공급 매장\n(축①② 수요 검증)")
stat(s, Inches(6.2), Inches(2.0), "95%", "가나초콜릿하우스\n구매전환율")
stat(s, Inches(8.9), Inches(2.0), "0", "4축 통합 + FWA\n경쟁 솔루션 수")

# GAP 시각화
rect(s, Inches(0.8), Inches(3.8), Inches(11.7), Inches(3.3), R_DARK2)
box(s, Inches(0.8), Inches(3.9), Inches(11.7), Inches(0.35), "메이아이가 못 채우는 공백 = RXR의 기회", sz=14, color=R_WHITE, bold=True, align=PP_ALIGN.CENTER)
# bars
max_w = Inches(8)
box(s, Inches(1.0), Inches(4.5), Inches(1.5), Inches(0.25), "메이아이", sz=11, color=R_CYAN, bold=True, align=PP_ALIGN.RIGHT)
bar(s, Inches(2.7), Inches(4.5), 40, max_w, Inches(0.3), R_CYAN, "유입 + 행동")

box(s, Inches(1.0), Inches(5.0), Inches(1.5), Inches(0.25), "GAP", sz=11, color=R_RED, bold=True, align=PP_ALIGN.RIGHT)
# dashed gap
gap_shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.7)+int(max_w*0.4), Inches(5.0), int(max_w*0.6), Inches(0.3))
gap_shape.fill.solid(); gap_shape.fill.fore_color.rgb = RGBColor(0x3A, 0x1A, 0x1A); gap_shape.line.fill.background()
box(s, Inches(2.7)+int(max_w*0.4)+Inches(0.2), Inches(5.02), Inches(3), Inches(0.25), "매출 + 인식 + 심리 = 비어있음", sz=10, color=R_RED)

box(s, Inches(1.0), Inches(5.5), Inches(1.5), Inches(0.25), "RXR", sz=11, color=R_RED, bold=True, align=PP_ALIGN.RIGHT)
bar(s, Inches(2.7), Inches(5.5), 100, max_w, Inches(0.3), R_RED, "전체 커버")

# ═══════════════════════════════════════
# SLIDE 14: CLOSING
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, R_DARK)

# 레드 라인
slide_shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.06))
slide_shape.fill.solid(); slide_shape.fill.fore_color.rgb = R_RED; slide_shape.line.fill.background()

box(s, Inches(1), Inches(1.5), Inches(11), Inches(0.5), "하나의 질문으로 정리하면", sz=16, color=R_GRAY, align=PP_ALIGN.CENTER)

# 기존 vs RXR
rect(s, Inches(1), Inches(2.3), Inches(5), Inches(3.0), R_DARK2)
box(s, Inches(1), Inches(2.4), Inches(5), Inches(0.3), "기존", sz=13, color=R_GRAY, bold=True, align=PP_ALIGN.CENTER)
box(s, Inches(1.2), Inches(2.9), Inches(4.6), Inches(1.8), "\"몇 명 왔어?\"\n\"얼마 팔렸어?\"\n\"좋았대?\"", sz=20, color=R_GRAY, align=PP_ALIGN.CENTER)
box(s, Inches(1.2), Inches(4.5), Inches(4.6), Inches(0.4), "조각난 데이터, 불완전한 답", sz=12, color=R_RED, align=PP_ALIGN.CENTER)

box(s, Inches(6.15), Inches(3.2), Inches(1), Inches(0.5), "VS", sz=24, color=R_GRAY2, bold=True, align=PP_ALIGN.CENTER)

rect(s, Inches(7.3), Inches(2.3), Inches(5), Inches(3.0), R_DARK2)
top_bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.3), Inches(2.3), Inches(5), Inches(0.06))
top_bar.fill.solid(); top_bar.fill.fore_color.rgb = R_RED; top_bar.line.fill.background()
box(s, Inches(7.3), Inches(2.4), Inches(5), Inches(0.3), "RXR", sz=13, color=R_RED, bold=True, align=PP_ALIGN.CENTER)
box(s, Inches(7.5), Inches(2.9), Inches(4.6), Inches(1.8), "\"봤다→들어왔다→머물렀다\n→샀다→공유했다\"\n\"그리고 진심이었다\"", sz=20, color=R_WHITE, align=PP_ALIGN.CENTER)
box(s, Inches(7.5), Inches(4.5), Inches(4.6), Inches(0.4), "완성된 스토리 + 심리적 근거", sz=12, color=R_GREEN, align=PP_ALIGN.CENTER)

# CTA
rect(s, Inches(2.5), Inches(5.7), Inches(8.3), Inches(1.3), R_RED)
box(s, Inches(2.5), Inches(5.8), Inches(8.3), Inches(0.5), "이것은 \"더 좋은 대시보드\"가 아닙니다", sz=15, color=R_WHITE, align=PP_ALIGN.CENTER)
box(s, Inches(2.5), Inches(6.3), Inches(8.3), Inches(0.5), "새로운 카테고리의 데이터 프로덕트입���다", sz=22, color=R_WHITE, bold=True, align=PP_ALIGN.CENTER)

# ─── 저장 ───
out = os.path.join(os.path.dirname(__file__), "rxr-4axis-business-pitch.pptx")
prs.save(out)
print(f"PPT saved: {out}")
