#!/usr/bin/env python3
"""파사드와 쇼윈도우의 비즈니스 심리학 - 강의자료 PPT 생성"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# === 색상 정의 ===
BG_DARK = RGBColor(0x0F, 0x0F, 0x18)
BG_CARD = RGBColor(0x1A, 0x1A, 0x28)
BG_CARD2 = RGBColor(0x22, 0x22, 0x34)
WHITE = RGBColor(0xF0, 0xF0, 0xF8)
TEXT_DIM = RGBColor(0x98, 0x98, 0xB0)
TEXT_MUTED = RGBColor(0x60, 0x60, 0x78)
ACCENT = RGBColor(0x7C, 0x6A, 0xEF)
ACCENT2 = RGBColor(0xA7, 0x8B, 0xFA)
GOLD = RGBColor(0xF0, 0xC0, 0x40)
GREEN = RGBColor(0x34, 0xD3, 0x99)
BLUE = RGBColor(0x60, 0xA5, 0xFA)
ROSE = RGBColor(0xF4, 0x72, 0xB6)
ORANGE = RGBColor(0xFB, 0x92, 0x3C)
CYAN = RGBColor(0x22, 0xD3, 0xEE)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Arial'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_paragraph(text_frame, text, font_size=14, color=WHITE, bold=False, space_before=0, space_after=0, alignment=PP_ALIGN.LEFT):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Arial'
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    p.alignment = alignment
    return p

def add_rounded_rect(slide, left, top, width, height, fill_color=BG_CARD, border_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_data_card(slide, left, top, value, label, color=GOLD, source=""):
    w, h = 2.7, 1.7
    add_rounded_rect(slide, left, top, w, h, BG_CARD, RGBColor(0x2A,0x2A,0x3A))
    add_textbox(slide, left, top+0.25, w, 0.6, value, font_size=32, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left, top+0.85, w, 0.5, label, font_size=11, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)
    if source:
        add_textbox(slide, left, top+1.3, w, 0.3, source, font_size=9, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 1: COVER
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 1.5, 0.6, 5, 0.4, "LECTURE MATERIAL", font_size=11, color=ACCENT, bold=True)
add_textbox(slide, 1.5, 1.2, 10, 1.2, "건물의 첫인상", font_size=52, color=WHITE, bold=True)
add_textbox(slide, 1.5, 2.5, 10, 0.8, "파사드와 쇼윈도우의 비즈니스 심리학", font_size=28, color=ACCENT2, bold=True)

tb = add_textbox(slide, 1.5, 3.6, 8, 1.0, "", font_size=16, color=TEXT_DIM)
tf = tb.text_frame
tf.paragraphs[0].text = '"예쁘니까"가 아니라 "돈이 되니까".'
tf.paragraphs[0].font.size = Pt(16)
tf.paragraphs[0].font.color.rgb = TEXT_DIM
add_paragraph(tf, "브랜드와 건물이 외벽에 투자해야 하는 정량적 이유와 심리학적 메커니즘.", 16, TEXT_DIM)

add_data_card(slide, 1.5, 5.0, "24%", "구매 결정이\n파사드에서 시작", GOLD)
add_data_card(slide, 4.5, 5.0, "+50%", "파사드 리노베이션 후\n매출 증가", GREEN)
add_data_card(slide, 7.5, 5.0, "$40M", "Glossier 매장 디자인\n획득 미디어 가치", ROSE)
add_data_card(slide, 10.5, 5.0, "+12.4%", "차별화 파사드\n임대료 프리미엄", BLUE)

# ============================================================
# SLIDE 2: 목차
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.8, 0.4, 3, 0.3, "AGENDA", font_size=11, color=ACCENT, bold=True)
add_textbox(slide, 0.8, 0.8, 8, 0.7, "오늘의 강의 구조", font_size=36, color=WHITE, bold=True)
add_textbox(slide, 0.8, 1.5, 10, 0.4, "왜 중요한가 → 얼마나 효과적인가 → 왜 작동하는가 → 누가 증명했는가 → 어떻게 진화했는가", font_size=14, color=TEXT_DIM)

agenda_items = [
    ("1", "왜 파사드인가 — 첫인상의 비즈니스 가치", "브랜드와 건물에 있어 '보여지는 것'이 매출과 직결되는 이유"),
    ("2", "숫자로 증명하는 파사드의 힘", "보행 전환, 매출, 미디어 가치, 자산 프리미엄 — 40+ 데이터"),
    ("3", "심리학이 설명하는 메커니즘", "10가지 행동경제학/인지심리학 이론으로 '왜 작동하는가'를 해부"),
    ("4", "케이스 스터디 — 돈으로 증명한 브랜드들", "Apple, Hermes, Glossier, Coex 파도 — 구체적 숫자와 전략"),
    ("5", "파사드의 진화 — 벽에서 미디어로", "7단계 진화: 기능 → 장식 → 스토리 → 건축 → 디지털 → AI"),
    ("6", "Takeaway — 파사드는 비용이 아니라 자산이다", "전통 광고와의 비교, ROI 공식, 6가지 핵심 인사이트"),
]

for i, (num, title, desc) in enumerate(agenda_items):
    y = 2.2 + i * 0.82
    circle = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(y), Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = BG_CARD
    circle.line.color.rgb = RGBColor(0x2A,0x2A,0x3A)
    circle.line.width = Pt(1)
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = ACCENT2
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.word_wrap = False

    add_textbox(slide, 1.5, y-0.02, 9, 0.3, title, font_size=16, color=WHITE, bold=True)
    add_textbox(slide, 1.5, y+0.32, 9, 0.3, desc, font_size=12, color=TEXT_DIM)


# ============================================================
# SLIDE 3: Chapter 1 - 왜 파사드인가
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.8, 0.4, 3, 0.3, "CHAPTER 1", font_size=11, color=ACCENT, bold=True)
add_textbox(slide, 0.8, 0.8, 10, 0.7, "왜 파사드인가", font_size=40, color=WHITE, bold=True)
add_textbox(slide, 0.8, 1.5, 10, 0.4, "첫인상의 비즈니스 가치", font_size=18, color=TEXT_DIM)

# Question box
q_shape = add_rounded_rect(slide, 0.8, 2.2, 11.5, 1.2, RGBColor(0x15,0x13,0x2A), ACCENT)
q_label = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(2.0), Inches(0.45), Inches(0.45))
q_label.fill.solid()
q_label.fill.fore_color.rgb = ACCENT
q_label.line.fill.background()
q_tf = q_label.text_frame
q_tf.paragraphs[0].text = "Q"
q_tf.paragraphs[0].font.size = Pt(14)
q_tf.paragraphs[0].font.color.rgb = WHITE
q_tf.paragraphs[0].font.bold = True
q_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

add_textbox(slide, 1.2, 2.4, 10.5, 0.8, "건물의 파사드가 단순히 '예쁘면 좋은 것'이 아니라\n비즈니스에 실질적 영향을 미친다고 말할 수 있는 근거는?", font_size=18, color=WHITE, bold=True)

# Answer
a_shape = add_rounded_rect(slide, 0.8, 3.7, 11.5, 3.3, BG_CARD, RGBColor(0x2A,0x2A,0x3A))
add_textbox(slide, 1.1, 3.85, 3, 0.3, "핵심 전제", font_size=11, color=ACCENT2, bold=True)

tb = add_textbox(slide, 1.1, 4.2, 10.8, 2.5, "", font_size=14, color=TEXT_DIM)
tf = tb.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "사람이 건물 앞을 지나가는 순간, 뇌는 0.05초 만에 첫인상을 형성합니다."
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.color.rgb = TEXT_DIM

lines = [
    "이 첫인상은 단순한 감상이 아닙니다.",
    '"들어갈 것인가, 지나칠 것인가"를 결정하는 무의식적 판단이며,',
    "이 판단이 곧 매출, 임대료, 브랜드 자산으로 직결됩니다.",
    "",
    "• 고객이 매장에 들어오기 전에 구매 의향이 결정됩니다.",
    "• 건물의 외관이 내부 상품의 가격 정당성을 만듭니다.",
    "• 파사드 하나가 연간 수십억 원의 무료 미디어를 생성합니다.",
    "• 차별화된 파사드는 같은 위치에서 임대료를 12%까지 높입니다.",
]
for line in lines:
    p = add_paragraph(tf, line, 14, TEXT_DIM if not line.startswith("•") else WHITE, line.startswith("•"), space_before=2)


# ============================================================
# SLIDE 4: 3가지 비즈니스 변수
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.8, 0.4, 3, 0.3, "CHAPTER 1", font_size=11, color=ACCENT, bold=True)
add_textbox(slide, 0.8, 0.8, 10, 0.6, "파사드가 결정하는 3가지 비즈니스 변수", font_size=30, color=WHITE, bold=True)

cards_data = [
    ("유입", "보행자가 멈추는가,\n지나치는가.", "쇼윈도 → 충동구매 62%", "shopPOPdisplays", GOLD),
    ("가격 인식", '"비싸도 괜찮다"는\n느낌이 드는가.', "환경 프라이밍 → WTP +25%", "Tversky & Kahneman", GREEN),
    ("자산 가치", "건물 자체의 가치가\n올라가는가.", "곡면 파사드 → 임대료 +12.4%", "Tandfonline", BLUE),
]

for i, (title, desc, evidence, source, color) in enumerate(cards_data):
    x = 0.8 + i * 4.1
    add_rounded_rect(slide, x, 1.7, 3.7, 3.5, BG_CARD, RGBColor(0x2A,0x2A,0x3A))
    add_textbox(slide, x+0.3, 1.9, 3.1, 0.5, title, font_size=24, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x+0.3, 2.5, 3.1, 0.8, desc, font_size=13, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x+0.3, 3.5, 3.1, 0.4, evidence, font_size=13, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x+0.3, 3.95, 3.1, 0.3, source, font_size=10, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

# Callout
add_rounded_rect(slide, 0.8, 5.5, 11.5, 1.5, RGBColor(0x1C,0x1A,0x10), RGBColor(0x3A,0x30,0x10))
add_textbox(slide, 1.1, 5.65, 3, 0.3, "강의 포인트", font_size=12, color=GOLD, bold=True)
tb = add_textbox(slide, 1.1, 6.0, 10.8, 0.9, "", font_size=14, color=WHITE)
tf = tb.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = '파사드의 가치는 "예쁘다/안 예쁘다"의 주관적 판단이 아닙니다.'
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.color.rgb = WHITE
add_paragraph(tf, '"유입률, 객단가, 임대료"라는 정량적 비즈니스 지표로 측정할 수 있습니다.', 14, WHITE)


# ============================================================
# SLIDE 5: Chapter 2 - 숫자로 증명 (유입 & 전환)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.8, 0.4, 3, 0.3, "CHAPTER 2", font_size=11, color=ACCENT, bold=True)
add_textbox(slide, 0.8, 0.8, 10, 0.6, "숫자로 증명하는 파사드의 힘", font_size=36, color=WHITE, bold=True)

add_textbox(slide, 0.8, 1.6, 5, 0.4, "● 유입과 전환", font_size=20, color=GOLD, bold=True)

cards = [
    ("76%", "구매 결정이\n매장 안에서 발생", "POPAI", GOLD),
    ("24%", "쇼윈도가\n구매 결정에 영향", "NPD Group", GOLD),
    ("62%", "쇼윈도에 의한\n충동구매 비율", "shopPOPdisplays", GOLD),
    ("+35%", "파사드 리노베이션 후\n유동인구 증가", "PRL Glass", GOLD),
]
for i, (v, l, s, c) in enumerate(cards):
    add_data_card(slide, 0.8 + i*3.1, 2.2, v, l, c, s)

add_textbox(slide, 0.8, 4.2, 5, 0.4, "● 매출 임팩트", font_size=20, color=GREEN, bold=True)

cards2 = [
    ("+50%", "파사드 리노베이션 후\n매출 증가", "PRL Glass (6개월 ROI)", GREEN),
    ("+44%", "리모델링 후\n신규 고객 매출", "Monash University", GREEN),
    ("+32%", "디지털 사이니지\n도입 매장 매출", "AIScreen (2025)", GREEN),
    ("300%", "비주얼 머천다이징\n투자 ROI", "Contra Vision", GREEN),
]
for i, (v, l, s, c) in enumerate(cards2):
    add_data_card(slide, 0.8 + i*3.1, 4.8, v, l, c, s)


# ============================================================
# SLIDE 6: Chapter 2 - 미디어 & 자산 + 해석
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.8, 0.4, 3, 0.3, "CHAPTER 2", font_size=11, color=ACCENT, bold=True)

add_textbox(slide, 0.8, 0.8, 5, 0.4, "● 미디어 가치 & 자산 프리미엄", font_size=20, color=ROSE, bold=True)

cards3 = [
    ("$40M", "Glossier 매장 디자인\nEarned Media", "520억원 상당", ROSE),
    ("1억+", "Coex 파도 일루전\n소셜 미디어 조회", "설치 수주 내", ROSE),
    ("+12.4%", "곡면 파사드 건물\n임대료 프리미엄", "Tandfonline", ORANGE),
    ("2~3년", "파사드 리모델링\n투자 회수 기간", "Monash Univ.", ORANGE),
]
for i, (v, l, s, c) in enumerate(cards3):
    add_data_card(slide, 0.8 + i*3.1, 1.4, v, l, c, s)

# Monash 해석
add_rounded_rect(slide, 0.8, 3.5, 11.5, 1.6, RGBColor(0x0F, 0x1A, 0x15), RGBColor(0x15, 0x3A, 0x28))
add_textbox(slide, 1.1, 3.6, 3, 0.3, "핵심 해석 — Monash University 연구", font_size=12, color=GREEN, bold=True)
tb = add_textbox(slide, 1.1, 3.95, 10.8, 1.0, "", font_size=13, color=WHITE)
tf = tb.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "신규 고객 매출 +44% vs 기존 고객 매출 +10%"
tf.paragraphs[0].font.size = Pt(13)
tf.paragraphs[0].font.color.rgb = WHITE
tf.paragraphs[0].font.bold = True
add_paragraph(tf, '→ 파사드 리모델링은 "기존 고객을 더 쓰게 하는 것"이 아니라 "새로운 고객을 끌어오는 것"', 13, TEXT_DIM, space_before=4)
add_paragraph(tf, '→ 리모델링 후 신규 고객 재방문율도 +16% (1년간 지속). 파사드 = "고객 획득 채널"', 13, TEXT_DIM, space_before=2)

# Glossier 해석
add_rounded_rect(slide, 0.8, 5.3, 11.5, 1.7, RGBColor(0x1A, 0x10, 0x18), RGBColor(0x3A, 0x18, 0x30))
add_textbox(slide, 1.1, 5.4, 3, 0.3, "핵심 해석 — Glossier", font_size=12, color=ROSE, bold=True)
tb = add_textbox(slide, 1.1, 5.75, 10.8, 1.0, "", font_size=13, color=WHITE)
tf = tb.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "광고비 0원. 대신 매장 자체를 '찍고 싶은 공간'으로 설계."
tf.paragraphs[0].font.size = Pt(13)
tf.paragraphs[0].font.color.rgb = WHITE
tf.paragraphs[0].font.bold = True
add_paragraph(tf, "고객이 자발적으로 280만 건의 인스타그램 포스팅 → $40M(520억) 미디어 가치.", 13, TEXT_DIM, space_before=4)
add_paragraph(tf, "파사드/매장 디자인이 곧 마케팅 비용을 대체. 현재 84% 리테일러가 이 모델을 따름.", 13, TEXT_DIM, space_before=2)


# ============================================================
# SLIDE 7: Chapter 3 - 심리학 도입
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.8, 0.4, 3, 0.3, "CHAPTER 3", font_size=11, color=ACCENT, bold=True)
add_textbox(slide, 0.8, 0.8, 10, 0.7, "심리학이 설명하는 메커니즘", font_size=36, color=WHITE, bold=True)
add_textbox(slide, 0.8, 1.5, 10, 0.5, '파사드가 효과적인 것은 "느낌" 때문이 아닙니다.\n인간의 인지 구조가 그렇게 설계되어 있기 때문입니다.', font_size=15, color=TEXT_DIM)

# Question
q_shape = add_rounded_rect(slide, 0.8, 2.4, 11.5, 1.0, RGBColor(0x15,0x13,0x2A), ACCENT)
add_textbox(slide, 1.2, 2.55, 10.5, 0.7, "왜 사람들은 에르메스 매장 앞에서 멈추고, 옆 건물은 그냥 지나칠까?", font_size=20, color=WHITE, bold=True)

# 6 theories summary grid
theories = [
    ("01", "단순노출효과", "반복 노출 → 무의식적 호감", "Zajonc, 1968", GOLD),
    ("02", "앵커링 효과", "첫인상이 가격 기준점 설정", "Kahneman, 1974", BLUE),
    ("03", "신호 이론", "비용 높은 신호 = 신뢰", "Spence, 1973", GREEN),
    ("04", "선택적 주의", "3D 파사드가 주의 필터 통과", "Broadbent, 1958", ORANGE),
    ("05", "미적-사용성", "아름다운 것 = 더 좋은 것", "Kurosu, 1995", ROSE),
    ("06", "Kaplan 매트릭스", "일관성×복잡성×신비감 균형", "Kaplan & Kaplan", CYAN),
]

for i, (num, name, desc, scholar, color) in enumerate(theories):
    col = i % 3
    row = i // 3
    x = 0.8 + col * 4.0
    y = 3.7 + row * 1.65
    add_rounded_rect(slide, x, y, 3.7, 1.4, BG_CARD, RGBColor(0x2A,0x2A,0x3A))
    add_textbox(slide, x+0.2, y+0.15, 0.5, 0.3, num, font_size=14, color=color, bold=True)
    add_textbox(slide, x+0.65, y+0.15, 2.8, 0.3, name, font_size=16, color=WHITE, bold=True)
    add_textbox(slide, x+0.2, y+0.55, 3.2, 0.4, desc, font_size=12, color=TEXT_DIM)
    add_textbox(slide, x+0.2, y+0.95, 3.2, 0.3, scholar, font_size=10, color=TEXT_MUTED)


# ============================================================
# SLIDE 8: 이론 상세 1 — 단순노출 + 앵커링
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.8, 0.4, 5, 0.3, "CHAPTER 3 — 핵심 이론 상세", font_size=11, color=ACCENT, bold=True)

# Theory 1
add_rounded_rect(slide, 0.8, 0.9, 11.5, 3.0, BG_CARD, RGBColor(0x2A,0x2A,0x3A))
num_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(0.75), Inches(0.5), Inches(0.4))
num_shape.fill.solid()
num_shape.fill.fore_color.rgb = ACCENT
num_shape.line.fill.background()
ntf = num_shape.text_frame
ntf.paragraphs[0].text = "1"
ntf.paragraphs[0].font.size = Pt(14)
ntf.paragraphs[0].font.color.rgb = WHITE
ntf.paragraphs[0].font.bold = True
ntf.paragraphs[0].alignment = PP_ALIGN.CENTER

add_textbox(slide, 1.1, 1.1, 5, 0.4, "단순노출효과  Mere Exposure Effect", font_size=18, color=WHITE, bold=True)
add_textbox(slide, 1.1, 1.5, 5, 0.3, "Robert Zajonc, 1968", font_size=11, color=TEXT_MUTED)

# Left column - What & Why
add_textbox(slide, 1.1, 1.9, 3, 0.25, "무엇인가", font_size=10, color=GOLD, bold=True)
add_textbox(slide, 1.1, 2.15, 4.8, 0.5, "어떤 대상에 반복적으로 노출되는 것만으로 호감이 증가.\n의식적 주의 불필요. 노출 사실을 기억 못 해도 효과 발생.", font_size=12, color=TEXT_DIM)

add_textbox(slide, 1.1, 2.7, 3, 0.25, "왜 파사드에 중요한가", font_size=10, color=GREEN, bold=True)
add_textbox(slide, 1.1, 2.95, 4.8, 0.6, "TV광고 15초, 디지털 광고는 스킵. 파사드는 수명 10~30년.\n매일 출퇴근 보행자에게 연간 수백 번 무의식적 노출.\n파사드는 세상에서 가장 긴 광고.", font_size=12, color=TEXT_DIM)

# Right column - Evidence
add_rounded_rect(slide, 6.5, 1.9, 5.4, 1.7, BG_CARD2)
add_textbox(slide, 6.7, 2.0, 1.5, 0.6, "60%", font_size=36, color=GOLD, bold=True)
add_textbox(slide, 8.2, 2.1, 3.5, 0.8, "단 1~5회 노출만으로\n선호도가 60% 증가합니다.\n파사드 앞을 5번 지나가면,\n이미 그 브랜드를 좋아하기 시작한 것.", font_size=12, color=TEXT_DIM)

# Theory 2
add_rounded_rect(slide, 0.8, 4.2, 11.5, 3.0, BG_CARD, RGBColor(0x2A,0x2A,0x3A))
num_shape2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(4.05), Inches(0.5), Inches(0.4))
num_shape2.fill.solid()
num_shape2.fill.fore_color.rgb = ACCENT
num_shape2.line.fill.background()
ntf2 = num_shape2.text_frame
ntf2.paragraphs[0].text = "2"
ntf2.paragraphs[0].font.size = Pt(14)
ntf2.paragraphs[0].font.color.rgb = WHITE
ntf2.paragraphs[0].font.bold = True
ntf2.paragraphs[0].alignment = PP_ALIGN.CENTER

add_textbox(slide, 1.1, 4.4, 5, 0.4, "앵커링 효과  Anchoring Effect", font_size=18, color=WHITE, bold=True)
add_textbox(slide, 1.1, 4.8, 5, 0.3, "Tversky & Kahneman, 1974", font_size=11, color=TEXT_MUTED)

add_textbox(slide, 1.1, 5.15, 3, 0.25, "무엇인가", font_size=10, color=GOLD, bold=True)
add_textbox(slide, 1.1, 5.4, 4.8, 0.4, "처음 접한 정보를 기준점(앵커)으로 삼아 이후 판단을 조정.\n첫 번째 인상이 모든 후속 판단의 '정상 범위'를 결정.", font_size=12, color=TEXT_DIM)

add_textbox(slide, 1.1, 5.85, 3, 0.25, "왜 파사드에 중요한가", font_size=10, color=GREEN, bold=True)
add_textbox(slide, 1.1, 6.1, 4.8, 0.7, "에르메스 긴자 메종의 유리 블록 파사드를 본 고객은\n문을 열기 전에 이미 '미술관급 브랜드'라는 앵커가 설정.\n이 앵커 위에서 버킨백 가격이 '당연한 것'이 됨.\n파사드 = 가격 정당화 장치.", font_size=12, color=TEXT_DIM)

add_rounded_rect(slide, 6.5, 5.15, 5.4, 1.7, BG_CARD2)
add_textbox(slide, 6.7, 5.25, 1.5, 0.6, "+25%", font_size=36, color=BLUE, bold=True)
add_textbox(slide, 8.2, 5.35, 3.5, 0.8, "환경적 프라이밍(파사드 포함)이\n지불의향(WTP)을 최대 25% 변화.\n같은 상품도 환경에 따라\n'비싸다' 또는 '합리적이다'가 달라짐.", font_size=12, color=TEXT_DIM)


# ============================================================
# SLIDE 9: 이론 상세 2 — 신호이론 + 선택적 주의
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.8, 0.4, 5, 0.3, "CHAPTER 3 — 핵심 이론 상세", font_size=11, color=ACCENT, bold=True)

# Theory 3
add_rounded_rect(slide, 0.8, 0.9, 11.5, 3.0, BG_CARD, RGBColor(0x2A,0x2A,0x3A))
num3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(0.75), Inches(0.5), Inches(0.4))
num3.fill.solid(); num3.fill.fore_color.rgb = ACCENT; num3.line.fill.background()
num3.text_frame.paragraphs[0].text = "3"
num3.text_frame.paragraphs[0].font.size = Pt(14); num3.text_frame.paragraphs[0].font.color.rgb = WHITE
num3.text_frame.paragraphs[0].font.bold = True; num3.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

add_textbox(slide, 1.1, 1.1, 5, 0.4, "신호 이론  Costly Signaling", font_size=18, color=WHITE, bold=True)
add_textbox(slide, 1.1, 1.5, 5, 0.3, "Michael Spence, 1973 / Zahavi 핸디캡 원리", font_size=11, color=TEXT_MUTED)

add_textbox(slide, 1.1, 1.9, 3, 0.25, "무엇인가", font_size=10, color=GOLD, bold=True)
add_textbox(slide, 1.1, 2.15, 4.8, 0.4, "비용이 많이 드는 신호일수록 신뢰도가 높음.\n가짜로 만들 수 없는 신호만이 진짜 정보를 전달.", font_size=12, color=TEXT_DIM)

add_textbox(slide, 1.1, 2.6, 3, 0.25, "왜 파사드에 중요한가", font_size=10, color=GREEN, bold=True)
add_textbox(slide, 1.1, 2.85, 4.8, 0.7, "인스타 광고 100만원이면 누구나 가능.\n건물 파사드에 수십억 투자는 위조 불가능.\nBergdorf 홀리데이 윈도우: 100명 장인 투입.\n에르메스 윈도우 1개: $500K(6.5억원).\n= '위조 불가능한 품질 신호'", font_size=12, color=TEXT_DIM)

add_rounded_rect(slide, 6.5, 1.9, 5.4, 1.7, BG_CARD2)
add_textbox(slide, 6.7, 2.0, 5, 0.3, "신호의 신뢰도 공식", font_size=12, color=GREEN, bold=True)
add_textbox(slide, 6.7, 2.4, 5, 0.5, "신뢰도 = 신호 비용  ÷  위조 가능성", font_size=20, color=WHITE, bold=True)
add_textbox(slide, 6.7, 2.95, 5, 0.5, "파사드: 비용 높음 + 위조 불가\n= 가장 신뢰도 높은 브랜드 신호", font_size=13, color=TEXT_DIM)

# Theory 4
add_rounded_rect(slide, 0.8, 4.2, 11.5, 3.0, BG_CARD, RGBColor(0x2A,0x2A,0x3A))
num4 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(4.05), Inches(0.5), Inches(0.4))
num4.fill.solid(); num4.fill.fore_color.rgb = ACCENT; num4.line.fill.background()
num4.text_frame.paragraphs[0].text = "4"
num4.text_frame.paragraphs[0].font.size = Pt(14); num4.text_frame.paragraphs[0].font.color.rgb = WHITE
num4.text_frame.paragraphs[0].font.bold = True; num4.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

add_textbox(slide, 1.1, 4.4, 5, 0.4, "선택적 주의 & 배너 맹시", font_size=18, color=WHITE, bold=True)
add_textbox(slide, 1.1, 4.8, 5, 0.3, "Broadbent (1958), Simons & Chabris (1999)", font_size=11, color=TEXT_MUTED)

add_textbox(slide, 1.1, 5.15, 3, 0.25, "핵심 역설", font_size=10, color=ORANGE, bold=True)
add_textbox(slide, 1.1, 5.4, 4.8, 0.7, "인간은 하루 4,000~10,000개 광고에 노출되지만\n대부분 무의식적으로 걸러냄 = 배너 맹시.\n디지털 배너 CTR = 0.05%\n\n하지만 파사드는 3차원 물리적 현저성으로\n주의 필터를 강제 통과. 보행자 정지율 20~40%.", font_size=12, color=TEXT_DIM)

add_rounded_rect(slide, 6.5, 5.15, 5.4, 1.7, BG_CARD2)
add_textbox(slide, 6.7, 5.25, 1.5, 0.6, "800x", font_size=36, color=ORANGE, bold=True)
add_textbox(slide, 8.2, 5.35, 3.5, 0.8, "파사드 정지율(20~40%)\nvs 배너 CTR(0.05%)\n\n파사드의 주의 포착력은\n디지털 광고의 최대 800배", font_size=12, color=TEXT_DIM)


# ============================================================
# SLIDE 10: 심리학 플로우
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.8, 0.4, 5, 0.3, "CHAPTER 3 — 메커니즘 요약", font_size=11, color=ACCENT, bold=True)
add_textbox(slide, 0.8, 0.8, 10, 0.6, "파사드가 작동하는 9단계 심리학적 플로우", font_size=28, color=WHITE, bold=True)

flow_steps = [
    ("1", "반복 노출", "단순노출 효과 → 무의식적 브랜드 선호", GOLD),
    ("2", "시각적 충격", "주의 경제학 + Kaplan 신비감 → 주의 획득", ORANGE),
    ("3", "심미적 판단", "미적-사용성 + 후광 → '좋은 브랜드' 기대", ROSE),
    ("4", "신뢰 형성", "신호 이론 → 건축 투자 = 품질 증명", GREEN),
    ("5", "진입 유인", "전망-피난처 + 가독성 → 본능적 안전감", CYAN),
    ("6", "가격 수용", "앵커링 → 프리미엄 가격 수용 범위 확장", BLUE),
    ("7", "체류 연장", "분위기 이론(PAD) → 체류시간+지출 증가", GOLD),
    ("8", "기억 각인", "피크-엔드 → 전체 경험의 기억 품질 결정", ACCENT2),
    ("9", "자발적 확산", "Social Proof → SNS 공유 → 자기강화 루프", ROSE),
]

for i, (num, trigger, desc, color) in enumerate(flow_steps):
    y = 1.6 + i * 0.62
    # Dot
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), Inches(y+0.08), Inches(0.18), Inches(0.18))
    dot.fill.solid(); dot.fill.fore_color.rgb = color; dot.line.fill.background()
    # Line
    if i < len(flow_steps) - 1:
        line_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.07), Inches(y+0.28), Inches(0.04), Inches(0.4))
        line_shape.fill.solid(); line_shape.fill.fore_color.rgb = RGBColor(0x2A,0x2A,0x3A); line_shape.line.fill.background()
    # Card
    add_rounded_rect(slide, 1.5, y-0.05, 10.5, 0.52, BG_CARD, RGBColor(0x2A,0x2A,0x3A))
    add_textbox(slide, 1.7, y, 1.5, 0.3, trigger, font_size=12, color=color, bold=True)
    add_textbox(slide, 3.3, y, 8.5, 0.3, desc, font_size=13, color=TEXT_DIM)


# ============================================================
# SLIDE 11: Case 1 — Apple
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.8, 0.4, 3, 0.3, "CHAPTER 4 — CASE 01", font_size=11, color=GREEN, bold=True)
add_textbox(slide, 0.8, 0.8, 10, 0.5, "Apple Store — 유리 한 장이 만든 세계 최고 매출", font_size=28, color=WHITE, bold=True)
add_textbox(slide, 0.8, 1.35, 10, 0.4, "투명성이라는 브랜드 철학을 건축으로 번역한 전략", font_size=14, color=TEXT_DIM)

# Data row
apple_data = [
    ("$7M", "초기 투자", GREEN), ("$15M+", "2019 리빌드", GREEN),
    ("$440M+", "5th Ave 연매출", GREEN), ("$5,546", "sqft당 매출 (1위)", GREEN),
    ("+15%", "인접 소매점 유입↑", BLUE),
]
for i, (v, l, c) in enumerate(apple_data):
    add_data_card(slide, 0.5 + i*2.55, 1.9, v, l, c)

# Narrative
add_rounded_rect(slide, 0.8, 3.9, 11.5, 3.2, BG_CARD, RGBColor(0x2A,0x2A,0x3A))
add_textbox(slide, 1.1, 4.0, 3, 0.25, "전략 해부", font_size=11, color=GREEN, bold=True)

tb = add_textbox(slide, 1.1, 4.3, 10.8, 2.5, "", font_size=13, color=TEXT_DIM)
tf = tb.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "Apple Store의 유리 파사드는 디자인 선택이 아니라 '우리는 숨길 것이 없다'는 철학의 물리적 선언."
tf.paragraphs[0].font.size = Pt(13); tf.paragraphs[0].font.color.rgb = TEXT_DIM
add_paragraph(tf, "", 6, TEXT_DIM)
add_paragraph(tf, "전망-피난처 이론: 유리 파사드 → 안을 들여다보고 싶은 본능(전망) 충족 + '안전하다'는 느낌(피난처) 충족 → 진입 장벽 극도로 낮아짐", 13, TEXT_DIM, space_before=4)
add_paragraph(tf, "", 6, TEXT_DIM)
add_paragraph(tf, "신호 이론: 유리 건축은 매우 비싸고 기술적으로 어려움 → '이 정도 기술력과 투자 여력'이라는 위조 불가능한 신호", 13, TEXT_DIM, space_before=4)
add_paragraph(tf, "", 6, TEXT_DIM)
add_paragraph(tf, "결과: $15M 투자 → 연 $440M 매출. ROI 약 2,900%. 건물 수명 동안 매년 반복.", 13, WHITE, True, space_before=4)


# ============================================================
# SLIDE 12: Case 2 — Hermes
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.8, 0.4, 5, 0.3, "CHAPTER 4 — CASE 02", font_size=11, color=GOLD, bold=True)
add_textbox(slide, 0.8, 0.8, 10, 0.5, "Hermes — 쇼윈도우 하나에 5억 원을 쓰는 이유", font_size=28, color=WHITE, bold=True)
add_textbox(slide, 0.8, 1.35, 10, 0.4, "광고 없이, 앰배서더 없이. 공간과 상품 자체로만 커뮤니케이션하는 전략.", font_size=14, color=TEXT_DIM)

h_data = [
    ("$500K+", "윈도우 1개 비용", GOLD), ("+30%", "보행 유입 프리미엄", GOLD),
    ("$93.7B", "브랜드 가치 (2025)", GOLD), ("0원", "광고비 / 앰배서더", GOLD),
]
for i, (v, l, c) in enumerate(h_data):
    add_data_card(slide, 0.8 + i*3.1, 1.9, v, l, c)

add_rounded_rect(slide, 0.8, 3.9, 11.5, 3.2, BG_CARD, RGBColor(0x2A,0x2A,0x3A))
add_textbox(slide, 1.1, 4.0, 3, 0.25, "전략 해부", font_size=11, color=GOLD, bold=True)

tb = add_textbox(slide, 1.1, 4.3, 10.8, 2.5, "", font_size=13, color=TEXT_DIM)
tf = tb.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = "에르메스 전략의 핵심은 'Counter-signaling'(역신호). 대부분의 럭셔리가 셀럽을 쓸 때, 에르메스는 상품과 공간 자체만으로 소통."
tf.paragraphs[0].font.size = Pt(13); tf.paragraphs[0].font.color.rgb = TEXT_DIM
add_paragraph(tf, "", 6, TEXT_DIM)
add_paragraph(tf, "신호 이론의 극단적 적용: 쇼윈도우 하나에 $500K(6.5억)을 투자하는 것은 '광고비'가 아니라 '우리 브랜드의 디테일과 장인정신'이라는 신뢰 신호.", 13, TEXT_DIM, space_before=4)
add_paragraph(tf, "", 6, TEXT_DIM)
add_paragraph(tf, "앵커링: 긴자 메종의 유리 블록 파사드(렌조 피아노, 2001)는 밤에 등불처럼 빛남. 고객이 문을 열기 전에 '미술관에 들어가는' 기대 → 버킨백 가격이 '당연한 것'이 됨.", 13, TEXT_DIM, space_before=4)
add_paragraph(tf, "", 6, TEXT_DIM)
add_paragraph(tf, "Vogue, AD 등 글로벌 미디어가 자발적으로 보도. 광고비 0원에 수십억 원의 미디어 가치.", 13, WHITE, True, space_before=4)


# ============================================================
# SLIDE 13: Case 3 — Glossier + Coex
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.8, 0.4, 5, 0.3, "CHAPTER 4 — CASE 03 & 04", font_size=11, color=ROSE, bold=True)

# Glossier
add_rounded_rect(slide, 0.8, 0.9, 5.5, 3.0, BG_CARD, RGBColor(0x2A,0x2A,0x3A))
add_textbox(slide, 1.1, 1.0, 4, 0.3, "GLOSSIER", font_size=11, color=ROSE, bold=True)
add_textbox(slide, 1.1, 1.3, 5, 0.4, "광고비 0원, 520억의 미디어", font_size=18, color=WHITE, bold=True)
add_textbox(slide, 1.1, 1.8, 4.8, 0.3, "280만+ 인스타그램 해시태그", font_size=24, color=ROSE, bold=True)
add_textbox(slide, 1.1, 2.2, 4.8, 0.3, "Earned Media Value  $40M (520억원)", font_size=14, color=WHITE, bold=True)
add_textbox(slide, 1.1, 2.6, 4.8, 0.8, "매장 자체를 '찍고 싶은 공간'으로 설계.\n고객이 자발적으로 콘텐츠를 생성하는 구조.\n전통 광고비 $40M은 수개월 소진.\n매장 디자인은 매년 가치가 재생산.", font_size=12, color=TEXT_DIM)

# Coex
add_rounded_rect(slide, 6.8, 0.9, 5.5, 3.0, BG_CARD, RGBColor(0x2A,0x2A,0x3A))
add_textbox(slide, 7.1, 1.0, 4, 0.3, "COEX SM TOWN 파도 일루전", font_size=11, color=CYAN, bold=True)
add_textbox(slide, 7.1, 1.3, 5, 0.4, "건물 벽 하나로 1억 뷰", font_size=18, color=WHITE, bold=True)
add_textbox(slide, 7.1, 1.8, 4.8, 0.3, "1억+ 소셜 미디어 조회", font_size=24, color=CYAN, bold=True)
add_textbox(slide, 7.1, 2.2, 4.8, 0.3, "설치 후 건물 임대 프리미엄 상승", font_size=14, color=WHITE, bold=True)
add_textbox(slide, 7.1, 2.6, 4.8, 0.8, "Kaplan '신비감(Mystery)' 극대화.\n익숙한 도시에서 예상 밖의 시각 경험.\n3D 빌보드 시청 시간: 일반 광고의 3~4배.\n68% 소비자: '프리미엄' 브랜드 인식.", font_size=12, color=TEXT_DIM)

# Korean studies
add_rounded_rect(slide, 0.8, 4.2, 11.5, 3.0, BG_CARD, RGBColor(0x2A,0x2A,0x3A))
add_textbox(slide, 1.1, 4.3, 5, 0.3, "한국 학술 연구 — 정량적 검증", font_size=14, color=WHITE, bold=True)

kr_studies = [
    ("쇼윈도 디스플레이와 구매의도 (KCI)", "컨셉/스토리 중심 > 시즌/코디네이션 중심 → 매장 진입+긍정 태도 형성에 더 효과적"),
    ("디지털 vs AI 미디어 파사드 (KCI)", "AI 미디어 파사드가 몰입성, 인지적 요소에서 유의미하게 높은 효과"),
    ("백화점 쇼윈도 현상학적 특성", "예술적 오브제 활용이 브랜드의 심리적 가치 향상에 기여"),
]
for i, (title, finding) in enumerate(kr_studies):
    y = 4.8 + i * 0.75
    add_textbox(slide, 1.1, y, 4, 0.3, title, font_size=12, color=ACCENT2, bold=True)
    add_textbox(slide, 5.2, y, 7, 0.4, finding, font_size=12, color=TEXT_DIM)


# ============================================================
# SLIDE 14: Chapter 5 - 진화 타임라인
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.8, 0.4, 3, 0.3, "CHAPTER 5", font_size=11, color=ACCENT, bold=True)
add_textbox(slide, 0.8, 0.8, 10, 0.5, "파사드의 진화 — 벽에서 미디어로", font_size=32, color=WHITE, bold=True)

phases = [
    ("Phase 1", "~1800s", "보호", "기능적 외벽. 돌/벽돌/목재. 르네상스: 장식+비례로 권위 표현.", TEXT_MUTED),
    ("Phase 2", "1800s", "상업", "1829 봉마르셰 최초 쇼윈도우. 판유리+전기 조명 → 야간 디스플레이.", TEXT_MUTED),
    ("Phase 3", "1920~60s", "극장", "Gene Moore(Tiffany's): 윈도우를 예술로. 마네킹, 조명, 내러티브.", TEXT_DIM),
    ("Phase 4", "1970~2000s", "건축", "Koolhaas 프라다($87M), Apple 유리 파사드. 스타 건축가 시대.", ACCENT2),
    ("Phase 5", "2000~15", "디지털", "Bloomberg LED 파사드, Burberry 디지털 플래그십(유입+50%).", ACCENT2),
    ("Phase 6", "2015~Now", "체험", "Coex 파도(1억뷰), Sphere($2.3B), Glossier(매장=콘텐츠).", CYAN),
    ("Phase 7", "Emerging", "AI", "개인 반응형, 에너지 생산, 디지털 트윈. AI 파사드 > 디지털 파사드.", CYAN),
]

for i, (phase, era, name, desc, color) in enumerate(phases):
    y = 1.6 + i * 0.78
    add_textbox(slide, 0.8, y, 1.0, 0.3, phase, font_size=11, color=color, bold=True)
    add_textbox(slide, 1.8, y, 0.8, 0.3, era, font_size=10, color=TEXT_MUTED)

    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.8), Inches(y+0.08), Inches(0.14), Inches(0.14))
    dot.fill.solid(); dot.fill.fore_color.rgb = color if color != TEXT_MUTED else RGBColor(0x44,0x44,0x55)
    dot.line.fill.background()
    if i < len(phases) - 1:
        ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.85), Inches(y+0.24), Inches(0.04), Inches(0.58))
        ln.fill.solid(); ln.fill.fore_color.rgb = RGBColor(0x2A,0x2A,0x3A); ln.line.fill.background()

    add_textbox(slide, 3.2, y-0.02, 1.2, 0.35, name, font_size=16, color=WHITE, bold=True)
    add_textbox(slide, 4.4, y, 8.2, 0.35, desc, font_size=12, color=TEXT_DIM)

# Bottom callout
add_rounded_rect(slide, 0.8, 7.15, 11.5, 0.0, BG_DARK)
add_textbox(slide, 0.8, 7.1, 11.5, 0.3, "보호(Shelter) → 장식 → 스토리 → 브랜드 → 미디어 → 경험 플랫폼 → AI", font_size=13, color=ACCENT2, bold=True, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 15: Takeaway - 비교표
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.8, 0.4, 3, 0.3, "CHAPTER 6 — TAKEAWAY", font_size=11, color=ACCENT, bold=True)
add_textbox(slide, 0.8, 0.8, 10, 0.5, "파사드는 비용이 아니라 자산이다", font_size=32, color=WHITE, bold=True)
add_textbox(slide, 0.8, 1.35, 10, 0.35, "전통 광고와 파사드, 같은 예산이라면 어디에 투자할 것인가", font_size=15, color=TEXT_DIM)

# Table
table_data = [
    ("비교 항목", "TV / 디지털 광고", "파사드 / 쇼윈도우"),
    ("수명", "15~30초 (단발)", "10~30년 (영구적)"),
    ("누적 노출", "캠페인 기간 한정", "24/7 × 365 × 수십 년"),
    ("주의 포착", "배너 맹시 CTR 0.05%", "보행자 정지율 20~40%"),
    ("감각 채널", "시각+청각 (2D)", "시각+촉각+공간감 (3D+다감각)"),
    ("신호 신뢰도", "누구나 집행 가능", "높은 투자 = 위조 불가능"),
    ("CPM 추이", "연 15~20% 상승", "시간이 갈수록 하락"),
    ("잔여 가치", "캠페인 종료 시 0", "자산 가치로 축적"),
]

table = slide.shapes.add_table(len(table_data), 3, Inches(0.8), Inches(1.9), Inches(11.5), Inches(4.0)).table
table.columns[0].width = Inches(2.5)
table.columns[1].width = Inches(4.5)
table.columns[2].width = Inches(4.5)

for r, row_data in enumerate(table_data):
    for c, cell_text in enumerate(row_data):
        cell = table.cell(r, c)
        cell.text = cell_text
        p = cell.text_frame.paragraphs[0]
        p.font.name = 'Arial'

        if r == 0:  # Header
            p.font.size = Pt(12)
            p.font.color.rgb = TEXT_MUTED
            p.font.bold = True
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_CARD2
        else:
            p.font.size = Pt(13)
            if c == 0:
                p.font.color.rgb = WHITE
                p.font.bold = True
            elif c == 1:
                p.font.color.rgb = TEXT_MUTED
            else:
                p.font.color.rgb = GREEN
                p.font.bold = True
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_CARD

# ROI formula
add_rounded_rect(slide, 0.8, 6.1, 11.5, 1.1, RGBColor(0x15,0x13,0x2A), ACCENT)
add_textbox(slide, 1.1, 6.2, 2, 0.25, "ROI FORMULA", font_size=10, color=ACCENT2, bold=True)
add_textbox(slide, 1.1, 6.45, 11, 0.5, "파사드 ROI = 매출 증가(+32~50%) + 유동인구 전환(+23~35%) + Earned Media($수천만) + 자산 프리미엄(+7~12%) + 브랜드 자산(장기)", font_size=14, color=WHITE, bold=True)


# ============================================================
# SLIDE 16: 6가지 인사이트
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 0.8, 0.4, 5, 0.3, "CHAPTER 6 — 6가지 핵심 인사이트", font_size=11, color=ACCENT, bold=True)
add_textbox(slide, 0.8, 0.8, 10, 0.5, "오늘 기억해야 할 6가지", font_size=32, color=WHITE, bold=True)

insights = [
    ("01", "영구적 미디어", "광고는 만료되지만, 파사드는 수십 년간 24/7\n인상을 생성. 누적 ROI는 어떤 매체보다 높다."),
    ("02", "인스타그램 효과", "포토제닉한 파사드 → 고객이 자발적으로\n콘텐츠 생산. Glossier: 520억원 미디어 가치."),
    ("03", "후광 효과", "인상적 파사드 → 내부 모든 것의 인식을 격상.\n같은 상품도 파사드에 따라 가격 평가가 달라진다."),
    ("04", "고객 획득 채널", "파사드 리모델링 → 신규 고객 매출 +44%,\n재방문율 +16%. 파사드는 '고객 획득' 도구.", ),
    ("05", "자산 프리미엄", "차별화 파사드 → 임대료 +10~12% 높이고,\n투자 회수는 2~3년. 비용이 아니라 자산 투자."),
    ("06", "주의 경제의 역설", "디지털 광고 피로 심화 → 물리적 파사드의\n주의 획득력 상대적 강화. 배너 맹시 시대의 기회."),
]

for i, (num, title, desc) in enumerate(insights):
    col = i % 2
    row = i // 2
    x = 0.8 + col * 6.0
    y = 1.6 + row * 1.85

    add_rounded_rect(slide, x, y, 5.6, 1.6, BG_CARD, RGBColor(0x2A,0x2A,0x3A))
    add_textbox(slide, x+0.2, y+0.15, 0.5, 0.4, num, font_size=24, color=RGBColor(0x2A,0x2A,0x50), bold=True)
    add_textbox(slide, x+0.7, y+0.2, 4.5, 0.3, title, font_size=16, color=WHITE, bold=True)
    add_textbox(slide, x+0.7, y+0.6, 4.6, 0.8, desc, font_size=12, color=TEXT_DIM)


# ============================================================
# SLIDE 17: 마무리
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, 1.5, 1.5, 10.3, 2.0,
    '"파사드에 1원을 투자하면,\n그것은 건축비가 아니라\n수십 년간 작동하는\n브랜드 자산에 대한 투자다."',
    font_size=36, color=ACCENT2, bold=True, alignment=PP_ALIGN.CENTER)

pills = [
    "단순노출 → 장기 투자",
    "앵커링 → 가격 정당화",
    "신호 이론 → 신뢰 증명",
    "주의 경제 → 주의 획득",
    "피크-엔드 → 기억 지배",
]
for i, pill_text in enumerate(pills):
    x = 1.5 + i * 2.2
    add_rounded_rect(slide, x, 4.5, 2.0, 0.5, BG_CARD, RGBColor(0x2A,0x2A,0x3A))
    add_textbox(slide, x, 4.55, 2.0, 0.4, pill_text, font_size=11, color=ACCENT2, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 1.5, 5.8, 10.3, 0.5, "Thank you.", font_size=18, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)


# ============================================================
# SAVE
# ============================================================
output_path = "/Users/choi_ai/do-better-workspace/30-knowledge/38-facade-lecture-material.pptx"
prs.save(output_path)
print(f"PPT saved: {output_path}")
