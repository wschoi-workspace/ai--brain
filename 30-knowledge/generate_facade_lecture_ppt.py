#!/usr/bin/env python3
"""40분 강의자료 PPT — 건물의 첫인상: 파사드와 쇼윈도우의 비즈니스 심리학"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# === 색상 ===
BG = RGBColor(0x0F, 0x0F, 0x18)
S1 = RGBColor(0x16, 0x16, 0x24)
S2 = RGBColor(0x1E, 0x1E, 0x30)
S3 = RGBColor(0x26, 0x26, 0x38)
BD = RGBColor(0x2A, 0x2A, 0x3A)
W = RGBColor(0xF0, 0xF0, 0xF8)
DIM = RGBColor(0x98, 0x98, 0xB0)
MUT = RGBColor(0x60, 0x60, 0x78)
AC = RGBColor(0x7C, 0x6A, 0xEF)
AC2 = RGBColor(0xA7, 0x8B, 0xFA)
GLD = RGBColor(0xF0, 0xC0, 0x40)
GRN = RGBColor(0x34, 0xD3, 0x99)
BLU = RGBColor(0x60, 0xA5, 0xFA)
RSE = RGBColor(0xF4, 0x72, 0xB6)
ORG = RGBColor(0xFB, 0x92, 0x3C)
CYN = RGBColor(0x22, 0xD3, 0xEE)
RED = RGBColor(0xEF, 0x44, 0x44)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

def tb(slide, l, t, w, h, txt, sz=14, c=W, b=False, al=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.text_frame.word_wrap = True
    p = box.text_frame.paragraphs[0]
    p.text = txt
    p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = b; p.font.name = 'Arial'
    p.alignment = al
    return box

def ap(tf, txt, sz=14, c=W, b=False, sb=0):
    p = tf.add_paragraph()
    p.text = txt; p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = b
    p.font.name = 'Arial'; p.space_before = Pt(sb)
    return p

def rect(slide, l, t, w, h, fc=S1, bc=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fc
    if bc: s.line.color.rgb = bc; s.line.width = Pt(1)
    else: s.line.fill.background()
    return s

def card(slide, l, t, val, lab, c=GLD, src=""):
    rect(slide, l, t, 2.65, 1.65, S1, BD)
    tb(slide, l, t+0.2, 2.65, 0.5, val, 30, c, True, PP_ALIGN.CENTER)
    tb(slide, l, t+0.75, 2.65, 0.5, lab, 11, DIM, False, PP_ALIGN.CENTER)
    if src: tb(slide, l, t+1.25, 2.65, 0.3, src, 9, MUT, False, PP_ALIGN.CENTER)

def img_placeholder(slide, l, t, w, h, label, num=""):
    r = rect(slide, l, t, w, h, S2, BD)
    tb(slide, l, t+h/2-0.25, w, 0.3, f"[IMAGE {num}]", 11, MUT, True, PP_ALIGN.CENTER)
    tb(slide, l, t+h/2+0.05, w, 0.35, label, 10, DIM, False, PP_ALIGN.CENTER)

def chapter_divider(title, subtitle, ch_num):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    tb(s, 0, 2.5, 13.333, 0.5, f"PART {ch_num}", 14, AC, True, PP_ALIGN.CENTER)
    tb(s, 0, 3.1, 13.333, 0.8, title, 40, W, True, PP_ALIGN.CENTER)
    tb(s, 0, 4.0, 13.333, 0.5, subtitle, 16, DIM, False, PP_ALIGN.CENTER)
    return s

# ============================================================
# S1: COVER
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tb(s, 1.2, 0.5, 5, 0.3, "40-MINUTE LECTURE", 11, AC, True)
tb(s, 1.2, 1.1, 10, 1.0, "건물의 첫인상", 54, W, True)
tb(s, 1.2, 2.4, 10, 0.6, "파사드와 쇼윈도우의 비즈니스 심리학", 26, AC2, True)
bx = tb(s, 1.2, 3.3, 8, 0.8, "", 16, DIM)
bx.text_frame.paragraphs[0].text = '"예쁘니까"가 아니라 "돈이 되니까".'
bx.text_frame.paragraphs[0].font.size = Pt(16); bx.text_frame.paragraphs[0].font.color.rgb = DIM
ap(bx.text_frame, "부동산 가치 × 브랜드 첫인상 × 커뮤니케이션 파급력", 15, DIM, sb=4)

card(s, 1.2, 4.8, "3.18x", "명동 1층 vs 2층\n임대료 격차", GLD)
card(s, 4.1, 4.8, "+50%", "파사드 리노베이션\n후 매출 증가", GRN)
card(s, 7.0, 4.8, "$40M", "Glossier 매장\nEarned Media", RSE)
card(s, 9.9, 4.8, "7,891억", "젠틀몬스터 그룹\n2024 매출", CYN)

# ============================================================
# S2: AGENDA
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tb(s, 0.8, 0.4, 3, 0.3, "AGENDA", 11, AC, True)
tb(s, 0.8, 0.8, 8, 0.6, "오늘의 강의 구조 (40분)", 32, W, True)
tb(s, 0.8, 1.4, 10, 0.4, "왜 중요한가 → 얼마나 효과적인가 → 누가 증명했는가 → 어떻게 작동하는가", 14, DIM)

items = [
    ("질문", "왜 1층 임대료가 2층의 3배인가?", "2분", MUT),
    ("Part 1", "부동산 가치 — 1층이 중요한 이유", "10분", GLD),
    ("Part 2", "브랜드 첫인상 — 파사드의 가치", "12분", GRN),
    ("Part 3", "커뮤니케이션 — 파급력의 메커니즘", "10분", RSE),
    ("Takeaway", "핵심 정리 + Q&A", "6분", AC2),
]
for i, (label, desc, dur, c) in enumerate(items):
    y = 2.1 + i * 1.0
    rect(s, 0.8, y, 11.5, 0.8, S1, BD)
    tb(s, 1.0, y+0.1, 1.2, 0.3, label, 14, c, True)
    tb(s, 2.4, y+0.1, 7, 0.3, desc, 15, W, True)
    tb(s, 10.8, y+0.1, 1.3, 0.3, dur, 13, MUT, False, PP_ALIGN.RIGHT)
    tb(s, 2.4, y+0.42, 7, 0.3, "", 11, DIM)

# ============================================================
# S3: 오프닝 질문
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tb(s, 0.8, 0.4, 3, 0.3, "OPENING", 11, AC, True)

qr = rect(s, 0.8, 1.0, 11.5, 2.0, RGBColor(0x14,0x12,0x28), AC)
ql = rect(s, 1.0, 0.82, 0.45, 0.45, AC)
ql.text_frame.paragraphs[0].text = "Q"
ql.text_frame.paragraphs[0].font.size = Pt(15)
ql.text_frame.paragraphs[0].font.color.rgb = W
ql.text_frame.paragraphs[0].font.bold = True
ql.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

tb(s, 1.3, 1.2, 10.5, 0.5, "명동 1층 임대료는 ㎡당 월 296,700원.", 22, W, True)
tb(s, 1.3, 1.75, 10.5, 0.5, "같은 건물 2층은 93,200원. 왜 1층이 3.18배 비싼 걸까요?", 22, W, True)
tb(s, 1.3, 2.3, 10.5, 0.4, '"사람이 많이 보니까"라고 답하면 50점입니다. 오늘 100점짜리 답을 찾아보겠습니다.', 14, DIM)

# 층별 효용비율 시각화
tb(s, 0.8, 3.4, 5, 0.4, "한국부동산원 층별 효용비율 (1층=100%)", 14, W, True)
floors = [("1층", 100, GLD), ("2층", 60.9, DIM), ("6~10층", 52.9, MUT), ("지하1층", 35.2, MUT)]
for i, (name, pct, c) in enumerate(floors):
    y = 4.0 + i * 0.6
    tb(s, 0.8, y, 1.2, 0.3, name, 13, W, True)
    bar_w = pct / 100 * 6
    rect(s, 2.2, y+0.02, bar_w, 0.3, c if c != MUT else S3)
    tb(s, 2.2 + bar_w + 0.2, y, 1.5, 0.3, f"{pct}%", 14, c, True)

tb(s, 7.5, 3.4, 5.5, 0.4, "글로벌 비교", 14, W, True)
rect(s, 7.5, 4.0, 5.3, 2.2, S1, BD)
tb(s, 7.8, 4.2, 4.8, 0.3, "NYC 5th Avenue", 12, GLD, True)
tb(s, 7.8, 4.55, 4.8, 0.3, "1층 $2,300+/sqft/년 (CBRE 2025)", 13, DIM)
tb(s, 7.8, 4.95, 4.8, 0.3, "NYC 1층 vs 지하층", 12, ORG, True)
tb(s, 7.8, 5.3, 4.8, 0.3, "$200 vs $40/sqft → 5배 차이 (The Real Deal)", 13, DIM)

# ============================================================
# PART 1 DIVIDER
# ============================================================
chapter_divider("1층이 중요한 이유", "부동산 가치 관점 — 1층은 건물 전체의 가치를 결정하는 앵커", "1")

# ============================================================
# S5: 1층 공실의 파괴력
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tb(s, 0.8, 0.4, 5, 0.3, "PART 1 — 부동산 가치", 11, GLD, True)
tb(s, 0.8, 0.8, 10, 0.5, "1층 공실이 건물 전체를 죽인다", 32, W, True)

card(s, 0.8, 1.7, "-20%", "공실 건물의\n인근 부동산 가치 하락", RED, "필라델피아 & 콜럼버스 연구")
card(s, 3.7, 1.7, "4.7조", "필라델피아 공실로\n총 가계 자산 감소", RED, "Lincoln Institute")
card(s, 6.6, 1.7, "-8.7%", "차압 매물 0.9km 이내\n부동산 가치 하락", RED, "Cleveland Fed")
card(s, 9.5, 1.7, "7x", "Active vs Dead\nFrontage 보행자 차이", GRN, "Jan Gehl")

rect(s, 0.8, 3.8, 11.5, 1.5, S1, BD)
tb(s, 1.1, 3.9, 3, 0.25, "Jan Gehl — Cities for People", 11, GLD, True)
bx = tb(s, 1.1, 4.2, 10.8, 1.0, "", 14, DIM)
bx.text_frame.paragraphs[0].text = "1층이 활성화된 거리 vs 텅 빈 벽으로 막힌 거리의 보행자 수 차이 = 7배"
bx.text_frame.paragraphs[0].font.size = Pt(14); bx.text_frame.paragraphs[0].font.color.rgb = W; bx.text_frame.paragraphs[0].font.bold = True
ap(bx.text_frame, '"죽은 전면(Dead Frontage)"은 보행자 통행량을 최대 40% 감소시킴', 13, DIM, sb=6)
ap(bx.text_frame, "→ 1층이 죽으면 거리가 죽고, 거리가 죽으면 건물이 죽는다", 13, DIM, sb=4)

# 1층 테넌트 = 건물 가치
rect(s, 0.8, 5.6, 11.5, 1.5, S1, BD)
tb(s, 1.1, 5.7, 4, 0.25, "반대로: 좋은 1층 테넌트 = 건물 전체 가치 상승", 11, GRN, True)
bx2 = tb(s, 1.1, 6.0, 10.8, 0.9, "", 13, DIM)
bx2.text_frame.paragraphs[0].text = "Whole Foods 1층 입점 → 상층 주거 임대료 +6.0% (RCLCO 2023)"
bx2.text_frame.paragraphs[0].font.size = Pt(13); bx2.text_frame.paragraphs[0].font.color.rgb = W
ap(bx2.text_frame, "복합용도 건물 내 리테일 임대료: 상권 대비 15~25% 높음 (CoStar/ICSC)", 13, DIM, sb=4)
ap(bx2.text_frame, "도산대로 럭셔리 1층 입점 후 건물 가치 40~60% 상승 (5년간)", 13, GRN, True, sb=4)

# ============================================================
# S6: 케이스 — 성수동
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tb(s, 0.8, 0.4, 5, 0.3, "PART 1 — CASE STUDY", 11, GLD, True)
tb(s, 0.8, 0.8, 10, 0.5, "성수동의 변신 — 파사드가 동네를 바꾸다", 28, W, True)

img_placeholder(s, 0.8, 1.6, 3.6, 2.2, "성수동 2018년 공장 거리", "1")
img_placeholder(s, 4.6, 1.6, 3.6, 2.2, "블루보틀 성수 유리 파사드", "2")
img_placeholder(s, 8.4, 1.6, 3.9, 2.2, "현재 성수동 플래그십 거리", "3")

card(s, 0.8, 4.2, "6x", "주말 방문객\n5만→30만+", GLD, "2018→2024")
card(s, 3.7, 4.2, "+50%", "1층 임대료\n5년간 상승", GLD, "서울 3배 속도")
card(s, 6.6, 4.2, "2,500만", "팝업 50평\n하루 최고 임대료", GLD, "연무장길")
card(s, 9.5, 4.2, "14층", "하우스노웨어\n2025 오픈", CYN, "젠틀몬스터 성수")

rect(s, 0.8, 6.2, 11.5, 0.9, RGBColor(0x1A, 0x18, 0x10), RGBColor(0x3A, 0x30, 0x10))
tb(s, 1.1, 6.3, 10.8, 0.6, '핵심: 성수동은 "건물이 좋아져서" 비싸진 게 아닙니다. "1층 매장의 파사드가 매력적으로 바뀌면서" 동네 전체가 바뀐 겁니다.', 14, W, True)

# ============================================================
# PART 2 DIVIDER
# ============================================================
chapter_divider("파사드가 가지는 첫인상의 가치", "브랜드 관점 — 0.05초의 판단이 매출, 가격, 경험을 결정한다", "2")

# ============================================================
# S8: 0.05초의 판단
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tb(s, 0.8, 0.4, 5, 0.3, "PART 2 — 브랜드 첫인상", 11, GRN, True)
tb(s, 0.8, 0.8, 10, 0.5, "0.05초의 판단 — 첫인상이 모든 것을 결정한다", 28, W, True)

data = [
    ("24%", "쇼윈도 → 구매 결정", "4명 중 1명", GLD),
    ("76%", "매장 내 구매 결정", "파사드가 입장권", GLD),
    ("62%", "쇼윈도 → 충동구매", "10명 중 6명", GLD),
    ("70%", "첫 방문 이유 = 외관", "FIT 연구", GLD),
    ("+35%", "파사드 리노베이션\n→ 유입 증가", "3개월 이내", GRN),
    ("+50%", "파사드 리노베이션\n→ 매출 증가", "6개월 내 ROI", GRN),
    ("+44%", "리모델링 후\n신규 고객 매출", "Monash Univ.", GRN),
    ("300%", "비주얼 머천다이징\n투자 ROI", "Contra Vision", GRN),
]
for i, (v, l, src, c) in enumerate(data):
    col = i % 4; row = i // 4
    x = 0.8 + col * 3.1; y = 1.6 + row * 2.1
    card(s, x, y, v, l, c, src)

rect(s, 0.8, 5.9, 11.5, 1.2, RGBColor(0x10, 0x1A, 0x14), RGBColor(0x15, 0x3A, 0x28))
tb(s, 1.1, 6.0, 10.8, 0.8, "파사드를 바꿨더니 유입 +35%, 매출 +50%. 마케팅 캠페인이 아니라 파사드 하나를 바꾼 효과.\nTV 광고 15초와 달리, 이 효과는 건물 수명인 10~30년 동안 지속됩니다.", 14, W, True)

# ============================================================
# S9: Case — Apple Store
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tb(s, 0.8, 0.4, 5, 0.3, "PART 2 — CASE: APPLE STORE", 11, GRN, True)
tb(s, 0.8, 0.8, 10, 0.5, "유리 한 장이 만든 세계 최고 매출", 28, W, True)

img_placeholder(s, 0.8, 1.5, 3.6, 2.5, "5th Ave 글래스 큐브 야경\n15장 유리 패널, 32ft 정육면체", "1")
img_placeholder(s, 4.6, 1.5, 3.6, 2.5, "Marina Bay Sands 수상 매장\n직경 30m 유리 구체, 물 위에 부유", "2")
img_placeholder(s, 8.4, 1.5, 3.9, 2.5, "내부에서 올려다본 유리 천장\n77,000sqft 매장에 자연광 유입", "3")

card(s, 0.8, 4.4, "$7M", "글래스 큐브\n초기 투자", GRN)
card(s, 3.4, 4.4, "$440M+", "5th Ave\n연매출", GRN)
card(s, 6.0, 4.4, "$5,546", "sqft당 매출\n세계 1위", GRN)
card(s, 8.6, 4.4, "+15%", "인접 소매점\n유입 증가", BLU)
card(s, 10.8, 4.4, "2,900%", "ROI\n(매년 반복)", GRN)

rect(s, 0.8, 6.3, 11.5, 0.9, S1, BD)
tb(s, 1.1, 6.35, 10.8, 0.7, '투명성 = 브랜드 철학의 물리적 선언. "전망-피난처 이론": 안이 보이면 안전하다고 느끼고, 진입 장벽 극도로 낮아짐.\n$15M 투자 → $440M 연매출. ROI 2,900%. 건물이 서 있는 한 매년 반복.', 13, DIM)

# ============================================================
# S10: Case — Hermès
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tb(s, 0.8, 0.4, 5, 0.3, "PART 2 — CASE: HERMÈS GINZA", 11, GRN, True)
tb(s, 0.8, 0.8, 10, 0.5, '쇼윈도우 하나에 6억 원을 쓰는 이유', 28, W, True)
tb(s, 0.8, 1.3, 10, 0.35, '렌조 피아노 설계 (2001) — 13,000개 맞춤 유리 블록, 광고 0원, 앰배서더 0명', 14, DIM)

img_placeholder(s, 0.8, 1.8, 3.6, 2.5, "긴자 메종 주간 전체 외관\n13,000개 유리 블록 = 스카프 직물 질감", "1")
img_placeholder(s, 4.6, 1.8, 3.6, 2.5, "야간 등불(Lantern) 효과\n유리 블록이 내부 빛을 은은하게 확산", "2")
img_placeholder(s, 8.4, 1.8, 3.9, 2.5, "유리 블록 클로즈업\n45cm×45cm, 하나하나 수공예 제작", "3")

card(s, 0.8, 4.7, "13,000", "수공예\n유리 블록", GLD)
card(s, 3.4, 4.7, "~6.5억", "쇼윈도 1개\n설치 비용", GLD)
card(s, 6.0, 4.7, "+30%", "아이코닉 파사드\n보행 유입", GLD)
card(s, 8.6, 4.7, "$93.7B", "브랜드 가치\n(2025)", GLD)

rect(s, 0.8, 6.5, 11.5, 0.7, S1, BD)
tb(s, 1.1, 6.55, 10.8, 0.5, '신호 이론: 13,000개 수공예 유리 블록 = 위조 불가능한 품질 신호. 앵커링: 미술관급 파사드 → 버킨백 가격이 "당연한 것"이 됨.', 13, DIM)

# ============================================================
# S11: Case — Gentle Monster
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tb(s, 0.8, 0.4, 5, 0.3, "PART 2 — CASE: GENTLE MONSTER", 11, GRN, True)
tb(s, 0.8, 0.8, 10, 0.5, '안경을 파는 미술관 — 매장의 85%가 아트', 28, W, True)

img_placeholder(s, 0.8, 1.5, 3.6, 2.3, "하우스도산 외관\n5층 브루탈리즘 콘크리트, 곡면 벽", "1")
img_placeholder(s, 4.6, 1.5, 3.6, 2.3, "1층 아트 인스톨레이션\n매 시즌 전체 교체, 키네틱 조각", "2")
img_placeholder(s, 8.4, 1.5, 3.9, 2.3, "상층부 아트 환경 + 안경 진열\n면적 85% 아트, 15% 제품", "3")

card(s, 0.5, 4.2, "7,891억", "아이아이컴바인드\n2024 매출", CYN)
card(s, 3.0, 4.2, "30%", "영업이익률\n2,338억원", CYN)
card(s, 5.5, 4.2, "60~70%", "외국인\n구매 비중", CYN)
card(s, 8.0, 4.2, "2x", "해외 매출\n1년간 성장", CYN)
card(s, 10.5, 4.2, "+60%", "도산대로 부동산\n5년간 상승", GLD)

rect(s, 0.8, 6.2, 11.5, 1.0, S1, BD)
tb(s, 1.1, 6.25, 10.8, 0.8, '"안경을 팔기 위해 안경을 보여주지 않는다." 매장 면적의 85%가 아트 인스톨레이션. 매 시즌 전체 인테리어를 교체하여\n"영구적 새로움" 창출. 결과: 하우스도산 하나가 도산대로 전체의 부동산 가치를 5년간 40~60% 올림.', 13, DIM)

# ============================================================
# S12: Case — Aesop
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tb(s, 0.8, 0.4, 5, 0.3, "PART 2 — CASE: AESOP", 11, GRN, True)
tb(s, 0.8, 0.8, 10, 0.5, '"No Two Stores Alike" — 같은 브랜드, 다른 얼굴', 28, W, True)

img_placeholder(s, 0.8, 1.5, 3.6, 2.3, "이솝 가로수길 (서울)\n한옥 영감 목재 + 미니멀 유리 전면", "1")
img_placeholder(s, 4.6, 1.5, 3.6, 2.3, "이솝 긴자 (도쿄)\n수백 개 재활용 신문지 튜브 파사드", "2")
img_placeholder(s, 8.4, 1.5, 3.9, 2.3, "시그니처 앰버 보틀 월\n어느 매장에나 공통 — 유일한 일관성", "3")

card(s, 0.8, 4.2, "914억", "이솝코리아\n2021 매출", ORG, "전년 대비 2배")
card(s, 3.5, 4.2, "75%+", "매출 중\n오프라인 비중", ORG)
card(s, 6.2, 4.2, "#2", "한국 =\n글로벌 2위 시장", ORG)
card(s, 8.9, 4.2, "3조원", "기업 가치\n(로레알 인수)", ORG)

rect(s, 0.8, 6.2, 11.5, 1.0, S1, BD)
tb(s, 1.1, 6.25, 10.8, 0.8, '매장마다 다른 건축가가 그 지역의 문화와 소재를 반영하여 설계. 호박색 보틀과 미니멀 사이니지만 일정하고 나머지는 전부 다름.\n각 매장이 "발견의 대상" → 여행지에서 이솝을 찾는 것 자체가 경험 → "골목 매장"이 3조원 기업가치를 만든 비결 = 공간이 곧 브랜딩', 13, DIM)

# ============================================================
# PART 3 DIVIDER
# ============================================================
chapter_divider("파급력의 메커니즘", "커뮤니케이션 관점 — 파사드는 24시간 자가증식하는 미디어다", "3")

# ============================================================
# S14: 전통 광고 vs 파사드
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tb(s, 0.8, 0.4, 5, 0.3, "PART 3 — 커뮤니케이션", 11, RSE, True)
tb(s, 0.8, 0.8, 10, 0.5, "파사드는 24시간 작동하는 미디어다", 28, W, True)

rows = [
    ("비교 항목", "TV / 디지털 광고", "파사드 / 쇼윈도우"),
    ("수명", "15~30초 (단발)", "10~30년 (영구적)"),
    ("누적 노출", "캠페인 기간 한정", "24/7 × 365 × 수십 년"),
    ("주의 포착", "배너 맹시 CTR 0.05%", "보행자 정지율 20~40%"),
    ("감각 채널", "시각+청각 (2D)", "시각+촉각+공간감 (3D)"),
    ("신호 신뢰도", "누구나 집행 가능", "높은 투자 = 위조 불가능"),
    ("CPM 추이", "연 15~20% 상승", "시간이 갈수록 하락"),
    ("잔여 가치", "캠페인 종료 시 0", "자산 가치로 축적"),
]
tbl = s.shapes.add_table(len(rows), 3, Inches(0.8), Inches(1.6), Inches(11.5), Inches(4.2)).table
tbl.columns[0].width = Inches(2.5); tbl.columns[1].width = Inches(4.5); tbl.columns[2].width = Inches(4.5)
for r, rd in enumerate(rows):
    for c, ct in enumerate(rd):
        cell = tbl.cell(r, c); cell.text = ct
        p = cell.text_frame.paragraphs[0]; p.font.name = 'Arial'
        if r == 0:
            p.font.size = Pt(12); p.font.color.rgb = MUT; p.font.bold = True
            cell.fill.solid(); cell.fill.fore_color.rgb = S2
        else:
            p.font.size = Pt(13)
            cell.fill.solid(); cell.fill.fore_color.rgb = S1
            if c == 0: p.font.color.rgb = W; p.font.bold = True
            elif c == 1: p.font.color.rgb = MUT
            else: p.font.color.rgb = GRN; p.font.bold = True

rect(s, 0.8, 6.1, 11.5, 1.0, RGBColor(0x1A, 0x10, 0x18), RGBColor(0x3A, 0x18, 0x30))
tb(s, 1.1, 6.2, 10.8, 0.7, "TV 광고 15초에 수억 원. 끝나면 가치 = 0. 파사드는 한 번 투자하면 10~30년 동안 24시간 작동.\n시간이 갈수록 노출당 비용 하락. 광고와 달리 건물 자산으로 축적.", 14, W, True)

# ============================================================
# S15: Case — Glossier
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tb(s, 0.8, 0.4, 5, 0.3, "PART 3 — CASE: GLOSSIER", 11, RSE, True)
tb(s, 0.8, 0.8, 10, 0.5, '광고비 0원, 520억의 미디어 가치', 28, W, True)

img_placeholder(s, 0.8, 1.5, 3.6, 2.3, 'NYC "Glossier Canyon"\n밀레니얼 핑크 곡선형 셀피 공간', "1")
img_placeholder(s, 4.6, 1.5, 3.6, 2.3, "LA 매장 핑크 타일 외관\n레트로 주유소 스타일", "2")
img_placeholder(s, 8.4, 1.5, 3.9, 2.3, "고객이 매장에서 셀피\n공간 자체가 콘텐츠 스튜디오", "3")

card(s, 0.8, 4.2, "280만+", "인스타그램\n해시태그 게시물", RSE)
card(s, 3.5, 4.2, "$40M", "Earned Media\nValue (520억)", RSE)
card(s, 6.2, 4.2, "0원", "전통\n광고비 지출", RSE)
card(s, 8.9, 4.2, "+40%", "인스타 매장\n첫방문 유입↑", RSE)

rect(s, 0.8, 6.2, 11.5, 1.0, S1, BD)
tb(s, 1.1, 6.25, 10.8, 0.8, '매장 자체를 "찍고 싶은 공간"으로 설계. 고객이 사진 찍고 → 올리고 → 그걸 본 사람이 방문하고 → 또 찍고.\n이 자기강화 루프(self-reinforcing loop)가 280만 건의 포스팅과 520억원의 미디어 가치를 만듦. 현재 84% 리테일러가 이 모델을 따름.', 13, DIM)

# ============================================================
# S16: Case — Coex 파도
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tb(s, 0.8, 0.4, 5, 0.3, "PART 3 — CASE: COEX WAVE", 11, RSE, True)
tb(s, 0.8, 0.8, 10, 0.5, '건물 벽 하나로 1억 뷰', 28, W, True)

img_placeholder(s, 0.8, 1.5, 3.6, 2.3, "정면 3D 파도 일루전\n1,620㎡ LED (80m×20m)", "1")
img_placeholder(s, 4.6, 1.5, 3.6, 2.3, "측면 → 평면 LED 확인\nanamorphic 기법 특정 각도 전용", "2")
img_placeholder(s, 8.4, 1.5, 3.9, 2.3, "야간 교차로 전경\n보행자들이 서서 촬영하는 모습", "3")

card(s, 0.8, 4.2, "1억+", "소셜 미디어\n조회수", CYN, "설치 수주 내")
card(s, 3.5, 4.2, "3~4x", "3D 빌보드\n시청 시간", CYN, "일반 광고 대비")
card(s, 6.2, 4.2, "68%", "프리미엄\n인식 소비자", CYN)
card(s, 8.9, 4.2, "58%", "SNS\n공유 의향", CYN)

rect(s, 0.8, 6.2, 11.5, 1.0, S1, BD)
tb(s, 1.1, 6.25, 10.8, 0.8, '건물 벽 LED 영상이 1억 뷰 = BTS 뮤직비디오 급. 광고비 0원, 사람들이 자발적으로 촬영 및 공유.\n설치 후 건물 임대 프리미엄 상승 → 미디어 파사드 = 건물 자산 가치를 높이는 투자. 제작: d\'strict, 4개월 소요.', 13, DIM)

# ============================================================
# S17: 파급력 숫자 종합
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tb(s, 0.8, 0.4, 5, 0.3, "PART 3 — 파급력 종합", 11, RSE, True)
tb(s, 0.8, 0.8, 10, 0.5, "커뮤니케이션 파급력 — 숫자 총정리", 28, W, True)

# 케이스 종합 테이블
cases = [
    ("Glossier", "280만 포스팅, EMV $40M", "광고비 0원"),
    ("Coex 파도", "1억+ 소셜 조회", "건물 임대 프리미엄 상승"),
    ("Gentle Monster", "연 EMV 50~80억원", "매 시즌 교체 → 영구적 새로움"),
    ("Selfridges", "연 EMV £15~20M", "문화 플랫폼 포지셔닝"),
    ("Apple 매장 오픈", "오픈당 EMV $5~10M", "인접 소매점 유입 +15%"),
]
tbl = s.shapes.add_table(len(cases)+1, 3, Inches(0.8), Inches(1.5), Inches(5.5), Inches(3.0)).table
tbl.columns[0].width = Inches(1.5); tbl.columns[1].width = Inches(2.3); tbl.columns[2].width = Inches(1.7)
hdr = [("브랜드", MUT), ("파급력 지표", MUT), ("비용 대비 효과", MUT)]
for c, (t, clr) in enumerate(hdr):
    cell = tbl.cell(0, c); cell.text = t
    cell.text_frame.paragraphs[0].font.size = Pt(11); cell.text_frame.paragraphs[0].font.color.rgb = clr
    cell.text_frame.paragraphs[0].font.bold = True; cell.text_frame.paragraphs[0].font.name = 'Arial'
    cell.fill.solid(); cell.fill.fore_color.rgb = S2
for r, (brand, metric, effect) in enumerate(cases):
    for c, t in enumerate([brand, metric, effect]):
        cell = tbl.cell(r+1, c); cell.text = t
        p = cell.text_frame.paragraphs[0]; p.font.size = Pt(11); p.font.name = 'Arial'
        p.font.color.rgb = W if c == 0 else (RSE if c == 1 else DIM)
        p.font.bold = (c == 0)
        cell.fill.solid(); cell.fill.fore_color.rgb = S1

# WOM 데이터
rect(s, 6.6, 1.5, 5.7, 5.5, S1, BD)
tb(s, 6.9, 1.6, 5, 0.3, "구전(WOM) & Earned Media 효과", 13, RSE, True)

wom = [
    ("5x", "WOM 판매 전환", "유료 미디어 대비"),
    ("200%", "WOM 고객 지출", "일반 고객 대비"),
    ("37%", "WOM 고객 유지", "더 오래 유지"),
    ("$5.50", "Earned Media ROI", "투자 $1당"),
    ("+29%", "UGC 활용 시", "웹 전환율 증가"),
    ("3x", "감정적 브랜드", "더 많은 WOM 생성"),
    ("4x", "OOH 브랜드 회상", "온라인 광고 대비"),
    ("+40%", "인스타 매장", "첫방문 유입 증가"),
]
for i, (v, l, s2_text) in enumerate(wom):
    y = 2.05 + i * 0.6
    tb(s, 6.9, y, 1.0, 0.3, v, 16, RSE, True)
    tb(s, 8.0, y, 2.3, 0.3, l, 12, W)
    tb(s, 10.3, y, 2.0, 0.3, s2_text, 10, MUT)

# ============================================================
# S18: 심리학 플로우
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tb(s, 0.8, 0.4, 5, 0.3, "PART 3 — 메커니즘", 11, RSE, True)
tb(s, 0.8, 0.8, 10, 0.5, "파사드가 작동하는 7단계 심리학적 플로우", 28, W, True)

steps = [
    ("반복 노출", "단순노출효과", "무의식적 호감 → 5회 노출로 선호도 60%↑", GLD),
    ("시각적 충격", "주의 경제학", "하루 10,000개 광고 중 '멈춤' → 정지율 800배", ORG),
    ("아름다운 파사드", "후광 효과", "내부 모든 것이 좋아 보임 → 사용성 50%↑ 인식", RSE),
    ("높은 투자", "신호 이론", "'진짜 좋은 브랜드' 신뢰 → 위조 불가능한 신호", GRN),
    ("가격 앵커", "앵커링", "'이 가격은 합리적' → WTP +25%", BLU),
    ("첫/마지막 인상", "피크-엔드", "전체 기억 품질 결정 → 경험을 bookend", AC2),
    ("SNS 공유", "사회적 증거", "자기강화 루프 → WOM 판매 전환 5배", RSE),
]

for i, (trigger, theory, result, color) in enumerate(steps):
    y = 1.5 + i * 0.78
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), Inches(y+0.08), Inches(0.16), Inches(0.16))
    dot.fill.solid(); dot.fill.fore_color.rgb = color; dot.line.fill.background()
    if i < len(steps) - 1:
        ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.06), Inches(y+0.26), Inches(0.04), Inches(0.56))
        ln.fill.solid(); ln.fill.fore_color.rgb = BD; ln.line.fill.background()
    rect(s, 1.4, y-0.03, 10.8, 0.56, S1, BD)
    tb(s, 1.6, y+0.02, 1.5, 0.25, trigger, 12, color, True)
    tb(s, 3.2, y+0.02, 1.8, 0.25, theory, 12, W, True)
    tb(s, 5.1, y+0.02, 7, 0.25, result, 11, DIM)

# ============================================================
# S19: Takeaway — 3관점 요약
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tb(s, 0.8, 0.4, 3, 0.3, "TAKEAWAY", 11, AC, True)
tb(s, 0.8, 0.8, 10, 0.5, "3가지 관점 요약", 32, W, True)

perspectives = [
    ("부동산", "1층이 건물 전체의\n가치를 결정한다", "공실 → -20% / 좋은 1층 → 전체 +20%\nWhole Foods 입점 → 상층 +6%", GLD),
    ("브랜드", "파사드는 0.05초의\n가격 정당화 장치", "리노베이션 → 매출 +50%\n환경 프라이밍 → WTP +25%", GRN),
    ("커뮤니케이션", "파사드는 24시간\n자가증식하는 미디어", "Glossier EMV $40M / Coex 1억뷰\nWOM 판매전환 5배 / EMV ROI $5.50", RSE),
]
for i, (title, msg, evidence, c) in enumerate(perspectives):
    x = 0.8 + i * 4.1
    rect(s, x, 1.6, 3.8, 4.0, S1, BD)
    tb(s, x+0.2, 1.75, 3.4, 0.4, title, 20, c, True, PP_ALIGN.CENTER)
    tb(s, x+0.2, 2.3, 3.4, 0.7, msg, 15, W, True, PP_ALIGN.CENTER)
    # Divider line
    rect(s, x+0.5, 3.2, 2.8, 0.02, BD)
    tb(s, x+0.2, 3.4, 3.4, 1.5, evidence, 12, DIM, False, PP_ALIGN.CENTER)

# ROI formula
rect(s, 0.8, 5.9, 11.5, 1.2, RGBColor(0x14,0x12,0x28), AC)
tb(s, 1.1, 5.95, 2, 0.25, "ROI FORMULA", 10, AC2, True)
tb(s, 1.1, 6.25, 11, 0.7, "파사드 ROI = 매출 증가(+32~50%) + 유동인구(+23~35%) + Earned Media($수천만~수백억) + 부동산 프리미엄(+7~12%) + 브랜드 자산(장기)\n투자 회수 기간: 2~3년 (리모델링 주기 7~10년 대비)", 14, W, True)

# ============================================================
# S20: 마무리
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)

tb(s, 0, 0.8, 13.333, 0.4, "처음에 질문했죠.", 16, DIM, False, PP_ALIGN.CENTER)
tb(s, 0, 1.2, 13.333, 0.5, '"왜 1층이 3배 비싼가."', 24, W, True, PP_ALIGN.CENTER)
tb(s, 0, 1.8, 13.333, 0.4, "100점짜리 답:", 16, DIM, False, PP_ALIGN.CENTER)

bx = tb(s, 1.5, 2.5, 10.3, 2.5, "", 17, W)
bx.text_frame.word_wrap = True
bx.text_frame.paragraphs[0].text = "1층의 파사드는"
bx.text_frame.paragraphs[0].font.size = Pt(17); bx.text_frame.paragraphs[0].font.color.rgb = DIM
bx.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
ap(bx.text_frame, "건물의 가치를 결정하고 (부동산)", 20, GLD, True, 12).alignment = PP_ALIGN.CENTER
ap(bx.text_frame, "브랜드의 가격을 정당화하며 (앵커링)", 20, GRN, True, 8).alignment = PP_ALIGN.CENTER
ap(bx.text_frame, "고객이 자발적으로 만드는 수백억의 미디어를 생성합니다 (커뮤니케이션)", 20, RSE, True, 8).alignment = PP_ALIGN.CENTER

tb(s, 0, 5.2, 13.333, 0.8, '"파사드에 1원을 투자하면,\n그것은 건축비가 아니라 수십 년간 작동하는 브랜드 자산에 대한 투자입니다."', 26, AC2, True, PP_ALIGN.CENTER)

tb(s, 0, 6.5, 13.333, 0.5, "Thank you.", 18, MUT, False, PP_ALIGN.CENTER)

# ============================================================
# S21: 이미지 소싱 가이드 (부록)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
tb(s, 0.8, 0.4, 3, 0.3, "APPENDIX", 11, MUT, True)
tb(s, 0.8, 0.8, 10, 0.5, "이미지 소싱 가이드 — 총 21컷", 24, W, True)

cases_img = [
    ("성수동", "2018 공장 거리", "블루보틀 유리 파사드", "현재 플래그십 거리"),
    ("Apple", "5th Ave 큐브 야경", "Marina Bay 수상매장", "내부 유리 천장"),
    ("Hermès", "긴자 주간 전체", "야간 등불 효과", "유리 블록 클로즈업"),
    ("Gentle Monster", "하우스도산 외관", "1층 설치 작품", "상층부 아트+제품"),
    ("Aesop", "가로수길 한옥+유리", "긴자 신문지 튜브", "앰버 보틀 월"),
    ("Glossier", "NYC 핑크 캐니언", "LA 핑크 외관", "고객 셀피 모습"),
    ("Coex 파도", "정면 3D 파도", "측면 평면 확인", "야간 교차로 전경"),
]
tbl = s.shapes.add_table(len(cases_img)+1, 4, Inches(0.8), Inches(1.5), Inches(11.5), Inches(4.0)).table
tbl.columns[0].width = Inches(2.0); tbl.columns[1].width = Inches(3.2); tbl.columns[2].width = Inches(3.2); tbl.columns[3].width = Inches(3.1)
hdrs = ["케이스", "이미지 1", "이미지 2", "이미지 3"]
for c, h in enumerate(hdrs):
    cell = tbl.cell(0, c); cell.text = h
    cell.text_frame.paragraphs[0].font.size = Pt(11); cell.text_frame.paragraphs[0].font.color.rgb = MUT
    cell.text_frame.paragraphs[0].font.bold = True; cell.text_frame.paragraphs[0].font.name = 'Arial'
    cell.fill.solid(); cell.fill.fore_color.rgb = S2
for r, row_data in enumerate(cases_img):
    for c, t in enumerate(row_data):
        cell = tbl.cell(r+1, c); cell.text = t
        p = cell.text_frame.paragraphs[0]; p.font.size = Pt(11); p.font.name = 'Arial'
        p.font.color.rgb = W if c == 0 else DIM
        p.font.bold = (c == 0)
        cell.fill.solid(); cell.fill.fore_color.rgb = S1

tb(s, 0.8, 5.8, 11.5, 0.8, "소싱 권장: ArchDaily, Dezeen, Archdaily Korea (고화질 건축 사진) / 브랜드 공식 프레스킷 / Pinterest 큐레이션", 12, MUT)

# ============================================================
# SAVE
# ============================================================
out = "/Users/choi_ai/do-better-workspace/30-knowledge/38-facade-lecture-40min.pptx"
prs.save(out)
print(f"PPT saved: {out}")
