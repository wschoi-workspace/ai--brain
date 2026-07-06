# -*- coding: utf-8 -*-
"""GUBI x 답십리 전시기획 제안서 → 편집 가능한 PPTX (텍스트/도형/표)"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

def C(h): return RGBColor.from_string(h)
BG=C('1A1A1A'); BG2=C('111111'); BG3=C('222222')
FG=C('F5F0EB'); FG2=C('C9C2BA'); FG3=C('7A7570')
ACCENT=C('6C5CE7'); ACCENTL=C('A29BFE')
LINE=C('3A3A3A'); AMBER=C('D9A34B'); BRASS=C('C8A86A')
ACARD=C('2A2540'); BRASSCARD=C('2C2719')
FONT='Pretendard'

prs=Presentation()
prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]
LM=0.9; TOTAL=13.333

def slide(bg=BG):
    s=prs.slides.add_slide(BLANK)
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb=bg; r.line.fill.background(); r.shadow.inherit=False
    return s

def _ea(run,font):
    rPr=run._r.get_or_add_rPr()
    ea=rPr.find(qn('a:ea'))
    if ea is None:
        ea=rPr.makeelement(qn('a:ea'),{}); rPr.append(ea)
    ea.set('typeface',font)

def run(p,text,size,color,bold=False,font=FONT):
    rn=p.add_run(); rn.text=text
    f=rn.font; f.size=Pt(size); f.color.rgb=color; f.bold=bold; f.name=font
    _ea(rn,font); return rn

def box(s,l,t,w,h,anchor=MSO_ANCHOR.TOP):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    return tf

def para(tf,first,sb=0,ls=None,align=None):
    p=tf.paragraphs[0] if first else tf.add_paragraph()
    if sb: p.space_before=Pt(sb)
    if ls: p.line_spacing=ls
    if align is not None: p.alignment=align
    return p

def lines(tf,text,size,color,bold=False,ls=1.15,first=True,sb=0,align=None):
    for i,ln in enumerate(text.split('\n')):
        p=para(tf, first and i==0, sb=(sb if i==0 else 0), ls=ls, align=align)
        run(p,ln,size,color,bold)
    return tf

def header(s,eyebrow,title,sub=None,tsize=27,eycol=ACCENTL):
    tf=box(s,LM,0.58,11.5,0.4); run(para(tf,True),eyebrow,12,eycol,True)
    tlines=title.count('\n')+1
    tf=box(s,LM,1.0,11.5,tlines*tsize*1.1/72+0.2)
    lines(tf,title,tsize,FG,True,ls=1.06)
    sub_top=1.0+tlines*tsize*1.1/72+0.12
    if sub:
        tf=box(s,LM,sub_top,10.9,1.0); lines(tf,sub,13,FG3,ls=1.3)
        sub_top+= (sub.count('\n')+1)*13*1.3/72+0.2
    return sub_top

def card(s,l,t,w,h,fill=BG3,ln=LINE,label=None,lcol=ACCENTL,title=None,tsize=15,body=None,bcol=FG3,bsize=11):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=fill
    sh.line.color.rgb=ln; sh.line.width=Pt(0.75); sh.shadow.inherit=False
    try: sh.adjustments[0]=0.05
    except: pass
    tf=sh.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.TOP
    tf.margin_left=Inches(0.2); tf.margin_right=Inches(0.2); tf.margin_top=Inches(0.16); tf.margin_bottom=Inches(0.14)
    first=True
    if label: lines(tf,label,10.5,lcol,True,first=first); first=False
    if title:
        for i,ln2 in enumerate(title.split('\n')):
            p=para(tf, first and i==0, sb=(7 if (first and i==0) and label else 0), ls=1.1)
            run(p,ln2,tsize,FG,True)
        first=False
    if body:
        for i,ln2 in enumerate(body.split('\n')):
            p=para(tf, first and i==0, sb=(8 if i==0 and not first else 0), ls=1.35)
            run(p,ln2,bsize,bcol)
        first=False
    return sh

def cols(n,top,h,gap=0.3):
    cw=(TOTAL-2*LM-(n-1)*gap)/n
    return [(LM+i*(cw+gap),top,cw,h) for i in range(n)]

def badge_pending(s):
    w,h=3.5,0.42; l=TOTAL-LM-w; t=0.5
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=C('2C2719'); sh.line.color.rgb=AMBER; sh.line.width=Pt(1); sh.shadow.inherit=False
    try: sh.adjustments[0]=0.5
    except: pass
    tf=sh.text_frame; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=para(tf,True,align=PP_ALIGN.CENTER); run(p,'🔄 UPDATE PENDING — 정보 업데이트 예정',10.5,AMBER,True)

def pagenum(s,n):
    tf=box(s,TOTAL-1.6,7.02,1.2,0.3); p=para(tf,True,align=PP_ALIGN.RIGHT); run(p,n,10,FG3)

# ───────────── s01 COVER ─────────────
s=slide()
tf=box(s,LM,0.7,11,0.4); run(para(tf,True),'EXHIBITION PROPOSAL — 2026',12,FG3,True)
tf=box(s,LM,3.0,11.5,0.5); run(para(tf,True),'GUBI × DAPSIMNI',16,ACCENTL,True)
tf=box(s,LM,3.5,11.5,2.0)
p=para(tf,True,ls=1.0); run(p,'Heritage',52,FG,True)
p=para(tf,False,ls=1.0); run(p,'Meets ',52,FG,True); run(p,'Heritage',52,ACCENT,True)
tf=box(s,LM,5.65,10.5,1.0)
lines(tf,'덴마크 디자인 헤리티지와 서울의 가장 독특한 로컬 헤리티지가 만나는 전시.\n가구를 전시하는 것이 아니라, 두 개의 헤리티지가 서로를 설명하게 만듭니다.',13,FG2,ls=1.4)
meta=[('For','GUBI'),('Concept','Heritage Reimagined'),('Date','v1 — 2026.06.12'),('By','Project Rent')]
for i,(k,v) in enumerate(meta):
    x=LM+i*2.9
    tf=box(s,x,6.7,2.7,0.6)
    p=para(tf,True); run(p,k.upper(),9,FG3,True)
    p=para(tf,False,sb=2); run(p,v,12,FG2)
pagenum(s,'01')

# ───────────── s02 Why Seoul ─────────────
s=slide()
header(s,'INTRO — WHY SEOUL, WHY NOW','새로움이 아니라,\n축적된 시간과 새로운 해석이 만나는 지점',
 '서울은 지금 전 세계에서 가장 빠르게 변화하는 문화 도시 중 하나입니다.\n그러나 진정한 문화적 매력은 오랜 역사가 새로운 세대의 손에서 다시 해석될 때 탄생합니다.')
for (l,t,w,h),(lab,lc,ti,bd,fill,tc) in zip(cols(3,4.0,2.6),[
  ('THE SHIFT',FG3,'전시의 정의가\n바뀌고 있다','사람들은 더 이상 "제품을 보러" 오지 않습니다. 경험하고 이야기할 수 있는 맥락을 찾아 움직입니다.',BG3,FG3),
  ('THE OPPORTUNITY',FG3,'서울이 주목받는\n이유의 변화','트렌디한 신상 거리가 아니라, 시간성이 살아있는 동네가 글로벌 크리에이터의 새 목적지가 됩니다.',BG3,FG3),
  ('OUR PROPOSAL',ACCENTL,'두 헤리티지의\n만남을 제안합니다','Project Rent는 GUBI의 덴마크 디자인 헤리티지와 답십리의 로컬 헤리티지가 교차하는 전시를 제안합니다.',ACARD,ACCENTL),
]):
    card(s,l,t,w,h,fill=fill,label=lab,lcol=lc,title=ti,tsize=17,body=bd)
pagenum(s,'02')

# ───────────── s03 About GUBI ─────────────
s=slide()
header(s,'ABOUT GUBI','GUBI는 가구를 파는 브랜드가 아닙니다',
 '1967년 코펜하겐에서 시작한 GUBI는, 잊혀진 20세기 디자인 아이콘을 아카이브에서 발굴하고\n현대의 감각으로 다시 소개하는 "헤리티지 재발견(Rediscovery)"의 브랜드입니다.')
card(s,LM,4.0,5.55,2.7,label='GUBI의 철학',title='"역사와 혁신의 섬세한 결합"',tsize=16,
 body='지난 100년의 디자이너·건축가·예술가 아카이브에서 잊혀진 명작을 발굴합니다. 단순 복각이 아니라, 새로운 동시대의 아이콘으로 되살려냅니다.')
# table-like card for icons
card(s,LM+5.95,4.0,5.55,2.7,label='발굴해 되살린 아이콘들',
 body='Beetle Chair  ·  미드센추리의 현대적 재해석\nBestlite  ·  1930년대 영국 모더니즘 조명\nMulti-Lite  ·  1970년대 덴마크 디자인\nGräshoppa  ·  북유럽 디자이너 아카이브 복원\n\n관능적이고 우아한 럭셔리 — 전통 스칸디 미니멀과 다른 GUBI의 미학',bsize=11.5)
pagenum(s,'03')

# ───────────── s04 Heritage Focus ─────────────
s=slide()
header(s,"GUBI'S HERITAGE FOCUS",'이번 전시가 향하는 곳 — 헤리티지',
 '이번 GUBI 전시는 신제품 진열이 아니라, 브랜드의 헤리티지 그 자체에 초점을 맞춥니다.\n질문은 하나입니다 — 이 헤리티지를 가장 설득력 있게 보여줄 장소는 어디인가.')
card(s,LM,3.95,11.53,1.3,fill=ACARD,ln=ACCENT,label='THE CORE QUESTION',
 body='헤리티지를 다루는 전시는, 헤리티지가 살아있는 장소에서 열릴 때 가장 강력합니다. 화이트큐브가 아니라 시간이 축적된 맥락 안에 놓일 때 GUBI의 철학은 비로소 증명됩니다.',bsize=13,bcol=FG2)
for (l,t,w,h),(lab,lc,ti,bd,fill) in zip(cols(3,5.45,1.55),[
  ('NOT THIS',FG3,'또 하나의 쇼룸','깨끗한 갤러리에 가구를 놓는 전시는 이미 많습니다. 맥락이 없으면 헤리티지는 "오래된 디자인"으로 소비됩니다.',BG3),
  ('BUT THIS',FG3,'맥락 속의 헤리티지','오래된 것과 새 해석이 공존하는 동네에서, GUBI의 발굴·재해석 철학은 장소 자체로 설명됩니다.',BG3),
  ('THE ANSWER',ACCENTL,'답십리','서울에서 "잊혀진 것이 다시 발견되는" 유일한 동네. GUBI 철학과 물리적으로 일치합니다.',ACARD),
]):
    card(s,l,t,w,h,fill=fill,label=lab,lcol=lc,title=ti,tsize=15,body=bd,bsize=10.5)
pagenum(s,'04')

# ───────────── s05 About Project Rent ─────────────
s=slide()
header(s,'ABOUT PROJECT RENT','We Build Brands, Not Just Spaces',
 'Project Rent는 단순 공간 운영사가 아닙니다. 매장을 "판매처"가 아니라 "미디어"로 설계하는\n리테일 미디어 플랫폼이자, 브랜드를 직접 만들고 운영하며 성장시키는 브랜드 빌더입니다.')
for (l,t,w,h),(ti,bd) in zip(cols(4,3.95,1.3),[
  ('브랜드 전략·컨설팅','브랜드의 본질을 정의하고 메시지를 설계'),
  ('공간 기획','공간을 콘텐츠이자 채널로 디자인'),
  ('콘텐츠 제작','에디토리얼 문법으로 스토리 전달'),
  ('실제 매장 운영','기획을 넘어 직접 운영·성과측정'),
]):
    card(s,l,t,w,h,fill=BG2,title=ti,tsize=13,body=bd,bsize=10.5)
for (l,t,w,h),(lab,bd) in zip(cols(3,5.45,1.55),[
  ('RETAIL AS MEDIA','공간이 곧 콘텐츠이자 채널. 매장을 미디어 자산으로 운용합니다.'),
  ('EDITORIAL DISCIPLINE','Dezeen·Wallpaper* 톤의 매거진 문법으로 브랜드를 편집합니다.'),
  ('MODULAR SCALABILITY','단발 팝업이 아닌 반복 가능한 모듈형 플랫폼으로 확장합니다.'),
]):
    card(s,l,t,w,h,fill=ACARD,ln=ACCENT,label=lab,body=bd,bsize=11)
pagenum(s,'05')

# ───────────── s06 성수동 ─────────────
s=slide()
header(s,'PROVEN THROUGH REAL PROJECTS — 01','성수동 팝업 문화를 만들어낸 팀',
 '"성수동을 어떻게 할까"가 아니라, 성수동을 팝업의 거리로 만든 시작점이 Project Rent였습니다.')
card(s,LM,3.7,7.2,1.7,fill=BG2,label='2019 — 분기점',title='가나초콜릿하우스',tsize=18,
 body='이 프로젝트 이전 성수동의 팝업은 분기당 3~4개. 이후 분기당 40~50개로 폭증하며 "성수동 팝업 시대"가 열렸습니다. Project Rent는 그 흐름을 만든 1세대 기획자입니다.')
tf=box(s,LM,5.55,7.2,1.2)
lines(tf,'대표 최원석은 KBS로부터 "성수동의 스티브 잡스"로 소개되었습니다. "성수동을 잡지 같은 거리로 만들겠다"는 비전 아래 작은 공간을 팝업의 성지로 바꿔왔습니다.',12,FG3,ls=1.4)
stats=[('250+','누적 기획·운영 팝업'),('60억','누적 계약 실적'),('특허 IP','RXR 성과측정 · 제10-2545502호')]
for i,(v,lab) in enumerate(stats):
    t=3.7+i*1.02
    sh=card(s,8.35,t,3.18,0.9,fill=BG3)
    tf=sh.text_frame
    p=tf.paragraphs[0]; run(p,v,24 if i<2 else 17,ACCENT,True)
    p=tf.add_paragraph(); p.space_before=Pt(2); run(p,lab,10,FG3,True)
tf=box(s,LM,6.75,11.5,0.4)
run(para(tf,True),'클라이언트: 현대자동차 · 매일유업 · 오비맥주 · 롯데웰푸드 · 배달의민족 · Apple TV+ (파친코2) · 한국관광공사 · KOCCA 더셀렉츠',10.5,FG3)
pagenum(s,'06')

# ───────────── s07 답십리 of ─────────────
s=slide()
header(s,'PROVEN THROUGH REAL PROJECTS — 02','그리고 지금, 답십리 — of',
 'Project Rent는 성수동에서 그치지 않았습니다. 답십리에서 of를 직접 기획·운영하며 그 동네의 새로운 분위기를 만들어내고 있습니다.')
card(s,LM,4.0,5.55,2.7,fill=ACARD,ln=ACCENT,label='of — 답십리 · @ofjec_t',
 title='론칭 3개월 만에\n서울의 로컬 핫스팟으로',tsize=18,
 body='Project Rent가 직접 기획·운영한 of는 론칭 후 단 3개월 만에 서울 디자인 업계·크리에이터 그룹·F&B 커뮤니티 사이에서 빠르게 입소문을 형성하며 답십리의 새 핫스팟으로 자리잡았습니다.')
card(s,LM+5.95,4.0,5.55,2.7,fill=BG3,label='이것이 의미하는 것',
 body='답십리가 뜨는 동네가 된 데에는 이유가 있습니다 — 그 분위기를 우리가 직접 설계하고 팔았기 때문입니다.\n\nGUBI에게 답십리를 제안하는 우리는 외부 에이전시가 아니라, 답십리의 지금을 만든 당사자입니다.',bsize=12)
pagenum(s,'07')

# ───────────── s08 Our Difference ─────────────
s=slide()
header(s,'OUR DIFFERENCE','우리는 제품을 진열하지 않습니다.\n맥락 속에 위치시킵니다.',tsize=30)
card(s,LM,3.9,5.55,1.9,fill=BG2,label='대부분의 전시',title='"무엇을 보여줄까"',tsize=20,
 body='제품을 공간에 배치합니다. 멋진 사진은 남지만, 왜 이곳이어야 했는지는 남지 않습니다.')
card(s,LM+5.95,3.9,5.55,1.9,fill=ACARD,ln=ACCENT,label='PROJECT RENT',
 title='"왜 이 장소에서\n이 브랜드가 이야기되는가"',tsize=18,
 body='브랜드와 장소가 서로를 설명하게 만듭니다. 전시가 끝나도 "그 동네의 그 경험"으로 기억됩니다.')
card(s,LM,6.05,11.53,0.85,fill=ACARD,ln=ACCENT,
 body="We don't build exhibitions.  We create context.",bsize=18,bcol=FG)
pagenum(s,'08')

# ───────────── s09 Why Dapsimni 1 ─────────────
s=slide()
header(s,'WHY DAPSIMNI — 01','서울의 숨겨진 헤리티지 디스트릭트',
 '답십리는 서울에서도 가장 독특한 역사를 가진 지역 — 오래된 물건들이 새로운 주인을 만나기 위해 모여드는 곳.')
for (l,t,w,h),(ti,bd) in zip(cols(4,3.85,1.4),[
  ('고가구','시간이 깃든 목가구가 모이는 거리'),
  ('골동품','서울 최대의 고미술·골동 상권'),
  ('수집 문화','가치를 알아보는 사람들의 네트워크'),
  ('장인 문화','복원하고 되살리는 손의 기술'),
]):
    card(s,l,t,w,h,fill=BG3,title=ti,tsize=15,body=bd,bsize=10.5)
    # brass title color override
card(s,LM,5.5,11.53,1.4,fill=ACARD,ln=ACCENT,label='THE HIDDEN PARALLEL',
 body='답십리는 서울에서 "잊혀진 것의 가치가 다시 발견되는" 거의 유일한 동네입니다. 오래된 물건이 버려지지 않고 누군가의 안목에 의해 다시 살아나는 곳 — 이것은 GUBI가 하는 일과 정확히 같습니다.',bsize=13,bcol=FG2)
pagenum(s,'09')

# ───────────── s10 Why Dapsimni 2 ─────────────
s=slide()
header(s,'WHY DAPSIMNI — 02','예상치 못한 크리에이티브 디스트릭트',
 '답십리는 과거를 보존하는 장소가 아니라, 과거를 현재의 문화로 재해석하는 장소로 진화하고 있습니다.')
for (l,t,w,h),(lab,lc,ti,bd,fill) in zip(cols(3,3.95,2.0),[
  ('NEW GENERATION',FG3,'젊은 크리에이터의 재발견','젊은 디자이너·크리에이터가 다시 주목하는 지역. 과거의 시간성과 새 문화가 공존합니다.',BG3),
  ('GLOBAL ATTENTION',FG3,'해외 크리에이터의 방문','이제 로컬만의 공간이 아닙니다. 해외 방문객과 크리에이터가 꾸준히 답십리를 찾습니다.',BG3),
  ('CULTURAL EXPERIMENT',ACCENTL,'of의 DJ 세션','유명 DJ가 of에서 직접 디제잉 세션을 진행하는 등 새로운 문화 실험이 이어집니다.',ACARD),
]):
    card(s,l,t,w,h,fill=fill,label=lab,lcol=lc,title=ti,tsize=15,body=bd)
card(s,LM,6.15,11.53,0.75,fill=ACARD,ln=ACCENT,
 body='답십리는 과거를 박제하는 곳이 아니라, 과거를 현재의 언어로 다시 쓰는 곳입니다.',bsize=15,bcol=FG)
pagenum(s,'10')

# ───────────── s11 WOW ─────────────
s=slide(BG2)
header(s,'WHY GUBI FITS HERE','Heritage Meets Heritage',
 'GUBI와 답십리는 같은 일을 합니다 — 잊혀진 것의 가치를 다시 발견하고, 현재의 감각으로 되살리는 것.',tsize=34,eycol=ACCENT)
card(s,LM,4.1,5.0,1.95,fill=BRASSCARD,ln=BRASS,label='GUBI',lcol=BRASS,
 title='Design Heritage\nReimagined',tsize=18,
 body='잊혀진 20세기 디자인 아이콘을 아카이브에서 발굴해 현대 감각으로 되살린다')
tf=box(s,LM+5.05,4.1,1.43,1.95,anchor=MSO_ANCHOR.MIDDLE)
p=para(tf,True,align=PP_ALIGN.CENTER); run(p,'MEETS',13,FG3,True)
card(s,LM+6.53,4.1,5.0,1.95,fill=ACARD,ln=ACCENT,label='DAPSIMNI',lcol=ACCENTL,
 title='Cultural Heritage\nReimagined',tsize=18,
 body='오래된 물건과 시간성을 새로운 세대의 안목으로 다시 발견하고 재해석한다')
card(s,LM,6.25,11.53,0.78,fill=ACARD,ln=ACCENT,
 body='GUBI는 답십리를 설명하고, 답십리는 GUBI를 증명합니다.',bsize=16,bcol=FG)
pagenum(s,'11')

# ───────────── s12 Concept ─────────────
s=slide()
header(s,'EXHIBITION CONCEPT','GUBI Across Dapsimni',
 '하나의 전시장에 모든 것을 담는 대신, 답십리 동네 전체를 하나의 전시 경험으로 구성합니다.')
card(s,LM,3.95,5.55,2.7,fill=BG3,label='관람이 아니라 탐험',
 body='방문객은 하나의 전시장을 둘러보는 것이 아니라, 답십리라는 동네 자체를 걸으며 GUBI의 세계관을 경험합니다. 골동품 거리를 지나 GUBI 헤리티지 컬렉션을 만나는 동선 자체가 전시의 내러티브가 됩니다.\n\n"전시장을 관람하는 것이 아니라, 동네를 탐험하며 브랜드를 만난다."',bsize=12)
# venue table
rows=[('거점','역할'),('두손갤러리','메인 전시 · 헤리티지 아카이브'),('고복희','디자인 오브제 · 스토리텔링'),
      ('of  @ofjec_t','라이프스타일 · 커뮤니티'),('호박','디자이너 인터뷰 · 브랜드 필름')]
tb=s.shapes.add_table(len(rows),2,Inches(LM+5.95),Inches(3.95),Inches(5.55),Inches(2.4)).table
tb.columns[0].width=Inches(2.1); tb.columns[1].width=Inches(3.45)
for ri,(a,b) in enumerate(rows):
    for ci,val in enumerate((a,b)):
        cell=tb.cell(ri,ci); cell.fill.solid()
        cell.fill.fore_color.rgb=BG2 if ri==0 else BG
        cell.margin_left=Inches(0.12); cell.margin_top=Inches(0.04); cell.margin_bottom=Inches(0.04)
        cell.vertical_anchor=MSO_ANCHOR.MIDDLE
        tf=cell.text_frame; tf.word_wrap=True
        p=tf.paragraphs[0]
        run(p,val,11 if ri else 10,(FG3 if ri==0 else (FG if ci==0 else FG3)),ri==0 or ci==0)
tf=box(s,LM+5.95,6.5,5.55,0.3); run(para(tf,True),'※ 각 거점 상세는 다음 페이지 — 현재 업데이트 예정',11,AMBER)
pagenum(s,'12')

# ───────────── s13~s16 placeholders ─────────────
venues=[
 ('MAIN EXHIBITION — VENUE 01','두손갤러리','Duson Gallery','약 200평 규모의 메인 전시 공간. GUBI 헤리티지의 중심 무대.',
  ['GUBI 메인 전시 — 대표 헤리티지 컬렉션','브랜드 히스토리 & 디자인 아카이브','메인 콘텐츠 프로그램의 중심 무대','약 200평 — 동선의 시작점이자 절정'],
  '· 정확한 면적·층고·공간 사진\n· 전시 구성안·동선 다이어그램\n· 전시 컬렉션 리스트\n· 프로그램 스케줄','13'),
 ('SPECIAL VENUE — 01','고복희','GOBOKI','디자인 오브제와 아카이브 아이템을 위한 스토리텔링 공간.',
  ['디자인 오브제 큐레이션','아카이브 아이템 전시','스토리텔링 중심 콘텐츠'],
  '· 공간 면적·사진·분위기\n· 전시 오브제 구성\n· 스토리텔링 시나리오','14'),
 ('SPECIAL VENUE — 02','of','@ofjec_t','Project Rent가 직접 운영하는 공간. 디자인과 라이프스타일, 커뮤니티가 만나는 곳.',
  ['디자인과 라이프스타일의 결합','크리에이티브 커뮤니티 프로그램','토크 & 소규모 이벤트','Project Rent 직영 — 운영 안정성'],
  '· 공간 면적·사진·운영 현황\n· 커뮤니티 프로그램 라인업\n· 토크/이벤트 스케줄\n· of 성과 데이터','15'),
 ('SPECIAL VENUE — 03','호박','KOBAK','디자이너 인터뷰와 브랜드 필름을 위한 특별 큐레이션 공간.',
  ['디자이너 인터뷰 콘텐츠','브랜드 필름 상영','특별 큐레이션 전시'],
  '· 공간 면적·사진\n· 영상 콘텐츠 구성\n· 큐레이션 테마','16'),
]
for ey,kr,en,sub,roles,pend,pg in venues:
    s=slide(); badge_pending(s)
    tf=box(s,LM,0.58,11,0.4); run(para(tf,True),ey,12,ACCENTL,True)
    tf=box(s,LM,1.02,11.5,0.8); p=para(tf,True); run(p,kr+'  ',28,FG,True); run(p,en,18,FG3,False)
    tf=box(s,LM,2.0,10.9,0.7); lines(tf,sub,13,FG3,ls=1.3)
    sh=card(s,LM,3.2,5.55,3.0,fill=BG3,label='이 공간의 역할')
    tf=sh.text_frame
    for r in roles:
        p=tf.add_paragraph(); p.space_before=Pt(6); p.line_spacing=1.3
        run(p,'•  '+r,12,FG3)
    card(s,LM+5.95,3.2,5.55,3.0,fill=BG2,ln=AMBER,label='🔄 업데이트 예정 항목',lcol=AMBER,
     body=pend+'\n\n자리표시용 페이지입니다. 확정 정보를 받아 업데이트합니다.',bsize=12)
    pagenum(s,pg)

# ───────────── s17 Journey ─────────────
s=slide()
header(s,'EXHIBITION JOURNEY','방문객의 여정 — Visitor Flow',
 '네 개의 거점이 하나의 이야기로 이어집니다. 동네를 걷는 것이 곧 전시를 경험하는 것입니다.')
flow=[('01','두손갤러리','메인 전시\n헤리티지 아카이브'),('02','고복희','디자인 오브제\n스토리텔링'),
      ('03','of','라이프스타일\n커뮤니티'),('04','호박','인터뷰\n브랜드 필름')]
cw=(TOTAL-2*LM-3*0.55)/4
for i,(n,ti,bd) in enumerate(flow):
    l=LM+i*(cw+0.55)
    sh=card(s,l,4.1,cw,1.8,fill=BG3)
    tf=sh.text_frame; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; run(p,n,13,ACCENT,True)
    p=tf.add_paragraph(); p.alignment=PP_ALIGN.CENTER; p.space_before=Pt(6); run(p,ti,16,FG,True)
    p=tf.add_paragraph(); p.alignment=PP_ALIGN.CENTER; p.space_before=Pt(4); p.line_spacing=1.2
    run(p,bd,10.5,FG3)
    if i<3:
        tf=box(s,l+cw+0.02,4.1,0.51,1.8,anchor=MSO_ANCHOR.MIDDLE)
        p=para(tf,True,align=PP_ALIGN.CENTER); run(p,'→',20,ACCENT,True)
card(s,LM,6.2,11.53,0.78,fill=ACARD,ln=ACCENT,
 body='전시장을 관람하는 것이 아니라, 답십리라는 동네 자체를 탐험하며 GUBI의 세계관을 경험합니다.',bsize=15,bcol=FG)
pagenum(s,'17')

# ───────────── s18 Beyond ─────────────
s=slide()
header(s,'BEYOND EXHIBITION','전시를 넘어, 문화 프로젝트로',
 '전시 기간 동안 다양한 문화 프로그램을 함께 운영하여, 단순 관람형 전시가 아닌 문화 프로젝트로 확장합니다.')
progs=[('Design Talk','디자인과 헤리티지 토크 세션',BG3,FG3),
 ('Curator Session','큐레이터와 함께하는 깊이 있는 관람',BG3,FG3),
 ('Community Gathering','크리에이터·디자인 커뮤니티 모임',BG3,FG3),
 ('Creative Networking','국내외 크리에이터를 잇는 네트워킹',BG3,FG3),
 ('Music Program','of의 검증된 DJ·음악 프로그램',BG3,FG3),
 ('+ Activation','GUBI 브랜드에 맞춰 추가 프로그램 설계',ACARD,ACCENTL)]
pos=cols(3,4.0,1.25)+cols(3,5.45,1.25)
for (l,t,w,h),(ti,bd,fill,lc) in zip(pos,progs):
    card(s,l,t,w,h,fill=fill,ln=(ACCENT if fill==ACARD else LINE),title=ti,tsize=15,body=bd,bsize=10.5)
pagenum(s,'18')

# ───────────── s19 Why PR ─────────────
s=slide()
header(s,'WHY PROJECT RENT','우리는 GUBI를 답십리에 배치하지 않습니다',
 '우리는 GUBI와 답십리가 서로를 설명하게 만드는 프로젝트를 제안합니다.')
for (l,t,w,h),(lab,ti,bd) in zip(cols(3,3.95,2.0),[
  ('공감 · 경험','','우리는 성수동 팝업 문화를 만들고, 답십리 of를 직접 운영합니다. 이 동네를 가장 잘 압니다.'),
  ('권위 · 성과','','250+ 팝업, 60억 누적 실적, RXR 성과측정 특허 — 기획에서 운영, 측정까지 책임집니다.'),
  ('관점 · 차별','','우리는 공간을 매장이 아니라 미디어로 설계합니다. 전시가 끝나도 이야기가 남습니다.'),
]):
    card(s,l,t,w,h,fill=BG2,label=lab,lcol=ACCENT,body=bd,bsize=13)
card(s,LM,6.15,11.53,0.78,fill=ACARD,ln=ACCENT,
 body='많은 에이전시는 전시를 만듭니다. Project Rent는 브랜드가 기억될 수 있는 맥락을 만듭니다.',bsize=15,bcol=FG)
pagenum(s,'19')

# ───────────── s20 Closing ─────────────
s=slide(BG2)
tf=box(s,0,1.4,13.333,0.5,anchor=MSO_ANCHOR.MIDDLE)
p=para(tf,True,align=PP_ALIGN.CENTER); run(p,'CLOSING',12,ACCENT,True)
tf=box(s,0,2.0,13.333,1.6)
p=para(tf,True,ls=1.2,align=PP_ALIGN.CENTER); run(p,'Heritage Is Not About',40,FG,True)
p=para(tf,False,ls=1.2,align=PP_ALIGN.CENTER); run(p,'Preserving The Past.',40,FG,True)
tf=box(s,0,4.0,13.333,0.6)
p=para(tf,True,align=PP_ALIGN.CENTER)
run(p,'It Is About Giving The Past ',22,FG2); run(p,'A New Future.',22,ACCENT,True)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(4.27),Inches(5.0),Inches(4.8),Inches(1.3))
sh.fill.solid(); sh.fill.fore_color.rgb=ACARD; sh.line.color.rgb=ACCENT; sh.line.width=Pt(1); sh.shadow.inherit=False
tf=sh.text_frame; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; run(p,'GUBI × DAPSIMNI',13,ACCENTL,True)
p=tf.add_paragraph(); p.alignment=PP_ALIGN.CENTER; p.space_before=Pt(6); run(p,'A Heritage Reimagined',24,FG,True)
tf=box(s,0,6.6,13.333,0.4); p=para(tf,True,align=PP_ALIGN.CENTER); run(p,'Project Rent · 2026',12,FG3)
pagenum(s,'20')

out='/Users/choi_ai/do-better-workspace/10-projects/31-gubi-dapsimni/reports/gubi-dapsimni-proposal-v1-2026-06-12.pptx'
prs.save(out)
print('SAVED:',out)
print('슬라이드 수:',len(prs.slides._sldIdLst))
