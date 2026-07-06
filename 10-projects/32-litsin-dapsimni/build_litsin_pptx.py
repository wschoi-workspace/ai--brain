# -*- coding: utf-8 -*-
"""LI-TSIN(리진) 답십리 브랜드 기획서 — 서사문체 최종본 → 편집 가능한 PPTX (텍스트/도형/표, 이미지 캡처 없음)"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

def C(h): return RGBColor.from_string(h)
BG=C('1A1A1A'); BG2=C('111111'); BG3=C('222222')
FG=C('F5F0EB'); FG2=C('C9C2BA'); FG3=C('7A7570')
ACCENT=C('6C5CE7'); ACCENTL=C('A29BFE')
LINE=C('333333'); LINE2=C('2A2A2A'); WARM=C('E9DCC8'); RUST=C('B86C4A'); GREEN=C('8FA37E'); BLUE=C('6F8AA3'); AMBER=C('D9A34B'); RED=C('E17055')
ACARD=C('1E1B2E')
FONT='Pretendard'

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]
LM=0.85; TOTAL=13.333; RM=TOTAL-LM; CW=RM-LM

def slide(bg=BG):
    s=prs.slides.add_slide(BLANK)
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb=bg; r.line.fill.background(); r.shadow.inherit=False
    return s

def run(p,text,size,color,bold=False,italic=False,font=FONT):
    rn=p.add_run(); rn.text=text
    f=rn.font; f.size=Pt(size); f.color.rgb=color; f.bold=bold; f.italic=italic; f.name=font
    rPr=rn._r.get_or_add_rPr(); ea=rPr.makeelement(qn('a:ea'),{}); ea.set('typeface',font); rPr.append(ea)
    return rn

def box(s,l,t,w,h,anchor=MSO_ANCHOR.TOP):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    return tf

def para(tf,first,sb=0,ls=None,align=None):
    p=tf.paragraphs[0] if first else tf.add_paragraph()
    if sb: p.space_before=Pt(sb)
    if ls: p.line_spacing=ls
    if align is not None: p.alignment=align
    return p

# **bold** 마크업 처리 (bold=ACCENTL 하이라이트)
def rich(p,text,size,base=FG2,hl=FG,bold_hl=True):
    for j,seg in enumerate(text.split('**')):
        if seg=='':continue
        run(p,seg,size,(hl if j%2==1 else base),(bold_hl and j%2==1))

def lines(tf,text,size,color,bold=False,ls=1.2,first=True,align=None):
    for i,ln in enumerate(text.split('\n')):
        p=para(tf,first and i==0,ls=ls,align=align); run(p,ln,size,color,bold)
    return tf

def rect(s,l,t,w,h,fill=None,ln=None,lw=0.75,rounded=False,radius=0.05):
    shp=MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sh=s.shapes.add_shape(shp,Inches(l),Inches(t),Inches(w),Inches(h)); sh.shadow.inherit=False
    if fill is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb=fill
    if ln is None: sh.line.fill.background()
    else: sh.line.color.rgb=ln; sh.line.width=Pt(lw)
    if rounded:
        try: sh.adjustments[0]=radius
        except: pass
    return sh

def hline(s,t,l=LM,r=RM,col=LINE2): rect(s,l,t,r-l,0.012,fill=col)

def header(s,eyebrow,page,title,sub=None,tsize=27):
    tf=box(s,LM,0.5,9.5,0.4); run(para(tf,True),eyebrow,11,FG3,True)
    tf=box(s,RM-1.2,0.5,1.2,0.4); run(para(tf,True,align=PP_ALIGN.RIGHT),page,11,FG3,True)
    hline(s,0.88)
    tl=title.count('\n')+1
    tf=box(s,LM,1.0,CW,tl*tsize*1.12/72+0.2)
    for i,ln in enumerate(title.split('\n')):
        p=para(tf,i==0,ls=1.06)
        rich(p,ln,tsize,FG,ACCENT,bold_hl=False)  # **부분은 accent
        for r in p.runs: r.font.bold=True
    y=1.0+tl*tsize*1.12/72+0.12
    if sub:
        tf=box(s,LM,y,11.3,1.6);
        for i,ln in enumerate(sub.split('\n')):
            p=para(tf,i==0,ls=1.34); rich(p,ln,12.5,FG2,FG)
        y+=(sub.count('\n')+1)*12.5*1.34/72+0.18
    return y

def footer(s,sec,pg='LI-TSIN'):
    hline(s,6.92)
    tf=box(s,LM,6.98,8,0.3); run(para(tf,True),sec,8.5,FG3,True)
    tf=box(s,RM-2,6.98,2,0.3); run(para(tf,True,align=PP_ALIGN.RIGHT),pg,8.5,FG3,True)

def callout(s,t,text,h=0.7,warn=False):
    col=AMBER if warn else ACCENT
    fill=C('221E0E') if warn else ACARD
    rect(s,LM,t,CW,h,fill=fill,ln=None); rect(s,LM,t,0.04,h,fill=col)
    tf=box(s,LM+0.28,t,CW-0.5,h,anchor=MSO_ANCHOR.MIDDLE); p=para(tf,True,ls=1.32)
    rich(p,text,12,FG2,(AMBER if warn else ACCENTL))

def note(s,t,text):
    tf=box(s,LM,t,CW,0.5); p=para(tf,True,ls=1.3); run(p,text,8.5,FG3)

# ════════════ 01 COVER ════════════
s=slide(BG2)
tf=box(s,LM,1.7,11,0.4); run(para(tf,True),'BRAND PROPOSAL · 답십리 고미술상가',12,FG3,True)
tf=box(s,LM,2.25,11,1.6); p=para(tf,True); run(p,'LI-TSIN',80,FG,True); run(p,'.',80,ACCENT,True)
tf=box(s,LM,4.1,11.5,0.5); run(para(tf,True),"La Fleur d'une Âme — 어느 세계에도 속하지 못한, 영혼의 꽃 한 송이",16,ACCENTL,False,True)
tf=box(s,LM,4.75,11.5,1.0);
p=para(tf,True,ls=1.4); run(p,'경계를 넘은 자가 차린, 답십리로 들어가는 관문.',17,FG2)
p=para(tf,False,ls=1.4); run(p,'가장 오래된 조선의 아름다움을, 가장 현대적인 혀와 코와 눈으로 다시 사는 문화 살롱.',17,FG2)
hline(s,6.55)
tf=box(s,LM,6.62,6,0.3); run(para(tf,True),'BRAND PROPOSAL',10,FG3,True)
tf=box(s,RM-4,6.62,4,0.3); run(para(tf,True,align=PP_ALIGN.RIGHT),'2026.06.15 · CONFIDENTIAL',10,FG3,True)

# ════════════ 02 THE TENSION ════════════
s=slide()
header(s,'01 — THE TENSION','02','서울에서 가장 깊은 콘텐츠, 가장 높은 **문턱**.',
 '답십리는 서울에서 고미술이 가장 깊게 모인 동네다. 그러나 깊이는 종종 문턱이 되고, 문턱은 끝내 벽이 된다. 골동과 장인이 우후죽순 모였으되, 정작 사람을 들이는 관문 하나가 없다. 그렇다면 묻자 — 이 깊은 동네에 왜 머물 이유가 없는가. **비어 있던 것은 콘텐츠가 아니라, 들어설 관문이었다.**')
# split 2칸
sw=CW/2; st=4.15; sh2=2.15
for i,(tag,tcol,big,bcol,h3,p) in enumerate([
    ('ASSET · 가진 것',GREEN,'No.1',GREEN,'서울 최대 고미술·공예 집적지','골동·복원·감정·장인 네트워크가 한 동네에 밀집. 콘텐츠 밀도는 국내 최강이자 복제 불가능한 도시 자산이다.'),
    ('FRICTION · 막힌 것',RED,'머물 이유',RED,'‘볼 것’은 많은데 ‘머물 이유’가 없다','흥정 문화, 전문가의 영역, 정서적 진입장벽. 젊은 방문객·관광객·문화 소비자가 들어와 머물 관문이 없다.')]):
    x=LM+i*sw
    rect(s,x,st,sw-(0.0 if i==1 else 0.0),sh2,fill=BG3,ln=LINE2)
    tf=box(s,x+0.28,st+0.22,sw-0.55,sh2-0.4)
    run(para(tf,True),tag,10.5,tcol,True)
    p2=para(tf,False,sb=6); run(p2,big,30,bcol,True)
    p2=para(tf,False,sb=6); run(p2,h3,15,FG,True)
    p2=para(tf,False,sb=6,ls=1.3); run(p2,p,10.5,FG2)
note(s,6.45,'출처 — 답십리 고미술상가 현황 관찰')
footer(s,'LI-TSIN — THE TENSION')

# ════════════ 03 THE WOMAN ════════════
s=slide(BG2)
header(s,'02 — THE WOMAN','03','왜 **리진**인가',
 '19세기 말, 조선이 가장 격렬하게 흔들리던 시대. 한 시대가 저물고 한 시대가 밀물처럼 들이치던 그 경계 위에 한 여자가 있었다. 동양과 서양, 전통과 현대 — **리진은 양쪽을 모두 통과했으나 어느 쪽에도 갇히지 않았다.** 경계에 선 자만이 두 세계를 옮길 수 있다. 우리가 빌리는 것은 그녀의 이름이 아니라, 그 번역의 태도다.')
rect(s,LM,4.15,0.035,1.5,fill=ACCENT)
tf=box(s,LM+0.25,4.15,6.0,1.6)
p=para(tf,True,ls=1.4); run(p,'조선의 궁중 무희로 태어나 파리의 연인이 되었으나, 결국 그 어떤 세계에도 속하지 못한 영혼의 꽃일 뿐이다.',17,FG,False,True)
p=para(tf,False,sb=10); run(p,"« …je ne suis que la fleur d'une âme qui n'appartient à aucun monde. »",11,ACCENTL,False,True)
pts=[('월경(越境)','경계를 넘은 여자. 조선 → 파리 → 조선. 두 세계를 통과하고 어디에도 갇히지 않은 인물.'),
     ('번역','그녀가 한 일이 곧 리진의 직무. 낯선 세계를 익숙한 감각으로 옮겨 사람을 들이는 통역사.'),
     ('신화','믿고 싶은 이야기. 브랜드는 사실이 아니라 서사로 소비된다 — 리진은 그 자체로 세계관.')]
ty=4.15
for n,t in pts:
    tf=box(s,7.3,ty,4.3,0.7)
    p=para(tf,True,ls=1.35); run(p,n+'  ',11,ACCENT,True); run(p,t,11,FG2)
    ty+=0.62
callout(s,6.0,'**서사 활용 원칙** — 리진(리심)은 실존 여부에 학계 논쟁이 있다. 카피는 역사적 사실로 단정하지 않고 전설·모티프의 톤으로 다룬다. 결격이 아니라 ‘신화적 여백’으로 쓴다.',h=0.62,warn=True)
footer(s,'LI-TSIN — THE WOMAN')

# ════════════ 04 THE PEAR BLOSSOM ════════════
s=slide()
header(s,'03 — THE NAME · LA FLEUR','04','이름에 담긴 향, **배꽃**',
 '**이화(梨花)에 월백(月白)하고.** 옛 시인은 배꽃 흰빛과 달빛 흰빛이 포개지는 봄밤을 그렇게 적었다. 배꽃은 요란하게 피지 않는다. 화려함은 눈을 끌고 곧 잊히지만, 청초함은 코끝에 남아 오래 간다. **리진의 미학은 후자다.** 화려하지 않으나 명징하다 — 그것이 곧 브랜드의 감각 코드다.')
pears=[('후각 · SCENT','배꽃의 첫 향','매장의 **시그니처 센트는 배꽃 향**. 문을 여는 첫 순간, 도심의 소음을 끄고 손님을 리진의 세계로 들인다.'),
       ('시각 · SYMBOL','달밤에 핀 흰 꽃','다섯 장 흰 꽃잎과 흰 달빛(**梨花月白**). 낮과 밤, 흑과 백의 경계에서 피는 꽃 — 경계를 넘은 리진의 상징.'),
       ('미각 · TASTE','맑고 깊은 단맛','화려한 겉모습 안의 **청초한 단맛**. 배의 맑은 향과 여백의 풍미로 번역된 조선의 미감 — 배 숙성 크림 디저트.')]
pw=(CW-2*0.3)/3; pt=4.0; ph=2.3
for i,(pk,h4,bd) in enumerate(pears):
    x=LM+i*(pw+0.3)
    rect(s,x,pt,pw,ph,fill=BG3,ln=LINE2); rect(s,x,pt,pw,0.035,fill=ACCENT)
    tf=box(s,x+0.24,pt+0.26,pw-0.45,ph-0.5)
    run(para(tf,True),pk,10,ACCENTL,True)
    p=para(tf,False,sb=8); run(p,h4,17,FG,True)
    p=para(tf,False,sb=8,ls=1.35); rich(p,bd,11,FG2,FG)
note(s,6.45,'모티프 출처 — 이조년(李兆年, 1268~1342)의 시조 「다정가」 첫 구 ‘이화에 월백하고’. 봄밤·배꽃·흰 달빛의 정서를 브랜드 모티프로 차용했다.')
footer(s,'LI-TSIN — THE PEAR BLOSSOM')

# ════════════ 05 IDENTITY SYSTEM ════════════
s=slide(BG2)
header(s,'04 — IDENTITY SYSTEM','05','로고와 **향으로 번역된 배꽃**',
 '배꽃이라는 모티프를 두 개의 실체로 번역한다. 하나는 눈으로 읽는 **워드마크**(LI-TSIN · 梨花 · La Fleur), 다른 하나는 코로 읽는 **시그니처 센트**(배꽃 → 묵향 → 침향). 로고는 보여야 닿지만, 향은 먼저 도착한다.')
# 좌 로고박스
rect(s,LM,3.95,3.5,2.5,fill=BG3,ln=LINE2)
tf=box(s,LM,4.55,3.5,1.4,anchor=MSO_ANCHOR.TOP)
p=para(tf,True,align=PP_ALIGN.CENTER); run(p,'LI-TSIN',32,FG,True)
p=para(tf,False,align=PP_ALIGN.CENTER,sb=10); run(p,'리진 · ',13,FG2); run(p,'梨花',13,ACCENTL,True)
p=para(tf,False,align=PP_ALIGN.CENTER,sb=10); run(p,"La Fleur d'une Âme",12,FG3,False,True)
# 우 scent 3행 + palette
scents=[('TOP · 첫 향',WARM,'배꽃 · 청배 · 새벽 이슬 — 문을 여는 순간 도심을 끄는 맑고 흰 첫인상'),
        ('HEART · 머묾',ACCENTL,'백목련 · 백차 · 은은한 묵향(墨) — 머무는 시간을 감싸는 정제된 깊이'),
        ('BASE · 여운',RUST,'침향 · 한지 · 머스크 — 떠난 뒤에도 옷깃에 남는 고요한 여운')]
sx=4.6; sy=3.95; sh3=0.62
for lab,col,desc in scents:
    rect(s,sx,sy,7.0,sh3,fill=BG3,ln=LINE2); rect(s,sx,sy,0.035,sh3,fill=col)
    tf=box(s,sx+0.2,sy,6.7,sh3,anchor=MSO_ANCHOR.MIDDLE)
    p=para(tf,True); run(p,lab+'   ',10.5,col,True); run(p,desc,11,FG2)
    sy+=sh3+0.12
# palette
pal=[('백자 크림','F5F0EB',C('333333')),('배꽃 웜','E9DCC8',C('333333')),('월백 퍼플','6C5CE7',FG),('묵·러스트','B86C4A',FG),('먹 차콜','1A1A1A',FG2)]
pw2=7.0/5; px=4.6; py=6.05
for nm,hexc,tc in pal:
    rect(s,px,py,pw2-0.06,0.34,fill=C(hexc),ln=LINE2)
    tf=box(s,px,py,pw2-0.06,0.34,anchor=MSO_ANCHOR.MIDDLE); run(para(tf,True,align=PP_ALIGN.CENTER),nm,8,tc,True)
    px+=pw2
footer(s,'LI-TSIN — IDENTITY SYSTEM')

# ════════════ 06 THE IDEA ════════════
s=slide()
header(s,'05 — THE IDEA','06','헤리티지 **트랜스레이션**',
 '우리는 전통을 보존하지 않는다. 보존된 것은 이미 멈춘 것이다. 문화는 박물관 유리장 안에서 박제될 때 죽고, 오늘의 사람들이 마시고·먹고·머물고·선물하는 경험 속으로 들어설 때 비로소 다시 숨을 쉰다. **그러므로 리진이 하는 일은 복고가 아니다. 번역이다.**')
# trans 2노드 + 화살표
nw=4.4; ny=4.25; nh=1.5
rect(s,LM+0.3,ny,nw,nh,fill=C('1E1714'),ln=RUST,lw=1)
tf=box(s,LM+0.3,ny,nw,nh,anchor=MSO_ANCHOR.MIDDLE)
p=para(tf,True,align=PP_ALIGN.CENTER); run(p,'SOURCE · 원형',10.5,RUST,True)
p=para(tf,False,align=PP_ALIGN.CENTER,sb=8); run(p,'가장 깊은 조선의 미학',20,FG,True)
tf=box(s,LM+0.3+nw,ny+nh/2-0.35,1.6,0.7,anchor=MSO_ANCHOR.MIDDLE)
p=para(tf,True,align=PP_ALIGN.CENTER); run(p,'→',34,ACCENT,True)
p=para(tf,False,align=PP_ALIGN.CENTER); run(p,'TRANSLATE',9,FG3,True)
rect(s,RM-nw-0.3,ny,nw,nh,fill=ACARD,ln=ACCENT,lw=1)
tf=box(s,RM-nw-0.3,ny,nw,nh,anchor=MSO_ANCHOR.MIDDLE)
p=para(tf,True,align=PP_ALIGN.CENTER); run(p,'OUTPUT · 번역',10.5,ACCENTL,True)
p=para(tf,False,align=PP_ALIGN.CENTER,sb=8); run(p,'가장 세련된 세계의 언어',20,FG,True)
callout(s,6.05,'“가장 깊은 조선의 미학을 가장 세련된 세계의 언어로 번역하다.” — **리진의 미션이자, 모든 제품·공간·프로그램을 관통하는 단 하나의 동사다.**')
footer(s,'LI-TSIN — THE IDEA')

# ════════════ 07 ONE SCENE ════════════
s=slide(BG2)
header(s,'06 — ONE SCENE','07','한 테이블 위, 동서고금',
 '키워드가 넷이면 소비자는 셋을 잊는다. 브랜드는 설명으로 기억되지 않는다. **한 장면**으로 각인된다. 백자 잔에 담긴 에스프레소, 그 곁의 고가구 한 점 — 동과 서, 과거와 현재가 한 상 위에 나란히 놓이는 찰나. 그 컷이 곧 리진의 시각 문법이다.')
rect(s,LM,4.0,CW,2.0,fill=BG3,ln=LINE2)
tf=box(s,LM+0.42,4.3,CW-0.8,1.0)
p=para(tf,True,ls=1.4); run(p,'파리에서 돌아온 리진이 차린 다실(茶室).',24,FG,True)
p=para(tf,False,ls=1.4); run(p,'잔은 조선 백자, 커피는 에스프레소, ',24,ACCENTL,True); run(p,'곁엔 답십리에서 발굴한 고가구 한 점.',24,FG,True)
tx=LM+0.42; ty2=5.55
for tag in ['越境 — 경계를 넘은 잔','고미술 + 현대 디저트 = 한 프레임','공간 = 사진 프레이밍 시스템','발견되는 한 컷']:
    w=0.4+len(tag)*0.125
    rect(s,tx,ty2,w,0.36,fill=BG2,ln=LINE2)
    tf=box(s,tx,ty2,w,0.36,anchor=MSO_ANCHOR.MIDDLE); run(para(tf,True,align=PP_ALIGN.CENTER),tag,10.5,FG2)
    tx+=w+0.14
note(s,6.45,'설계 원칙 — 모든 좌석·집기는 ‘고미술 + 현대 F&B’가 한 프레임에 잡히도록 배치한다. 인테리어가 아니라 미디어를 설계한다.')
footer(s,'LI-TSIN — ONE SCENE')

# ════════════ 08 THREE PILLARS ════════════
s=slide()
header(s,'07 — THREE PILLARS','08','세 가지 경험, **하나의 주연**',
 '리진은 세 개의 기둥으로 선다. 그러나 셋의 무게는 같지 않다. **F&B는 사람을 들이는 미끼, 공예 큐레이션이 가닿아야 할 목적지다.** 주연은 언제나 공예가 잡는다 — 그래야 ‘예쁜 한옥 카페’ 한 줄로 수렴하지 않는다.')
pil=[('PILLAR 01 · THE HERITAGE','★ 주연 — 목적지','고미술과 공예','답십리의 가장 강력한 자산. 고미술·현대 공예를 큐레이션하는 살아있는 갤러리. 커피를 마시러 왔다가 한국의 아름다움을 발견한다.','브랜드 무게중심 · 신뢰성 · 큐레이션 커머스의 원천',True),
     ('PILLAR 02 · THE FUSION','미끼 — 게이트','재해석한 디저트','약과를 그대로 내지 않는다. 조선의 맛을 현대 파티세리로 번역 — 오미자 무스, 막걸리 가나슈, 흑임자 밀푀유, 배 숙성 크림.','대중적 접근성 · 시각적 극대화 · F&B 기프팅',False),
     ('PILLAR 03 · THE MODERN','미끼 — 앵커','다도(茶道) 커피','추출·서빙에 궁중 다례의 리추얼을 입힌다. 백자 선의 잔, 공예가의 드리퍼. 커피 한 잔이 숨을 고르는 의식이 된다.','일상 반복 구매 · 앵커 콘텐츠 · 객단가 기반',False)]
widths=[CW*0.40,CW*0.30,CW*0.30]; px=LM; pt=3.95; ph=2.55
for i,(pl,pr,h4,p,role,lead) in enumerate(pil):
    w=widths[i]-(0.16 if i<2 else 0)
    rect(s,px,pt,w,ph,fill=(ACARD if lead else BG3),ln=(ACCENT if lead else LINE2),lw=(1 if lead else 0.75))
    tf=box(s,px+0.24,pt+0.22,w-0.45,ph-0.4)
    run(para(tf,True),pl,10,(ACCENTL if lead else FG3),True)
    p2=para(tf,False,sb=5); run(p2,pr,11,ACCENT,True)
    p2=para(tf,False,sb=7); run(p2,h4,18,FG,True)
    p2=para(tf,False,sb=8,ls=1.3); run(p2,p,10.5,FG2)
    p2=para(tf,False,sb=8,ls=1.25); run(p2,'역할 · ',9.5,FG3,True); run(p2,role,9.5,FG3)
    px+=widths[i]
footer(s,'LI-TSIN — THREE PILLARS')

# ════════════ 09 GATEWAY MODEL ════════════
s=slide(BG2)
header(s,'08 — BUSINESS MODEL','09','게이트웨이 **3층 깔때기**',
 '정체성은 하나로 고정한다 — 동네로 들어가는 게이트웨이. 수익은 그 위에 **F&B → 큐레이션 커머스 → 살롱** 3층으로 쌓는다. **미끼는 커피지만, 깔때기의 끝은 답십리다.**')
fl=[('1층 · GATE','다도 커피 · 디저트','**사람을 들인다.** 직영 F&B로 일상 방문·회전·SNS 인지. 객단가 1.5–2.5만, 시즌 디저트로 재방문.','회전 · 인지',GREEN),
    ('2층 · COMMERCE','공예 큐레이션','**곁의 공예를 산다.** 점주·장인과 위탁 큐레이션. QR로 원점포 연결 = 게이트웨이 실증. 재고 부담 낮음.','위탁 수수료 · 5만~수백만',ACCENT),
    ('3층 · SALON','멤버십 · 프로그램','**단골·컬렉터가 된다.** 다도 클래스, 도슨트 투어, 작가 드롭, 살롱 멤버십으로 락인.','락인 · 프리미엄',AMBER)]
fy=4.0; fh=0.78
for fn,fsub,fd,fm,col in fl:
    rect(s,LM,fy,CW,fh,fill=BG3,ln=LINE2); rect(s,LM,fy,0.04,fh,fill=col)
    tf=box(s,LM+0.3,fy,2.6,fh,anchor=MSO_ANCHOR.MIDDLE)
    p=para(tf,True); run(p,fn,11,ACCENT,True)
    p=para(tf,False,sb=3); run(p,fsub,15,FG,True)
    tf=box(s,LM+3.1,fy,6.6,fh,anchor=MSO_ANCHOR.MIDDLE); p=para(tf,True,ls=1.3); rich(p,fd,11,FG2,FG)
    tf=box(s,RM-2.4,fy,2.3,fh,anchor=MSO_ANCHOR.MIDDLE); run(para(tf,True,align=PP_ALIGN.RIGHT),fm,10,ACCENTL,True)
    fy+=fh+0.16
footer(s,'LI-TSIN — GATEWAY MODEL')

# ════════════ 10 ANCHOR THESIS ════════════
s=slide()
header(s,'09 — ANCHOR THESIS','10','리진은 카페가 아니라 **동네의 로비**',
 '리진의 첫 목표는 제 매출이 아니다. **답십리 전체를 깨우는 앵커가 되는 것이다.** 사람들은 커피를 마시러 온다. 그러다 고미술을 발견한다. 그리고 결국 동네 전체를 걷는다.')
ast=[('+α','유동인구','젊은 방문객·외국인·문화 소비자를 동네로 유입시키는 첫 관문'),
     ('QR','전환','공예 코너 → 원점포 유입률로 ‘동네로 들어갔는가’를 실측'),
     ('↑','자산가치','동네의 문화 권위와 자산가치 상승을 가장 먼저 만드는 거점')]
ay=4.0
for av,al,ad in ast:
    tf=box(s,LM,ay,4.5,0.7)
    p=para(tf,True); run(p,av+'  ',26,ACCENT,True); run(p,al,13,FG2)
    p=para(tf,False,sb=2,ls=1.25); run(p,ad,10.5,FG3)
    ay+=0.84
tf=box(s,5.9,3.95,5.7,1.4)
p=para(tf,True,ls=1.34); run(p,'우리가 만드는 것은 카페도, 편집샵도, 갤러리도 아니다. ',16,FG,True); run(p,'답십리 전체를 잇는 하나의 문화적 로비다.',16,ACCENTL,True)
tf=box(s,5.9,5.35,5.7,0.9)
p=para(tf,True,ls=1.4); run(p,'1년차의 성패는 객단가나 회전이 아니라 ① 동네 트래픽 ② 공예 QR→상가 유입률 ③ 멤버십 사전신청 수로 본다. 성공은 답십리가 현대 문화 지도에 처음 이름을 올리는 것으로 측정한다.',11,FG2)
footer(s,'LI-TSIN — ANCHOR THESIS')

# ════════════ 11 EVOLUTION ════════════
s=slide(BG2)
header(s,'10 — EVOLUTION','11','게이트웨이에서 **플랫폼**으로',
 '리진은 한 점포에서 멈추지 않는다. 장면이 SNS로 퍼지고, 방문이 공예의 발견·구매로 이어지고, 단골이 멤버가 되고, 작가·점주가 합류해 큐레이션이 깊어진다. **이 바퀴가 돌수록 리진은 답십리를 만든 게이트웨이 플랫폼이 된다.**')
evo=[('초기 · 0–1년','게이트 살롱 오픈','F&B + 공예 코너 + 동네 연계. 인지·SNS 장면 자산 축적. 게이트웨이 가설 실증.','F&B 60 / 커머스 30 / 프로그램 10',True),
     ('중기 · 1–3년','커머스 · 멤버십 정착','큐레이션 커머스 강화, 멤버십·도슨트 정착. 컬렉터·작가 네트워크 형성.','F&B 40 / 커머스 35 / 프로그램 25',False),
     ('후기 · 3년~','답십리 게이트웨이 플랫폼','트래픽 분배 · 브랜드 라이선스 · 헤리티지 컨설팅. 답십리 = 리진이 만든 동네.','플랫폼 · B2B 상승 (분산)',False)]
ew=CW/3; et=4.0; eh=2.5
for i,(stage,h4,p,mix,now) in enumerate(evo):
    x=LM+i*ew
    rect(s,x,et,ew-0.02,eh,fill=(ACARD if now else BG3),ln=LINE2)
    tf=box(s,x+0.24,et+0.24,ew-0.45,eh-0.5)
    run(para(tf,True),stage,10,(ACCENTL if now else FG3),True)
    p2=para(tf,False,sb=7); run(p2,h4,16,FG,True)
    p2=para(tf,False,sb=8,ls=1.3); run(p2,p,10.5,FG2)
    p2=para(tf,False,sb=10,ls=1.25); run(p2,'수익 · ',9,FG3,True); run(p2,mix,9,ACCENT if now else FG3,now)
footer(s,'LI-TSIN — EVOLUTION')

# ════════════ 12 CLOSING ════════════
s=slide(BG2)
tf=box(s,LM,2.0,11.5,0.4); run(para(tf,True),'CLOSING',12,FG3,True)
tf=box(s,LM,2.6,11.8,2.0)
p=para(tf,True,ls=1.04); run(p,'리진이,',60,FG,True)
p=para(tf,False,ls=1.04); run(p,'답십리로 돌아왔다.',60,ACCENT,True)
tf=box(s,LM,5.0,11.5,1.3)
p=para(tf,True,ls=1.45); run(p,'리진이 차린 것은 카페가 아니라 관문이다. 사람들은 커피 한 잔을 들고 들어와, 조선의 가장 오래된 아름다움과 오늘의 가장 새로운 감각을 한 자리에서 마주한다.',16,FG2)
p=para(tf,False,ls=1.45,sb=8); run(p,'시간은 흐르고 아름다움은 쉽게 박제되지만, 리진은 그 박제를 풀어 멈춰 있던 것을 다시 살게 한다.',16,FG2)
hline(s,6.6)
tf=box(s,LM,6.66,7,0.3); run(para(tf,True),"LI-TSIN · La Fleur d'une Âme",10,FG3,True)
tf=box(s,RM-5,6.66,5,0.3); run(para(tf,True,align=PP_ALIGN.RIGHT),'다음 단계 · 공간 스펙 TBD · 예산 TBD · 오픈 TBD',10,FG3,True)

out='/Users/choi_ai/do-better-workspace/10-projects/32-litsin-dapsimni/litsin-brand-proposal-2026-06-15.pptx'
prs.save(out)
print('saved:', out, '| slides:', len(prs.slides._sldIdLst))
