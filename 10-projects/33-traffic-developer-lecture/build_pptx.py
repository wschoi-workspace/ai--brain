#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Traffic Developer 강의 — 편집 가능한 PPTX 생성 (텍스트/도형/표 객체)
폰트: 화면 HTML 대비 약 +30% 강의용 크기
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

# ───────── 색상 ─────────
BG   = RGBColor(0x1A,0x1A,0x1A); BG2 = RGBColor(0x11,0x11,0x11); BG3 = RGBColor(0x22,0x22,0x22)
FG   = RGBColor(0xF5,0xF0,0xEB); FG2 = RGBColor(0xC9,0xC2,0xBA); FG3 = RGBColor(0x7A,0x75,0x70)
ACC  = RGBColor(0x6C,0x5C,0xE7); ACCL= RGBColor(0xA2,0x9B,0xFE)
GREEN= RGBColor(0x8F,0xA3,0x7E); RED = RGBColor(0xE1,0x70,0x55); AMBER=RGBColor(0xD9,0xA3,0x4B); BLUE=RGBColor(0x6F,0x8A,0xA3)
LINE = RGBColor(0x33,0x33,0x33); LINE2=RGBColor(0x2A,0x2A,0x2A)
ACCBOX = RGBColor(0x24,0x20,0x3A)  # accent 박스 배경(짙은 보라)
ENDBOX = RGBColor(0x26,0x1F,0x3D)
REDBOX = RGBColor(0x2A,0x1C,0x18); GREENBOX=RGBColor(0x1E,0x24,0x1B); AMBERBOX=RGBColor(0x2A,0x23,0x14); BLUEBOX=RGBColor(0x17,0x1F,0x26)
FONT = "Pretendard"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
MW = 0.62                 # 좌우 마진
CW = 13.333 - 2*MW        # 콘텐츠 폭

def _set_font(run, size, color, bold=False, spacing=None):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    # 한글 폰트(ea) 지정
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {}); rPr.append(ea)
    ea.set('typeface', FONT)
    cs = rPr.find(qn('a:cs'))
    if cs is None:
        cs = rPr.makeelement(qn('a:cs'), {}); rPr.append(cs)
    cs.set('typeface', FONT)

def textbox(s, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, ls=1.0, sa=0.0):
    """runs: list of paragraphs; each paragraph = list of (text,size,color,bold[,spacing])"""
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if ls: p.line_spacing = ls
        if sa: p.space_after = Pt(sa)
        for seg in para:
            txt, size, color, bold = seg[0], seg[1], seg[2], seg[3]
            sp = seg[4] if len(seg) > 4 else None
            r = p.add_run(); r.text = txt
            _set_font(r, size, color, bold, sp)
    return tb

def _no_shadow(sp):
    spPr = sp._element.spPr
    el = spPr.makeelement(qn('a:effectLst'), {})
    spPr.append(el)

def box(s, l, t, w, h, fill, line, lw=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06):
    sp = s.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(lw)
    try:
        sp.adjustments[0] = radius
    except Exception:
        pass
    _no_shadow(sp)
    sp.text_frame.word_wrap = True
    sp.text_frame.margin_left = Inches(0.08); sp.text_frame.margin_right = Inches(0.08)
    sp.text_frame.margin_top = Inches(0.05); sp.text_frame.margin_bottom = Inches(0.05)
    return sp

def box_text(sp, paras, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, ls=1.05):
    tf = sp.text_frame; tf.vertical_anchor = anchor
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if ls: p.line_spacing = ls
        for seg in para:
            txt, size, color, bold = seg[0], seg[1], seg[2], seg[3]
            r = p.add_run(); r.text = txt
            _set_font(r, size, color, bold)

def arrow(s, l, t, w, ch="→", color=ACC, size=18):
    textbox(s, l, t, w, 0.5, [[(ch, size, color, False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ───────── 슬라이드 공통 ─────────
def new_slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    r = box(s, -0.1, -0.1, 13.533, 7.7, bg, None, shape=MSO_SHAPE.RECTANGLE, radius=0)
    r.shadow.inherit = False
    return s

def footer(s, left, right):
    textbox(s, MW, 7.06, CW*0.6, 0.3, [[(left.upper(), 8, FG3, False)]], align=PP_ALIGN.LEFT)
    textbox(s, MW+CW*0.4, 7.06, CW*0.6, 0.3, [[(right.upper(), 8, FG3, False)]], align=PP_ALIGN.RIGHT)
    box(s, MW, 7.0, CW, 0.012, LINE2, None, shape=MSO_SHAPE.RECTANGLE, radius=0)

def head_block(s, label, head_runs, sub_runs=None, y=0.55):
    textbox(s, MW, y, CW, 0.3, [[(label.upper(), 11, FG3, False)]])
    textbox(s, MW, y+0.32, CW, 0.9, [head_runs], ls=1.05)
    box(s, MW, y+1.18, 0.55, 0.028, ACC, None, shape=MSO_SHAPE.RECTANGLE, radius=0)
    if sub_runs:
        textbox(s, MW, y+1.3, CW, 0.8, [sub_runs], ls=1.3)

# ───────── 인포그래픽 헬퍼 ─────────
def chain_row(s, items, y, h=0.85, total=None, start_x=None, arr="→"):
    """items: list of (paras, kind) kind: normal/start/end"""
    total = total or CW
    start_x = MW if start_x is None else start_x
    n = len(items); aw = 0.5
    bw = (total - aw*(n-1)) / n
    x = start_x
    for i,(paras,kind) in enumerate(items):
        if kind=="end": fill,line,tc = ENDBOX, ACC, FG
        elif kind=="start": fill,line,tc = BG3, LINE, FG
        else: fill,line,tc = BG3, LINE2, FG2
        sp = box(s, x, y, bw, h, fill, line)
        box_text(sp, paras)
        x += bw
        if i < n-1:
            arrow(s, x, y, aw, arr, ACC if kind!="start" or True else FG3, 18); x += aw
    return y+h

def flow_col(s, items, cx, top, bw=3.6, bh=0.62, gap=0.34, arr="↓", color=ACC):
    y = top
    for i,(paras,kind) in enumerate(items):
        if kind=="end": fill,line = ENDBOX, ACC
        elif kind=="red": fill,line = REDBOX, RED
        elif kind=="green": fill,line = GREENBOX, GREEN
        else: fill,line = BG3, LINE2
        sp = box(s, cx-bw/2, y, bw, bh, fill, line); box_text(sp, paras)
        y += bh
        if i < len(items)-1:
            arrow(s, cx-0.25, y-0.02, 0.5, arr, color, 16); y += gap
    return y

def formula_row(s, steps, y, h=1.0):
    """steps: list of (num, title, final_bool)"""
    n=len(steps); aw=0.42; bw=(CW - aw*(n-1))/n; x=MW
    for i,(num,title,final) in enumerate(steps):
        fill,line = (ENDBOX,ACC) if final else (BG3,LINE2)
        tcol = ACCL if final else FG
        sp = box(s, x, y, bw, h, fill, line)
        box_text(sp, [[(num, 9, FG3, False)], [(title, 13.5, tcol, True)]], ls=1.1)
        x += bw
        if i<n-1: arrow(s, x, y, aw, "→", ACC, 16); x+=aw

def kpi_box(s, l, t, w, h, label, value, vunit, subv, vcolor=FG):
    sp = box(s, l, t, w, h, BG3, LINE2)
    box_text(sp, [[(label.upper(), 10, FG3, False)],
                  [(value, 26, vcolor, False), (vunit, 12, FG2, False)],
                  [(subv, 9.5, FG3, False)]], ls=1.15)

def hero_box(s, l, t, w, h, value, unit, desc, src, vcolor=ACC, fill=BG3, line=LINE2):
    sp = box(s, l, t, w, h, fill, line)
    paras = [[(value, 40, vcolor, False), (unit, 18, FG2, False)],
             [(desc, 12, FG2, False)]]
    if src: paras.append([(src, 9, FG3, False)])
    box_text(sp, paras, ls=1.15)

def card(s, l, t, w, h, title=None, body=None, tag=None, tagcolor=ACCL, fill=BG3, line=LINE2, src=None, anchor=MSO_ANCHOR.TOP):
    sp = box(s, l, t, w, h, fill, line)
    paras=[]
    if tag: paras.append([(tag, 10, tagcolor, False)])
    if title: paras.append([(title, 15, FG, True)])
    if body: paras.append([(body, 12, FG2, False)])
    if src: paras.append([(src, 9, FG3, False)])
    box_text(sp, paras, align=PP_ALIGN.LEFT, anchor=anchor, ls=1.25)

def vs_pair(s, left, right, y=2.7, h=3.4):
    halfw=(CW-1.0)/2
    lsp=box(s, MW, y, halfw, h, left.get('fill',BG3), left.get('line',LINE2))
    box_text(lsp, left['paras'], align=left.get('align',PP_ALIGN.LEFT), anchor=MSO_ANCHOR.TOP, ls=1.3)
    rsp=box(s, MW+halfw+1.0, y, halfw, h, right.get('fill',ACCBOX), right.get('line',ACC))
    box_text(rsp, right['paras'], align=right.get('align',PP_ALIGN.LEFT), anchor=MSO_ANCHOR.TOP, ls=1.3)
    c=box(s, MW+halfw+0.32, y+h/2-0.36, 0.72, 0.72, BG2, LINE, shape=MSO_SHAPE.OVAL)
    box_text(c, [[("VS", 12, FG3, False)]])

# ════════════════════════════════════════════════════════════
# 슬라이드 정의
# ════════════════════════════════════════════════════════════

# 1. 커버
s=new_slide(BG2)
textbox(s, MW, 2.0, CW, 0.4, [[("REAL ESTATE DEVELOPMENT · THE NEXT 10 YEARS", 12, FG3, False)]], align=PP_ALIGN.CENTER)
textbox(s, MW, 2.5, CW, 1.7, [[("부동산 개발의 다음 10년", 44, FG, False)],
                              [("건물에서 트래픽으로", 44, ACCL, False)]], align=PP_ALIGN.CENTER, ls=1.1)
box(s, 13.333/2-0.3, 4.35, 0.6, 0.03, ACC, None, shape=MSO_SHAPE.RECTANGLE, radius=0)
textbox(s, MW, 4.6, CW, 0.9, [[("사람들은 건물을 찾아가지 않는다.", 19, FG2, False)],
                              [("이유를 찾아간다.", 19, FG2, False)]], align=PP_ALIGN.CENTER, ls=1.3)
textbox(s, MW, 5.8, CW, 0.3, [[("BY PROJECT RENT", 12, FG3, False)]], align=PP_ALIGN.CENTER)
footer(s, "Building → Traffic", "by Project Rent · 2026")

# 2. 강의 목표
s=new_slide(BG2)
head_block(s, "Why This Lecture",
           [("이 강의는 ", 28, FG, True), ("부동산 강의가 아니다", 28, ACCL, True)],
           [("가치가 어떻게 만들어지는가 — 그리고 왜 앞으로 부동산 개발은 건설업이 아니라 ", 15, FG2, False),
            ("운영업", 15, FG, True), ("에 가까워질 것인가를 설명하는 강의다.", 15, FG2, False)])
halfw=(CW-0.4)/2
card(s, MW, 2.65, halfw, 1.9, "가치는 어떻게 만들어지는가", "땅값과 건물값은 결과일 뿐이다. 그 값을 움직이는 진짜 동력이 무엇인지 거슬러 올라간다.", "QUESTION 1", FG3)
card(s, MW+halfw+0.4, 2.65, halfw, 1.9, "왜 개발이 건설업에서 운영업이 되는가", "짓고 파는 사업에서, 사람이 찾아올 이유를 만들고 운영하는 사업으로. 패러다임 전환을 추적한다.", "QUESTION 2", ACCL, fill=ACCBOX, line=ACC)
chain_row(s, [([[("토지의 시대",13,FG,False)]],"start"),([[("건물의 시대",13,FG2,False)]],"normal"),([[("트래픽의 시대",13,FG,True)]],"end")], 4.9, h=0.75)
footer(s, "Lecture Goal", "02")

# ── PART 간지 헬퍼 ──
def divider_slide(num, kicker, q_lines, note, foot):
    s=new_slide(BG2)
    textbox(s, MW, 5.2, CW, 3.0, [[(num, 300, RGBColor(0x26,0x24,0x22), False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, MW, 2.0, CW, 0.4, [[(kicker.upper(), 13, ACCL, False)]], align=PP_ALIGN.CENTER)
    textbox(s, MW, 2.7, CW, 1.8, [[(seg if isinstance(seg,tuple) else (seg,36,FG,False)) for seg in line] for line in q_lines], align=PP_ALIGN.CENTER, ls=1.15)
    textbox(s, MW, 5.0, CW, 0.5, [[(note, 16, FG2, False)]], align=PP_ALIGN.CENTER)
    footer(s, foot, "")
    return s

# 3. PART 1
s=new_slide(BG2)
textbox(s, MW, 1.0, CW, 4.0, [[("1", 300, RGBColor(0x26,0x24,0x22), False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
textbox(s, MW, 2.0, CW, 0.4, [[("PART 1", 13, ACCL, False)]], align=PP_ALIGN.CENTER)
textbox(s, MW, 2.7, CW, 1.7, [[("부동산은 현재 가치일까,", 36, FG, False)],[("미래 가치", 36, ACCL, False),("일까?", 36, FG, False)]], align=PP_ALIGN.CENTER, ls=1.15)
textbox(s, MW, 5.0, CW, 0.5, [[("우리가 사는 것은 지금의 입지인가, 앞으로의 기대인가", 16, FG2, False)]], align=PP_ALIGN.CENTER)
footer(s, "Part 1 — 현재 가치 vs 미래 가치", "03")

# 4. P1 타임라인
s=new_slide(BG)
head_block(s, "P1 · The Nature of Value",
           [("부동산은 현재 가치일까, 미래 가치일까?", 28, FG, True)],
           [("사람들은 강남이 좋아서 샀을까, 좋아질 것이라 믿어서 샀을까. 부동산의 본질은 현재가 아니라 ", 15, FG2, False),("미래 가치의 거래", 15, FG, True),("에 있다.", 15, FG2, False)])
tl=[("1973","강남","영동지구 개발촉진지구 지정\n정책·인프라가 만든 계획도시", False),
    ("2009","판교","신도시 입주·2012 테크노밸리\nIT 산업이 만든 자족도시", False),
    ("2010s","성수","서울숲(2000s)→산업·문화 집적\n자생적으로 떠오른 거리", False),
    ("Now","다음은?","미래 기대가 먼저,\n가격은 뒤따른다", True)]
nw=CW/4; y=3.4
box(s, MW+nw*0.5, y+0.16, CW-nw, 0.02, LINE, None, shape=MSO_SHAPE.RECTANGLE, radius=0)
for i,(yr,pl,ds,now) in enumerate(tl):
    cx=MW+nw*i+nw/2
    dot=box(s, cx-0.1, y, 0.2, 0.2, ACC if now else BG3, ACC if now else FG3, 2, shape=MSO_SHAPE.OVAL)
    textbox(s, cx-nw/2, y+0.4, nw, 0.5, [[(yr, 22, FG, False)]], align=PP_ALIGN.CENTER)
    textbox(s, cx-nw/2, y+0.95, nw, 0.4, [[(pl, 14, ACCL if now else FG, True)]], align=PP_ALIGN.CENTER)
    textbox(s, cx-nw/2, y+1.35, nw, 0.9, [[(ds, 11, FG3, False)]], align=PP_ALIGN.CENTER, ls=1.25)
textbox(s, MW, 6.35, CW, 0.3, [[("출처: 서울정책아카이브(강남), 판교테크노밸리(판교), 건국대 글로컬문화전략연구소(성수)", 9, FG3, False)]], align=PP_ALIGN.CENTER)
footer(s, "Building → Traffic", "04")

# 5. P2
s=new_slide(BG)
head_block(s, "P2 · Buying the Future",
           [("사람들은 왜 미래의 강남을 사는가", 28, FG, True)],
           [("부동산 가격은 현재 가치의 반영이 아니다. 사람들이 믿는 ", 15, FG2, False),("미래의 크기", 15, FG, True),("를 거래하는 시장이다.", 15, FG2, False)])
flow_col(s, [([[("미래 기대",13,FG,True)]],"end"),([[("사람과 자본의 유입",13,FG2,False)]],"normal"),([[("가격 형성",13,FG2,False)]],"normal")], cx=MW+3.4, top=2.75, bw=3.4)
cx2=MW+CW-3.5
for i,(t,b) in enumerate([("판교","아직 없는 산업단지를 미리 샀다 — 테크노밸리 기대"),("마곡","기업 이전이 '예정'된 시점에 가격이 먼저 움직였다"),("한남동","'앞으로 더 귀해질 희소성'에 대한 베팅")]):
    card(s, cx2, 2.75+i*1.25, 4.4, 1.1, t, b, fill=BLUEBOX, line=BLUE)
footer(s, "Building → Traffic", "05")

# 6. P3
s=new_slide(BG)
head_block(s, "P3 · Order of Cause",
           [("살기 좋아서 비싼가, 비싸서 살기 좋은가", 28, FG, True)],
           [("두 가지 모두 맞다. 하지만 중요한 것은 ", 15, FG2, False),("순서", 15, FG, True),("다. 좋은 이유가 먼저인가, 가격이 먼저인가.", 15, FG2, False)])
chain_row(s, [([[("좋은 이유",13,FG,False)]],"start"),([[("사람이 모임",13,FG2,False)]],"normal"),([[("가격 상승",13,FG2,False)]],"normal"),([[("인프라 유입",13,FG2,False)]],"normal"),([[("더 좋은 이유",13,FG,True)]],"end")], 3.3, h=0.95)
textbox(s, MW, 4.9, CW, 0.9, [[("가치는 직선이 아니라 ", 15, FG2, False),("순환",15,FG,True),("한다. 한 바퀴를 먼저 돌리는 '이유'가 어디서 오는지가 핵심이다.",15,FG2,False)]], align=PP_ALIGN.CENTER, ls=1.4)
footer(s, "Building → Traffic", "06")

# 7. P4
s=new_slide(BG2)
head_block(s, "P4 · Definition",
           [("부동산 가격은 ", 28, FG, True),("사람들이 믿는 미래의 크기", 28, ACC, True),("다", 28, FG, True)],
           [("결국 사람들은 미래를 사고 있다. 그렇다면 사람들은 무엇을 보고 미래를 판단할까.", 15, FG2, False)])
chain_row(s, [([[("입지",13,FG2,False)]],"normal"),([[("학군",13,FG2,False)]],"normal"),([[("개발 호재",13,FG2,False)]],"normal"),([[("미래 기대 = 가격",13,FG,True)]],"end")], 3.4, h=0.95)
textbox(s, MW, 5.0, CW, 0.9, [[("지금까지 '미래의 근거'는 입지·학군·개발 호재였다. 그런데 이 근거들이 흔들리기 시작했다 — 그것이 다음 이야기다.", 15, FG2, False)]], align=PP_ALIGN.CENTER, ls=1.4)
footer(s, "Building → Traffic", "07")

# 8. PART 2
s=new_slide(BG2)
textbox(s, MW, 1.0, CW, 4.0, [[("2", 300, RGBColor(0x26,0x24,0x22), False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
textbox(s, MW, 2.0, CW, 0.4, [[("PART 2", 13, ACCL, False)]], align=PP_ALIGN.CENTER)
textbox(s, MW, 2.7, CW, 1.7, [[("왜 지금", 36, FG, False)],[("이 이야기가 중요한가", 36, FG, False)]], align=PP_ALIGN.CENTER, ls=1.15)
textbox(s, MW, 5.0, CW, 0.5, [[("공급의 시대가 끝나고 있다", 16, FG2, False)]], align=PP_ALIGN.CENTER)
footer(s, "Part 2 — 왜 지금인가", "08")

# 9. P5
s=new_slide(BG)
head_block(s, "P5 · PF Crisis",
           [("PF가 어려운 진짜 이유", 28, FG, True)],
           [("PF(프로젝트 파이낸싱) 위기를 흔히 금리 문제로 본다. 하지만 더 큰 원인은 ", 15, FG2, False),("공급 중심 모델의 한계", 15, FG, True),("다.", 15, FG2, False)])
chain_row(s, [([[("인구 증가",12.5,FG2,False)]],"start"),([[("공간 부족",12.5,FG2,False)]],"normal"),([[("개발",12.5,FG2,False)]],"normal"),([[("분양",12.5,FG2,False)]],"normal"),([[("성공",12.5,GREEN,True)]],"normal")], 3.3, h=0.9, arr="→")
textbox(s, MW, 4.9, CW, 0.9, [[("이 공식은 '공급이 부족하던 시대'의 공식이었다. ",15,AMBER,False),("분양만 되면 끝나던 모델이, 사람이 안 오면 멈춘다.",15,FG2,False)]], align=PP_ALIGN.CENTER, ls=1.4)
footer(s, "Building → Traffic", "09")

# 10. P6
s=new_slide(BG)
head_block(s, "P6 · The Shortage Era",
           [("우리는 너무 오래 ", 28, FG, True),("공급 부족 시대", 28, ACCL, True),("에 살았다", 28, FG, True)],
           [("짓기만 하면 가치가 올라가던 시대가 있었다. 그래서 우리는 개발을 '건물 공급'으로만 이해하게 되었다.", 15, FG2, False)])
chain_row(s, [([[("공간 부족",13,FG2,False)]],"start"),([[("건물 공급",13,FG2,False)]],"normal"),([[("가치 상승",13,FG,True)]],"end")], 3.4, h=0.95, total=9, start_x=MW+(CW-9)/2)
textbox(s, MW, 5.0, CW, 0.9, [[("이 단순한 등식이 수십 년간 작동했다. 문제는 이 등식의 ",15,FG2,False),("첫 칸(공간 부족)",15,FG,True),("이 사라지고 있다는 것이다.",15,FG2,False)]], align=PP_ALIGN.CENTER, ls=1.4)
footer(s, "Building → Traffic", "10")

# 11. P7
s=new_slide(BG2)
textbox(s, MW, 1.4, CW, 0.4, [[("P7 · THE NEW QUESTION", 11, FG3, False)]], align=PP_ALIGN.CENTER)
textbox(s, MW, 2.0, CW, 0.6, [[("이제 건물을 지으면 사람이 오는 시대가 아니다", 22, FG2, False)]], align=PP_ALIGN.CENTER)
textbox(s, MW, 3.0, CW, 1.2, [[("건물", 44, FG, False),("   ≠   ", 50, ACC, False),("트래픽", 44, FG, False)]], align=PP_ALIGN.CENTER)
textbox(s, MW, 4.7, CW, 0.6, [[("사람이 올 ", 18, FG2, False),("이유", 18, FG, True),("를 만들어야 하는 시대다", 18, FG2, False)]], align=PP_ALIGN.CENTER)
footer(s, "Building → Traffic", "11")

# 12. PART 3
s=new_slide(BG2)
textbox(s, MW, 1.0, CW, 4.0, [[("3", 300, RGBColor(0x26,0x24,0x22), False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
textbox(s, MW, 2.0, CW, 0.4, [[("PART 3", 13, ACCL, False)]], align=PP_ALIGN.CENTER)
textbox(s, MW, 2.7, CW, 1.7, [[("일본은 이미", 36, FG, False)],[("미래를 경험", 36, ACCL, False),("했다", 36, FG, False)]], align=PP_ALIGN.CENTER, ls=1.15)
textbox(s, MW, 5.0, CW, 0.5, [[("한국보다 먼저 늙은 나라가 찾은 답", 16, FG2, False)]], align=PP_ALIGN.CENTER)
footer(s, "Part 3 — 일본의 선행 경험", "12")

# 13. P8
s=new_slide(BG)
head_block(s, "P8 · Japan Aged First",
           [("일본은 먼저 늙었다", 28, FG, True)],
           [("일본은 한국보다 먼저 인구 감소·고령화·저성장을 경험했다. 그 결과 부동산은 '짓는 것'에서 ", 15, FG2, False),("'운영하는 것'", 15, FG, True),("으로 무게가 옮겨갔다.", 15, FG2, False)])
tw=(CW-0.8)/3; ty=2.9; th=2.9
hero_box(s, MW, ty, tw, th, "13.8", "%", "전국 빈집률 (역대 최고)\n총무성 주택·토지통계조사, 2023", None, vcolor=RED, fill=REDBOX, line=RED)
hero_box(s, MW+tw+0.4, ty, tw, th, "900", "만 호", "전국 빈집 총수\n총주택 6,502만 호 (총무성, 2023)", None, vcolor=FG)
card(s, MW+(tw+0.4)*2, ty, tw, th, "그런데 도쿄 도심 오피스는 2%대 만실", "빈집은 넘치는데 좋은 입지·운영되는 공간은 만실. 부동산이 양극화되고 있다는 신호.", None, fill=BLUEBOX, line=BLUE, src="도쿄 도심5구 평균 공실률, 미키상사 2025~26", anchor=MSO_ANCHOR.MIDDLE)
footer(s, "Building → Traffic", "13")

# 14. P9 vs
s=new_slide(BG)
head_block(s, "P9 · Development Redefined",
           [("일본의 개발은 왜 달라졌는가", 28, FG, True)],
           [("건물 완공이 사업의 '종료'가 아니라 '시작'이 되었다. 분양 후 손 떼는 모델에서, 운영하며 가치를 키우는 모델로.", 15, FG2, False)])
y=2.85; h=3.5; halfw=(CW-1.0)/2
lsp=box(s, MW, y, halfw, h, REDBOX, RED); box_text(lsp,[[("과거 · 한국형",11,RED,False)],[("개발",13,FG2,False)],[("↓",14,RED,False)],[("분양",13,FG2,False)],[("↓",14,RED,False)],[("종료 — 손을 뗀다",13,FG,True)]],ls=1.35)
rsp=box(s, MW+halfw+1.0, y, halfw, h, GREENBOX, GREEN); box_text(rsp,[[("현재 · 일본형",11,GREEN,False)],[("개발",13,FG2,False)],[("↓",14,GREEN,False)],[("운영",13,FG2,False)],[("↓",14,GREEN,False)],[("가치 상승 — 계속 키운다",13,FG,True)]],ls=1.35)
c=box(s, MW+halfw+0.34, y+h/2-0.36, 0.72, 0.72, BG2, LINE, shape=MSO_SHAPE.OVAL); box_text(c,[[("VS",12,FG3,False)]])
footer(s, "Building → Traffic", "14")

# 15. P10 시모키타
s=new_slide(BG)
head_block(s, "P10 · Case · Shimokitazawa",
           [("시모키타자와 — 크게 짓지 않고, 작게 잇다", 26, FG, True)],
           [("오다큐 전철은 지하화로 생긴 철도 상부 부지를 대형 쇼핑몰이 아니라 ", 15, FG2, False),("작은 브랜드와 콘텐츠", 15, FG, True),("로 채웠다.", 15, FG2, False)])
flow_col(s, [([[("철도 상부 부지 (약 1.7km)",12.5,FG2,False)]],"normal"),([[("13개 소규모 시설로 분산",12.5,FG2,False)]],"normal"),([[("작은 브랜드·가게의 집합",12.5,FG2,False)]],"normal"),([[("체류 → 재방문",12.5,FG,True)]],"end")], cx=MW+3.6, top=2.7, bw=4.2, bh=0.58, gap=0.28)
kpi_box(s, MW+CW-4.3, 2.7, 4.3, 1.2, "전면 개장", "2022", ".5", "下北線路街 (시모키타 라인로가)")
card(s, MW+CW-4.3, 4.05, 4.3, 1.85, None, "'지원형 개발(支援型開発)' — 사업자가 주도하지 않고 지역을 뒤에서 받친다. BONUS TRACK(2020)·reload(2021) 단계적 조성.", fill=BG3, line=LINE2, src="출처: 오다큐 그룹 공식, 2022", anchor=MSO_ANCHOR.MIDDLE)
footer(s, "Building → Traffic", "15")

# 16. P11 다이칸야마
s=new_slide(BG)
head_block(s, "P11 · Case · Daikanyama T-SITE",
           [("다이칸야마 츠타야 — 책이 아니라 취향을 팔다", 26, FG, True)],
           [("서점을 판 것이 아니다. ", 15, FG2, False),("취향과 체류", 15, FG, True),("를 판매했다. '숲속의 도서관'이라는 라이프스타일 제안 공간.", 15, FG2, False)])
chain_row(s, [([[("책",13,FG,False)]],"start"),([[("취향",13,FG2,False)]],"normal"),([[("체류",13,FG2,False)]],"normal"),([[("커뮤니티",13,FG,True)]],"end")], 2.85, h=0.8, total=7.4)
textbox(s, MW, 3.85, 7.4, 0.9, [[("'무엇을 파는 공간인가'가 아니라 '어떤 사람이 머무는 공간인가'로 질문이 바뀌었다.", 14, FG2, False)]], ls=1.4)
kpi_box(s, MW+CW-4.3, 2.7, 4.3, 1.2, "개관", "2011", "", "운영: CCC · 설계: Klein Dytham")
card(s, MW+CW-4.3, 4.05, 4.3, 1.85, None, "타깃은 '프리미엄 에이지'. 책·음악·영화·여행을 한 공간에서 큐레이션해 머무는 시간 자체를 상품화했다.", fill=AMBERBOX, line=AMBER, src="출처: CCC 공식 · Klein Dytham", anchor=MSO_ANCHOR.MIDDLE)
footer(s, "Building → Traffic", "16")

# 17. P12 모리빌딩
s=new_slide(BG)
head_block(s, "P12 · Case · Mori Building",
           [("모리빌딩 — 건물이 아니라 지역을 운영한다", 26, FG, True)],
           [("하나의 건물이 아니라 ", 15, FG2, False),("거리와 지역 전체", 15, FG, True),("를 장기 보유·운영한다. 자산 소유보다 '도시 운영'에 가깝다.", 15, FG2, False)])
chain_row(s, [([[("건물",13,FG,False)]],"start"),([[("거리",13,FG2,False)]],"normal"),([[("지역",13,FG2,False)]],"normal"),([[("도시",13,FG,True)]],"end")], 2.85, h=0.8, total=7.4)
textbox(s, MW, 3.85, 7.4, 1.0, [[("에어리어 매니지먼트",14,FG,True),(" — 개발이 끝이 아니라 그 지역을 수십 년간 운영하며 가치를 관리한다.",14,FG2,False)]], ls=1.4)
kpi_box(s, MW+CW-4.3, 2.7, 4.3, 1.15, "롯폰기힐스", "2003", "", "타운 매니지먼트 도입")
kpi_box(s, MW+CW-4.3, 4.0, 4.3, 1.15, "아자부다이힐스", "2023", "", "약 35년·권리자 300명 협의")
textbox(s, MW+CW-4.3, 5.25, 4.3, 0.3, [[("출처: 모리빌딩 공식", 9, FG3, False)]])
footer(s, "Building → Traffic", "17")

# 18. PART 4
s=new_slide(BG2)
textbox(s, MW, 1.0, CW, 4.0, [[("4", 300, RGBColor(0x26,0x24,0x22), False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
textbox(s, MW, 2.0, CW, 0.4, [[("PART 4", 13, ACCL, False)]], align=PP_ALIGN.CENTER)
textbox(s, MW, 2.7, CW, 1.7, [[("성수동은", 36, FG, False)],[("왜 올랐는가", 36, ACCL, False)]], align=PP_ALIGN.CENTER, ls=1.15)
textbox(s, MW, 5.0, CW, 0.5, [[("국가사업도, 신도시도 아니었는데", 16, FG2, False)]], align=PP_ALIGN.CENTER)
footer(s, "Part 4 — 성수동의 공식", "18")

# 19. P13
s=new_slide(BG)
head_block(s, "P13 · Not an Exception",
           [("성수동은 특수사례가 아니다", 28, FG, True)],
           [("국가사업도 신도시도 아니었다. 대형 개발 없이 서울 최고 수준의 가치 상승을 만들었다. 우연이 아니라 ", 15, FG2, False),("공식", 15, FG, True),("이다.", 15, FG2, False)])
hero_box(s, MW, 2.9, 5.0, 2.9, "+110", "%", "성수 외국인 관광객 증가\n(2025 1분기, 전년 동기 대비)\nCBRE 코리아", None, vcolor=GREEN, fill=GREENBOX, line=GREEN)
card(s, MW+5.4, 2.9, CW-5.4, 2.9, "평당 월매출이 강남을 추월", "국내 패션 브랜드 기준 — 성수 634만원 vs 강남 167만원(약 4배). 글로벌 스포츠웨어도 성수 876 vs 강남 509만원.", None, fill=BG3, line=LINE2, src="출처: CBRE 코리아 리테일 인사이트, 2025~26", anchor=MSO_ANCHOR.MIDDLE)
footer(s, "Building → Traffic", "19")

# 20. P14
s=new_slide(BG)
head_block(s, "P14 · The First Movers",
           [("가장 먼저 온 사람들", 28, FG, True)],
           [("성수동에 가장 먼저 온 것은 자본이 아니었다. ", 15, FG2, False),("창작자", 15, FG, True),("가 먼저 왔다. 디자이너, 작은 브랜드, 기획자들.", 15, FG2, False)])
cw3=(CW-0.8)/3
for i,(t,b) in enumerate([("디자이너","싼 임대료의 옛 공장·창고를 스튜디오로. 작업 공간을 찾아 모였다."),("작은 브랜드","쇼룸과 플래그십을 실험. '여기서만 볼 수 있는 것'을 만들었다."),("기획자","전시·팝업·콘텐츠를 기획. 거리에 일어날 사건을 설계했다.")]):
    card(s, MW+(cw3+0.4)*i, 2.75, cw3, 1.8, t, b)
chain_row(s, [([[("디자이너 + 브랜드 + 기획자",13,FG2,False)]],"normal"),([[("생산자 밀도",13,FG,True)]],"end")], 4.85, h=0.75, total=8, start_x=MW+(CW-8)/2)
footer(s, "Building → Traffic", "20")

# 21. P15
s=new_slide(BG)
head_block(s, "P15 · Density → Content",
           [("생산자가 모이면 ", 28, FG, True),("콘텐츠", 28, ACCL, True),("가 생긴다", 28, FG, True)],
           [("생산자 밀도가 임계점을 넘으면 브랜드·전시·팝업이 생기고, 거리에 '문화'가 만들어진다. 그리고 문화는 SNS를 타고 퍼진다.", 15, FG2, False)])
chain_row(s, [([[("생산자 밀도",13,FG,False)]],"start"),([[("브랜드·전시·팝업",13,FG2,False)]],"normal"),([[("거리의 문화",13,FG2,False)]],"normal"),([[("SNS 확산",13,FG,True)]],"end")], 3.3, h=0.95)
textbox(s, MW, 4.9, CW, 0.9, [[("콘텐츠는 '광고'가 아니라 '사건'이다. 사람들은 광고를 보러 오지 않지만, 사건은 보러 온다.", 15, FG2, False)]], align=PP_ALIGN.CENTER, ls=1.4)
footer(s, "Building → Traffic", "21")

# 22. P16
s=new_slide(BG2)
textbox(s, MW, 1.5, CW, 0.4, [[("P16 · CONTENT PULLS PEOPLE", 11, FG3, False)]], align=PP_ALIGN.CENTER)
textbox(s, MW, 2.1, CW, 0.6, [[("사람들은 ", 26, FG, False),("콘텐츠", 26, ACC, True),("를 따라온다", 26, FG, False)]], align=PP_ALIGN.CENTER)
textbox(s, MW, 3.2, CW, 0.8, [[("콘텐츠", 36, FG, False),("   →   ", 36, ACC, False),("트래픽", 36, FG, False)]], align=PP_ALIGN.CENTER)
textbox(s, MW, 4.7, CW, 0.9, [[("사람들은 성수를 찾아간 것이 아니다. 그곳에서 벌어지는 ", 17, FG2, False),("경험", 17, FG, True),("을 찾아갔다.", 17, FG2, False)]], align=PP_ALIGN.CENTER, ls=1.3)
footer(s, "Building → Traffic", "22")

# 23. P17
s=new_slide(BG)
head_block(s, "P17 · Capital Follows People",
           [("자본은 사람을 따라온다", 28, FG, True)],
           [("브랜드와 자본은 사람들이 모인 곳으로 이동한다. 트래픽이 먼저고 투자가 나중이다. ", 15, FG2, False),("순서가 뒤집히지 않는다.", 15, FG, True)])
chain_row(s, [([[("트래픽",13,FG,False)]],"start"),([[("브랜드 입점",13,FG2,False)]],"normal"),([[("투자 유입",13,FG,True)]],"end")], 3.4, h=0.95, total=9, start_x=MW+(CW-9)/2)
textbox(s, MW, 5.0, CW, 0.9, [[("대기업 플래그십과 투자금은 '뜰 곳'을 미리 아는 게 아니다. 사람이 모인 것을 ",15,FG2,False),("확인하고",15,FG,True),(" 들어온다.",15,FG2,False)]], align=PP_ALIGN.CENTER, ls=1.4)
footer(s, "Building → Traffic", "23")

# 24. P18 공식
s=new_slide(BG2)
textbox(s, MW, 0.7, CW, 0.35, [[("P18 · THE SEONGSU FORMULA", 11, FG3, False)]], align=PP_ALIGN.CENTER)
textbox(s, MW, 1.1, CW, 0.7, [[("성수동 공식", 28, FG, True)]], align=PP_ALIGN.CENTER)
textbox(s, MW, 1.85, CW, 0.5, [[("성수동은 우연이 아니다. ", 15, FG2, False),("가치 창조 공식", 15, FG, True),("이 순서대로 작동한 결과다.", 15, FG2, False)]], align=PP_ALIGN.CENTER)
formula_row(s, [("01","생산자",False),("02","콘텐츠",False),("03","트래픽",False),("04","브랜드",False),("05","자본",False),("06","부동산 가치",True)], 3.1, h=1.1)
textbox(s, MW, 4.7, CW, 0.9, [[("이 공식의 출발점은 건물도 자본도 아닌 ",15,FG2,False),("생산자",15,FG,True),("다. 종착점인 부동산 가치는 가장 마지막에 따라온다.",15,FG2,False)]], align=PP_ALIGN.CENTER, ls=1.4)
footer(s, "Building → Traffic", "24")

# 25. PART 5
s=new_slide(BG2)
textbox(s, MW, 1.0, CW, 4.0, [[("5", 300, RGBColor(0x26,0x24,0x22), False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
textbox(s, MW, 2.0, CW, 0.4, [[("PART 5", 13, ACCL, False)]], align=PP_ALIGN.CENTER)
textbox(s, MW, 2.7, CW, 1.7, [[("왜 어떤 곳은", 36, FG, False)],[("죽는가", 36, RGBColor(0xE8,0x9B,0x85), False)]], align=PP_ALIGN.CENTER, ls=1.15)
textbox(s, MW, 5.0, CW, 0.5, [[("공식이 멈추면, 거리도 멈춘다", 16, FG2, False)]], align=PP_ALIGN.CENTER)
footer(s, "Part 5 — 쇠퇴의 메커니즘", "25")

# 26. P19 가로수길
s=new_slide(BG)
head_block(s, "P19 · Case · Garosu-gil",
           [("가로수길 — 성공 이후, 이유를 잃다", 27, FG, True)],
           [("뜬 뒤 임대료가 치솟자 창작자가 떠났고, 콘텐츠가 사라졌다. ", 15, FG2, False),("공식이 역방향으로", 15, FG, True),(" 돌기 시작했다.", 15, FG2, False)])
flow_col(s, [([[("창작자 감소",12.5,FG2,False)]],"red"),([[("콘텐츠 감소",12.5,FG2,False)]],"red"),([[("트래픽 감소",12.5,FG2,False)]],"red"),([[("가치 감소",12.5,FG,True)]],"red")], cx=MW+3.6, top=2.7, bw=4.2, bh=0.56, gap=0.26, color=RED)
hero_box(s, MW+CW-4.6, 2.7, 4.6, 3.1, "41.6", "%", "가로수길 공실률 (2025 1분기)\n전년 동기 대비 +0.4%p · 타 상권 대비 최대 약 12배.\n임대료는 오히려 우상향 — 공실과 임대료가 함께 오르는 역설.", "출처: 국민일보·뉴스1 (원자료 한국부동산원)", vcolor=RED, fill=REDBOX, line=RED)
footer(s, "Building → Traffic", "26")

# 27. P20 을지로
s=new_slide(BG)
head_block(s, "P20 · Case · Euljiro",
           [("을지로 — 콘텐츠는 있었지만", 28, FG, True)],
           [("을지로에는 콘텐츠가 있었다. 그러나 그것이 ", 15, FG2, False),("확장 가능한 플랫폼", 15, FG, True),("으로 이어지지 못했다. 트래픽이 자본과 연결되지 않았다.", 15, FG2, False)])
chain_row(s, [([[("콘텐츠",13,FG,False)]],"start"),([[("트래픽",13,FG2,False)]],"normal"),([[("자본 연결 ✕",13,FG3,False)]],"normal")], 3.4, h=0.95, total=9, start_x=MW+(CW-9)/2)
textbox(s, MW, 5.0, CW, 0.9, [[("콘텐츠와 트래픽이 있어도 ",15,FG2,False),("지속 가능한 구조",15,FG,True),("로 묶지 못하면 거리는 한순간의 유행으로 끝난다.",15,FG2,False)]], align=PP_ALIGN.CENTER, ls=1.4)
footer(s, "Building → Traffic", "27")

# 28. P21 임대료
s=new_slide(BG2)
head_block(s, "P21 · Rent is a Thermometer",
           [("임대료는 ", 28, FG, True),("체온계", 28, ACC, True),("다", 28, FG, True)],
           [("임대료는 원인이 아니라 결과다. 거리의 건강 상태를 보여주는 지표일 뿐, 임대료를 올린다고 거리가 건강해지지 않는다.", 15, FG2, False)])
chain_row(s, [([[("콘텐츠",13,FG,False)]],"start"),([[("트래픽",13,FG2,False)]],"normal"),([[("가치",13,FG2,False)]],"normal"),([[("임대료",13,FG,True)]],"end")], 3.4, h=0.95)
textbox(s, MW, 5.0, CW, 0.9, [[("체온계 숫자를 손으로 올린다고 열이 나지 않는다. 임대료만 먼저 올리면 ",15,FG2,False),("앞 단계가 무너진다(가로수길).",15,FG,True)]], align=PP_ALIGN.CENTER, ls=1.4)
footer(s, "Building → Traffic", "28")

# 29. PART 6
s=new_slide(BG2)
textbox(s, MW, 1.0, CW, 4.0, [[("6", 300, RGBColor(0x26,0x24,0x22), False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
textbox(s, MW, 2.0, CW, 0.4, [[("PART 6", 13, ACCL, False)]], align=PP_ALIGN.CENTER)
textbox(s, MW, 2.7, CW, 1.7, [[("트래픽은", 36, FG, False)],[("관리 가능한가", 36, ACCL, False)]], align=PP_ALIGN.CENTER, ls=1.15)
textbox(s, MW, 5.0, CW, 0.5, [[("우연인가, 아니면 설계할 수 있는가", 16, FG2, False)]], align=PP_ALIGN.CENTER)
footer(s, "Part 6 — 트래픽의 측정과 설계", "29")

# 30. P22 vs
s=new_slide(BG)
head_block(s, "P22 · The Key Question",
           [("가장 중요한 질문", 28, FG, True)],
           [("지금까지의 모든 사례가 향하는 단 하나의 질문 — 트래픽은 ", 15, FG2, False),("우연", 15, FG, True),("인가, 아니면 ", 15, FG2, False),("설계와 운영", 15, FG, True),("의 결과인가.", 15, FG2, False)])
y=2.85; h=3.4; halfw=(CW-1.0)/2
lsp=box(s, MW, y, halfw, h, REDBOX, RED); box_text(lsp,[[("OPTION A",11,RED,False)],[("우연",26,FG,False)],[("운이 좋아 떴고, 운이 다하면 진다 — 통제 불가능한 영역",12,FG2,False)]],ls=1.3)
rsp=box(s, MW+halfw+1.0, y, halfw, h, ACCBOX, ACC); box_text(rsp,[[("OPTION B",11,ACCL,False)],[("운영",26,ACCL,False)],[("측정하고 설계하고 키울 수 있는 것 — 개발사의 새로운 역량",12,FG2,False)]],ls=1.3)
c=box(s, MW+halfw+0.34, y+h/2-0.36, 0.72, 0.72, BG2, LINE, shape=MSO_SHAPE.OVAL); box_text(c,[[("VS",12,FG3,False)]])
footer(s, "Building → Traffic", "30")

# 31. P23 측정 가능
s=new_slide(BG)
head_block(s, "P23 · Traffic is Measurable",
           [("트래픽은 측정 가능하다", 28, FG, True)],
           [("좋은 공간은 단순 방문자 수가 아니라 ", 15, FG2, False),("체류시간과 재방문율", 15, FG, True),("로 판단된다. 측정할 수 있다면, 관리할 수 있다.", 15, FG2, False)])
chain_row(s, [([[("방문",13,FG,False)]],"start"),([[("체류",13,FG2,False)]],"normal"),([[("재방문",13,FG2,False)]],"normal"),([[("전환",13,FG,True)]],"end")], 2.9, h=0.8)
kw=(CW-0.9)/4
for i,(l,v,sb,col) in enumerate([("VISIT","방문","얼마나 오는가",FG),("DWELL","체류","얼마나 머무는가",FG),("RETURN","재방문","다시 오는가",FG),("CONVERT","전환","가치로 이어지는가",ACCL)]):
    sp=box(s, MW+(kw+0.3)*i, 4.0, kw, 1.5, BG3, LINE2)
    box_text(sp,[[(l,10,FG3,False)],[(v,18,col,False)],[(sb,10,FG3,False)]],ls=1.2)
footer(s, "Building → Traffic", "31")

# 32. P24 선행지표
s=new_slide(BG)
head_block(s, "P24 · Leading Indicators",
           [("미래 가치의 ", 28, FG, True),("선행지표", 28, ACCL, True)],
           [("가격보다 먼저 움직이는 신호들이 있다. 공시지가·임대료가 오른 뒤엔 이미 늦다 — 그 전에 움직이는 것을 봐야 한다.", 15, FG2, False)])
halfw=(CW-0.5)/2
ls=box(s, MW, 2.9, halfw, 2.9, BG3, LINE2)
box_text(ls,[[("후행지표 · 결과",11,FG3,False)],[("· 공시지가",13,FG3,False)],[("· 임대료",13,FG3,False)],[("· 거래가",13,FG3,False)],[("이미 모두가 아는 신호 — 투자의 영역",11,FG3,False)]],align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,ls=1.5)
rs=box(s, MW+halfw+0.5, 2.9, halfw, 2.9, ACCBOX, ACC)
box_text(rs,[[("선행지표 · 원인",11,ACCL,False)],[("· 생산자 밀도",13,FG,False)],[("· 콘텐츠 생성량",13,FG,False)],[("· 체류시간",13,FG,False)],[("· 브랜드 유입",13,FG,False)],[("먼저 움직이는 신호 — 개발의 영역",11,ACCL,False)]],align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,ls=1.4)
footer(s, "Building → Traffic", "32")

# 33. P25 다음 성수
s=new_slide(BG)
head_block(s, "P25 · Finding the Next Seongsu",
           [("미래의 성수를 찾는 방법", 28, FG, True)],
           [("이미 유명한 곳을 사는 것은 ", 15, FG2, False),("투자", 15, FG, True),("다. 새로운 생산자가 모이는 곳을 먼저 찾는 것은 ", 15, FG2, False),("개발", 15, FG, True),("이다.", 15, FG2, False)])
y=3.0; bw=4.6; h=1.4
l=box(s, MW+1.0, y, bw, h, BG3, LINE); box_text(l,[[("현재의 성수",15,FG,True)],[("이미 비싸다 · 투자",11,FG3,False)]],ls=1.3)
arrow(s, MW+1.0+bw+0.2, y+h/2-0.25, 0.8, "→", ACC, 22)
r=box(s, MW+1.0+bw+1.2, y, bw, h, ENDBOX, ACC); box_text(r,[[("다음 성수  ",15,FG,True),("?",20,ACC,True)],[("생산자가 모이는 곳 · 개발",11,ACCL,False)]],ls=1.3)
textbox(s, MW, 5.1, CW, 0.9, [[("선행지표(생산자 밀도·콘텐츠·체류)를 읽을 수 있다면, 다음 성수는 ",15,FG2,False),("예측이 아니라 발견",15,FG,True),("의 대상이 된다.",15,FG2,False)]], align=PP_ALIGN.CENTER, ls=1.4)
footer(s, "Building → Traffic", "33")

# 34. PART 7
s=new_slide(BG2)
textbox(s, MW, 1.0, CW, 4.0, [[("7", 300, RGBColor(0x26,0x24,0x22), False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
textbox(s, MW, 2.0, CW, 0.4, [[("PART 7", 13, ACCL, False)]], align=PP_ALIGN.CENTER)
textbox(s, MW, 2.7, CW, 1.7, [[("새로운 개발사의", 36, FG, False)],[("등장", 36, ACCL, False)]], align=PP_ALIGN.CENTER, ls=1.15)
textbox(s, MW, 5.0, CW, 0.5, [[("Builder에서 Operator로", 16, FG2, False)]], align=PP_ALIGN.CENTER)
footer(s, "Part 7 — Traffic Developer", "34")

# 35. P26 vs
s=new_slide(BG)
head_block(s, "P26 · Builder → Operator",
           [("Builder에서 ", 28, FG, True),("Operator", 28, ACCL, True),("로", 28, FG, True)],
           [("건설 경쟁력보다 운영 경쟁력이 중요해진다. 잘 짓는 회사가 아니라, 사람이 계속 찾아오게 ", 15, FG2, False),("운영하는 회사", 15, FG, True),("가 이긴다.", 15, FG2, False)])
y=2.85; h=3.4; halfw=(CW-1.0)/2
ls=box(s, MW, y, halfw, h, BG3, LINE2); box_text(ls,[[("과거 · BUILDER",11,FG3,False)],[("잘 짓는 경쟁",15,FG,True)],[("· 토지 확보 → 인허가 → 시공",12,FG2,False)],[("· 분양·매각으로 사업 종료",12,FG2,False)],[("· 핵심 역량: 건설·금융",12,FG2,False)]],align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,ls=1.45)
rs=box(s, MW+halfw+1.0, y, halfw, h, ACCBOX, ACC); box_text(rs,[[("미래 · OPERATOR",11,ACCL,False)],[("계속 오게 만드는 경쟁",15,FG,True)],[("· 콘텐츠·브랜드·커뮤니티 설계",12,FG,False)],[("· 준공 이후가 사업의 시작",12,FG,False)],[("· 핵심 역량: 운영·기획",12,FG,False)]],align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,ls=1.45)
arrow(s, MW+halfw+0.32, y+h/2-0.3, 0.72, "→", ACC, 22)
footer(s, "Building → Traffic", "35")

# 36. P27 스타벅스
s=new_slide(BG)
head_block(s, "P27 · Case · Starbucks",
           [("카페도 같다", 28, FG, True)],
           [("사람들은 커피를 사는 것이 아니라 ", 15, FG2, False),("머물 이유", 15, FG, True),("를 산다. 스타벅스가 판 것은 음료가 아니라 '제3의 공간'이었다.", 15, FG2, False)])
chain_row(s, [([[("커피",13,FG,False)]],"start"),([[("체류",13,FG2,False)]],"normal"),([[("커뮤니티",13,FG,True)]],"end")], 2.95, h=0.85, total=7.0)
card(s, MW+7.4, 2.7, CW-7.4, 2.9, None, "'제3의 공간(Third Place)' — 집(제1)과 직장(제2) 사이, 머무는 공간. 사회학자 레이 올든버그가 제시하고 스타벅스가 브랜드 비전으로 차용했다.", fill=AMBERBOX, line=AMBER, src="출처: Ray Oldenburg, 『The Great Good Place』(1989)", anchor=MSO_ANCHOR.MIDDLE)
footer(s, "Building → Traffic", "36")

# 37. P28 애플
s=new_slide(BG)
head_block(s, "P28 · Case · Apple",
           [("리테일도 같다", 28, FG, True)],
           [("온라인이 상품 판매를 대체했다. 오프라인에 남은 것은 ", 15, FG2, False),("경험", 15, FG, True),("이다. 애플스토어는 물건을 파는 곳이 아니라 브랜드를 체험하는 곳이다.", 15, FG2, False)])
chain_row(s, [([[("상품",13,FG,False)]],"start"),([[("경험",13,FG2,False)]],"normal"),([[("팬덤",13,FG,True)]],"end")], 2.95, h=0.85, total=7.0)
card(s, MW+7.4, 2.7, CW-7.4, 2.9, None, "애플은 매장을 '타운스퀘어(town square)'로 재정의했다. \"매장이 우리가 만드는 가장 큰 제품\"이라며 판매보다 워크숍·커뮤니티·체험에 무게를 뒀다.", fill=BLUEBOX, line=BLUE, src="출처: Angela Ahrendts, Fortune (2016)", anchor=MSO_ANCHOR.MIDDLE)
footer(s, "Building → Traffic", "37")

# 38. P29 Traffic Developer
s=new_slide(BG2)
textbox(s, MW, 1.1, CW, 0.35, [[("P29 · TRAFFIC DEVELOPER", 11, FG3, False)]], align=PP_ALIGN.CENTER)
textbox(s, MW, 1.5, CW, 0.7, [[("Traffic Developer", 30, FG, True)]], align=PP_ALIGN.CENTER)
textbox(s, MW, 2.3, CW, 0.6, [[("미래의 개발사는 건물을 만드는 사람이 아니다. 사람들이 ", 15, FG2, False),("찾아올 이유", 15, FG, True),("를 만드는 사람이다.", 15, FG2, False)]], align=PP_ALIGN.CENTER)
chain_row(s, [([[("브랜드",12.5,FG2,False)]],"normal"),([[("콘텐츠",12.5,FG2,False)]],"normal"),([[("커뮤니티",12.5,FG2,False)]],"normal"),([[("운영",12.5,FG2,False)]],"normal"),([[("Traffic Developer",12.5,FG,True)]],"end")], 3.5, h=0.85, arr="→")
textbox(s, MW, 4.9, CW, 0.9, [[("건물 짓는 역량 위에, 사람을 모으고 머물게 하는 역량을 더한다. 이것이 다음 10년 개발사의 정체성이다.", 15, FG2, False)]], align=PP_ALIGN.CENTER, ls=1.4)
footer(s, "Building → Traffic", "38")

# 39. P30 결론
s=new_slide(BG2)
textbox(s, MW, 0.85, CW, 0.35, [[("P30 · CONCLUSION", 11, FG3, False)]], align=PP_ALIGN.CENTER)
textbox(s, MW, 1.25, CW, 0.7, [[("사람들은 건물을 찾아가지 않는다", 28, FG, True)]], align=PP_ALIGN.CENTER)
textbox(s, MW, 2.05, CW, 0.7, [[("과거의 가치는 토지에서, 산업화 시대 가치는 건물에서 나왔다. 앞으로의 가치는 사람들이 ", 14, FG2, False),("반복적으로 찾아올 이유", 14, FG, True),("에서 나온다.", 14, FG2, False)]], align=PP_ALIGN.CENTER, ls=1.3)
# 3단 공식
steps=[("PAST","토지 가치",FG3,BG3,LINE2),("INDUSTRIAL","건물 가치",FG2,BG3,LINE2),("NEXT 10 YEARS","트래픽 가치",ACCL,ENDBOX,ACC)]
bw=3.4; aw=0.5; tx=MW+(CW-(bw*3+aw*2))/2; y=3.3
for i,(n,t,tc,fl,ln) in enumerate(steps):
    sp=box(s, tx, y, bw, 1.2, fl, ln); box_text(sp,[[(n,9,FG3,False)],[(t,16,tc,True)]],ls=1.15)
    tx+=bw
    if i<2: arrow(s, tx, y, aw, "→", ACC, 18); tx+=aw
textbox(s, MW, 4.9, CW, 0.9, [[("토지 가치 · 건물 가치 · 트래픽 가치", 17, FG2, False)],[("＝ 미래 자산 가치", 18, ACCL, True)]], align=PP_ALIGN.CENTER, ls=1.4)
footer(s, "Building → Traffic", "39")

# 40. 엔딩
s=new_slide(BG2)
textbox(s, MW, 1.4, CW, 0.35, [[("THE NEXT 10 YEARS", 11, FG3, False)]], align=PP_ALIGN.CENTER)
textbox(s, MW, 2.0, CW, 2.2, [[("부동산 개발의 다음 10년은", 24, FG, False)],
                              [("건물을 얼마나 잘 짓느냐의 경쟁이 아니다.", 24, FG, False)],
                              [("사람들이 얼마나 ", 24, FG, False),("오래, 자주, 반복적으로", 24, ACCL, True)],
                              [("찾아오게 만드느냐의 경쟁이다.", 24, FG, False)]], align=PP_ALIGN.CENTER, ls=1.35)
box(s, 13.333/2-0.3, 4.55, 0.6, 0.03, ACC, None, shape=MSO_SHAPE.RECTANGLE, radius=0)
textbox(s, MW, 4.8, CW, 0.9, [[("프로젝트렌트는 공간을 기획하는 회사가 아니라,", 15, FG2, False)],
                              [("미래 자산가치를 설계하는 ", 15, FG2, False),("트래픽 디벨로퍼", 15, FG, True),("입니다.", 15, FG2, False)]], align=PP_ALIGN.CENTER, ls=1.35)
textbox(s, MW, 6.0, CW, 0.3, [[("BY PROJECT RENT", 12, FG3, False)]], align=PP_ALIGN.CENTER)
footer(s, "Building → Traffic", "by Project Rent · 2026")

out="/Users/choi_ai/do-better-workspace/10-projects/33-traffic-developer-lecture/traffic-developer-lecture.pptx"
prs.save(out)
print("saved:", out, "slides:", len(prs.slides._sldIdLst))
