# -*- coding: utf-8 -*-
"""PROJECT RENT IR V10 — Physical Ad-Tech Network → 편집 가능한 PPTX (텍스트/도형/표, 이미지 캡처 없음)"""
import math
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def C(h): return RGBColor.from_string(h)
BG=C('1A1A1A'); BG2=C('111111'); BG3=C('222222')
FG=C('F5F0EB'); FG2=C('C9C2BA'); FG3=C('7A7570')
ACCENT=C('6C5CE7'); ACCENTL=C('A29BFE')
LINE=C('333333'); LINE2=C('2A2A2A'); AMBER=C('D9A34B'); GREEN=C('8FA37E'); RUST=C('B86C4A'); RED=C('E17055'); BLUE=C('6F8AA3')
ACARD=C('221E33')
FONT='Pretendard'

prs=Presentation()
prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]
LM=0.9; TOTAL=13.333; RM=TOTAL-LM; CW=RM-LM

def slide(bg=BG):
    s=prs.slides.add_slide(BLANK)
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb=bg; r.line.fill.background(); r.shadow.inherit=False
    return s

def run(p,text,size,color,bold=False,font=FONT,spacing=None):
    rn=p.add_run(); rn.text=text
    f=rn.font; f.size=Pt(size); f.color.rgb=color; f.bold=bold; f.name=font
    from pptx.oxml.ns import qn
    rPr=rn._r.get_or_add_rPr(); ea=rPr.makeelement(qn('a:ea'),{}); ea.set('typeface',font); rPr.append(ea)
    return rn

def box(s,l,t,w,h,anchor=MSO_ANCHOR.TOP):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    return tf

def para(tf,first,sb=0,ls=None,align=None):
    p=tf.paragraphs[0] if first else tf.add_paragraph()
    if sb: p.space_before=Pt(sb)
    if ls: p.line_spacing=ls
    if align is not None: p.alignment=align
    return p

def lines(tf,text,size,color,bold=False,ls=1.2,first=True,sb=0,align=None):
    for i,ln in enumerate(text.split('\n')):
        p=para(tf, first and i==0, sb=(sb if i==0 else 0), ls=ls, align=align)
        run(p,ln,size,color,bold)
    return tf

def rect(s,l,t,w,h,fill=None,ln=None,lw=0.75,rounded=False,radius=0.06):
    shp=MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sh=s.shapes.add_shape(shp,Inches(l),Inches(t),Inches(w),Inches(h))
    sh.shadow.inherit=False
    if fill is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb=fill
    if ln is None: sh.line.fill.background()
    else: sh.line.color.rgb=ln; sh.line.width=Pt(lw)
    if rounded:
        try: sh.adjustments[0]=radius
        except: pass
    return sh

def hline(s,t,l=LM,r=RM,col=LINE2):
    rect(s,l,t,r-l,0.012,fill=col)

# 헤더: eyebrow 좌/페이지 우 + 라인 + bold 제목 + 서브카피
def header(s,eyebrow,page,title,sub=None,tsize=28):
    tf=box(s,LM,0.52,9.5,0.4); run(para(tf,True),eyebrow,11,FG3,True)
    tf=box(s,RM-1.2,0.52,1.2,0.4); run(para(tf,True,align=PP_ALIGN.RIGHT),page,11,FG3,True)
    hline(s,0.92)
    tlines=title.count('\n')+1
    tf=box(s,LM,1.06,CW,tlines*tsize*1.12/72+0.2); lines(tf,title,tsize,FG,True,ls=1.08)
    y=1.06+tlines*tsize*1.12/72+0.14
    if sub:
        tf=box(s,LM,y,11.0,1.0);
        for i,ln in enumerate(sub.split('\n')):
            p=para(tf,i==0,ls=1.32);
            # bold 부분 처리: '**'로 감싼 텍스트는 FG bold
            seg=ln.split('**')
            for j,t in enumerate(seg):
                if t=='':continue
                run(p,t,13.5,(FG if j%2==1 else FG2),(j%2==1))
        y+=(sub.count('\n')+1)*13.5*1.32/72+0.22
    return y

def footer(s,pg):
    hline(s,6.92)
    tf=box(s,LM,6.98,8,0.3); run(para(tf,True),'PHYSICAL AD-TECH NETWORK · BY PROJECT RENT',8.5,FG3,False)
    tf=box(s,RM-2,6.98,2,0.3); p=para(tf,True,align=PP_ALIGN.RIGHT); run(p,pg+' / 20',9,ACCENT,True)

def cols(n,top,h,gap=0.28,l=LM,total_w=CW):
    cw=(total_w-(n-1)*gap)/n
    return [(l+i*(cw+gap),top,cw,h) for i in range(n)]

def card(s,l,t,w,h,label=None,lcol=ACCENTL,title=None,tsize=14,body=None,bcol=FG3,bsize=10.5,fill=BG3,accent_top=False):
    sh=rect(s,l,t,w,h,fill=fill,ln=LINE2,rounded=True,radius=0.045)
    if accent_top:
        rect(s,l,t,w,0.035,fill=ACCENT)
    tf=sh.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.TOP
    tf.margin_left=Inches(0.18); tf.margin_right=Inches(0.16); tf.margin_top=Inches(0.16); tf.margin_bottom=Inches(0.12)
    first=True
    if label: lines(tf,label,9.5,lcol,True,first=first); first=False
    if title:
        for i,ln in enumerate(title.split('\n')):
            p=para(tf,first and i==0,sb=(7 if first and i==0 and label else 0),ls=1.1); run(p,ln,tsize,FG,True)
        first=False
    if body:
        for i,ln in enumerate(body.split('\n')):
            p=para(tf,first and i==0,sb=(8 if i==0 and not first else 0),ls=1.3);
            seg=ln.split('**')
            for j,t in enumerate(seg):
                if t=='':continue
                run(p,t,bsize,(ACCENTL if j%2==1 else bcol),(j%2==1))
        first=False
    return sh

def kpi(s,l,t,w,h,label,value,unit='',sub=None,vcol=ACCENT):
    rect(s,l,t,w,h,fill=BG3,ln=LINE2,rounded=True,radius=0.06)
    tf=box(s,l,t+0.16,w,0.3,anchor=MSO_ANCHOR.TOP); p=para(tf,True,align=PP_ALIGN.CENTER); run(p,label,9,FG3,True)
    tf=box(s,l,t+h*0.34,w,0.6,anchor=MSO_ANCHOR.TOP); p=para(tf,True,align=PP_ALIGN.CENTER)
    run(p,value,30,vcol,False);
    if unit: run(p,unit,15,FG2,False)
    if sub:
        tf=box(s,l,t+h-0.42,w,0.3,anchor=MSO_ANCHOR.TOP); p=para(tf,True,align=PP_ALIGN.CENTER); run(p,sub,8.5,FG3)

# ir-flow: 가로 박스 n개
def flow(s,t,h,steps):
    n=len(steps); gap=0.0; cw=CW/n
    for i,(num,title,desc) in enumerate(steps):
        l=LM+i*cw
        rect(s,l,t,cw-0.04,h,fill=BG3,ln=LINE2)
        tf=box(s,l+0.18,t+h*0.30,cw-0.32,h*0.6,anchor=MSO_ANCHOR.TOP)
        p=para(tf,True); run(p,num,18,ACCENT,False)
        p=para(tf,False,sb=6,ls=1.1); run(p,title,12.5,FG,True)
        p=para(tf,False,sb=5,ls=1.25); run(p,desc,9,FG3)
        if i<n-1:
            tf=box(s,l+cw-0.18,t+h/2-0.15,0.3,0.3,anchor=MSO_ANCHOR.MIDDLE); run(para(tf,True,align=PP_ALIGN.CENTER),'→',14,ACCENT,True)

def callout(s,t,text,h=0.62):
    rect(s,LM,t,CW,h,fill=C('1E1B2E'),ln=None)
    rect(s,LM,t,0.04,h,fill=ACCENT)
    tf=box(s,LM+0.28,t,CW-0.5,h,anchor=MSO_ANCHOR.MIDDLE); p=para(tf,True,ls=1.3)
    seg=text.split('**')
    for j,t2 in enumerate(seg):
        if t2=='':continue
        run(p,t2,12.5,(ACCENTL if j%2==1 else FG2),(j%2==1))

def note(s,t,text):
    tf=box(s,LM,t,CW,0.4); p=para(tf,True,ls=1.3)
    seg=text.split('**')
    for j,t2 in enumerate(seg):
        if t2=='':continue
        run(p,t2,8.5,(FG2 if j%2==1 else FG3),(j%2==1))

# ═══════════════ P01 COVER ═══════════════
s=slide(BG2)
tf=box(s,LM,1.5,11,0.4); run(para(tf,True),'PROJECT RENT · IR DECK 2026 · CONFIDENTIAL',11,FG3,True)
tf=box(s,LM,2.0,11.5,1.7)
p=para(tf,True,ls=1.04); run(p,'Building the ',40,FG,True)
p=para(tf,False,ls=1.04); run(p,'Physical Ad-Tech',40,ACCENT,True); run(p,' Network',40,FG,True)
tf=box(s,LM,3.85,11,0.4); p=para(tf,True); run(p,'광고 · 행동 · 구매 데이터를 연결하는 ',15,FG2); run(p,'오프라인 행동 데이터 네트워크',15,FG,True)
rect(s,LM,4.45,0.03,1.5,fill=LINE)
tf=box(s,LM+0.28,4.5,11,1.5)
p=para(tf,True,ls=1.4); run(p,'Google은 사람들이 ',16,FG2); run(p,'무엇을 검색하는지',16,FG3); run(p,' 이해했다.',16,FG2)
p=para(tf,False,ls=1.4); run(p,'Amazon은 사람들이 ',16,FG2); run(p,'무엇을 구매하는지',16,FG3); run(p,' 이해했다.',16,FG2)
p=para(tf,False,ls=1.4); run(p,'Project Rent는 사람들이 ',16,FG,True); run(p,'왜 움직이는지',16,ACCENT,True); run(p,'를 이해한다.',16,FG,True)
tf=box(s,LM,6.4,11.5,0.4); run(para(tf,True),'PHYSICAL AD-TECH NETWORK   ·   HIGH RESOLUTION CONSUMER DATA   ·   BY PROJECT RENT',10,FG3,True)

# ═══════════════ P02 PROBLEM ═══════════════
s=slide()
header(s,'PART 1 · OPENING — PROBLEM','02','Advertising is Broken',
 '**광고는 넘치지만, 사람은 이해되지 않는다.** 쿠키가 닫히고 비용이 오르며, 온라인 데이터의 정교함이 구조적으로 하락하고 있다 — 광고 산업 전체가 흔들리는 중이다.')
cs=cols(4,3.7,2.4)
data=[('01 · 데이터','쿠키리스 가속','Safari·Firefox 이미 차단 — **웹 절반이 쿠키리스**. iOS ATT로 앱 추적 제한.'),
      ('02 · 비용','CAC 상승','이커머스 CAC **’23→’25 40~60%↑**. 유료 트래픽만으론 성장 정체.'),
      ('03 · 효율','Meta·Google 저하','Meta CPM 2020 대비 **+89%**, Google CPC ’25 **+33.7%**.'),
      ('04 · 한계','온라인 데이터의 한계','클릭·노출까지만 안다. **‘왜’는 모른다.**')]
for (l,t,w,h),(lb,ti,bd) in zip(cs,data):
    card(s,l,t,w,h,label=lb,title=ti,tsize=15,body=bd,accent_top=True)
note(s,6.3,'**출처:** 이커머스 CAC·CPM — Shopify·industry benchmarks(2024~2025). 쿠키 — Safari·Firefox 차단 + Apple ATT. ※ 구글 Chrome은 2024.7 서드파티 쿠키 폐지 철회.')
footer(s,'02')

# ═══════════════ P03 HIDDEN OPPORTUNITY ═══════════════
s=slide()
header(s,'PART 1 · OPENING — HIDDEN OPPORTUNITY','03','80% of Consumption Happens Offline',
 '**소비의 대부분은 오프라인에서 일어나지만, 데이터의 대부분은 온라인에서 수집된다.** 가장 큰 소비 현장이 가장 측정되지 않는다 — 거대한 공백이 비어 있다.')
rect(s,LM,3.8,5.0,2.5,fill=BG2,ln=LINE2,rounded=True,radius=0.04)
tf=box(s,LM+0.4,4.4,2.2,1.0); p=para(tf,True); run(p,'80',46,ACCENT); run(p,'%',20,FG2)
tf=box(s,LM+0.4,5.5,2.2,0.4); run(para(tf,True),'소비 지출이 오프라인에서',9.5,FG3)
tf=box(s,LM+2.7,4.4,2.2,1.0); p=para(tf,True); run(p,'~0',46,RED); run(p,'%',20,FG2)
tf=box(s,LM+2.7,5.5,2.2,0.4); run(para(tf,True),'행동 단위 측정 가능 비율',9.5,FG3)
tf=box(s,6.3,3.9,5.1,0.4); run(para(tf,True),'측정되지 않는 시장',11,FG3,True)
bl=[('국내 오프라인 마케팅 10조+',' — OOH 4.3조 + 카페·리테일 6조'),
    ('국내 디지털 광고 11.3조',' (’25, KOBACO) — 측정은 여기 편중'),
    ('이 거대한 소비 현장이 ‘감’으로',' 집행된다 — 최적화 불가')]
tf=box(s,6.3,4.4,5.1,2.0)
for i,(b,r) in enumerate(bl):
    p=para(tf,i==0,sb=(0 if i==0 else 11),ls=1.3); run(p,'•  ',12,ACCENT,True); run(p,b,12.5,FG,True); run(p,r,12.5,FG2)
note(s,6.4,'**출처:** 국내 디지털광고 11.3조 — KOBACO(2025). 오프라인 마케팅·측정비율 — 내부 추정.')
footer(s,'03')

# ═══════════════ P04 INSIGHT ═══════════════
s=slide(BG2)
header(s,'PART 1 · OPENING — INSIGHT','04','Behavior is the New Data',
 '**검색어는 속일 수 있지만, 행동은 속일 수 없다.** 미래 경쟁은 데이터의 양이 아니라 해상도(High Resolution)에서 갈린다.')
cardL=rect(s,LM,3.7,5.55,2.6,fill=BG3,ln=LINE2,rounded=True,radius=0.04)
tf=box(s,LM+0.3,3.95,5,2.2)
run(para(tf,True),'온라인 데이터',10,FG3,True)
p=para(tf,False,sb=6); run(p,'클릭 = 의도',24,FG2,True)
for t in ['무엇을 클릭/검색했나','표면적 관심 (저해상)','쿠키·추적에 의존','맥락 없는 단편']:
    p=para(tf,False,sb=8,ls=1.2); run(p,'•  ',11,ACCENT); run(p,t,11.5,FG2)
cardR=rect(s,6.78,3.7,5.55,2.6,fill=ACARD,ln=ACCENT,lw=1,rounded=True,radius=0.04)
tf=box(s,6.78+0.3,3.95,5,2.2)
run(para(tf,True),'오프라인 행동 데이터',10,ACCENTL,True)
p=para(tf,False,sb=6); run(p,'행동 = 진실',24,FG,True)
for t,b in [('방문·체류·반응·구매',False),('실제 행동의 맥락 (고해상)',True),('비식별 퍼스트파티',False),('‘왜 움직였는가’',False)]:
    p=para(tf,False,sb=8,ls=1.2); run(p,'•  ',11,ACCENT); run(p,t,11.5,(FG if b else FG2),b)
footer(s,'04')

# ═══════════════ P05 SOLUTION ═══════════════
s=slide()
header(s,'PART 2 · SOLUTION','05','Physical Ad-Tech Network',
 '**오프라인 공간을 광고 → 행동 → 구매 데이터가 연결되는 네트워크로 전환한다.** 지금까지 끊겨 있던 광고와 커머스를 하나의 측정 가능한 흐름으로 잇는다.')
flow(s,3.7,2.0,[('01','광고','브랜드 캠페인이 거점에서 시작'),('02','방문','고객이 현장에 진입'),('03','체류','머무름·동선이 데이터로'),
                ('04','반응','관심·선택을 비식별 기록'),('05','구매','전환 발생·매출'),('06','재방문','관계·리타겟으로 순환')])
callout(s,6.0,'모든 단계가 **행동 데이터**를 남긴다 — 광고가 곧 데이터 수집 장치가 된다.')
footer(s,'05')

# ═══════════════ P06 WHY DIFFERENT ═══════════════
s=slide()
header(s,'PART 2 · SOLUTION — WHY DIFFERENT','06','Most Companies Analyze Data. We Generate Data.',
 '**대부분은 데이터를 분석하지만, 우리는 데이터가 생성되는 현장을 직접 운영한다.** 데이터를 사오는 게 아니라 만든다 — 이것이 복제할 수 없는 진입장벽이다.', tsize=25)
rect(s,LM,3.9,4.3,2.4,fill=BG2,ln=LINE2,rounded=True,radius=0.04)
tf=box(s,LM+0.3,4.15,3.7,0.4); run(para(tf,True),'데이터가 생성되는 현장',10,FG3,True)
xx=LM+0.3
for chip in ['팝업','F&B','공간','운영','판매']:
    w=0.55+len(chip)*0.12
    rect(s,xx,4.65,w,0.38,fill=ACARD,ln=ACCENT,lw=0.75,rounded=True,radius=0.3)
    tf=box(s,xx,4.65,w,0.38,anchor=MSO_ANCHOR.MIDDLE); run(para(tf,True,align=PP_ALIGN.CENTER),chip,10.5,ACCENTL,True); xx+=w+0.12
tf=box(s,LM+0.3,5.4,3.7,0.8); p=para(tf,True,ls=1.35); run(p,'직영 ',12,FG2); run(p,'8거점',12,FG,True); run(p,' · ',12,FG2); run(p,'350+ 팝업·40+ 계약',12,FG,True); run(p,' — 데이터 소스를 ',12,FG2); run(p,'소유',12,ACCENTL,True); run(p,'한다.',12,FG2)
bl=[('데이터를 사오지 않는다 — 만든다.',' 현장을 직접 운영해 행동 데이터를 원천 생성'),
    ('풀스택 동시 실행.',' 기획·공간·F&B·운영·판매·데이터를 한 몸으로'),
    ('분석만 하는 회사와 다르다.',' 데이터가 태어나는 곳을 가졌다')]
tf=box(s,5.5,3.95,5.9,2.3)
for i,(b,r) in enumerate(bl):
    p=para(tf,i==0,sb=(0 if i==0 else 13),ls=1.3); run(p,'•  ',12,ACCENT,True); run(p,b,12.5,FG,True); run(p,r,12.5,FG2)
footer(s,'06')

# ═══════════════ P07 ENGINE ═══════════════
s=slide()
header(s,'PART 2 · SOLUTION — PRODUCT ARCHITECTURE','07','The Engine Behind the Network',
 '**광고를 측정하는 것이 아니라 행동을 측정한다.** Bracket → Physical AI → RXR → Attribution → Commerce로 이어지는 엔진. 기술은 목적이 아니라 네트워크의 심장이다.')
flow(s,3.9,1.9,[('01','Bracket','유휴 공간을 지능형 미디어 노드로'),('02','Physical AI','LAS 실시간 취향 분류 (32종+, 비식별)'),
                ('03','RXR','행동 분석 엔진 — 네트워크의 심장'),('04','Attribution','오프라인 경험 → 온라인 매출 증명'),('05','Commerce','전환·재방문으로 닫는 루프')])
note(s,6.1,'**IP:** 특허 제10-2545502호(라이프스타일 AI 세그멘테이션) · LAS 32종+ 비식별 처리 · CTO 전 퓨리오사 출신 · Alpha/Beta 구축 중.')
footer(s,'07')

# ═══════════════ P08 HIGH RES ═══════════════
s=slide(BG2)
header(s,'PART 2 · SOLUTION — CORE ASSET','08','The Data No One Else Can Build',
 '**이 데이터는 사올 수 없다 — 350개 이상의 라이브 리테일 환경을 직접 운영해야만 쌓인다.** 측정사는 외부 데이터에 의존하지만, 우리는 데이터가 생성되는 현장을 소유한다. 양이 아니라 1st-party의 질이 자산이다.')
cs=cols(4,3.9,2.3)
for (l,t,w,h),(num,ti,bd) in zip(cs,[('01','방문','누가, 언제, 어떤 경로로'),('02','체류','어디서, 얼마나 머물렀나'),('03','반응','무엇 앞에서 멈추고 선택'),('04','구매','경험이 전환으로 이어졌나')]):
    sh=rect(s,l,t,w,h,fill=BG3,ln=LINE2,rounded=True,radius=0.045); rect(s,l,t,w,0.035,fill=ACCENT)
    tf=box(s,l+0.22,t+0.3,w-0.4,h-0.5)
    p=para(tf,True); run(p,num,30,ACCENT)
    p=para(tf,False,sb=8); run(p,ti,16,FG,True)
    p=para(tf,False,sb=6,ls=1.25); run(p,bd,10.5,FG3)
note(s,6.35,'**350+ 팝업 · 40+ 기업 계약 리서치**에서 축적한 1st-party 오프라인 행동 데이터. ※ 고유 프로필·이벤트 실측치는 집계 후 갱신 [확인 필요].')
footer(s,'08')

# ═══════════════ P09 TRACTION ═══════════════
s=slide()
header(s,'PART 3 · TRACTION','09','Already Proven',
 '**우리는 이 구조를 설명하는 회사가 아니라, 이미 운영해본 회사다.** 외부 자본 없이 6년 흑자로 쌓은 실행력 위에 데이터 레이어를 얹는다.')
# bar chart 좌
tf=box(s,LM,3.7,6,0.3); run(para(tf,True),'매출 추이 (억원)',10,FG3,True)
bars=[('19',3.7,True),('20',11.5,True),('21',27,True),('22',27,True),('23',30,True),('24',40,False)]
bw=0.62; gap=0.32; bx=LM+0.1; base=6.2; maxh=2.0; mx=40
for lab,val,dim in bars:
    bh=val/mx*maxh
    rect(s,bx,base-bh,bw,bh,fill=(LINE if dim else ACCENT))
    tf=box(s,bx-0.1,base-bh-0.28,bw+0.2,0.25); run(para(tf,True,align=PP_ALIGN.CENTER),str(val),9.5,(FG if not dim else FG2),not dim)
    tf=box(s,bx-0.1,base+0.05,bw+0.2,0.25); run(para(tf,True,align=PP_ALIGN.CENTER),"'"+lab,9,(ACCENT if not dim else FG3))
    bx+=bw+gap
# KPI 우 2x2
ks=[('2024 매출','40','억'),('영업이익률','21.6','%'),("CAGR '19~'24",'183','%'),('연속 흑자','6','년')]
kc=cols(2,3.75,1.15,gap=0.25,l=7.0,total_w=RM-7.0)
kc2=cols(2,5.05,1.15,gap=0.25,l=7.0,total_w=RM-7.0)
for (l,t,w,h),(lb,v,u) in zip(kc,ks[:2]): kpi(s,l,t,w,h,lb,v,u)
for (l,t,w,h),(lb,v,u) in zip(kc2,ks[2:]): kpi(s,l,t,w,h,lb,v,u)
footer(s,'09')

# ═══════════════ P10 CASE ═══════════════
s=slide()
header(s,'PART 3 · TRACTION — CASE STUDY','10','Advertising That Generates Assets',
 '**광고비는 사라지는 비용이 아니라, 축적되는 자산이 될 수 있다.** 각 캠페인이 성과를 내는 동시에 행동 데이터로 축적된다 — ‘비용 → 자산’ 전환의 실증.')
cs=cols(4,3.9,2.2)
for (l,t,w,h),(lb,v,u,sub) in zip(cs,[('가나초콜릿하우스','4,500','%','마케팅 효율'),('Apple TV+ 파친코','88','%','시청 희망률'),('비욘드','100','%','재방문 희망률'),('삼성스토어','+44','%','방문객 활성화')]):
    kpi(s,l,t,w,h,lb,v,u,sub)
note(s,6.35,'**계약 기반 1st-party 데이터 프로젝트** — 현대 오픈디자인랩 · 삼성 바스켓스토어 · 애플TV · 가나초콜릿하우스 외. 고객 350+ 팝업·40+ 계약. ※수치 측정정의는 부록.')
footer(s,'10')

# ═══════════════ P11 WHY NOW ═══════════════
s=slide()
header(s,'PART 3 · TRACTION — WHY NOW','11','The Perfect Timing',
 '**네 변곡점이 동시에 열렸다.** 쿠키리스로 퍼스트파티 데이터 가치가 오르고, 리테일미디어가 성장하며, 팝업이 폭발하고, AI 비용이 급락한 지금이 카테고리 선점의 창이다.')
cs=cols(4,3.7,2.4)
data=[('01 · 데이터','쿠키리스','웹 ~50% 쿠키리스 + ATT → **퍼스트파티·오프라인 가치↑**'),
      ('02 · 시장','Retail Media','글로벌 **2030 $312B · CAGR 11%**'),
      ('03 · 행동','팝업 폭발','국내 팝업 **3,077개 (+80% YoY)**'),
      ('04 · 기술','AI 발전','행동 데이터 해석 비용 급락')]
for (l,t,w,h),(lb,ti,bd) in zip(cs,data):
    card(s,l,t,w,h,label=lb,title=ti,tsize=15,body=bd,accent_top=True)
note(s,6.3,'**출처:** Retail Media — Forrester(2025) $184B→$312B (Forrester 고유 정의 · 타 기관 추정 $37~57B). 팝업 3,077개 — 업계 집계(2024).')
footer(s,'11')

# ═══════════════ P12 BUSINESS MODEL ═══════════════
s=slide()
header(s,'PART 4 · BUSINESS MODEL','12','How We Make Money',
 '**데이터는 수익원이 아니라, 수익을 강화하는 자산이다.** 같은 데이터가 광고 · 커머스 · 인텔리전스 3중으로 수익화된다. 현 매출 40억은 광고 수익을 중심으로 발생 중.')
cs=cols(3,3.9,2.3)
bms=[('01 · AD REVENUE','광고 수익','광고비\n프로모션\n브랜드 캠페인',ACCENTL,False),
     ('02 · COMMERCE REVENUE','커머스 수익','판매 수수료\n특약매입\n위탁판매',GREEN,False),
     ('03 · INTELLIGENCE REVENUE','데이터 수익','데이터\n리포트\nSaaS 구독',ACCENTL,True)]
for (l,t,w,h),(lb,ti,bd,lc,acc) in zip(cs,bms):
    sh=rect(s,l,t,w,h,fill=(ACARD if acc else BG3),ln=(ACCENT if acc else LINE2),lw=(1 if acc else 0.75),rounded=True,radius=0.045)
    tf=sh.text_frame; tf.word_wrap=True; tf.margin_left=Inches(0.2); tf.margin_top=Inches(0.18)
    run(para(tf,True),lb,9.5,lc,True)
    p=para(tf,False,sb=8); run(p,ti,17,FG,True)
    for i,x in enumerate(bd.split('\n')):
        p=para(tf,False,sb=(9 if i==0 else 5),ls=1.2); run(p,'•  ',11,ACCENT); run(p,x,11.5,FG2)
footer(s,'12')

# ═══════════════ P13 J-CURVE ═══════════════
s=slide()
header(s,'PART 4 · BUSINESS MODEL — J-CURVE','13','From Advertising to Intelligence',
 '**광고에서 시작해 지능으로 끝난다.** 데이터가 쌓일수록 가속하는 비선형 성장. (SaaS 단위지표 NRR·payback은 [확인 필요 · Beta 후])')
cs=cols(3,3.9,2.3)
ph=[('PHASE 1 · 현재 40억','Physical Ad Inventory','브랜드가 거점을 광고 지면으로 구매 → 광고·공간 수익. 동시에 행동 데이터 축적.',True),
    ('PHASE 2 · ’26~’27','Behavior Data SaaS','취향·행동 데이터를 구독 모델로 — Recurring Revenue.',False),
    ('PHASE 3 · ’27~','Commerce Platform','방문자 전용 결제·추천 — Transaction Revenue.',False)]
for (l,t,w,h),(lb,ti,bd,now) in zip(cs,ph):
    rect(s,l,t,w,h,fill=BG3,ln=LINE2,rounded=True,radius=0.045)
    # phase badge
    bw=2.2 if now else 1.6
    rect(s,l+0.2,t+0.2,bw,0.34,fill=(ACARD if now else C('1A1A1A')),ln=(ACCENT if now else LINE2),lw=0.75,rounded=True,radius=0.3)
    tf=box(s,l+0.2,t+0.2,bw,0.34,anchor=MSO_ANCHOR.MIDDLE); run(para(tf,True,align=PP_ALIGN.CENTER),lb,9,(ACCENTL if now else FG3),True)
    tf=box(s,l+0.2,t+0.75,w-0.4,1.4)
    p=para(tf,True); run(p,ti,17,FG,True)
    p=para(tf,False,sb=8,ls=1.3); run(p,bd,11,FG2)
callout(s,6.05,'Phase 2는 새 사업이 아니다 — 기업은 **이미 계약형으로 우리 데이터를 구매 중**(현대·삼성·애플TV·가나). 이를 구독으로 제품화한다.')
footer(s,'13')

# ═══════════════ P14 FLYWHEEL ═══════════════
s=slide(BG2)
header(s,'PART 4 · BUSINESS MODEL — FLYWHEEL','14','The Network Effect',
 '**더 많은 거점이 더 많은 데이터를 만들고, 더 많은 데이터가 더 많은 브랜드를 부른다.** 선형 성장이 아니라 시간이 갈수록 벌어지는 데이터 네트워크 효과 — 그래서 플랫폼이다.')
# flywheel 좌측
cx,cy,r=3.6,4.95,1.45; nd=1.25
nodes=[('01','거점','방문이 시작'),('02','행동 데이터','방문·체류 수집'),('03','AI','RXR 해석'),('04','광고 효율','ROAS 증명'),('05','브랜드 유입','→ 더 많은 거점')]
for i,(num,ti,sub) in enumerate(nodes):
    a=math.radians(-90+i*72); nx=cx+r*math.cos(a); ny=cy+r*math.sin(a)
    o=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(nx-nd/2),Inches(ny-nd/2),Inches(nd),Inches(nd))
    o.fill.solid(); o.fill.fore_color.rgb=BG3; o.line.color.rgb=(ACCENT if i==0 else LINE); o.line.width=Pt(1.25 if i==0 else 0.75); o.shadow.inherit=False
    tf=o.text_frame; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=para(tf,True,align=PP_ALIGN.CENTER); run(p,num,9,ACCENT,True)
    p=para(tf,False,align=PP_ALIGN.CENTER,sb=2); run(p,ti,11,FG,True)
    p=para(tf,False,align=PP_ALIGN.CENTER,sb=1); run(p,sub,7.5,FG3)
tf=box(s,cx-1.0,cy-0.25,2.0,0.6,anchor=MSO_ANCHOR.MIDDLE)
p=para(tf,True,align=PP_ALIGN.CENTER); run(p,'Network Effect',12,ACCENT,True)
p=para(tf,False,align=PP_ALIGN.CENTER); run(p,'쓸수록 강해진다',9,FG3)
# 우측 텍스트
tf=box(s,6.6,4.4,5.0,1.8)
p=para(tf,True,ls=1.4); run(p,'더 많은 거점이 더 많은 데이터를 만들고,',15,FG,True)
p=para(tf,False,ls=1.4); run(p,'더 많은 데이터가 ',15,FG,True); run(p,'더 많은 브랜드를 부른다.',15,ACCENT,True)
callout(s,6.0,'선형 성장(매장×매출)이 아니라 **비선형 성장(데이터 네트워크 효과)**. 시간이 갈수록 벌어지는 해자.')
footer(s,'14')

# ═══════════════ P15 MARKET EXPANSION ═══════════════
s=slide()
header(s,'PART 4 · BUSINESS MODEL — MARKET EXPANSION','15','The Future of Distribution',
 '**우리는 광고 플랫폼이 아니라, 브랜드가 발견되고 성장하는 인프라를 만든다.** 한국에서 검증한 네트워크를 글로벌 유통 인프라로 — $312B 시장으로 확장하는 10배 스토리.')
flow(s,3.9,1.9,[('01','한국','성수 직영 거점에서 검증 — 350+ 팝업'),('02','일본','기존 글로벌 실적 (더셀렉츠 도쿄)'),
                ('03','미국','Apple TV+ 파친코 등 글로벌 IP 협업'),('04','글로벌 네트워크','발견·성장·유통이 일어나는 인프라')])
footer(s,'15')

# ═══════════════ P16 COMPETITION & WEDGE ═══════════════
s=slide()
header(s,'PART 5 · WHY WE WIN — COMPETITION','16','Everyone Measures. We Operate.',
 '**경쟁자는 발자국을 측정한다. 우리는 발자국이 생기는 현장을 운영한다.** 측정만 하는 회사, 자사 채널에 갇힌 RMN과 달리 — 우리는 350+ 현장을 직접 소유·운영하며 1st-party 데이터를 생성한다.')
rows=[('구분','Cosmose · 측정류','대기업 RMN','in-store 스크린','Project Rent'),
      ('방식','행동 측정만','자사 채널 광고','노출만','현장 소유·운영'),
      ('데이터 소유','외부 의존','자사 한정','없음','350+ 현장 생성'),
      ('1st-party','제한적','자사몰','없음','오프라인 1st-party'),
      ('현장','미소유','자사 매장','임대 면적','직영 8 + 350+')]
tb=s.shapes.add_table(len(rows),5,Inches(LM),Inches(3.7),Inches(CW),Inches(2.3)).table
tb.columns[0].width=Inches(1.6)
for ci in range(1,4): tb.columns[ci].width=Inches(2.3)
tb.columns[4].width=Inches(CW-1.6-3*2.3)
for ri,row in enumerate(rows):
    tb.rows[ri].height=Inches(0.44)
    for ci,val in enumerate(row):
        cell=tb.cell(ri,ci); cell.fill.solid()
        own=(ci==4)
        if ri==0: cell.fill.fore_color.rgb=(C('1E1B2E') if own else BG2)
        else: cell.fill.fore_color.rgb=(C('1E1B2E') if own else C('1E1E1E'))
        cell.margin_left=Inches(0.13); cell.margin_top=Inches(0.03); cell.margin_bottom=Inches(0.03); cell.vertical_anchor=MSO_ANCHOR.MIDDLE
        tf=cell.text_frame; tf.word_wrap=True; p=tf.paragraphs[0]
        col=(ACCENTL if (ri==0 and own) else FG3 if ri==0 else (FG if own else FG2))
        run(p,val,9.5,col,(ri==0 or own))
callout(s,6.18,'카테고리에 **현장을 소유·운영하는 네트워크형 승자는 아직 없다** — 우리가 그 자리를 만든다.',h=0.5)
footer(s,'16')

# ═══════════════ P17 WHY PROJECT RENT ═══════════════
s=slide()
header(s,'PART 5 · WHY WE WIN','17','Why Project Rent',
 '**비전을 가진 회사는 많지만, 우리는 이미 비전의 일부를 운영하고 있다.** 기획·공간·F&B·운영·판매·데이터를 한 몸으로 실행하는 풀스택 — 분석만, 중개만 하는 곳과 다른 희소함.')
rect(s,LM,3.95,CW,1.0,fill=BG2,ln=LINE2,rounded=True,radius=0.04)
tf=box(s,LM+0.3,4.12,4,0.35); run(para(tf,True),'통합 역량 (풀스택)',10,FG3,True)
xx=LM+0.3
for chip in ['기획','공간','F&B','운영','판매','데이터']:
    w=0.55+len(chip)*0.13
    rect(s,xx,4.5,w,0.38,fill=ACARD,ln=ACCENT,lw=0.75,rounded=True,radius=0.3)
    tf=box(s,xx,4.5,w,0.38,anchor=MSO_ANCHOR.MIDDLE); run(para(tf,True,align=PP_ALIGN.CENTER),chip,11,ACCENTL,True); xx+=w+0.14
bl=[('풀스택 동시 실행은 희소하다.',' 기획부터 데이터까지 한 팀이 한 흐름으로'),
    ('분석만 하는 곳과 다르다.',' 데이터가 태어나는 현장을 소유'),
    ('중개만 하는 곳과 다르다.',' 공간을 운영하며 행동·구매까지 연결')]
tf=box(s,LM,5.3,CW,1.4)
for i,(b,r) in enumerate(bl):
    p=para(tf,i==0,sb=(0 if i==0 else 12),ls=1.3); run(p,'•  ',12,ACCENT,True); run(p,b,13,FG,True); run(p,r,13,FG2)
footer(s,'17')

# ═══════════════ P18 TEAM ═══════════════
s=slide()
header(s,'PART 5 · WHY WE WIN — TEAM','18','Built by Operators',
 '**우리는 컨설턴트가 아니라 운영자다.** 350+ 팝업·40+ 기업 계약으로 오프라인 리테일 시장을 직접 만든 팀 — 비전을 말하는 게 아니라 운영해 왔다.')
rect(s,LM,3.9,4.3,2.4,fill=BG2,ln=LINE2,rounded=True,radius=0.04)
tf=box(s,LM+0.3,4.15,3.7,2.0)
run(para(tf,True),'FOUNDER · CEO',10,ACCENTL,True)
p=para(tf,False,sb=8); run(p,'최원석',22,FG,True)
p=para(tf,False,sb=8,ls=1.35); run(p,'산업디자이너 → 대기업 브랜더 → 필라멘트앤코 → 프로젝트렌트. 공간·브랜드·비즈니스를 잇는 ',12,FG2); run(p,'비즈니스 디자이너',12,ACCENTL,True); run(p,'. "성수동의 스티브 잡스"(KBS).',12,FG2)
tcards=[('CTO','전 퓨리오사 출신','Physical AI 엔진 리드'),('ADVISORY','행동과학 R&D','행동·취향 모델링')]
for (l,t,w,h),(lb,ti,bd) in zip(cols(2,3.9,1.15,gap=0.25,l=5.5,total_w=RM-5.5),tcards):
    card(s,l,t,w,h,label=lb,title=ti,tsize=14,body=bd,bsize=10.5)
card(s,5.5,5.15,RM-5.5,1.15,label='조직',title='6명 · 전원 AX · 인당 6.7억',tsize=14,body='미디어: KBS · 메종코리아 · 롱블랙. 소수 정예 운영 조직.',bsize=10.5)
footer(s,'18')

# ═══════════════ P19 VALUATION BRIDGE ═══════════════
s=slide(BG2)
header(s,'PART 5 · WHY WE WIN — VALUATION','19',"What You're Actually Buying",
 '**에이전시 경제가 데이터 자산을 키운다. 이 라운드는 그 자산이 SaaS가 되는 순간을 산다.** 40억 에이전시 가격이 아니라, 흑자로 굴러가며 1st-party 데이터를 쌓는 회사의 옵션 가치다.')
cs=cols(3,3.9,2.3)
vb=[('① 하방 — 닫혔다','40억 흑자 사업','6년 연속 흑자·영업이익률 21.6%가 데이터 자산을 자가 생성. 다운사이드가 닫혀 있다.',False),
    ('② 자산 — 독점','1st-party 데이터','350+ 현장 · 40+ 계약에서 나오는, 아무도 못 만드는 오프라인 행동 데이터.',False),
    ('③ 상방 — 열렸다','데이터 SaaS','이미 계약형으로 팔리는 데이터를 구독 제품화 — 여기가 당신이 사는 업사이드.',True)]
for (l,t,w,h),(lb,ti,bd,acc) in zip(cs,vb):
    sh=rect(s,l,t,w,h,fill=(ACARD if acc else BG3),ln=(ACCENT if acc else LINE2),lw=(1 if acc else 0.75),rounded=True,radius=0.045)
    tf=sh.text_frame; tf.word_wrap=True; tf.margin_left=Inches(0.2); tf.margin_top=Inches(0.18); tf.margin_right=Inches(0.16)
    run(para(tf,True),lb,9.5,(ACCENTL if acc else FG3),True)
    p=para(tf,False,sb=8); run(p,ti,16,FG,True)
    p=para(tf,False,sb=8,ls=1.3); run(p,bd,10.5,FG2)
callout(s,6.2,'Pre 200억 = 하방 닫힌 흑자 사업 + 상방 열린 데이터 SaaS의 **비대칭 베팅**.',h=0.5)
footer(s,'19')

# ═══════════════ P20 INVESTMENT & CLOSING ═══════════════
s=slide(BG2)
header(s,'PART 5 · WHY WE WIN — INVESTMENT & CLOSING','20','Building the Offline Data Network',
 '**Google은 웹을 인덱싱했다. 우리는 현실 세계의 소비 행동을 인덱싱한다.** 우리는 광고를 측정하는 회사가 아니라, 사람의 행동을 측정하는 회사입니다.')
# budget 좌
tf=box(s,LM,3.8,5.5,0.3); run(para(tf,True),'THE ASK · 자금 사용처',10,FG3,True)
bud=[('RXR 엔진 · 데이터 인프라',40,'8억 (40%)',ACCENT),('거점 · 커머스 확장',25,'5억 (25%)',GREEN),('인력 (CTO·데이터)',15,'3억 (15%)',AMBER),('거점 자산화',10,'2억 (10%)',RUST),('운영자금 (24M)',10,'2억 (10%)',FG3)]
by=4.2
for lab,pct,val,col in bud:
    tf=box(s,LM,by,2.3,0.3,anchor=MSO_ANCHOR.MIDDLE); run(para(tf,True),lab,10.5,FG,True)
    rect(s,LM+2.4,by+0.02,2.3,0.24,fill=BG3,ln=LINE)
    rect(s,LM+2.4,by+0.02,2.3*pct/100,0.24,fill=col)
    tf=box(s,LM+4.8,by,0.85,0.3,anchor=MSO_ANCHOR.MIDDLE); run(para(tf,True,align=PP_ALIGN.RIGHT),val,10.5,col,True)
    by+=0.46
# KPI 우 3
for (l,t,w,h),(lb,v,u) in zip(cols(3,3.85,1.2,gap=0.2,l=6.6,total_w=RM-6.6),[('Pre-money','200','억'),('Raising','20','억'),('Runway','24','M')]):
    kpi(s,l,t,w,h,lb,v,u)
tf=box(s,6.6,5.25,5.0,0.3); run(para(tf,True),'미래 비전 — 하나의 흐름',10,FG3,True)
xx=6.6
for chip in ['광고','행동','구매','데이터','지능','유통']:
    w=0.5+len(chip)*0.13
    rect(s,xx,5.65,w,0.38,fill=ACARD,ln=ACCENT,lw=0.75,rounded=True,radius=0.3)
    tf=box(s,xx,5.65,w,0.38,anchor=MSO_ANCHOR.MIDDLE); run(para(tf,True,align=PP_ALIGN.CENTER),chip,10,ACCENTL,True); xx+=w+0.1
tf=box(s,6.6,6.2,5.0,0.3); run(para(tf,True),'문의: ws.choi@project-rent.com · 010-2870-0561',10,FG3)
footer(s,'20')

out='/Users/choi_ai/do-better-workspace/20-operations/R_define/프로젝트렌트IR_V10-adtech.pptx'
prs.save(out)
print('saved:', out, '| slides:', len(prs.slides._sldIdLst))
